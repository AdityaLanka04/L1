"""
Comprehensive Slide Analysis System
Batches multiple slides per AI call to avoid rate limits.
"""

import os
import json
import logging
import time
from typing import List, Dict, Any, Optional
from pathlib import Path
import PyPDF2
from groq import Groq
from sqlalchemy.orm import Session
from datetime import datetime, timezone
from activity_logger import log_ai_tokens
from ai_usage import extract_usage_from_openai_like

logger = logging.getLogger(__name__)

groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

BATCH_SIZE = 4
BATCH_DELAY = 3.0


class ComprehensiveSlideAnalyzer:
    def __init__(self, db: Session):
        self.db = db
        self.model = "llama-3.3-70b-versatile"
        self.current_user_id = None

    def call_ai(self, prompt: str, max_tokens: int = 3000, temperature: float = 0.3, retries: int = 3) -> str:
        for attempt in range(retries):
            try:
                response = groq_client.chat.completions.create(
                    model=self.model,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=max_tokens,
                    temperature=temperature
                )
                self._log_usage(response)
                return response.choices[0].message.content
            except Exception as e:
                err = str(e)
                logger.error(f"AI call error (attempt {attempt + 1}/{retries}): {err}")
                if "429" in err or "rate" in err.lower():
                    wait = (attempt + 1) * 5
                    logger.info(f"Rate limited — waiting {wait}s before retry")
                    time.sleep(wait)
                elif attempt < retries - 1:
                    time.sleep(2)
                else:
                    return ""
        return ""

    def _log_usage(self, response):
        if not self.current_user_id:
            return
        usage = extract_usage_from_openai_like(response)
        if not usage:
            return
        try:
            log_ai_tokens(
                user_id=self.current_user_id,
                tool_name="slide_explorer_ai",
                prompt_tokens=usage.get("prompt_tokens", 0),
                completion_tokens=usage.get("completion_tokens", 0),
                total_tokens=usage.get("total_tokens", 0),
                model=self.model,
                metadata={"provider": "groq", "source": "slide_analysis"}
            )
        except Exception:
            pass

    def extract_slide_content(self, file_path: Path, file_type: str) -> List[Dict[str, Any]]:
        slides_data = []

        if file_type == 'pdf':
            try:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        text = page.extract_text() or ""
                        slides_data.append({
                            "slide_number": page_num,
                            "content": text.strip(),
                            "title": f"Page {page_num}"
                        })
            except Exception as e:
                logger.error(f"Error reading PDF: {e}")
                raise

        elif file_type in ['ppt', 'pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                for slide_num, ppt_slide in enumerate(prs.slides, 1):
                    slide_text = []
                    slide_title = f"Slide {slide_num}"
                    for shape in ppt_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                                if shape.placeholder_format.type == 1:
                                    slide_title = shape.text.strip()
                            slide_text.append(shape.text.strip())
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = [cell.text for cell in row.cells if cell.text]
                                if row_text:
                                    slide_text.append(" | ".join(row_text))
                    slides_data.append({
                        "slide_number": slide_num,
                        "content": "\n".join(slide_text),
                        "title": slide_title
                    })
            except ImportError:
                logger.error("python-pptx not installed")
                raise
            except Exception as e:
                logger.error(f"Error reading PowerPoint: {e}")
                raise

        return slides_data

    def _fallback_analysis(self, slide_number: int, title: str) -> Dict[str, Any]:
        return {
            "slide_number": slide_number,
            "title": title,
            "detailed_explanation": f"This slide covers {title}. Review the surrounding slides for full context.",
            "key_concepts": [f"Main topic: {title}"],
            "definitions": {},
            "exam_questions": [{"question": f"What are the main points in {title}?", "type": "short-answer", "difficulty": "medium", "answer_hint": "Review the slide carefully"}],
            "practical_applications": [],
            "common_misconceptions": [],
            "study_tips": ["Re-read the slide and take notes"],
            "cross_references": [],
            "difficulty_level": "intermediate",
            "estimated_study_time": "3 minutes"
        }

    def analyze_batch(self, batch: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Analyze a batch of slides in a single AI call."""
        slides_block = ""
        for s in batch:
            content = s["content"][:800] if s["content"] else "(no text content)"
            slides_block += f"\n---\nSLIDE {s['slide_number']}: {s['title']}\n{content}\n"

        prompt = f"""You are an expert educational content analyzer. Analyze the following {len(batch)} presentation slides and return a JSON array with one object per slide.

SLIDES:
{slides_block}

For each slide return this JSON object:
{{
  "slide_number": <number>,
  "title": "<title>",
  "detailed_explanation": "## Overview\n<one paragraph intro>\n\n## Key Points\n- <point 1>\n- <point 2>\n- <point 3>\n\n## Details\n<one paragraph deeper explanation>",
  "key_concepts": ["<concept1>", "<concept2>", "<concept3>"],
  "definitions": {{"<term>": "<definition>"}},
  "exam_questions": [{{"question": "<q>", "type": "short-answer|essay|multiple-choice", "difficulty": "easy|medium|hard", "answer_hint": "<hint>"}}],
  "practical_applications": ["<application>"],
  "common_misconceptions": ["<misconception>"],
  "study_tips": ["<tip>"],
  "cross_references": ["<reference>"],
  "difficulty_level": "introductory|intermediate|advanced",
  "estimated_study_time": "<X> minutes"
}}

Return ONLY a valid JSON array of {len(batch)} objects. No extra text."""

        raw = self.call_ai(prompt, max_tokens=3500, temperature=0.3)

        import re
        results = []

        match = re.search(r'\[[\s\S]*\]', raw)
        if match:
            try:
                cleaned = re.sub(r'[\x00-\x1f\x7f-\x9f]', ' ', match.group())
                parsed = json.loads(cleaned)
                if isinstance(parsed, list):
                    for item in parsed:
                        item["slide_number"] = int(item.get("slide_number", 0))
                    return parsed
            except json.JSONDecodeError as e:
                logger.error(f"Batch JSON decode error: {e}")

        logger.warning("Batch parse failed — falling back to per-slide analysis")
        for s in batch:
            results.append(self._analyze_single_slide(s))
        return results

    def _analyze_single_slide(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """Single-slide fallback used only when batch parse fails."""
        slide_number = slide_data["slide_number"]
        title = slide_data["title"]
        content = slide_data["content"]

        if not content or len(content.strip()) < 20:
            return self._fallback_analysis(slide_number, title)

        prompt = f"""Analyze this slide and return ONLY a single JSON object. The detailed_explanation must use markdown: start with ## Overview, then ## Key Points with bullet list, then ## Details paragraph.

Slide {slide_number}: {title}
{content[:1200]}

Return:
{{"slide_number":{slide_number},"title":"{title}","detailed_explanation":"## Overview\\n<intro para>\\n\\n## Key Points\\n- <point1>\\n- <point2>\\n- <point3>\\n\\n## Details\\n<deeper para>","key_concepts":["c1","c2","c3"],"definitions":{{}},"exam_questions":[{{"question":"Q","type":"short-answer","difficulty":"medium","answer_hint":"hint"}}],"practical_applications":[],"common_misconceptions":[],"study_tips":["tip"],"cross_references":[],"difficulty_level":"intermediate","estimated_study_time":"5 minutes"}}"""

        raw = self.call_ai(prompt, max_tokens=900, temperature=0.3)
        import re
        match = re.search(r'\{[\s\S]*\}', raw)
        if match:
            try:
                cleaned = re.sub(r'[\x00-\x1f\x7f-\x9f]', ' ', match.group())
                result = json.loads(cleaned)
                result["slide_number"] = slide_number
                result["title"] = title
                return result
            except Exception:
                pass
        return self._fallback_analysis(slide_number, title)

    def analyze_presentation(
        self,
        slide_id: int,
        file_path: Path,
        file_type: str,
        force_reanalyze: bool = False
    ) -> Dict[str, Any]:
        from models import SlideAnalysis, UploadedSlide

        try:
            slide = self.db.query(UploadedSlide).filter(UploadedSlide.id == slide_id).first()
            if slide:
                self.current_user_id = slide.user_id
        except Exception:
            self.current_user_id = None

        if not force_reanalyze:
            existing = self.db.query(SlideAnalysis).filter(SlideAnalysis.slide_id == slide_id).first()
            if existing and existing.analysis_data:
                logger.info(f"Returning cached analysis for slide_id {slide_id}")
                return json.loads(existing.analysis_data)

        logger.info(f"Extracting content from {file_path}")
        slides_data = self.extract_slide_content(file_path, file_type)
        total_slides = len(slides_data)
        logger.info(f"Analyzing {total_slides} slides in batches of {BATCH_SIZE}")

        analyzed_slides = []
        batches = [slides_data[i:i + BATCH_SIZE] for i in range(0, total_slides, BATCH_SIZE)]

        for batch_idx, batch in enumerate(batches):
            slide_nums = [s["slide_number"] for s in batch]
            logger.info(f"Analyzing batch {batch_idx + 1}/{len(batches)}: slides {slide_nums}")
            batch_results = self.analyze_batch(batch)

            slide_map = {s["slide_number"]: s for s in batch}
            for result in batch_results:
                analyzed_slides.append(result)

            if batch_idx < len(batches) - 1:
                logger.info(f"Waiting {BATCH_DELAY}s between batches")
                time.sleep(BATCH_DELAY)

        analyzed_slides.sort(key=lambda x: x.get("slide_number", 0))

        presentation_summary = self._generate_presentation_summary(analyzed_slides)
        result = {
            "slide_id": slide_id,
            "total_slides": total_slides,
            "analyzed_at": datetime.now(timezone.utc).isoformat(),
            "presentation_summary": presentation_summary,
            "slides": analyzed_slides
        }

        self._store_analysis(slide_id, result)
        return result

    def _generate_presentation_summary(self, analyzed_slides: List[Dict[str, Any]]) -> Dict[str, Any]:
        all_concepts = []
        all_questions = []
        difficulty_counts = {"introductory": 0, "intermediate": 0, "advanced": 0}
        total_study_time = 0

        for slide in analyzed_slides:
            all_concepts.extend(slide.get("key_concepts", []))
            all_questions.extend(slide.get("exam_questions", []))
            difficulty = slide.get("difficulty_level", "intermediate")
            difficulty_counts[difficulty] = difficulty_counts.get(difficulty, 0) + 1
            time_str = slide.get("estimated_study_time", "5 minutes")
            try:
                total_study_time += int(''.join(filter(str.isdigit, time_str)))
            except Exception:
                total_study_time += 5

        return {
            "total_concepts": len(set(all_concepts)),
            "total_exam_questions": len(all_questions),
            "difficulty_distribution": difficulty_counts,
            "estimated_total_study_time": f"{total_study_time} minutes",
            "recommended_review_sessions": max(1, total_study_time // 30)
        }

    def _store_analysis(self, slide_id: int, analysis_data: Dict[str, Any]):
        from models import SlideAnalysis
        try:
            existing = self.db.query(SlideAnalysis).filter(SlideAnalysis.slide_id == slide_id).first()
            if existing:
                existing.analysis_data = json.dumps(analysis_data)
                existing.analyzed_at = datetime.now(timezone.utc)
            else:
                self.db.add(SlideAnalysis(
                    slide_id=slide_id,
                    analysis_data=json.dumps(analysis_data),
                    analyzed_at=datetime.now(timezone.utc)
                ))
            self.db.commit()
            logger.info(f"Stored analysis for slide_id {slide_id}")
        except Exception as e:
            logger.error(f"Error storing analysis: {e}")
            self.db.rollback()


def get_or_create_analysis(
    slide_id: int,
    file_path: Path,
    file_type: str,
    db: Session,
    force_reanalyze: bool = False
) -> Dict[str, Any]:
    analyzer = ComprehensiveSlideAnalyzer(db)
    return analyzer.analyze_presentation(slide_id, file_path, file_type, force_reanalyze)
