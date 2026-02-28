import os
import io
import json
import logging
import tempfile
import traceback
import subprocess
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional

import PyPDF2
from fastapi import APIRouter, Body, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, Response
from sqlalchemy.orm import Session

import models
from database import get_db
from deps import call_ai, get_current_user, get_user_by_email, get_user_by_username

try:
    from ai_media_processor import ai_media_processor
except ImportError:
    ai_media_processor = None

try:
    from math_processor import process_math_in_response
except ImportError:
    process_math_in_response = None

logger = logging.getLogger(__name__)

UPLOAD_DIR = Path("uploads/slides")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

router = APIRouter(prefix="/api", tags=["media"])

@router.post("/transcribe_audio/")
async def transcribe_audio(
    audio_file: UploadFile = File(...),
    user_id: str = Form(...),
    db: Session = Depends(get_db),
):
    temp_audio_path = None
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        audio_content = await audio_file.read()
        if len(audio_content) == 0:
            raise HTTPException(status_code=400, detail="Empty audio file")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".webm", mode="wb") as temp_file:
            temp_file.write(audio_content)
            temp_audio_path = temp_file.name

        try:
            from groq import Groq
            groq_api_key = os.getenv("GROQ_API_KEY")
            groq_client = Groq(api_key=groq_api_key) if groq_api_key else None
            if not groq_client:
                raise HTTPException(status_code=500, detail="Transcription service not configured")

            with open(temp_audio_path, "rb") as f:
                transcription = groq_client.audio.transcriptions.create(
                    file=f,
                    model="whisper-large-v3-turbo",
                    response_format="json",
                    language="en",
                )

            transcript_text = transcription.text
            return {
                "status": "success",
                "transcript": transcript_text,
                "length": len(transcript_text),
                "model_used": "whisper-large-v3-turbo",
            }

        except HTTPException:
            raise
        except Exception as groq_error:
            logger.error(f"Groq API error: {str(groq_error)}")
            raise HTTPException(status_code=500, detail=f"Transcription failed: {str(groq_error)}")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in transcribe_audio: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to transcribe audio: {str(e)}")
    finally:
        if temp_audio_path and os.path.exists(temp_audio_path):
            try:
                os.remove(temp_audio_path)
            except Exception as cleanup_error:
                logger.warning(f"Failed to cleanup temp file: {cleanup_error}")

@router.get("/test_transcribe")
def test_transcribe_endpoint():
    groq_api_key = os.getenv("GROQ_API_KEY")
    return {
        "status": "endpoint exists",
        "message": "Transcribe audio endpoint is registered",
        "groq_available": groq_api_key is not None,
    }

@router.post("/transcribe_audio_test/")
async def transcribe_audio_test(
    user_id: str = Form(...),
    db: Session = Depends(get_db),
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        groq_api_key = os.getenv("GROQ_API_KEY")
        return {
            "status": "success",
            "message": "Endpoint working, user found",
            "user_id": user.id,
            "groq_configured": groq_api_key is not None,
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Test endpoint error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/media/process")
async def process_media(
    user_id: str = Form(...),
    file: UploadFile = File(None),
    youtube_url: str = Form(None),
    note_style: str = Form("detailed"),
    difficulty: str = Form("intermediate"),
    subject: str = Form("general"),
    custom_instructions: str = Form(None),
    generate_flashcards: bool = Form(False),
    generate_quiz: bool = Form(False),
    db: Session = Depends(get_db),
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        if not ai_media_processor:
            raise HTTPException(status_code=500, detail="Media processor not available")

        transcript_data = None
        analysis_data = None
        source_type = None
        filename = None
        file_path = None

        logger.info(f"Received - File: {file.filename if file else None}, YouTube URL: '{youtube_url}'")

        if youtube_url and youtube_url.strip():
            logger.info(f"Processing YouTube URL: {youtube_url}")
            result = await ai_media_processor.process_youtube_video(youtube_url.strip())

            if not result.get("success"):
                raise HTTPException(status_code=400, detail=result.get("error", "Failed to process YouTube video"))

            transcript_data = result
            source_type = "youtube"
            filename = result["video_info"]["title"]

        elif file:
            logger.info(f"Processing uploaded file: {file.filename}")

            upload_dir = "backend/uploads/media"
            os.makedirs(upload_dir, exist_ok=True)

            safe_name = "".join(c for c in (file.filename or "upload") if c.isalnum() or c in (".", "_", "-"))
            file_path = os.path.join(
                upload_dir,
                f"{user.id}_{int(datetime.now().timestamp())}_{safe_name}",
            )

            with open(file_path, "wb") as f:
                content = await file.read()
                f.write(content)

            result = await ai_media_processor.transcribe_audio_groq(file_path)

            if not result.get("success"):
                raise HTTPException(status_code=400, detail=result.get("error", "Failed to transcribe audio"))

            transcript_data = result
            source_type = "upload"
            filename = file.filename

        else:
            raise HTTPException(status_code=400, detail="Please provide either a file or YouTube URL")

        logger.info("Performing AI analysis...")
        analysis_result = await ai_media_processor.analyze_transcript_ai(
            transcript_data["transcript"],
            {"subject": subject, "difficulty": difficulty},
            user_id=user.id,
        )

        if not analysis_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"AI analysis failed: {analysis_result.get('error', 'Unknown error')}",
            )

        analysis_data = analysis_result.get("analysis", {})

        logger.info(f"Generating {note_style} notes...")
        notes_result = await ai_media_processor.generate_notes_ai(
            transcript_data["transcript"],
            analysis_data,
            note_style,
            {
                "difficulty": difficulty,
                "subject": subject,
                "custom_instructions": custom_instructions,
            },
            user_id=user.id,
        )

        if not notes_result.get("success"):
            raise HTTPException(status_code=500, detail="Note generation failed")

        flashcards = []
        if generate_flashcards:
            logger.info("Generating flashcards...")
            flashcard_result = await ai_media_processor.generate_flashcards_ai(
                transcript_data["transcript"],
                analysis_data,
                count=10,
                user_id=user.id,
            )
            if flashcard_result.get("success"):
                flashcards = flashcard_result.get("flashcards", [])
            else:
                logger.warning(f"Flashcard generation failed: {flashcard_result.get('error', 'Unknown')}")

        quiz_questions = []
        if generate_quiz:
            logger.info("Generating quiz...")
            quiz_result = await ai_media_processor.generate_quiz_ai(
                transcript_data["transcript"],
                analysis_data,
                count=10,
                user_id=user.id,
            )
            if quiz_result.get("success"):
                quiz_questions = quiz_result.get("questions", [])
            else:
                logger.warning(f"Quiz generation failed: {quiz_result.get('error', 'Unknown')}")

        key_moments = []
        if transcript_data.get("has_timestamps") and transcript_data.get("segments"):
            key_moments = await ai_media_processor.extract_key_moments(
                transcript_data["segments"],
                analysis_data,
            )

        response = {
            "success": True,
            "source_type": source_type,
            "filename": filename,
            "transcript": transcript_data["transcript"],
            "language": transcript_data.get("language", "en"),
            "language_name": ai_media_processor.get_language_name(transcript_data.get("language", "en")),
            "duration": transcript_data.get("duration", 0),
            "has_timestamps": transcript_data.get("has_timestamps", False),
            "segments": transcript_data.get("segments", [])[:50],
            "analysis": analysis_data,
            "notes": {
                "content": notes_result["content"],
                "style": note_style,
            },
            "flashcards": flashcards,
            "quiz_questions": quiz_questions,
            "key_moments": key_moments,
            "video_info": transcript_data.get("video_info") if source_type == "youtube" else None,
        }

        return response

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Media processing error: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

@router.post("/media/save-notes")
async def save_media_notes(
    user_id: str = Body(...),
    title: str = Body(...),
    content: str = Body(...),
    transcript: str = Body(None),
    analysis: dict = Body(None),
    flashcards: list = Body(None),
    quiz_questions: list = Body(None),
    key_moments: list = Body(None),
    db: Session = Depends(get_db),
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        new_note = models.Note(
            user_id=user.id,
            title=title,
            content=content,
            custom_font="__MEDIA_NOTE__",
            transcript=transcript,
            analysis=json.dumps(analysis) if analysis else None,
            flashcards=json.dumps(flashcards) if flashcards else None,
            quiz_questions=json.dumps(quiz_questions) if quiz_questions else None,
            key_moments=json.dumps(key_moments) if key_moments else None,
            created_at=datetime.now(timezone.utc),
            updated_at=datetime.now(timezone.utc),
        )

        db.add(new_note)
        db.commit()
        db.refresh(new_note)

        if flashcards and len(flashcards) > 0:
            flashcard_set = models.FlashcardSet(
                user_id=user.id,
                title=f"Flashcards: {title}",
                source_type="media",
                source_id=new_note.id,
            )
            db.add(flashcard_set)
            db.commit()
            db.refresh(flashcard_set)

            for fc in flashcards:
                flashcard = models.Flashcard(
                    set_id=flashcard_set.id,
                    question=fc.get("question", ""),
                    answer=fc.get("answer", ""),
                    difficulty=fc.get("difficulty", "medium"),
                )
                db.add(flashcard)

            db.commit()

        from gamification_system import award_points
        award_points(db, user.id, "note_created")

        return {
            "success": True,
            "note_id": new_note.id,
            "message": "Notes saved successfully",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error saving notes: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to save notes: {str(e)}")

@router.post("/media/regenerate-notes")
async def regenerate_notes(
    transcript: str = Body(...),
    analysis: dict = Body(...),
    note_style: str = Body("detailed"),
    difficulty: str = Body("intermediate"),
    subject: str = Body("general"),
    custom_instructions: str = Body(None),
    current_user: models.User = Depends(get_current_user),
):
    try:
        if not ai_media_processor:
            raise HTTPException(status_code=500, detail="Media processor not available")

        notes_result = await ai_media_processor.generate_notes_ai(
            transcript,
            analysis,
            note_style,
            {
                "difficulty": difficulty,
                "subject": subject,
                "custom_instructions": custom_instructions,
            },
            user_id=current_user.id,
        )

        if not notes_result.get("success"):
            raise HTTPException(status_code=500, detail="Note regeneration failed")

        return {
            "success": True,
            "content": notes_result["content"],
            "style": note_style,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error regenerating notes: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to regenerate notes: {str(e)}")

@router.get("/media/estimate-cost")
async def estimate_processing_cost(
    duration_seconds: int = Query(...),
    file_size_mb: float = Query(...),
):
    if not ai_media_processor:
        raise HTTPException(status_code=500, detail="Media processor not available")
    cost_info = ai_media_processor.estimate_processing_cost(duration_seconds, file_size_mb)
    return cost_info

@router.post("/media/generate-title")
async def generate_smart_title(
    transcript: str = Body(...),
    key_concepts: list = Body([]),
    summary: str = Body(""),
):
    try:
        prompt = f"""Generate a concise 3-4 word title for these notes.

Summary: {summary[:200]}
Key Concepts: {', '.join(key_concepts[:5])}
Content: {transcript[:500]}

Return ONLY the title, nothing else. Make it descriptive and catchy."""

        title = call_ai(prompt, max_tokens=20, temperature=0.7)
        title = title.strip().strip('"').strip("'")
        words = title.split()
        if len(words) > 4:
            title = " ".join(words[:4])

        return {"title": title}

    except Exception as e:
        logger.error(f"Title generation error: {str(e)}")
        return {"title": "Media Notes"}

@router.get("/get_note/{note_id}")
async def get_single_note(
    note_id: int,
    db: Session = Depends(get_db),
):
    try:
        note = db.query(models.Note).filter(
            models.Note.id == note_id,
            models.Note.is_deleted == False,
        ).first()

        if not note:
            raise HTTPException(status_code=404, detail="Note not found")

        analysis = {}
        flashcards = []
        quiz_questions = []
        key_moments = []

        try:
            if note.analysis:
                analysis = json.loads(note.analysis)
        except Exception:
            pass

        try:
            if note.flashcards:
                flashcards = json.loads(note.flashcards)
        except Exception:
            pass

        try:
            if note.quiz_questions:
                quiz_questions = json.loads(note.quiz_questions)
        except Exception:
            pass

        try:
            if note.key_moments:
                key_moments = json.loads(note.key_moments)
        except Exception:
            pass

        return {
            "id": note.id,
            "title": note.title,
            "content": note.content,
            "created_at": note.created_at.isoformat(),
            "updated_at": note.updated_at.isoformat(),
            "folder_id": note.folder_id,
            "is_favorite": note.is_favorite,
            "custom_font": note.custom_font,
            "transcript": note.transcript or "",
            "analysis": analysis,
            "flashcards": flashcards,
            "quiz_questions": quiz_questions,
            "key_moments": key_moments,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching note {note_id}: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/media/history")
async def get_media_history(
    user_id: str = Query(...),
    db: Session = Depends(get_db),
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        all_notes = (
            db.query(models.Note)
            .filter(
                models.Note.user_id == user.id,
                models.Note.is_deleted == False,
            )
            .order_by(models.Note.created_at.desc())
            .limit(50)
            .all()
        )

        history = []
        for note in all_notes:
            is_media_note = note.custom_font == "__MEDIA_NOTE__"

            if not is_media_note:
                flashcard_set = (
                    db.query(models.FlashcardSet)
                    .filter(
                        models.FlashcardSet.source_type == "media",
                        models.FlashcardSet.source_id == note.id,
                    )
                    .first()
                )
                if flashcard_set:
                    is_media_note = True

            if is_media_note:
                history.append({
                    "id": note.id,
                    "title": note.title,
                    "created_at": note.created_at.isoformat(),
                    "preview": note.content[:200] if note.content else "",
                })

        logger.info(f"Found {len(history)} media-generated notes for user {user_id}")
        return {"history": history}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"History fetch error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/upload_slides")
async def upload_slides(
    user_id: str = Form(...),
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db),
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        uploaded_slides = []

        for file in files:
            if not file.filename.lower().endswith((".pdf", ".ppt", ".pptx")):
                continue

            file_content = await file.read()
            file_size = len(file_content)

            timestamp = int(datetime.now(timezone.utc).timestamp())
            clean_name = "".join(c for c in (file.filename or "upload") if c.isalnum() or c in (".", "_", "-"))
            safe_filename = f"{user.id}_{timestamp}_{clean_name}"
            file_path = UPLOAD_DIR / safe_filename

            with open(file_path, "wb") as f:
                f.write(file_content)

            page_count = 0
            extracted_text = ""

            if file.filename.lower().endswith(".pdf"):
                try:
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    page_count = len(pdf_reader.pages)

                    for page in pdf_reader.pages[:10]:
                        extracted_text += page.extract_text() + "\n"

                    extracted_text = extracted_text[:10000]
                except Exception as e:
                    logger.error(f"Error extracting PDF text: {str(e)}")

            elif file.filename.lower().endswith((".ppt", ".pptx")):
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(file_content))
                    page_count = len(prs.slides)

                    for slide_idx, ppt_slide in enumerate(prs.slides):
                        if slide_idx >= 10:
                            break
                        for shape in ppt_slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                extracted_text += shape.text.strip() + "\n"

                    extracted_text = extracted_text[:10000]
                except Exception as e:
                    logger.error(f"Error extracting PowerPoint content: {str(e)}")

            slide = models.UploadedSlide(
                user_id=user.id,
                filename=safe_filename,
                original_filename=file.filename,
                file_path=str(file_path),
                file_size=file_size,
                file_type=file.content_type or "application/pdf",
                page_count=page_count,
                extracted_text=extracted_text,
                processing_status="completed",
                uploaded_at=datetime.now(timezone.utc),
                processed_at=datetime.now(timezone.utc),
            )

            db.add(slide)
            uploaded_slides.append(slide)

        db.commit()

        return {
            "status": "success",
            "uploaded_count": len(uploaded_slides),
            "message": f"Successfully uploaded {len(uploaded_slides)} slide(s)",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading slides: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to upload slides: {str(e)}")

@router.get("/get_uploaded_slides")
def get_uploaded_slides(
    user_id: str = Query(...),
    db: Session = Depends(get_db),
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        slides = (
            db.query(models.UploadedSlide)
            .filter(models.UploadedSlide.user_id == user.id)
            .order_by(models.UploadedSlide.uploaded_at.desc())
            .all()
        )

        return {
            "slides": [
                {
                    "id": slide.id,
                    "filename": slide.original_filename,
                    "file_size": slide.file_size,
                    "file_type": slide.file_type,
                    "page_count": slide.page_count,
                    "uploaded_at": slide.uploaded_at.isoformat() + "Z",
                    "preview_url": slide.preview_url,
                    "processing_status": slide.processing_status,
                }
                for slide in slides
            ]
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting slides: {str(e)}")
        return {"slides": []}

@router.delete("/delete_slide/{slide_id}")
def delete_slide(
    slide_id: int,
    db: Session = Depends(get_db),
):
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        if Path(slide.file_path).exists():
            Path(slide.file_path).unlink()

        db.delete(slide)
        db.commit()

        return {
            "status": "success",
            "message": "Slide deleted successfully",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting slide: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/analyze_slide/{slide_id}")
async def analyze_slide(
    slide_id: int,
    force_reanalyze: bool = Query(False),
    db: Session = Depends(get_db),
):
    try:
        from comprehensive_slide_analyzer import get_or_create_analysis

        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        file_path = Path(slide.file_path)
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")

        file_ext = slide.original_filename.lower().split(".")[-1]
        if file_ext not in ["pdf", "ppt", "pptx"]:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        analysis_result = get_or_create_analysis(
            slide_id=slide_id,
            file_path=file_path,
            file_type=file_ext,
            db=db,
            force_reanalyze=force_reanalyze,
        )

        return {
            "status": "success",
            "filename": slide.original_filename,
            "total_slides": analysis_result["total_slides"],
            "presentation_summary": analysis_result["presentation_summary"],
            "slides": analysis_result["slides"],
            "analyzed_at": analysis_result["analyzed_at"],
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing slide: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error analyzing slide: {str(e)}")

@router.get("/slide_file/{slide_id}")
async def get_slide_file(
    slide_id: int,
    db: Session = Depends(get_db),
):
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        file_path = Path(slide.file_path)
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")

        filename = slide.original_filename.lower()
        if filename.endswith(".pdf"):
            media_type = "application/pdf"
        elif filename.endswith(".pptx"):
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif filename.endswith(".ppt"):
            media_type = "application/vnd.ms-powerpoint"
        else:
            media_type = "application/octet-stream"

        return FileResponse(
            path=str(file_path),
            media_type=media_type,
            filename=slide.original_filename,
            headers={
                "Content-Disposition": f'inline; filename="{slide.original_filename}"'
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error serving slide file: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error serving slide file: {str(e)}")

@router.get("/slide_image/{slide_id}/{page_number}")
async def get_slide_image(
    slide_id: int,
    page_number: int,
    db: Session = Depends(get_db),
):
    import base64

    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        file_path = Path(slide.file_path)
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")

        if slide.original_filename.lower().endswith(".pdf"):
            try:
                import fitz
                doc = fitz.open(str(file_path))

                if page_number < 1 or page_number > len(doc):
                    doc.close()
                    raise HTTPException(
                        status_code=400,
                        detail=f"Invalid page number. File has {len(doc)} pages.",
                    )

                page = doc[page_number - 1]
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                doc.close()

                return Response(content=img_bytes, media_type="image/png")

            except ImportError:
                raise HTTPException(status_code=500, detail="PyMuPDF not installed for PDF rendering")
            except HTTPException:
                raise
            except Exception as e:
                logger.error(f"Error rendering PDF page: {e}")
                raise HTTPException(status_code=500, detail=f"Error rendering PDF: {str(e)}")

        elif slide.original_filename.lower().endswith((".ppt", ".pptx")):
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        result = subprocess.run(
                            [
                                "soffice",
                                "--headless",
                                "--convert-to",
                                "pdf",
                                "--outdir",
                                temp_dir,
                                str(file_path),
                            ],
                            capture_output=True,
                            timeout=60,
                        )

                        converted_files = [f for f in os.listdir(temp_dir) if f.endswith(".pdf")]
                        if converted_files:
                            pdf_path = os.path.join(temp_dir, converted_files[0])

                            import fitz
                            doc = fitz.open(pdf_path)

                            if page_number < 1 or page_number > len(doc):
                                doc.close()
                                raise HTTPException(status_code=400, detail="Invalid page number")

                            page = doc[page_number - 1]
                            mat = fitz.Matrix(2.0, 2.0)
                            pix = page.get_pixmap(matrix=mat)
                            img_bytes = pix.tobytes("png")
                            doc.close()

                            return Response(content=img_bytes, media_type="image/png")
                    except (subprocess.TimeoutExpired, FileNotFoundError):
                        pass

                try:
                    from pptx import Presentation as PptxPresentation
                    from pptx.util import Inches, Pt
                    from PIL import Image, ImageDraw, ImageFont

                    prs = PptxPresentation(str(file_path))

                    if page_number < 1 or page_number > len(prs.slides):
                        raise HTTPException(
                            status_code=400,
                            detail=f"Invalid page number. File has {len(prs.slides)} slides.",
                        )

                    ppt_slide = prs.slides[page_number - 1]

                    width, height = 1280, 720
                    img = Image.new("RGB", (width, height), color="#ffffff")
                    draw = ImageDraw.Draw(img)

                    try:
                        title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 36)
                        body_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
                        small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
                    except Exception:
                        try:
                            title_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36
                            )
                            body_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20
                            )
                            small_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16
                            )
                        except Exception:
                            title_font = ImageFont.load_default()
                            body_font = ImageFont.load_default()
                            small_font = ImageFont.load_default()

                    y_position = 40
                    texts_extracted = []

                    for shape in ppt_slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts_extracted.append(shape.text.strip())

                    if texts_extracted:
                        title = texts_extracted[0][:100]
                        words = title.split()
                        lines = []
                        current_line = []
                        for word in words:
                            current_line.append(word)
                            test_line = " ".join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=title_font)
                            if bbox[2] > width - 80:
                                current_line.pop()
                                if current_line:
                                    lines.append(" ".join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(" ".join(current_line))

                        for line in lines[:2]:
                            draw.text((40, y_position), line, fill="#1a1a2e", font=title_font)
                            y_position += 45

                        y_position += 20

                    for text in texts_extracted[1:]:
                        if y_position > height - 60:
                            break

                        words = text.split()
                        lines = []
                        current_line = []
                        for word in words:
                            current_line.append(word)
                            test_line = " ".join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=body_font)
                            if bbox[2] > width - 100:
                                current_line.pop()
                                if current_line:
                                    lines.append(" ".join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(" ".join(current_line))

                        for line in lines[:4]:
                            if y_position > height - 60:
                                break
                            draw.text((50, y_position), f"• {line}", fill="#333333", font=body_font)
                            y_position += 28

                        y_position += 15

                    draw.text((width - 60, height - 30), f"{page_number}", fill="#888888", font=small_font)

                    if not texts_extracted:
                        draw.text(
                            (width // 2, height // 2 - 20),
                            f"Slide {page_number}",
                            fill="#1a1a2e",
                            font=title_font,
                            anchor="mm",
                        )
                        draw.text(
                            (width // 2, height // 2 + 30),
                            "(Content preview not available)",
                            fill="#888888",
                            font=body_font,
                            anchor="mm",
                        )

                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format="PNG", quality=95)
                    img_bytes = img_buffer.getvalue()

                    return Response(content=img_bytes, media_type="image/png")

                except ImportError as e:
                    logger.error(f"python-pptx not available: {e}")
                    from PIL import Image, ImageDraw

                    img = Image.new("RGB", (800, 600), color="#1a1a2e")
                    draw = ImageDraw.Draw(img)
                    draw.text((400, 280), f"Slide {page_number}", fill="#d7b38c", anchor="mm")
                    draw.text((400, 320), "Preview not available", fill="#888888", anchor="mm")

                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format="PNG")
                    img_bytes = img_buffer.getvalue()

                    return Response(content=img_bytes, media_type="image/png")

            except HTTPException:
                raise
            except Exception as e:
                logger.error(f"Error rendering PowerPoint: {e}")
                raise HTTPException(status_code=500, detail=f"Error rendering PowerPoint: {str(e)}")

        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting slide image: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error getting slide image: {str(e)}")
