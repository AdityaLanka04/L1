import os
import sys
import re
import json
import logging
import warnings
from datetime import datetime, timezone, timedelta
from typing import Dict, List, Optional, Any
import asyncio
import math  
import tempfile
from fastapi import FastAPI, Form, Depends, HTTPException, status, Query, Request, File, UploadFile, Body
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordRequestForm, HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel,EmailStr 
from jose import JWTError, jwt
from sqlalchemy.orm import Session
from sqlalchemy import func, and_
from sqlalchemy import text 
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests
import requests
from dotenv import load_dotenv
from groq import Groq

# Suppress warnings
warnings.filterwarnings('ignore', category=FutureWarning)
warnings.filterwarnings('ignore', category=DeprecationWarning)
warnings.filterwarnings('ignore', category=RuntimeWarning)

# Configure minimal logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s: %(message)s'
)

# Set specific loggers to WARNING or ERROR
logging.getLogger('uvicorn.access').setLevel(logging.WARNING)
logging.getLogger('sqlalchemy').setLevel(logging.WARNING)
logging.getLogger('knowledge_graph').setLevel(logging.ERROR)

# Set up logger
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# Media Processing Imports (optional - will be imported when needed)
try:
    from youtube_transcript_api import YouTubeTranscriptApi
    YOUTUBE_AVAILABLE = True
except ImportError:
    YOUTUBE_AVAILABLE = False
    logger.warning("youtube-transcript-api not installed")

try:
    from pytube import YouTube
    PYTUBE_AVAILABLE = True
except ImportError:
    PYTUBE_AVAILABLE = False
    logger.warning("pytube not installed")

try:
    from langdetect import detect, LangDetectException
    LANGDETECT_AVAILABLE = True
except ImportError:
    LANGDETECT_AVAILABLE = False
    logger.warning("langdetect not installed")

try:
    import pycountry
    PYCOUNTRY_AVAILABLE = True
except ImportError:
    PYCOUNTRY_AVAILABLE = False
    logger.warning("pycountry not installed")

try:
    from pydub import AudioSegment
    PYDUB_AVAILABLE = True
except ImportError:
    PYDUB_AVAILABLE = False
    logger.warning("pydub not installed")

import models
from database import SessionLocal, engine, get_db
from models import get_db

DATABASE_URL = os.getenv("DATABASE_URL", "")


from pathlib import Path
import PyPDF2
import io
from argon2 import PasswordHasher
from argon2.exceptions import VerifyMismatchError

sys.path.append(os.path.dirname(os.path.abspath(__file__)))
load_dotenv()

from ai_personality import PersonalityEngine, AdaptiveLearningModel
from neural_adaptation import get_rl_agent, ConversationContextAnalyzer
import advanced_prompting
from flashcard_api_minimal import register_flashcard_api_minimal
from question_bank_enhanced import register_question_bank_api
from proactive_ai_system import get_proactive_ai_engine
from adaptive_learning_api import register_adaptive_learning_api
from math_processor import process_math_in_response, enhance_display_math

# LangGraph Agent System
from agents.setup import register_agent_routes

from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
import os
from pathlib import Path
from fastapi.middleware.cors import CORSMiddleware

from fastapi import WebSocket, WebSocketDisconnect
from websocket_manager import (
    manager, 
    notify_battle_challenge, 
    notify_battle_accepted, 
    notify_battle_declined,
    notify_battle_started,
    notify_battle_completed
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

models.Base.metadata.create_all(bind=engine)


def sync_sequences():
    if "postgres" not in DATABASE_URL:
        return

    sequences = [
        ("chat_messages_id_seq", "chat_messages"),
        ("chat_sessions_id_seq", "chat_sessions"),
        ("activities_id_seq", "activities"),
    ]

    with engine.connect() as connection:
        for sequence_name, table_name in sequences:
            sequence_exists = connection.execute(
                text("SELECT to_regclass(:sequence_name)"),
                {"sequence_name": sequence_name}
            ).scalar()

            if not sequence_exists:
                continue

            connection.execute(
                text(
                    f"SELECT setval(:sequence_name, COALESCE((SELECT MAX(id) FROM {table_name}), 0) + 1, false)"
                ),
                {"sequence_name": sequence_name}
            )

        connection.commit()


sync_sequences()





app = FastAPI(title="Brainwave Backend API", version="3.0.0")  #  Keep this, remove duplicate

# ==================== CORS CONFIGURATION ====================



app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:8000",
        "https://cerbyl.com",
        "https://www.cerbyl.com",
        "https://ceryl.onrender.com",
        "https://l1.vercel.app", 
        "https://l1-theta.vercel.app",
        "https://l1-7i4bnhcn1-asphar0057s-projects.vercel.app"
    ],
    allow_origin_regex=r"https://(l1-.*\.vercel\.app|.*cerbyl\.com)$",  # Allow all Vercel preview deploys and cerbyl.com
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# CORS Configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "https://*.onrender.com",
        "https://*.vercel.app",
        "https://brainwave-ai.vercel.app"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

ph = PasswordHasher()
security = HTTPBearer()

GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID")
SECRET_KEY = os.getenv("SECRET_KEY", "your-super-secret-key-change-this-in-production")
ALGORITHM = "HS256"

# AI API Configuration - Gemini first (free tier), Groq as fallback
GEMINI_API_KEY = os.getenv("GOOGLE_GENERATIVE_AI_KEY") or os.getenv("GEMINI_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_MODEL = "llama-3.3-70b-versatile"
GEMINI_MODEL = "gemini-2.0-flash"  # Gemini 2.0 Flash (stable, 10 req/min, 4M tokens/min, 1500 req/day)

# Initialize AI clients
groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

try:
    import google.generativeai as genai
    if GEMINI_API_KEY:
        genai.configure(api_key=GEMINI_API_KEY)
        # Don't create model here, create it in ai_utils with correct name
        gemini_client = genai  # Pass the genai module itself
    else:
        gemini_client = None
        logger.warning("  No GEMINI_API_KEY found in .env")
except ImportError as e:
    gemini_client = None
    logger.error(f"  google-generativeai not installed: {e}")
    logger.warning("Install with: pip install google-generativeai")
except Exception as e:
    gemini_client = None
    logger.error(f"  Error initializing Gemini: {e}")

if groq_client:
    pass  # Groq configured
else:
    logger.error(" NO GROQ_API_KEY found")

# Determine which AI to use
if gemini_client:
    logger.info(" Using GEMINI as primary AI (free tier)")
    primary_ai = "gemini"
elif groq_client:
    logger.info(" Using GROQ as primary AI")
    primary_ai = "groq"
else:
    logger.error(" NO AI API KEYS CONFIGURED!")
    primary_ai = None

# Initialize unified AI client
from ai_utils import UnifiedAIClient
# Use Gemini as primary, Groq as fallback
unified_ai = UnifiedAIClient(gemini_client, groq_client, GEMINI_MODEL, GROQ_MODEL, GEMINI_API_KEY)

# Unified AI call function - uses the unified_ai client
def call_ai(prompt: str, max_tokens: int = 2000, temperature: float = 0.7) -> str:
    """
    Call AI with Gemini as primary, Groq as fallback
    Used throughout the entire app for all AI calls
    Automatically processes math notation in responses
    """
    response = unified_ai.generate(prompt, max_tokens, temperature)
    # Post-process math notation to ensure proper LaTeX wrapping
    response = process_math_in_response(response)
    response = enhance_display_math(response)
    return response

register_flashcard_api_minimal(app)

register_question_bank_api(app, unified_ai, get_db)

register_adaptive_learning_api(app, unified_ai)

# Register AI Chat Agent
from ai_chat_integration import register_ai_chat_agent
register_ai_chat_agent(app)
logger.info(" AI Chat Agent registered successfully")

# Register LangGraph Agent System Routes
register_agent_routes(app)
logger.info(" LangGraph Agent routes registered")

class Token(BaseModel):
    access_token: str
    token_type: str

class UserCreate(BaseModel):
    first_name: str
    last_name: str
    email: str
    username: str
    password: str
    age: Optional[int] = None
    field_of_study: Optional[str] = None
    learning_style: Optional[str] = None
    school_university: Optional[str] = None

class GoogleAuth(BaseModel):
    token: str

class ChatSessionCreate(BaseModel):
    user_id: str
    title: str = "New Chat"

class ChatMessageSave(BaseModel):
    chat_id: int
    user_message: str
    ai_response: str

class NoteCreate(BaseModel):
    user_id: str
    title: str = "New Note"
    content: str = ""

class NoteUpdate(BaseModel):
    note_id: int
    title: str
    content: str


class UserProfileUpdate(BaseModel):
    user_id: str
    firstName: Optional[str] = None
    lastName: Optional[str] = None
    email: Optional[str] = None
    age: Optional[int] = None
    fieldOfStudy: Optional[str] = None
    learningStyle: Optional[str] = None
    schoolUniversity: Optional[str] = None
    preferredSubjects: Optional[List[str]] = []
    difficultyLevel: Optional[str] = "intermediate"
    studySchedule: Optional[str] = "flexible"
    learningPace: Optional[str] = "moderate"
    motivationFactors: Optional[List[str]] = []
    weakAreas: Optional[List[str]] = []
    strongAreas: Optional[List[str]] = []
    careerGoals: Optional[str] = None
    studyGoals: Optional[str] = None
    timeZone: Optional[str] = None
    studyEnvironment: Optional[str] = "quiet"
    preferredLanguage: Optional[str] = "english"
    preferredSessionLength: Optional[int] = None  
    bestStudyTimes: Optional[List[str]] = [] 

class FolderCreate(BaseModel):
    user_id: str
    name: str
    color: Optional[str] = "#D7B38C"
    parent_id: Optional[int] = None

class FolderUpdate(BaseModel):
    folder_id: int
    name: str
    color: Optional[str] = None

class NoteUpdateFolder(BaseModel):
    note_id: int
    folder_id: Optional[int] = None

class NoteFavorite(BaseModel):
    note_id: int
    is_favorite: bool

class ChatFolderCreate(BaseModel):
    user_id: str
    name: str
    color: Optional[str] = "#D7B38C"
    parent_id: Optional[int] = None

class ChatUpdateFolder(BaseModel):
    chat_id: int
    folder_id: Optional[int] = None

class GenerateChatTitleRequest(BaseModel):
    chat_id: int
    user_id: str

class AIWritingAssistRequest(BaseModel):
    user_id: str
    content: str
    action: str  # "continue", "improve", "simplify", "expand", "tone_change"
    tone: Optional[str] = "professional"

class ShareContentRequest(BaseModel):
    content_type: str  # 'chat' or 'note'
    content_id: int
    friend_ids: List[int]
    message: Optional[str] = None
    permission: str = 'view'  # 'view' or 'edit'

class RemoveSharedAccessRequest(BaseModel):
    share_id: int



def verify_password(plain_password, hashed_password):
    """Verify password with Argon2"""
    try:
        ph.verify(hashed_password, plain_password)
        return True
    except VerifyMismatchError:
        return False

def get_password_hash(password):
    """Hash password with Argon2 (no length limit)"""
    return ph.hash(password)

def create_access_token(data: dict, expires_delta: timedelta = None):
    to_encode = data.copy()
    # Token expires in 30 days by default
    expire = datetime.utcnow() + (expires_delta or timedelta(days=30))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_user_by_username(db, username: str):
    return db.query(models.User).filter(models.User.username == username).first()

def get_user_by_email(db, email: str):
    return db.query(models.User).filter(models.User.email == email).first()

def get_comprehensive_profile_safe(db: Session, user_id: int):
    """Safely get comprehensive profile, returns None if schema mismatch"""
    try:
        return db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user_id
        ).first()
    except Exception as e:
        logger.warning(f"Could not load comprehensive profile: {e}")
        return None

def authenticate_user(db, username: str, password: str):
    user = get_user_by_username(db, username)
    if not user:
        user = get_user_by_email(db, username)
    if not user or not verify_password(password, user.hashed_password):
        return False
    return user

def verify_google_token(token: str):
    try:
        idinfo = id_token.verify_oauth2_token(token, google_requests.Request(), GOOGLE_CLIENT_ID)
        if idinfo['iss'] not in ['accounts.google.com', 'https://accounts.google.com']:
            raise ValueError('Wrong issuer.')
        return idinfo
    except ValueError as e:
        raise HTTPException(status_code=400, detail=f"Invalid token: {str(e)}")

def verify_token(credentials: HTTPAuthorizationCredentials = Depends(security)):
    try:
        payload = jwt.decode(credentials.credentials, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise HTTPException(status_code=401, detail="Invalid authentication credentials")
        return username
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid authentication credentials")
def calculate_day_streak(db: Session, user_id: int) -> int:
    today = datetime.now(timezone.utc).date()
    
    recent_days = db.query(models.DailyLearningMetrics.date).filter(
        models.DailyLearningMetrics.user_id == user_id,
        models.DailyLearningMetrics.questions_answered > 0
    ).order_by(models.DailyLearningMetrics.date.desc()).all()
    
    if not recent_days:
        return 0
    
    recent_dates = [day[0] for day in recent_days]
    
    if today not in recent_dates and (today - timedelta(days=1)) not in recent_dates:
        return 0
    
    streak = 0
    check_date = today if today in recent_dates else today - timedelta(days=1)
    
    while check_date in recent_dates:
        streak += 1
        check_date -= timedelta(days=1)
    
    return streak

def get_current_user(db: Session = Depends(get_db), token: str = Depends(verify_token)):
    user = get_user_by_username(db, token)
    if not user:
        user = get_user_by_email(db, token)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    return user

# Playlist API endpoints are defined below in main.py

def build_user_profile_dict(user, comprehensive_profile=None) -> Dict[str, Any]:
    profile = {
        "user_id": getattr(user, "id", "unknown"),
        "first_name": getattr(user, "first_name", "Student"),
        "last_name": getattr(user, "last_name", ""),
        "field_of_study": getattr(user, "field_of_study", "General Studies"),
        "learning_style": getattr(user, "learning_style", "Mixed"),
        "school_university": getattr(user, "school_university", "Student"),
        "age": getattr(user, "age", None)
    }
    
    if comprehensive_profile:
        profile.update({
            "difficulty_level": getattr(comprehensive_profile, "difficulty_level", "intermediate"),
            "learning_pace": getattr(comprehensive_profile, "learning_pace", "moderate"),
            "study_environment": getattr(comprehensive_profile, "study_environment", "quiet"),
            "preferred_language": getattr(comprehensive_profile, "preferred_language", "english"),
            "study_goals": getattr(comprehensive_profile, "study_goals", None),
            "career_goals": getattr(comprehensive_profile, "career_goals", None),
            "primary_archetype": getattr(comprehensive_profile, "primary_archetype", ""),
            "secondary_archetype": getattr(comprehensive_profile, "secondary_archetype", ""),
            "archetype_description": getattr(comprehensive_profile, "archetype_description", "")
        })
    
    return profile

MATHEMATICAL_FORMATTING_INSTRUCTIONS = """
Use LaTeX: inline $x$, display $$x$$
"""

# Then define the function ONCE (around line 648)
async def generate_ai_response(prompt: str, user_profile: Dict[str, Any]) -> str:
    try:
        system_prompt = f"""You are an expert AI tutor helping {user_profile.get('first_name', 'a student')} who is studying {user_profile.get('field_of_study', 'various subjects')}.
Learning Style: {user_profile.get('learning_style', 'Mixed')}
Level: {user_profile.get('difficulty_level', 'intermediate')}
Pace: {user_profile.get('learning_pace', 'moderate')}

Provide clear, educational responses tailored to the student's profile.

{MATHEMATICAL_FORMATTING_INSTRUCTIONS}"""

        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=2048,
            top_p=0.9,
        )
        
        response = chat_completion.choices[0].message.content
        # Post-process math notation
        response = process_math_in_response(response)
        response = enhance_display_math(response)
        return response
    except Exception as e:
        logger.error(f"Groq API error: {e}")
        return f"I apologize, but I encountered an error processing your request. Please try again."

def extract_topic_keywords(text: str) -> List[str]:
    stop_words = {
        'what', 'when', 'where', 'which', 'who', 'how', 'does', 'can', 'will', 
        'should', 'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been',
        'have', 'has', 'had', 'do', 'does', 'did', 'explain', 'tell', 'help',
        'understand', 'know', 'about', 'this', 'that', 'with', 'from', 'for'
    }
    
    words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
    keywords = [w for w in words if w not in stop_words]
    
    return keywords[:10]


def get_relevant_past_conversations(db: Session, user_id: int, current_topic: str, limit: int = 5) -> List[Dict]:
    try:
        recent_sessions = db.query(models.ChatSession).filter(
            models.ChatSession.user_id == user_id
        ).order_by(models.ChatSession.updated_at.desc()).limit(20).all()
        
        relevant_conversations = []
        
        for session in recent_sessions:
            messages = db.query(models.ChatMessage).filter(
                models.ChatMessage.chat_session_id == session.id
            ).order_by(models.ChatMessage.timestamp.asc()).limit(3).all()
            
            if messages:
                session_content = " ".join([
                    f"{m.user_message} {m.ai_response}" for m in messages
                ])
                
                topic_keywords = current_topic.lower().split()
                relevance_score = sum(
                    1 for keyword in topic_keywords 
                    if len(keyword) > 3 and keyword in session_content.lower()
                )
                
                if relevance_score > 0:
                    relevant_conversations.append({
                        'session_id': session.id,
                        'title': session.title,
                        'relevance_score': relevance_score,
                        'last_updated': session.updated_at,
                        'preview': session_content[:200]
                    })
        
        relevant_conversations.sort(
            key=lambda x: (x['relevance_score'], x['last_updated']), 
            reverse=True
        )
        
        return relevant_conversations[:limit]
        
    except Exception as e:
        logger.error(f"Error getting relevant conversations: {str(e)}")
        return []


def get_user_learning_context(db: Session, user_id: int) -> Dict[str, Any]:
    try:
        recent_activities = db.query(models.Activity).filter(
            models.Activity.user_id == user_id
        ).order_by(models.Activity.timestamp.desc()).limit(20).all()
        
        topics_studied = {}
        for activity in recent_activities:
            topic = activity.topic or "General"
            topics_studied[topic] = topics_studied.get(topic, 0) + 1
        
        daily_metrics = db.query(models.DailyLearningMetrics).filter(
            models.DailyLearningMetrics.user_id == user_id
        ).order_by(models.DailyLearningMetrics.date.desc()).limit(7).all()
        
        avg_questions_per_day = (
            sum(m.questions_answered for m in daily_metrics) / len(daily_metrics)
            if daily_metrics else 0
        )
        
        avg_accuracy = (
            sum(m.correct_answers / max(m.questions_answered, 1) for m in daily_metrics) 
            / len(daily_metrics) * 100
            if daily_metrics else 0
        )
        
        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user_id
        ).first()
        
        weak_areas = []
        strong_areas = []
        
        if comprehensive_profile:
            try:
                weak_areas = json.loads(comprehensive_profile.weak_areas or "[]")
                strong_areas = json.loads(comprehensive_profile.strong_areas or "[]")
            except:
                pass
        
        return {
            'topics_studied': topics_studied,
            'most_frequent_topics': sorted(topics_studied.items(), key=lambda x: x[1], reverse=True)[:3],
            'avg_questions_per_day': round(avg_questions_per_day, 1),
            'avg_accuracy': round(avg_accuracy, 1),
            'weak_areas': weak_areas,
            'strong_areas': strong_areas,
            'total_learning_days': len(daily_metrics)
        }
        
    except Exception as e:
        logger.error(f"Error building learning context: {str(e)}")
        return {}


def build_user_profile_dict(user, comprehensive_profile=None) -> Dict[str, Any]:
    profile = {
        "user_id": getattr(user, "id", "unknown"),
        "first_name": getattr(user, "first_name", "Student"),
        "last_name": getattr(user, "last_name", ""),
        "field_of_study": getattr(user, "field_of_study", "General Studies"),
        "learning_style": getattr(user, "learning_style", "Mixed"),
        "school_university": getattr(user, "school_university", "Student"),
        "age": getattr(user, "age", None)
    }
    
    if comprehensive_profile:
        profile.update({
            "difficulty_level": getattr(comprehensive_profile, "difficulty_level", "intermediate"),
            "learning_pace": getattr(comprehensive_profile, "learning_pace", "moderate"),
            "study_environment": getattr(comprehensive_profile, "study_environment", "quiet"),
            "preferred_language": getattr(comprehensive_profile, "preferred_language", "english"),
            "study_goals": getattr(comprehensive_profile, "study_goals", None),
            "career_goals": getattr(comprehensive_profile, "career_goals", None),
            "primary_archetype": getattr(comprehensive_profile, "primary_archetype", ""),
            "secondary_archetype": getattr(comprehensive_profile, "secondary_archetype", ""),
            "archetype_description": getattr(comprehensive_profile, "archetype_description", "")
        })
    
    return profile

#@app.get("/")
#async def root():
 #   return {
 #       "message": "Brainwave Backend API v3.0.0",
  ##     "ai_provider": "Groq",
   #     "api_version": "3.0.0"
   # }

from datetime import datetime, timezone
import os

@app.get("/api/health")
def health_check():
    return {
        "status": "healthy",
        "message": "Brainwave API is running",
        "ai_provider": "Groq",
        "frontend": "https://cerbyl.com",
        "timestamp": datetime.now(timezone.utc).isoformat() + 'Z'
    }

@app.on_event("startup")
async def fix_database_sequences():
    """Automatically fix PostgreSQL sequences and initialize agent system on startup"""
    # Only run for PostgreSQL, not SQLite
    if "postgres" not in DATABASE_URL:
        pass  # SQLite doesn't need sequence fixes
    else:
        logger.info(" Checking and fixing PostgreSQL database sequences...")
        
        # Run migrations on startup (production)
        try:
            logger.info("🔄 Running database migrations...")
            from migration import run_migration
            run_migration()
        except Exception as e:
            logger.error(f" Migration error: {e}")
        
        try:
            db = SessionLocal()
            
            tables_to_fix = [
                'users',
                'chat_sessions',
                'notes',
                'activities',
                'flashcard_sets',
                'daily_learning_metrics',
                'user_stats',
                'folders',
                'question_sets',
                'learning_reviews',
                'uploaded_slides'
            ]
            
            for table in tables_to_fix:
                try:
                    # Get max ID
                    max_id_query = text(f"SELECT COALESCE(MAX(id), 0) FROM {table}")
                    max_id = db.execute(max_id_query).scalar()
                    
                    # Fix the sequence (PostgreSQL only)
                    fix_query = text(f"SELECT setval('{table}_id_seq', :next_id)")
                    db.execute(fix_query, {"next_id": max_id + 1})
                    
                    logger.info(f" Fixed sequence for {table}: next_id = {max_id + 1}")
                    
                except Exception as e:
                    logger.warning(f" Could not fix {table}: {str(e)}")
            
            db.commit()
            db.close()
            
        except Exception as e:
            logger.error(f" Sequence fix failed: {str(e)}")
    
    # Initialize LangGraph Agent System
    try:
        from agents.setup import setup_agent_system
        await setup_agent_system(
            app, 
            unified_ai, 
            enable_knowledge_graph=True,
            db_session_factory=SessionLocal
        )
    except Exception as e:
        logger.warning(f" Agent system initialization failed (non-critical): {e}")
        logger.info("   App will continue without agent system")

@app.get("/api/fix-reminder-timezones")
async def fix_reminder_timezones(user_id: str = Query(...), db: Session = Depends(get_db)):
    """
    Fix existing reminders that were stored as UTC time.
    This converts early morning times (likely UTC) to reasonable local times.
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get user's reminders with dates
        reminders = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.reminder_date != None
        ).all()
        
        fixed_reminders = []
        
        for reminder in reminders:
            if not reminder.reminder_date:
                continue
                
            original_time = reminder.reminder_date
            hour = original_time.hour
            
            # If it's between midnight and 6 AM, it might be a UTC time
            # that was meant to be a reasonable local time
            if 0 <= hour <= 6:
                # For very early times like 3:05 AM, assume it should be 9:05 AM
                if hour <= 3:
                    adjusted_time = original_time + timedelta(hours=6)
                # For 4-6 AM, assume it should be 9-11 AM  
                else:
                    adjusted_time = original_time + timedelta(hours=5)
                
                fixed_reminders.append({
                    "id": reminder.id,
                    "title": reminder.title,
                    "original": original_time.strftime('%Y-%m-%d %H:%M'),
                    "fixed": adjusted_time.strftime('%Y-%m-%d %H:%M')
                })
                
                reminder.reminder_date = adjusted_time
        
        if fixed_reminders:
            db.commit()
            return {
                "status": "success",
                "message": f"Fixed {len(fixed_reminders)} reminders",
                "fixed_reminders": fixed_reminders
            }
        else:
            return {
                "status": "success", 
                "message": "No reminders needed fixing",
                "fixed_reminders": []
            }
            
    except Exception as e:
        db.rollback()
        return {
            "status": "error",
            "message": str(e)
        }

@app.get("/api/fix-sequences")
async def fix_sequences_now(db: Session = Depends(get_db)):
    """
    Emergency sequence fix endpoint
    Call this once: https://ceryl.onrender.com/api/fix-sequences
    """
    try:
        tables = [
            'users', 'chat_sessions', 'notes', 'activities',
            'flashcard_sets', 'daily_learning_metrics', 'user_stats',
            'folders', 'question_sets', 'learning_reviews'
        ]
        
        fixed = []
        errors = []
        
        for table in tables:
            try:
                # Get max ID
                max_id_query = text(f"SELECT COALESCE(MAX(id), 0) FROM {table}")
                max_id = db.execute(max_id_query).scalar()
                
                # Fix sequence
                fix_query = text(f"SELECT setval('{table}_id_seq', :next_id)")
                db.execute(fix_query, {"next_id": max_id + 1})
                
                fixed.append(f"{table}: next_id={max_id + 1}")
                
            except Exception as e:
                errors.append(f"{table}: {str(e)}")
        
        db.commit()
        
        return {
            "status": "success",
            "fixed": fixed,
            "errors": errors if errors else None,
            "message": f"Fixed {len(fixed)} sequences"
        }
        
    except Exception as e:
        db.rollback()
        return {
            "status": "error",
            "message": str(e)
        }

@app.get("/api/get_daily_goal_progress")
def get_daily_goal_progress(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        today = datetime.now(timezone.utc).date()
        
        questions_today = db.query(models.Activity).filter(
            models.Activity.user_id == user.id,
            func.date(models.Activity.timestamp) == today
        ).count()
        
        daily_goal = 20
        
        activities = db.query(func.date(models.Activity.timestamp)).filter(
            models.Activity.user_id == user.id
        ).distinct().order_by(func.date(models.Activity.timestamp).desc()).all()
        
        streak = 0
        if activities:
            check_date = datetime.now(timezone.utc).date()
            for activity_date in [a[0] for a in activities]:
                if activity_date == check_date or activity_date == check_date - timedelta(days=1):
                    streak += 1
                    check_date = activity_date - timedelta(days=1)
                else:
                    break
        
        return {
            "questions_today": questions_today,
            "daily_goal": daily_goal,
            "percentage": min(int((questions_today / daily_goal) * 100), 100),
            "streak": streak
        }
    except Exception as e:
        logger.error(f"Error getting daily goal: {str(e)}")
        return {"questions_today": 0, "daily_goal": 20, "percentage": 0, "streak": 0}



# -----------------------------
# 1️⃣ Pydantic model (JSON body)
# -----------------------------
class RegisterPayload(BaseModel):
    first_name: str
    last_name: str
    email: EmailStr
    username: str
    password: str
    age: int | None = None
    field_of_study: str | None = None
    learning_style: str | None = None
    school_university: str | None = None


# -----------------------------
# 2️⃣ Register endpoint (JSON)
# -----------------------------
from sqlalchemy import text  # Make sure this is imported

@app.post("/api/register")
async def register(payload: RegisterPayload, db: Session = Depends(get_db)):
    logger.info(f"Registering user: {payload.username}")

    # --- validation ---
    if len(payload.password) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters long")

    if get_user_by_username(db, payload.username):
        raise HTTPException(status_code=400, detail="Username already registered")

    if get_user_by_email(db, payload.email):
        raise HTTPException(status_code=400, detail="Email already registered")

    # --- password hashing ---
    hashed_password = get_password_hash(payload.password)

    # --- create user record with auto-retry ---
    max_retries = 2
    for attempt in range(max_retries):
        try:
            db_user = models.User(
                first_name=payload.first_name,
                last_name=payload.last_name,
                email=payload.email,
                username=payload.username,
                hashed_password=hashed_password,
                age=payload.age,
                field_of_study=payload.field_of_study,
                learning_style=payload.learning_style,
                school_university=payload.school_university,
                google_user=False,
            )
            db.add(db_user)
            db.commit()
            db.refresh(db_user)

            # --- initialize stats ---
            user_stats = models.UserStats(user_id=db_user.id)
            db.add(user_stats)
            db.commit()

            logger.info(f" User {payload.username} registered successfully")
            return {"message": "User registered successfully"}
            
        except Exception as e:
            db.rollback()
            error_msg = str(e).lower()
            
            # Check if it's a sequence/duplicate key error
            if "duplicate key" in error_msg and "pkey" in error_msg and attempt < max_retries - 1:
                logger.warning(f" Sequence error detected, fixing... (attempt {attempt + 1})")
                
                try:
                    # Fix the sequence using text()
                    max_id_query = text("SELECT COALESCE(MAX(id), 0) FROM users")
                    max_id = db.execute(max_id_query).scalar()
                    
                    fix_query = text("SELECT setval('users_id_seq', :next_id)")
                    db.execute(fix_query, {"next_id": max_id + 1})
                    db.commit()
                    
                    logger.info(f" Sequence fixed to {max_id + 1}, retrying registration...")
                    continue  # Retry registration
                    
                except Exception as fix_error:
                    logger.error(f" Failed to fix sequence: {str(fix_error)}")
            
            # If not a sequence error or retry failed, raise the error
            logger.error(f" Registration failed: {str(e)}")
            raise HTTPException(status_code=500, detail="Registration failed. Please try again.")
    
    # If we get here, all retries failed
    raise HTTPException(status_code=500, detail="Registration failed after retries")

@app.post("/api/token", response_model=Token)
def login(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = authenticate_user(db, form_data.username, form_data.password)
    if not user:
        raise HTTPException(status_code=400, detail="Incorrect username or password")
    
    # Update last_login for session tracking (use naive datetime for SQLite)
    user.last_login = datetime.utcnow()
    db.commit()
    
    access_token = create_access_token(data={"sub": user.username})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/api/token_form")
async def login_form(username: str = Form(...), password: str = Form(...), db: Session = Depends(get_db)):
    user = authenticate_user(db, username, password)
    if not user:
        raise HTTPException(status_code=401, detail="Incorrect username or password")
    
    # Update last_login for session tracking (use naive datetime for SQLite)
    user.last_login = datetime.utcnow()
    db.commit()
    
    access_token = create_access_token(data={"sub": user.username, "user_id": user.id})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/api/google-auth")
def google_auth(auth_data: GoogleAuth, db: Session = Depends(get_db)):
    try:
        try:
            url = f"https://oauth2.googleapis.com/tokeninfo?id_token={auth_data.token}"
            response = requests.get(url)
            user_info = response.json() if response.status_code == 200 else verify_google_token(auth_data.token)
        except:
            user_info = verify_google_token(auth_data.token)
        
        email = user_info.get('email')
        if not email:
            raise HTTPException(status_code=400, detail="Email not found")
        
        user = get_user_by_email(db, email)
        
        if not user:
            user = models.User(
                first_name=user_info.get('given_name', ''),
                last_name=user_info.get('family_name', ''),
                email=email,
                username=email,
                hashed_password=get_password_hash("google_oauth"),
                picture_url=user_info.get('picture', ''),
                google_user=True
            )
            db.add(user)
            db.commit()
            db.refresh(user)
            
            user_stats = models.UserStats(user_id=user.id)
            db.add(user_stats)
            db.commit()
        else:
            # Update last_login for existing user (use naive datetime for SQLite)
            user.last_login = datetime.utcnow()
            db.commit()
        
        access_token = create_access_token(data={"sub": user.username})
        
        return {
            "access_token": access_token,
            "token_type": "bearer",
            "user_info": {
                "email": email,
                "given_name": user_info.get('given_name'),
                "family_name": user_info.get('family_name'),
                "picture": user_info.get('picture'),
                "google_user": True
            }
        }
    except Exception as e:
        logger.error(f"Google auth error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/me")
async def get_current_user_info(current_user: models.User = Depends(get_current_user)):
    return {
        "id": current_user.id,
        "first_name": current_user.first_name,
        "last_name": current_user.last_name,
        "email": current_user.email,
        "username": current_user.username,
        "age": current_user.age,
        "field_of_study": current_user.field_of_study,
        "learning_style": current_user.learning_style,
        "school_university": current_user.school_university,
        "picture_url": current_user.picture_url,
        "google_user": current_user.google_user
    }

#!/usr/bin/env python3
# Enhanced AI Chat System with Conversation Memory and Context Awareness

import json
import re
from datetime import datetime, timezone, timedelta
from typing import Dict, List, Optional, Any
import asyncio
from sqlalchemy import and_, desc, func
from groq import Groq


def get_user_learning_context(db, user_id: int) -> Dict[str, Any]:
    """
    Build comprehensive learning context about the user
    """
    try:
        # Get user's recent topics
        recent_activities = db.query(models.Activity).filter(
            models.Activity.user_id == user_id
        ).order_by(models.Activity.timestamp.desc()).limit(20).all()
        
        topics_studied = {}
        for activity in recent_activities:
            topic = activity.topic or "General"
            topics_studied[topic] = topics_studied.get(topic, 0) + 1
        
        # Get learning patterns
        daily_metrics = db.query(models.DailyLearningMetrics).filter(
            models.DailyLearningMetrics.user_id == user_id
        ).order_by(models.DailyLearningMetrics.date.desc()).limit(7).all()
        
        avg_questions_per_day = (
            sum(m.questions_answered for m in daily_metrics) / len(daily_metrics)
            if daily_metrics else 0
        )
        
        avg_accuracy = (
            sum(m.correct_answers / max(m.questions_answered, 1) for m in daily_metrics) 
            / len(daily_metrics) * 100
            if daily_metrics else 0
        )
        
        # Get user's weak and strong areas from profile
        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user_id
        ).first()
        
        weak_areas = []
        strong_areas = []
        
        if comprehensive_profile:
            try:
                weak_areas = json.loads(comprehensive_profile.weak_areas or "[]")
                strong_areas = json.loads(comprehensive_profile.strong_areas or "[]")
            except:
                pass
        
        return {
            'topics_studied': topics_studied,
            'most_frequent_topics': sorted(topics_studied.items(), key=lambda x: x[1], reverse=True)[:3],
            'avg_questions_per_day': round(avg_questions_per_day, 1),
            'avg_accuracy': round(avg_accuracy, 1),
            'weak_areas': weak_areas,
            'strong_areas': strong_areas,
            'total_learning_days': len(daily_metrics)
        }
        
    except Exception as e:
        logger.error(f"Error building learning context: {str(e)}")
        return {}


@app.post("/api/ask/")
async def ask_ai_enhanced(
    user_id: str = Form(...),
    question: str = Form(...),
    chat_id: Optional[str] = Form(None),
    db: Session = Depends(get_db)
):
    logger.info(f"🤖 Processing AI query for user {user_id}")
    
    try:
        chat_id_int = int(chat_id) if chat_id else None
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get or create daily metric
        from datetime import datetime, timezone
        today = datetime.now(timezone.utc).date()
        daily_metric = db.query(models.DailyLearningMetrics).filter(
            and_(
                models.DailyLearningMetrics.user_id == user.id,
                models.DailyLearningMetrics.date == today
            )
        ).first()
        
        if not daily_metric:
            daily_metric = models.DailyLearningMetrics(
                user_id=user.id,
                date=today,
                questions_answered=0,
                correct_answers=0,
                sessions_completed=0,
                time_spent_minutes=0.0,
                accuracy_rate=0.0,
                engagement_score=0.0
            )
            db.add(daily_metric)
            db.commit()
            db.refresh(daily_metric)
        
        if chat_id_int:
            chat_session = db.query(models.ChatSession).filter(
                models.ChatSession.id == chat_id_int,
                models.ChatSession.user_id == user.id
            ).first()
            if not chat_session:
                raise HTTPException(status_code=404, detail="Chat session not found")
        
        comprehensive_profile = get_comprehensive_profile_safe(db, user.id)
        
        user_profile = {
            "user_id": user.id,
            "first_name": user.first_name or "Student",
            "last_name": user.last_name or "",
            "field_of_study": user.field_of_study or "General Studies",
            "learning_style": user.learning_style or "Mixed",
            "difficulty_level": "intermediate",
            "learning_pace": "moderate",
            "primary_archetype": "",
            "secondary_archetype": "",
            "brainwave_goal": "",
            "preferred_subjects": []
        }
        
        if comprehensive_profile:
            user_profile.update({
                "difficulty_level": comprehensive_profile.difficulty_level or "intermediate",
                "learning_pace": comprehensive_profile.learning_pace or "moderate",
                "primary_archetype": comprehensive_profile.primary_archetype or "",
                "secondary_archetype": comprehensive_profile.secondary_archetype or "",
                "brainwave_goal": comprehensive_profile.brainwave_goal or ""
            })
            
            try:
                if comprehensive_profile.preferred_subjects:
                    subjects_data = comprehensive_profile.preferred_subjects
                    if isinstance(subjects_data, str):
                        user_profile["preferred_subjects"] = json.loads(subjects_data)
                    else:
                        user_profile["preferred_subjects"] = subjects_data
            except Exception as e:
                logger.error(f"Error parsing subjects: {e}")
                user_profile["preferred_subjects"] = []
        
        logger.info(f"📊 Profile: {user_profile['primary_archetype']} | {user_profile['field_of_study']}")
        
        conversation_history = []
        if chat_id_int:
            recent_messages = db.query(models.ChatMessage).filter(
                models.ChatMessage.chat_session_id == chat_id_int
            ).order_by(models.ChatMessage.timestamp.desc()).limit(10).all()
            
            for msg in reversed(recent_messages):
                conversation_history.append({
                    'user_message': msg.user_message,
                    'ai_response': msg.ai_response,
                    'timestamp': msg.timestamp
                })
        
        learning_context = advanced_prompting.get_user_learning_context(db, user.id)
        
        topic_keywords = " ".join(advanced_prompting.extract_topic_keywords(question))
        relevant_past_chats = advanced_prompting.get_relevant_past_conversations(
            db, user.id, topic_keywords, limit=3
        )
        
        rl_agent = get_rl_agent(db, user.id)
        
        context = {
            'message_length': len(question),
            'question_complexity': ConversationContextAnalyzer.calculate_complexity(question),
            'archetype': user_profile.get('primary_archetype', ''),
            'time_of_day': datetime.now(timezone.utc).hour,
            'session_length': len(conversation_history),
            'previous_ratings': [
                msg.get('userRating', 3) for msg in conversation_history
                if isinstance(msg, dict) and msg.get('userRating')
            ],
            'topic_keywords': advanced_prompting.extract_topic_keywords(question),
            'sentiment': ConversationContextAnalyzer.analyze_sentiment(question),
            'formality_level': 0.5
        }
        
        state = rl_agent.encode_state(context)
        response_adjustments = rl_agent.get_response_adjustment()
        
        logger.info(f"🧠 Neural adjustments: {response_adjustments}")
        
        # Use unified call_ai function (Gemini primary, Groq fallback)
        response = await advanced_prompting.generate_enhanced_ai_response(
            question,
            user_profile,
            learning_context,
            conversation_history,
            relevant_past_chats,
            db,
            call_ai,  # Pass the unified AI function
            GROQ_MODEL
        )
        
        known_mistake = rl_agent.check_for_known_mistakes(response)
        if known_mistake:
            logger.info(f" Correcting known mistake")
            response = known_mistake
        
        # All tracking (messages, activities, metrics, points) is handled by save_chat_message endpoint
        # This endpoint ONLY generates the AI response
        
        topic = advanced_prompting.get_topic_from_question(question)
        
        db.commit()
        
        action = rl_agent.network.predict(state)
        next_state = rl_agent.encode_state({**context, 'session_length': len(conversation_history) + 1})
        rl_agent.remember(state, action, 0.0, next_state)
        
        advanced_prompting.save_conversation_memory(db, user.id, question, response, context['topic_keywords'])
        
        return {
            "answer": response,
            "ai_confidence": 0.92,
            "misconception_detected": False,
            "should_request_feedback": False,
            "topics_discussed": [topic],
            "query_type": "conversational_learning",
            "model_used": GROQ_MODEL,
            "ai_provider": "Groq",
            "archetype_used": user_profile.get('primary_archetype', 'None'),
            "questions_today": daily_metric.questions_answered,
            "learning_context_used": bool(learning_context),
            "relevant_past_chats": len(relevant_past_chats),
            "neural_network_active": True,
            "response_adjustments": response_adjustments
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f" ERROR IN ASK_AI: {str(e)}")
        logger.error(f" FULL TRACEBACK:\n{error_details}")
        print(f"\n{'='*80}")
        print(f" ERROR IN /api/ask/")
        print(f" Error: {str(e)}")
        print(f" Type: {type(e).__name__}")
        print(f" Traceback:\n{error_details}")
        print(f"{'='*80}\n")
        return {
            "answer": "I apologize, but I encountered an error. Could you please rephrase your question?",
            "ai_confidence": 0.3,
            "misconception_detected": False,
            "should_request_feedback": False,
            "topics_discussed": ["error"],
            "query_type": "error",
            "model_used": "error_handler",
            "ai_provider": "Groq",
            "error_debug": str(e)  # Add error for debugging
        }
    
@app.post("/api/test_ai_simple")
async def test_ai_simple(question: str = Form(...)):
    """Simple test endpoint that bypasses all complexity"""
    try:
        response = call_ai(f"Answer this question in one sentence: {question}", max_tokens=200, temperature=0.7)
        return {"answer": response, "status": "success"}
    except Exception as e:
        import traceback
        return {"answer": f"Error: {str(e)}", "status": "error", "traceback": traceback.format_exc()}

@app.post("/api/ask_simple/")
async def ask_simple(
    user_id: str = Form(...),
    question: str = Form(...),
    chat_id: Optional[str] = Form(None),
    db: Session = Depends(get_db)
):
    """Simplified /ask endpoint that bypasses all complex logic"""
    print(f"\n🔥 ASK_SIMPLE CALLED 🔥")
    print(f"🔥 User: {user_id}, Question: {question[:50]}...")
    
    try:
        # Get user
        print(f" Looking up user: {user_id}")
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            print(f" User not found in database: {user_id}")
            return {
                "answer": "Please log in again - your session may have expired.",
                "ai_confidence": 0.0,
                "misconception_detected": False,
                "should_request_feedback": False,
                "topics_discussed": ["error"],
                "query_type": "error",
                "model_used": "error",
                "ai_provider": "Error"
            }
        
        # Verify chat belongs to user if chat_id provided
        chat_id_int = None
        if chat_id:
            try:
                chat_id_int = int(chat_id)
                chat_session = db.query(models.ChatSession).filter(
                    models.ChatSession.id == chat_id_int,
                    models.ChatSession.user_id == user.id
                ).first()
                if not chat_session:
                    print(f" Chat session {chat_id_int} not found or doesn't belong to user {user.id}")
                    # Don't create a new session - the frontend should have created it
                    # Just use the provided chat_id and let the message save handle it
                    print(f" Will use provided chat_id: {chat_id_int}")
            except ValueError as e:
                print(f" Invalid chat_id format: {str(e)}")
                chat_id_int = None
        
        # Load chat history for context
        chat_history = ""
        if chat_id_int:
            try:
                recent_messages = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == chat_id_int
                ).order_by(models.ChatMessage.timestamp.desc()).limit(10).all()
                
                if recent_messages:
                    # Reverse to get chronological order
                    recent_messages = list(reversed(recent_messages))
                    chat_history = "\n\nPrevious conversation:\n"
                    for msg in recent_messages:
                        # Each message has both user_message and ai_response
                        chat_history += f"Student: {msg.user_message}\n"
                        chat_history += f"You: {msg.ai_response}\n"
                    print(f"📜 Loaded {len(recent_messages)} previous message pairs for context")
            except Exception as e:
                print(f" Error loading chat history: {str(e)}")
        
        # Build personalized prompt
        first_name = user.first_name or "there"
        field_of_study = user.field_of_study or "your studies"
        
        prompt = f"""You are a helpful AI tutor assisting {first_name}, who is studying {field_of_study}.

Be warm, encouraging, and personalized. Address them by name when appropriate.
Provide clear, educational responses tailored to their level.

CRITICAL - Mathematical Notation Rules:
- ONLY use LaTeX for actual mathematical formulas and equations
- Use $...$ for inline math: "The derivative $f'(x) = 2x$ shows..."
- Use $$...$$ for display equations on their own line
- Regular text should NOT be in LaTeX - only the math parts
- Examples:
  * CORRECT: "The function $f(x) = e^x$ is its own derivative"
  * WRONG: "$The function f(x) = e^x is its own derivative$"
  * CORRECT: "For the series $$\\sum_{{n=1}}^{{\\infty}} \\frac{{1}}{{n^2}} = \\frac{{\\pi^2}}{{6}}$$"
  * CORRECT: "When $x^2 + y^2 = r^2$, we have a circle of radius $r$"

{chat_history}

Current Question: {question}"""
        
        # Generate response using simple AI call
        print(f"🔥 Calling AI for {first_name}...")
        response = call_ai(prompt, max_tokens=2000, temperature=0.7)
        print(f"🔥 AI response received: {len(response)} chars")
        
        # Save message to database and award points
        if chat_id_int:
            try:
                print(f" Attempting to save message for chat_id: {chat_id_int}, user_id: {user.id}")
                
                # Check if this exact message already exists (prevent double-saving)
                existing_message = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == chat_id_int,
                    models.ChatMessage.user_message == question,
                    models.ChatMessage.ai_response == response
                ).first()
                
                if not existing_message:
                    # Save the message
                    chat_message = models.ChatMessage(
                        chat_session_id=chat_id_int,
                        user_id=user.id,
                        user_message=question,
                        ai_response=response,
                        is_user=True
                    )
                    db.add(chat_message)
                    print(f" Message object created and added to session")
                    
                    # Update chat session timestamp
                    chat_session = db.query(models.ChatSession).filter(
                        models.ChatSession.id == chat_id_int
                    ).first()
                    if chat_session:
                        chat_session.updated_at = datetime.now(timezone.utc)
                        print(f" Chat session timestamp updated")
                    
                    # Award points for AI chat (only once per message)
                    try:
                        from gamification_system import award_points
                        result = award_points(db, user.id, "ai_chat")
                        print(f" Points awarded: {result}")
                    except Exception as gam_error:
                        print(f" Failed to award AI chat points: {gam_error}")
                        import traceback
                        traceback.print_exc()
                    
                    db.commit()
                    print(f" Database committed - message saved!")
                else:
                    print(f" Message already exists, skipping save and point award")
            except Exception as save_error:
                print(f" Error saving message: {str(save_error)}")
                import traceback
                traceback.print_exc()
                db.rollback()
        else:
            print(f" No chat_id_int, cannot save message")
        
        return {
            "answer": response,
            "ai_confidence": 0.9,
            "misconception_detected": False,
            "should_request_feedback": False,
            "topics_discussed": ["general"],
            "query_type": "simple",
            "model_used": "groq",
            "ai_provider": "Groq",
            "chat_id": chat_id_int  # Return the actual chat_id used (in case it was created/changed)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"\n ERROR IN ASK_SIMPLE: {str(e)}")
        print(f" Traceback:\n{error_details}\n")
        return {
            "answer": f"Error: {str(e)}",
            "ai_confidence": 0.3,
            "misconception_detected": False,
            "should_request_feedback": False,
            "topics_discussed": ["error"],
            "query_type": "error",
            "model_used": "error",
            "ai_provider": "Error"
        }

@app.post("/api/ask_with_files/")
async def ask_with_files(
    user_id: str = Form(...),
    question: str = Form(...),
    chat_id: Optional[str] = Form(None),
    files: List[UploadFile] = File(default=[]),
    db: Session = Depends(get_db)
):
    """Handle questions with file attachments"""
    print(f"\n🔥 ASK_WITH_FILES CALLED 🔥")
    print(f"🔥 User: {user_id}, Question: {question[:50]}..., Files: {len(files)}")
    
    try:
        # Get user
        print(f" Looking up user: {user_id}")
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            print(f" User not found in database: {user_id}")
            return {
                "answer": "Please log in again - your session may have expired.",
                "ai_confidence": 0.0,
                "misconception_detected": False,
                "should_request_feedback": False,
                "topics_discussed": ["error"],
                "query_type": "error",
                "model_used": "error",
                "ai_provider": "Error",
                "files_processed": 0,
                "file_summaries": [],
                "has_file_context": False
            }
        
        # Verify chat belongs to user if chat_id provided
        chat_id_int = None
        if chat_id:
            try:
                chat_id_int = int(chat_id)
                chat_session = db.query(models.ChatSession).filter(
                    models.ChatSession.id == chat_id_int,
                    models.ChatSession.user_id == user.id
                ).first()
                if not chat_session:
                    print(f" Chat session {chat_id_int} not found or doesn't belong to user {user.id}")
                    # Create a new chat session instead of failing
                    print(f" Creating new chat session for user {user.id}")
                    new_chat = models.ChatSession(
                        user_id=user.id,
                        title="New Chat with Files"
                    )
                    db.add(new_chat)
                    db.commit()
                    db.refresh(new_chat)
                    chat_id_int = new_chat.id
                    print(f" Created new chat session: {chat_id_int}")
            except Exception as e:
                print(f" Error validating chat_id: {str(e)}")
                # Create new chat on error too
                try:
                    new_chat = models.ChatSession(
                        user_id=user.id,
                        title="New Chat with Files"
                    )
                    db.add(new_chat)
                    db.commit()
                    db.refresh(new_chat)
                    chat_id_int = new_chat.id
                    print(f" Created new chat session after error: {chat_id_int}")
                except:
                    chat_id_int = None
        
        # Process files if provided
        file_summaries = []
        file_context = ""
        
        if files and len(files) > 0:
            print(f"🔥 Processing {len(files)} files...")
            for file in files:
                if file.filename:
                    file_content = ""
                    
                    try:
                        # Read file content
                        content = await file.read()
                        
                        # Process based on file type
                        if file.content_type == 'application/pdf' or file.filename.lower().endswith('.pdf'):
                            # Extract text from PDF
                            try:
                                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                                pdf_text = ""
                                for page in pdf_reader.pages:
                                    pdf_text += page.extract_text() + "\n"
                                file_content = pdf_text[:4000]  # Limit to 4000 chars
                                print(f"📄 Extracted {len(file_content)} chars from PDF: {file.filename}")
                            except Exception as pdf_error:
                                print(f" PDF extraction error: {pdf_error}")
                                file_content = f"[Could not extract text from PDF: {file.filename}]"
                        
                        elif file.content_type and file.content_type.startswith('image/'):
                            # For images, we'll describe them in context
                            file_content = f"[Image file attached: {file.filename}]"
                            print(f"🖼️ Image detected: {file.filename}")
                        
                        else:
                            # Try to read as text
                            try:
                                file_content = content.decode('utf-8')[:4000]
                                print(f" Text file read: {file.filename}")
                            except:
                                file_content = f"[Binary file: {file.filename}]"
                        
                        file_summaries.append({
                            "file_name": file.filename,
                            "file_type": file.content_type,
                            "status": "processed",
                            "content_length": len(file_content)
                        })
                        
                        # Add to context
                        file_context += f"\n\n=== Content from {file.filename} ===\n{file_content}\n"
                        
                    except Exception as e:
                        print(f" Error processing file {file.filename}: {str(e)}")
                        file_summaries.append({
                            "file_name": file.filename,
                            "file_type": file.content_type,
                            "status": "error",
                            "error": str(e)
                        })
                        file_context += f"\n[Error reading file: {file.filename}]"
        
        # Load chat history for context
        chat_history = ""
        if chat_id_int:
            try:
                recent_messages = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == chat_id_int
                ).order_by(models.ChatMessage.timestamp.desc()).limit(10).all()
                
                if recent_messages:
                    # Reverse to get chronological order
                    recent_messages = list(reversed(recent_messages))
                    chat_history = "\n\nPrevious conversation:\n"
                    for msg in recent_messages:
                        # Each message has both user_message and ai_response
                        chat_history += f"Student: {msg.user_message}\n"
                        chat_history += f"You: {msg.ai_response}\n"
                    print(f"📜 Loaded {len(recent_messages)} previous message pairs for context")
            except Exception as e:
                print(f" Error loading chat history: {str(e)}")
        
        # Build personalized prompt
        first_name = user.first_name or "there"
        field_of_study = user.field_of_study or "your studies"
        
        # Build prompt with file context if available
        if file_context:
            prompt = f"""You are a helpful AI tutor assisting {first_name}, who is studying {field_of_study}.

The user has uploaded file(s) with the following content:
{file_context}
{chat_history}

Based on the file content above, the conversation history, and the user's question, provide a clear, educational response.
Be warm, encouraging, and personalized. Address them by name when appropriate.

CRITICAL - Mathematical Notation Rules:
- ONLY use LaTeX for actual mathematical formulas and equations
- Use $...$ for inline math: "The derivative $f'(x) = 2x$ shows..."
- Use $$...$$ for display equations on their own line
- Regular text should NOT be in LaTeX - only the math parts
- Examples:
  * CORRECT: "The function $f(x) = e^x$ is its own derivative"
  * WRONG: "$The function f(x) = e^x is its own derivative$"
  * CORRECT: "For the series $$\\sum_{{n=1}}^{{\\infty}} \\frac{{1}}{{n^2}} = \\frac{{\\pi^2}}{{6}}$$"
  * CORRECT: "When $x^2 + y^2 = r^2$, we have a circle of radius $r$"

Current Question: {question}"""
        else:
            prompt = f"""You are a helpful AI tutor assisting {first_name}, who is studying {field_of_study}.
        
Be warm, encouraging, and personalized. Address them by name when appropriate.
Provide clear, educational responses tailored to their level.

CRITICAL - Mathematical Notation Rules:
- ONLY use LaTeX for actual mathematical formulas and equations
- Use $...$ for inline math: "The derivative $f'(x) = 2x$ shows..."
- Use $$...$$ for display equations on their own line
- Regular text should NOT be in LaTeX - only the math parts
- Examples:
  * CORRECT: "The function $f(x) = e^x$ is its own derivative"
  * WRONG: "$The function f(x) = e^x is its own derivative$"
  * CORRECT: "For the series $$\\sum_{{n=1}}^{{\\infty}} \\frac{{1}}{{n^2}} = \\frac{{\\pi^2}}{{6}}$$"
  * CORRECT: "When $x^2 + y^2 = r^2$, we have a circle of radius $r$"

{chat_history}

Current Question: {question}"""
        
        # Generate response using simple AI call
        print(f"🔥 Calling AI for {first_name}...")
        response = call_ai(prompt, max_tokens=2000, temperature=0.7)
        print(f"🔥 AI response received: {len(response)} chars")
        
        # Save message to database and award points
        if chat_id_int:
            try:
                print(f" Attempting to save message for chat_id: {chat_id_int}, user_id: {user.id}")
                
                # Check if this exact message already exists (prevent double-saving)
                existing_message = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == chat_id_int,
                    models.ChatMessage.user_message == question,
                    models.ChatMessage.ai_response == response
                ).first()
                
                if not existing_message:
                    # Save the message
                    chat_message = models.ChatMessage(
                        chat_session_id=chat_id_int,
                        user_id=user.id,
                        user_message=question,
                        ai_response=response,
                        is_user=True
                    )
                    db.add(chat_message)
                    print(f" Message object created and added to session")
                    
                    # Update chat session timestamp
                    chat_session = db.query(models.ChatSession).filter(
                        models.ChatSession.id == chat_id_int
                    ).first()
                    if chat_session:
                        chat_session.updated_at = datetime.now(timezone.utc)
                        print(f" Chat session timestamp updated")
                    
                    # Award points for AI chat (only once per message)
                    try:
                        from gamification_system import award_points
                        result = award_points(db, user.id, "ai_chat")
                        print(f" Points awarded: {result}")
                    except Exception as gam_error:
                        print(f" Failed to award AI chat points: {gam_error}")
                        import traceback
                        traceback.print_exc()
                    
                    db.commit()
                    print(f" Database committed - message saved!")
                else:
                    print(f" Message already exists, skipping save and point award")
            except Exception as save_error:
                print(f" Error saving message: {str(save_error)}")
                import traceback
                traceback.print_exc()
                db.rollback()
        else:
            print(f" No chat_id_int, cannot save message")
        
        return {
            "answer": response,
            "ai_confidence": 0.9,
            "misconception_detected": False,
            "should_request_feedback": False,
            "topics_discussed": ["general"],
            "query_type": "with_files",
            "model_used": "groq",
            "ai_provider": "Groq",
            "files_processed": len(file_summaries),
            "file_summaries": file_summaries,
            "has_file_context": len(file_summaries) > 0,
            "chat_id": chat_id_int  # Return the actual chat_id used
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"\n ERROR IN ASK_WITH_FILES: {str(e)}")
        print(f" Traceback:\n{error_details}\n")
        return {
            "answer": f"Error processing your request: {str(e)}",
            "ai_confidence": 0.3,
            "misconception_detected": False,
            "should_request_feedback": False,
            "topics_discussed": ["error"],
            "query_type": "error",
            "model_used": "error",
            "ai_provider": "Error",
            "files_processed": 0,
            "file_summaries": [],
            "has_file_context": False
        }


# ==================== Enhanced AI Chat Agent Endpoint ====================

@app.post("/api/ask_agent/")
async def ask_agent(
    user_id: str = Form(...),
    question: str = Form(...),
    chat_id: Optional[str] = Form(None),
    chat_mode: Optional[str] = Form(None),  # tutoring, socratic, explanation, etc.
    response_style: Optional[str] = Form(None),  # concise, detailed, step_by_step, etc.
    db: Session = Depends(get_db)
):
    """
    Enhanced AI chat endpoint using the LangGraph-based Chat Agent.
    Provides intelligent tutoring with:
    - Emotional intelligence (detects confusion, frustration, curiosity)
    - Adaptive response styles
    - Memory-aware context
    - Self-reflection and quality improvement
    - Suggested follow-up questions
    """
    import time
    start_time = time.time()
    
    try:
        # Get user
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            print(f" User not found: {user_id}")
            return {
                "answer": "Please log in again - your session may have expired.",
                "ai_confidence": 0.0,
                "chat_mode": "error",
                "emotional_state": "neutral",
                "suggested_questions": [],
                "topics_discussed": ["error"]
            }
        
        # Validate chat_id
        chat_id_int = None
        if chat_id:
            try:
                chat_id_int = int(chat_id)
            except ValueError:
                chat_id_int = None
        
        # Try to use the LangGraph Chat Agent
        try:
            from agents.agent_api import get_chat_agent
            chat_agent = get_chat_agent()
            
            # Build state for the agent
            session_id = f"chat_{user.id}_{chat_id_int or 'new'}"
            
            state = {
                "user_id": str(user.id),
                "user_input": question,
                "session_id": session_id,
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "user_preferences": {
                    "learning_style": user.learning_style or "mixed",
                    "difficulty_level": "intermediate",
                    "field_of_study": user.field_of_study or "general"
                }
            }
            
            # Add optional mode/style
            if chat_mode:
                state["chat_mode"] = chat_mode
            if response_style:
                state["response_style"] = response_style
            
            # Invoke the chat agent
            result = await chat_agent.invoke(state)
            
            response = result.response
            metadata = result.metadata
            
        except Exception as agent_error:
            print(f" Chat agent failed, falling back to simple AI: {agent_error}")
            # Fallback to simple AI call
            first_name = user.first_name or "there"
            field_of_study = user.field_of_study or "your studies"
            
            prompt = f"""You are a helpful AI tutor assisting {first_name}, who is studying {field_of_study}.
Be warm, encouraging, and personalized. Provide clear, educational responses.

Question: {question}"""
            
            response = call_ai(prompt, max_tokens=2000, temperature=0.7)
            metadata = {
                "chat_mode": "tutoring",
                "response_style": "conversational",
                "emotional_state": "neutral",
                "quality_score": 0.7,
                "concepts_discussed": [],
                "suggested_questions": [],
                "fallback": True
            }
        
        # Save message to database
        if chat_id_int:
            try:
                existing_message = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == chat_id_int,
                    models.ChatMessage.user_message == question,
                    models.ChatMessage.ai_response == response
                ).first()
                
                if not existing_message:
                    chat_message = models.ChatMessage(
                        chat_session_id=chat_id_int,
                        user_id=user.id,
                        user_message=question,
                        ai_response=response,
                        is_user=True
                    )
                    db.add(chat_message)
                    
                    # Update chat session timestamp
                    chat_session = db.query(models.ChatSession).filter(
                        models.ChatSession.id == chat_id_int
                    ).first()
                    if chat_session:
                        chat_session.updated_at = datetime.now(timezone.utc)
                    
                    # Award points
                    try:
                        from gamification_system import award_points
                        award_points(db, user.id, "ai_chat")
                    except Exception as gam_error:
                        print(f" Points award failed: {gam_error}")
                    
                    db.commit()
                    print(f" Message saved to chat {chat_id_int}")
            except Exception as save_error:
                print(f" Save error: {save_error}")
                db.rollback()
        
        execution_time = (time.time() - start_time) * 1000
        
        return {
            "answer": response,
            "ai_confidence": metadata.get("quality_score", 0.8),
            "misconception_detected": False,
            "should_request_feedback": metadata.get("quality_score", 0.8) < 0.6,
            "topics_discussed": metadata.get("concepts_discussed", []),
            "query_type": "agent",
            "model_used": "langgraph_chat_agent",
            "ai_provider": "LangGraph",
            "chat_id": chat_id_int,
            # Enhanced agent metadata
            "chat_mode": metadata.get("chat_mode", "tutoring"),
            "response_style": metadata.get("response_style", "conversational"),
            "emotional_state": metadata.get("emotional_state", "neutral"),
            "quality_score": metadata.get("quality_score", 0.7),
            "suggested_questions": metadata.get("suggested_questions", []),
            "learning_actions": metadata.get("learning_actions", []),
            "execution_time_ms": execution_time
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"\n ERROR IN ASK_AGENT: {str(e)}")
        print(f" Traceback:\n{error_details}\n")
        return {
            "answer": f"Error: {str(e)}",
            "ai_confidence": 0.3,
            "chat_mode": "error",
            "emotional_state": "neutral",
            "suggested_questions": [],
            "topics_discussed": ["error"]
        }


# ==================== Flashcard Agent Endpoint ====================

@app.post("/api/flashcard_agent/")
async def flashcard_agent_endpoint(
    user_id: str = Form(...),
    action: str = Form(...),  # generate, review, analyze, recommend, explain
    topic: Optional[str] = Form(None),
    content: Optional[str] = Form(None),
    card_count: int = Form(10),
    difficulty: str = Form("medium"),
    depth_level: str = Form("standard"),  # surface, standard, deep
    review_results: Optional[str] = Form(None),  # JSON string of review results
    is_public: bool = Form(False),
    db: Session = Depends(get_db)
):
    """
    Flashcard Agent endpoint for intelligent flashcard operations.
    Supports:
    - generate: Create flashcards from topic or content
    - review: Process review session with spaced repetition
    - analyze: Analyze performance and identify weak areas
    - recommend: Get personalized study recommendations
    - explain: Get detailed explanation of a concept
    """
    import time
    start_time = time.time()
    
    print(f"\n FLASHCARD_AGENT CALLED ")
    print(f" User: {user_id}, Action: {action}, Topic: {topic}")
    
    try:
        # Get user
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            print(f" User not found: {user_id}")
            return {
                "success": False,
                "error": "User not found",
                "action": action
            }
        
        # Try to use the LangGraph Flashcard Agent
        try:
            from agents.agent_api import get_flashcard_agent
            flashcard_agent = get_flashcard_agent()
            
            # Build state for the agent
            session_id = f"flashcard_{user.id}_{action}"
            
            state = {
                "user_id": str(user.id),
                "action": action,
                "session_id": session_id,
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "card_count": card_count,  # Direct state value
                "difficulty": difficulty,   # Direct state value
                "depth_level": depth_level,  # surface, standard, deep
                "action_params": {
                    "topic": topic,
                    "card_count": card_count,
                    "difficulty": difficulty,
                    "depth_level": depth_level
                }
            }
            
            # Add action-specific data
            if topic:
                state["topic"] = topic
                if content:
                    state["user_input"] = f"Generate flashcards from this content about {topic}"
                else:
                    state["user_input"] = f"Generate flashcards about {topic}"
            elif content:
                state["user_input"] = "Generate flashcards from this content"
            
            if content:
                state["source_content"] = content
            if review_results:
                try:
                    state["review_results"] = json.loads(review_results)
                except json.JSONDecodeError:
                    state["review_results"] = []
            
            # Invoke the flashcard agent
            result = await flashcard_agent.invoke(state)
            
            response_data = result.metadata.get("response_data", {})
            execution_time = (time.time() - start_time) * 1000
            
            # Save generated cards to database if action is generate
            if action == "generate" and response_data.get("cards"):
                cards = response_data["cards"]
                set_title = f"AI Generated: {topic or 'Study Session'}"
                
                # Generate unique share code
                import random
                import string
                share_code = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
                
                # Create flashcard set
                flashcard_set = models.FlashcardSet(
                    user_id=user.id,
                    title=set_title,
                    description=f"Generated {len(cards)} cards via AI Agent",
                    source_type="ai_agent",
                    share_code=share_code,
                    is_public=is_public
                )
                db.add(flashcard_set)
                db.commit()
                db.refresh(flashcard_set)
                
                logger.info(f"Agent flashcard set created: id={flashcard_set.id}, is_public={is_public}")
                
                # Add cards to set
                for card in cards:
                    db_card = models.Flashcard(
                        set_id=flashcard_set.id,
                        question=card.get("question", ""),
                        answer=card.get("answer", ""),
                        difficulty=card.get("difficulty", difficulty)
                    )
                    db.add(db_card)
                
                db.commit()
                
                # Award points
                try:
                    from gamification_system import award_points
                    award_points(db, user.id, "flashcard_create")
                except Exception as gam_error:
                    print(f" Points award failed: {gam_error}")
                
                return {
                    "success": True,
                    "action": action,
                    "response": result.response,
                    "cards": cards,
                    "set_id": flashcard_set.id,
                    "share_code": share_code,
                    "set_title": set_title,
                    "card_count": len(cards),
                    "execution_time_ms": execution_time
                }
            
            # Ensure cards are at top level for generate action
            response = {
                "success": result.success,
                "action": action,
                "response": result.response,
                "data": response_data,
                "execution_time_ms": execution_time
            }
            
            # If generate action, ensure cards are at top level
            if action == "generate" and response_data.get("cards"):
                response["cards"] = response_data["cards"]
                response["card_count"] = len(response_data["cards"])
            
            return response
            
        except Exception as agent_error:
            print(f" Flashcard agent failed, falling back to simple generation: {agent_error}")
            
            # Fallback to simple flashcard generation
            if action == "generate" and (topic or content):
                from flashcard_minimal import generate_flashcards_minimal
                
                # Use content if provided, otherwise use topic
                if content:
                    cards = generate_flashcards_minimal(
                        unified_ai,
                        content,
                        card_count,
                        difficulty,
                        is_topic=False  # Content mode
                    )
                    set_title = topic or "Chat Study Cards"
                else:
                    cards = generate_flashcards_minimal(
                        unified_ai,
                        topic,
                        card_count,
                        difficulty,
                        is_topic=True
                    )
                    set_title = f"Flashcards: {topic}"
                
                if cards:
                    # Save to database
                    import random
                    import string
                    share_code = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
                    
                    flashcard_set = models.FlashcardSet(
                        user_id=user.id,
                        title=set_title,
                        description=f"Generated {len(cards)} cards",
                        source_type="ai_generated",
                        share_code=share_code,
                        is_public=is_public
                    )
                    db.add(flashcard_set)
                    db.commit()
                    db.refresh(flashcard_set)
                    
                    for card in cards:
                        db_card = models.Flashcard(
                            set_id=flashcard_set.id,
                            question=card.get("question", ""),
                            answer=card.get("answer", ""),
                            difficulty=card.get("difficulty", difficulty)
                        )
                        db.add(db_card)
                    
                    db.commit()
                    
                    return {
                        "success": True,
                        "action": action,
                        "response": f"Generated {len(cards)} flashcards",
                        "cards": cards,
                        "set_id": flashcard_set.id,
                        "share_code": share_code,
                        "set_title": set_title,
                        "fallback": True
                    }
            
            return {
                "success": False,
                "action": action,
                "error": str(agent_error),
                "fallback": True
            }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"\n ERROR IN FLASHCARD_AGENT: {str(e)}")
        print(f" Traceback:\n{error_details}\n")
        return {
            "success": False,
            "action": action,
            "error": str(e)
        }


@app.post("/api/create_chat_session")
def create_chat_session(
    session_data: ChatSessionCreate,
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, session_data.user_id) or get_user_by_email(db, session_data.user_id)
        if not user:
            logger.error(f"User not found: {session_data.user_id}")
            raise HTTPException(status_code=404, detail="User not found")
        
        chat_session = models.ChatSession(
            user_id=user.id,
            title=session_data.title
        )
        db.add(chat_session)
        db.commit()
        db.refresh(chat_session)
        
        logger.info(f" Chat session created: ID={chat_session.id} for user {user.email}")
        
        return {
            "id": chat_session.id,
            "session_id": chat_session.id,
            "title": chat_session.title,
            "created_at": chat_session.created_at.isoformat() + 'Z',
            "updated_at": chat_session.updated_at.isoformat() + 'Z',
            "status": "success"
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error creating chat session: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to create chat session: {str(e)}")

@app.put("/api/rename_chat_session")
def rename_chat_session(
    data: dict,
    db: Session = Depends(get_db)
):
    try:
        chat_id = data.get('chat_id')
        new_title = data.get('new_title')
        
        if not chat_id or not new_title:
            raise HTTPException(status_code=400, detail="chat_id and new_title are required")
        
        chat_session = db.query(models.ChatSession).filter(
            models.ChatSession.id == chat_id
        ).first()
        
        if not chat_session:
            raise HTTPException(status_code=404, detail="Chat session not found")
        
        chat_session.title = new_title
        chat_session.updated_at = datetime.now(timezone.utc)
        db.commit()
        
        logger.info(f" Chat session renamed: ID={chat_id} to '{new_title}'")
        
        return {
            "status": "success",
            "chat_id": chat_id,
            "new_title": new_title
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error renaming chat session: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to rename chat session: {str(e)}")

@app.get("/api/get_user_achievements")
def get_user_achievements(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get user's achievements
        user_achievements = db.query(models.UserAchievement).filter(
            models.UserAchievement.user_id == user.id
        ).order_by(models.UserAchievement.earned_at.desc()).all()
        
        achievements = []
        for ua in user_achievements:
            achievement = db.query(models.Achievement).filter(
                models.Achievement.id == ua.achievement_id
            ).first()
            
            if achievement:
                achievements.append({
                    "id": achievement.id,
                    "name": achievement.name,
                    "description": achievement.description,
                    "icon": achievement.icon,
                    "points": achievement.points,
                    "category": achievement.category,
                    "rarity": achievement.rarity,
                    "earned_at": ua.earned_at.isoformat() + 'Z'
                })
        
        return {
            "achievements": achievements,
            "total_points": sum(a["points"] for a in achievements)
        }
        
    except Exception as e:
        logger.error(f"Error getting achievements: {str(e)}")
        return {
            "achievements": [],
            "total_points": 0
        }

@app.get("/api/get_chat_sessions")
def get_chat_sessions(user_id: str = Query(...), db: Session = Depends(get_db)):
    user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    sessions = db.query(models.ChatSession).filter(
        models.ChatSession.user_id == user.id
    ).order_by(models.ChatSession.updated_at.desc()).all()
    
    return {
        "sessions": [
            {
                "id": session.id,
                "title": session.title,
                "folder_id": session.folder_id,
                "created_at": session.created_at.isoformat() + 'Z' if session.created_at else None,
                "updated_at": session.updated_at.isoformat() + 'Z' if session.updated_at else None
            }
            for session in sessions
        ]
    }

@app.get("/api/get_chat_messages")
def get_chat_messages(chat_id: int = Query(...), db: Session = Depends(get_db)):
    try:
        logger.info(f"📥 Loading messages for chat_id: {chat_id}")
        
        messages = db.query(models.ChatMessage).filter(
            models.ChatMessage.chat_session_id == chat_id
        ).order_by(models.ChatMessage.timestamp.asc()).all()
        
        logger.info(f"Found {len(messages)} message pairs")
        
        result = []
        
        for msg in messages:
            # Add user message
            result.append({
                "id": f"user_{msg.id}",
                "type": "user",
                "content": msg.user_message,
                "timestamp": msg.timestamp.isoformat() + 'Z'
            })
            
            # Add AI message
            result.append({
                "id": f"ai_{msg.id}",
                "type": "ai",
                "content": msg.ai_response,
                "timestamp": msg.timestamp.isoformat() + 'Z',
                "aiConfidence": 0.85,
                "shouldRequestFeedback": False
            })
        
        logger.info(f"📤 Returning {len(result)} individual messages")
        return result
        
    except Exception as e:
        logger.error(f"Error in get_chat_messages: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
    
@app.post("/api/submit_response_feedback")
async def submit_response_feedback(
    user_id: str = Form(...),
    rating: int = Form(...),
    message_context: str = Form(None),
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        from ai_personality import AdaptiveLearningModel
        
        adaptive_model = AdaptiveLearningModel()
        adaptive_model.load_model_state(db, user.id)
        adaptive_model.update_from_feedback(rating, 'explanation_depth')
        adaptive_model.save_model_state(db, user.id)
        
        feedback = models.UserFeedback(
            user_id=user.id,
            feedback_type="rating",
            rating=rating,
            topic_context=message_context,
            is_processed=False
        )
        db.add(feedback)
        db.commit()
        
        return {
            "status": "success",
            "message": "Feedback recorded and model updated"
        }
        
    except Exception as e:
        logger.error(f"Error submitting feedback: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
@app.post("/api/submit_advanced_feedback")
async def submit_advanced_feedback(
    user_id: str = Form(...),
    rating: int = Form(...),
    feedback_text: str = Form(None),
    improvement_suggestion: str = Form(None),
    message_content: str = Form(None),
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        from ai_personality import AdaptiveLearningModel
        from neural_adaptation import get_rl_agent, ConversationContextAnalyzer
        
        adaptive_model = AdaptiveLearningModel()
        adaptive_model.load_model_state(db, user.id)
        adaptive_model.update_from_feedback(rating, 'explanation_depth')
        adaptive_model.save_model_state(db, user.id)
        
        rl_agent = get_rl_agent(db, user.id)
        
        reward = rl_agent.get_reward(rating, feedback_text)
        
        context = {
            'message_length': len(message_content or ''),
            'question_complexity': 0.5,
            'archetype': '',
            'time_of_day': datetime.now(timezone.utc).hour,
            'session_length': 0,
            'previous_ratings': [],
            'topic_keywords': [],
            'sentiment': ConversationContextAnalyzer.analyze_sentiment(feedback_text or ''),
            'formality_level': 0.5
        }
        
        state = rl_agent.encode_state(context)
        next_state = state
        action = rl_agent.network.predict(state)
        
        rl_agent.remember(state, action, reward, next_state)
        rl_agent.learn_from_experience(batch_size=16)
        
        if rating >= 4:
            rl_agent.learn_from_positive_feedback(message_content or '', rating)
        elif rating <= 2 and feedback_text:
            rl_agent.learn_from_negative_feedback(
                message_content or '', rating, feedback_text
            )
        
        if improvement_suggestion:
            correction_intent = ConversationContextAnalyzer.extract_correction_intent(
                improvement_suggestion
            )
            if correction_intent:
                rl_agent.save_correction(
                    message_content or '',
                    improvement_suggestion,
                    context
                )
        
        rl_agent.save_model()
        
        feedback = models.UserFeedback(
            user_id=user.id,
            feedback_type="rating_with_learning",
            rating=rating,
            feedback_text=feedback_text,
            topic_context=message_content,
            is_processed=True,
            resulted_in_improvement=True
        )
        db.add(feedback)
        db.commit()
        
        return {
            "status": "success",
            "message": "Feedback processed with neural network",
            "reward": float(reward),
            "model_updated": True,
            "corrections_learned": len(rl_agent.user_corrections)
        }
        
    except Exception as e:
        logger.error(f"Error in advanced feedback: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_chat_history/{session_id}")
async def get_chat_history(session_id: str, db: Session = Depends(get_db)):
    try:
        session_id_int = int(session_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid session ID")
    
    messages = db.query(models.ChatMessage).filter(
        models.ChatMessage.chat_session_id == session_id_int
    ).order_by(models.ChatMessage.timestamp.asc()).all()
    
    return {
        "session_id": session_id,
        "messages": [
            {
                "user_message": msg.user_message,
                "ai_response": msg.ai_response,
                "timestamp": msg.timestamp.isoformat() + 'Z'
            }
            for msg in messages
        ]
    }

@app.post("/api/save_chat_message")
def save_chat_message(message_data: ChatMessageSave, db: Session = Depends(get_db)):
    chat_session = db.query(models.ChatSession).filter(
        models.ChatSession.id == message_data.chat_id
    ).first()
    if not chat_session:
        raise HTTPException(status_code=404, detail="Chat session not found")
    
    # Check if this exact message already exists (to prevent double-saving and double-counting)
    existing_message = db.query(models.ChatMessage).filter(
        models.ChatMessage.chat_session_id == message_data.chat_id,
        models.ChatMessage.user_message == message_data.user_message,
        models.ChatMessage.ai_response == message_data.ai_response
    ).first()
    
    if existing_message:
        # Message already saved, don't award points again
        return {"status": "success", "message": "Message already exists"}
    
    chat_message = models.ChatMessage(
        chat_session_id=message_data.chat_id,
        user_message=message_data.user_message,
        ai_response=message_data.ai_response,
        is_user=True  # This represents a user message
    )
    db.add(chat_message)
    
    chat_session.updated_at = datetime.now(timezone.utc)
    
    message_count = db.query(models.ChatMessage).filter(
        models.ChatMessage.chat_session_id == message_data.chat_id
    ).count()
    
    if chat_session.title == "New Chat" and message_count == 0:
        words = message_data.user_message.strip().split()
        new_title = " ".join(words[:4]) + ("..." if len(words) > 4 else "")
        chat_session.title = new_title[:50] if new_title else "New Chat"
    
    # Only award points for USER messages, not AI responses
    # This endpoint is called once per user message (which includes the AI response)
    # So we award points once per call
    try:
        from gamification_system import award_points
        award_points(db, chat_session.user_id, "ai_chat")
    except Exception as gam_error:
        logger.warning(f"Failed to award AI chat points: {gam_error}")
    
    db.commit()
    return {"status": "success", "message": "Message saved successfully"}

@app.delete("/api/delete_chat_session/{session_id}")
def delete_chat_session(
    session_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    chat_session = db.query(models.ChatSession).filter(
        models.ChatSession.id == session_id,
        models.ChatSession.user_id == current_user.id
    ).first()
    
    if not chat_session:
        raise HTTPException(status_code=404, detail="Chat session not found")
    
    db.query(models.ChatMessage).filter(
        models.ChatMessage.chat_session_id == session_id
    ).delete()
    
    db.delete(chat_session)
    db.commit()
    
    return {"message": "Chat session deleted successfully"}

# ==================== NOTES ENDPOINTS ====================
# Add this helper function BEFORE the notes endpoints


def to_utc_iso(dt):
    """Convert datetime to ISO format with UTC timezone indicator"""
    if dt is None:
        return None
    return dt.isoformat() + 'Z' if not dt.isoformat() + 'Z'.endswith('Z') else dt.isoformat() + 'Z'

def clean_conversational_elements(text: str) -> str:
    """Remove conversational phrases from AI-generated content"""
    conversational_patterns = [
        r'^(Hi|Hello|Hey|Greetings)[,!.\s]+.*?\n',
        r'^(Sure|Of course|Certainly|Absolutely)[,!.\s]+.*?\n',
        r'^(Here\'s|Here is|Let me)[,\s]+.*?\n',
        r'^(I\'ll|I will|I can)[,\s]+.*?\n',
        r'Hope this helps.*$',
        r'Feel free to.*$',
        r'Let me know.*$',
    ]
    
    cleaned_text = text
    for pattern in conversational_patterns:
        cleaned_text = re.sub(pattern, '', cleaned_text, flags=re.MULTILINE | re.IGNORECASE)
    
    # Remove multiple newlines and clean up
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
    return cleaned_text.lstrip('\n\r\t ')


# ==================== NOTES ENDPOINTS ====================

@app.get("/api/get_notes")
def get_notes(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all notes for a user (excluding deleted)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # CRITICAL: Query with explicit filter for non-deleted notes
        notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            models.Note.is_deleted == False  # This line filters out deleted notes
        ).order_by(models.Note.updated_at.desc()).all()
        
        logger.info(f"Retrieved {len(notes)} active notes for user {user.email}")
        
        # Double-check: filter again in Python to be absolutely sure
        result = []
        for note in notes:
            if not note.is_deleted:  # Extra safety check
                result.append({
                    "id": note.id,
                    "title": note.title,
                    "content": note.content,
                    "created_at": note.created_at.isoformat() + 'Z',
                    "updated_at": note.updated_at.isoformat() + 'Z',
                    "is_favorite": getattr(note, 'is_favorite', False),
                    "folder_id": getattr(note, 'folder_id', None),
                    "custom_font": getattr(note, 'custom_font', 'Inter'),
                    "is_deleted": False
                })
        
        logger.info(f"Returning {len(result)} notes after filtering")
        return result
        
    except Exception as e:
        logger.error(f"Error getting notes: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))
# Add these endpoints to your main.py file
@app.get("/api/debug_notes/{user_id}")
def debug_notes(user_id: str, db: Session = Depends(get_db)):
    """Debug endpoint to see ALL notes"""
    user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
    if not user:
        return {"error": "User not found"}
    
    notes = db.query(models.Note).filter(
        models.Note.user_id == user.id
    ).all()
    
    result = []
    for note in notes:
        result.append({
            "id": note.id,
            "title": note.title,
            "is_deleted": note.is_deleted,
            "deleted_at": str(note.deleted_at) if note.deleted_at else None,
        })
    
    return {"total_notes": len(notes), "notes": result}

@app.post("/api/transcribe_audio/")
async def transcribe_audio(
    audio_file: UploadFile = File(...),
    user_id: str = Form(...),
    db: Session = Depends(get_db)
):
    temp_audio_path = None
    try:
        logger.info(f"🎤 Transcribing audio for user: {user_id}")
        logger.info(f"📁 File info - Name: {audio_file.filename}, Type: {audio_file.content_type}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        audio_content = await audio_file.read()
        logger.info(f"📊 Audio size: {len(audio_content)} bytes")
        
        if len(audio_content) == 0:
            raise HTTPException(status_code=400, detail="Empty audio file")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".webm", mode='wb') as temp_file:
            temp_file.write(audio_content)
            temp_audio_path = temp_file.name
        
        logger.info(f" Saved to: {temp_audio_path}")
        
        try:
            with open(temp_audio_path, "rb") as f:
                logger.info(" Calling Groq Whisper API...")
                transcription = groq_client.audio.transcriptions.create(
                    file=f,
                    model="whisper-large-v3-turbo",
                    response_format="json",
                    language="en"
                )
            
            transcript_text = transcription.text
            logger.info(f" Transcription successful: '{transcript_text[:100]}...'")
            
            return {
                "status": "success",
                "transcript": transcript_text,
                "length": len(transcript_text),
                "model_used": "whisper-large-v3-turbo"
            }
            
        except Exception as groq_error:
            logger.error(f" Groq API error: {str(groq_error)}")
            logger.error(f"Error type: {type(groq_error).__name__}")
            raise HTTPException(
                status_code=500, 
                detail=f"Transcription failed: {str(groq_error)}"
            )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f" Unexpected error: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to transcribe audio: {str(e)}"
        )
    finally:
        if temp_audio_path and os.path.exists(temp_audio_path):
            try:
                os.remove(temp_audio_path)
                logger.info(f" Cleaned up temp file: {temp_audio_path}")
            except Exception as cleanup_error:
                logger.warning(f" Failed to cleanup temp file: {cleanup_error}")

@app.get("/api/test_transcribe")
def test_transcribe_endpoint():
    """Test endpoint to verify transcription route exists"""
    return {
        "status": "endpoint exists",
        "message": "Transcribe audio endpoint is registered",
        "groq_available": GROQ_API_KEY is not None,
        "groq_key_prefix": GROQ_API_KEY[:10] if GROQ_API_KEY else "None"
    }


@app.post("/api/transcribe_audio_test/")
async def transcribe_audio_test(
    user_id: str = Form(...),
    db: Session = Depends(get_db)
):
    """Simplified test version without actual audio file"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        return {
            "status": "success",
            "message": "Endpoint working, user found",
            "user_id": user.id,
            "groq_configured": GROQ_API_KEY is not None
        }
    except Exception as e:
        logger.error(f"Test endpoint error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/fix_all_notes")
def fix_all_notes(db: Session = Depends(get_db)):
    """Emergency fix - set all NULL is_deleted to False and ensure data integrity"""
    try:
        # Get all notes
        all_notes = db.query(models.Note).all()
        
        fixed_count = 0
        verified_count = 0
        error_count = 0
        
        for note in all_notes:
            try:
                # Fix NULL or None is_deleted values
                if note.is_deleted is None:
                    note.is_deleted = False
                    fixed_count += 1
                
                # Fix any integer 1 that should be True (for SQLite compatibility)
                elif note.is_deleted == 1:
                    note.is_deleted = True
                    verified_count += 1
                
                # Fix any integer 0 that should be False
                elif note.is_deleted == 0:
                    note.is_deleted = False
                    verified_count += 1
                
                # If marked as deleted but no deleted_at timestamp, add one
                if note.is_deleted and note.deleted_at is None:
                    note.deleted_at = datetime.now(timezone.utc)
                    logger.info(f"Added deleted_at timestamp to note {note.id}")
                
                # If not deleted but has deleted_at timestamp, clear it
                if not note.is_deleted and note.deleted_at is not None:
                    note.deleted_at = None
                    logger.info(f"Cleared deleted_at timestamp from note {note.id}")
                    
            except Exception as e:
                error_count += 1
                logger.error(f"Error fixing note {note.id}: {str(e)}")
                continue
        
        # Commit all changes
        db.commit()
        
        logger.info(f" Note fix completed: {fixed_count} fixed, {verified_count} verified, {error_count} errors")
        
        return {
            "status": "success",
            "total_notes": len(all_notes),
            "fixed_count": fixed_count,
            "verified_count": verified_count,
            "error_count": error_count,
            "message": f"Successfully fixed {fixed_count} notes with NULL is_deleted values"
        }
    
    except Exception as e:
        logger.error(f"Critical error in fix_all_notes: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to fix notes: {str(e)}"
        )
  

@app.get("/api/get_folders")
def get_folders(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all folders for a user with note counts"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        folders = db.query(models.Folder).filter(
            models.Folder.user_id == user.id
        ).order_by(models.Folder.name.asc()).all()
        
        result = []
        for folder in folders:
            #  Count only non-deleted notes
            note_count = db.query(models.Note).filter(
                models.Note.folder_id == folder.id,
                models.Note.is_deleted == False  #  Exclude deleted notes from count
            ).count()
            
            result.append({
                "id": folder.id,
                "name": folder.name,
                "color": folder.color,
                "parent_id": folder.parent_id,
                "note_count": note_count,
                "created_at": folder.created_at.isoformat() + 'Z'
            })
        
        return {"folders": result}
    except Exception as e:
        logger.error(f"Error getting folders: {str(e)}")
        return {"folders": []}


@app.get("/api/get_trash")
def get_trash(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all deleted notes (trash)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get notes deleted within last 30 days (timezone-aware)
        thirty_days_ago = datetime.now(timezone.utc) - timedelta(days=30)
        
        notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            models.Note.is_deleted == True,
            models.Note.deleted_at != None,
            models.Note.deleted_at >= thirty_days_ago
        ).order_by(models.Note.deleted_at.desc()).all()
        
        logger.info(f"📥 Found {len(notes)} deleted notes for user {user.email}")
        
        result = []
        for note in notes:
            # Handle timezone-aware datetime safely
            if note.deleted_at:
                if note.deleted_at.tzinfo is None:
                    # If naive, make it aware
                    deleted_at_aware = note.deleted_at.replace(tzinfo=timezone.utc)
                else:
                    deleted_at_aware = note.deleted_at
                
                days_since_deletion = (datetime.now(timezone.utc) - deleted_at_aware).days
                days_remaining = max(0, 30 - days_since_deletion)
            else:
                days_remaining = 30
            
            result.append({
                "id": note.id,
                "title": note.title,
                "content": note.content,
                "deleted_at": note.deleted_at.isoformat() + 'Z' if note.deleted_at else None,
                "days_remaining": days_remaining
            })
        
        return {"trash": result}
        
    except Exception as e:
        logger.error(f"Error getting trash: {str(e)}", exc_info=True)
        return {"trash": []}


# ==================== ACTIVITY TIMELINE ENDPOINTS ====================

@app.get("/api/get_flashcards")
def get_flashcards(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all flashcards for a user (for activity timeline)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get all flashcard sets for the user
        flashcard_sets = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.user_id == user.id
        ).order_by(models.FlashcardSet.created_at.desc()).all()
        
        result = []
        for fs in flashcard_sets:
            # Get flashcards in this set
            flashcards = db.query(models.Flashcard).filter(
                models.Flashcard.set_id == fs.id
            ).all()
            
            for card in flashcards:
                result.append({
                    "id": card.id,
                    "set_id": fs.id,
                    "set_title": fs.title,
                    "question": card.question,
                    "answer": card.answer,
                    "difficulty": card.difficulty,
                    "created_at": (card.created_at.isoformat() + 'Z' if card.created_at else fs.created_at.isoformat() + 'Z'),
                    "updated_at": (card.updated_at.isoformat() + 'Z' if card.updated_at else fs.created_at.isoformat() + 'Z')
                })
        
        logger.info(f"Retrieved {len(result)} flashcards for user {user.email}")
        return result
        
    except Exception as e:
        logger.error(f"Error getting flashcards: {str(e)}", exc_info=True)
        return []


@app.get("/get_flashcards_in_set")
@app.get("/api/get_flashcards_in_set")
def get_flashcards_in_set(set_id: int = Query(...), db: Session = Depends(get_db)):
    """Get all flashcards in a specific set"""
    try:
        flashcard_set = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.id == set_id
        ).first()
        
        if not flashcard_set:
            raise HTTPException(status_code=404, detail="Flashcard set not found")
        
        flashcards = db.query(models.Flashcard).filter(
            models.Flashcard.set_id == set_id
        ).all()
        
        return {
            "set_id": flashcard_set.id,
            "set_title": flashcard_set.title,
            "share_code": getattr(flashcard_set, 'share_code', None),
            "description": flashcard_set.description or "",
            "flashcards": [
                {
                    "id": card.id,
                    "question": card.question,
                    "answer": card.answer,
                    "difficulty": card.difficulty or "medium",
                    "times_reviewed": card.times_reviewed or 0,
                    "correct_count": card.correct_count or 0,
                    "marked_for_review": card.marked_for_review if hasattr(card, 'marked_for_review') else False
                }
                for card in flashcards
            ]
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting flashcards in set: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/get_flashcard_history")
@app.get("/api/get_flashcard_history")
def get_flashcard_history(user_id: str = Query(...), limit: int = Query(50), db: Session = Depends(get_db)):
    """Get flashcard sets with mastery/accuracy for a user"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get all flashcard sets for the user
        flashcard_sets = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.user_id == user.id
        ).order_by(models.FlashcardSet.created_at.desc()).limit(limit).all()
        
        result = []
        for fs in flashcard_sets:
            # Get all cards in this set
            cards = db.query(models.Flashcard).filter(
                models.Flashcard.set_id == fs.id
            ).all()
            
            # Calculate accuracy percentage based on card reviews
            total_reviews = sum(c.times_reviewed or 0 for c in cards)
            total_correct = sum(c.correct_count or 0 for c in cards)
            
            if total_reviews > 0:
                accuracy = (total_correct / total_reviews) * 100
            else:
                accuracy = 0.0
            
            result.append({
                "id": fs.id,
                "share_code": getattr(fs, 'share_code', None),
                "title": fs.title,
                "description": fs.description or "",
                "card_count": len(cards),
                "accuracy_percentage": round(accuracy, 1),
                "source_type": fs.source_type or "manual",
                "is_public": fs.is_public if hasattr(fs, 'is_public') else False,
                "created_at": fs.created_at.isoformat() + 'Z' if fs.created_at else None,
                "updated_at": fs.updated_at.isoformat() + 'Z' if fs.updated_at else None
            })
        
        logger.info(f"Retrieved {len(result)} flashcard sets for user {user.email}")
        return {"flashcard_history": result}
        
    except Exception as e:
        logger.error(f"Error getting flashcard history: {str(e)}", exc_info=True)
        return {"flashcard_history": []}


@app.get("/get_flashcard_statistics")
@app.get("/api/get_flashcard_statistics")
def get_flashcard_statistics(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get flashcard statistics for a user"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"total_sets": 0, "total_cards": 0, "cards_mastered": 0, "average_accuracy": 0}
        
        # Count sets
        total_sets = db.query(func.count(models.FlashcardSet.id)).filter(
            models.FlashcardSet.user_id == user.id
        ).scalar() or 0
        
        # Count cards
        total_cards = db.query(func.count(models.Flashcard.id)).join(
            models.FlashcardSet, models.Flashcard.set_id == models.FlashcardSet.id
        ).filter(
            models.FlashcardSet.user_id == user.id
        ).scalar() or 0
        
        # Count mastered cards (correct_count >= 3)
        cards_mastered = db.query(func.count(models.Flashcard.id)).join(
            models.FlashcardSet, models.Flashcard.set_id == models.FlashcardSet.id
        ).filter(
            models.FlashcardSet.user_id == user.id,
            models.Flashcard.correct_count >= 3
        ).scalar() or 0
        
        # Calculate average accuracy
        cards = db.query(models.Flashcard).join(
            models.FlashcardSet, models.Flashcard.set_id == models.FlashcardSet.id
        ).filter(
            models.FlashcardSet.user_id == user.id
        ).all()
        
        total_reviews = sum(c.times_reviewed or 0 for c in cards)
        total_correct = sum(c.correct_count or 0 for c in cards)
        average_accuracy = (total_correct / total_reviews * 100) if total_reviews > 0 else 0
        
        return {
            "total_sets": total_sets,
            "total_cards": total_cards,
            "cards_mastered": cards_mastered,
            "average_accuracy": round(average_accuracy, 1)
        }
        
    except Exception as e:
        logger.error(f"Error getting flashcard statistics: {str(e)}", exc_info=True)
        return {"total_sets": 0, "total_cards": 0, "cards_mastered": 0, "average_accuracy": 0}


@app.get("/api/get_flashcards_for_review")
@app.post("/api/get_flashcards_for_review")
def get_flashcards_for_review(user_id: str = Query(None), db: Session = Depends(get_db)):
    """Get flashcards that need review"""
    try:
        if not user_id:
            return {"cards": [], "count": 0}
            
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"cards": [], "count": 0}
        
        # Get cards marked for review or with low accuracy
        cards = db.query(models.Flashcard).join(
            models.FlashcardSet, models.Flashcard.set_id == models.FlashcardSet.id
        ).filter(
            models.FlashcardSet.user_id == user.id,
            models.Flashcard.marked_for_review == True
        ).all()
        
        result = []
        for card in cards:
            result.append({
                "id": card.id,
                "set_id": card.set_id,
                "question": card.question,
                "answer": card.answer,
                "times_reviewed": card.times_reviewed or 0,
                "correct_count": card.correct_count or 0
            })
        
        return {"cards": result, "count": len(result)}
        
    except Exception as e:
        logger.error(f"Error getting flashcards for review: {str(e)}", exc_info=True)
        return {"cards": [], "count": 0}


@app.post("/generate_flashcards")
@app.post("/api/generate_flashcards")
async def generate_flashcards_endpoint(
    user_id: str = Form(...),
    topic: str = Form(None),
    generation_type: str = Form("topic"),
    chat_data: str = Form(None),
    card_count: int = Form(10),
    difficulty: str = Form("medium"),
    set_title: str = Form(None),
    is_public: bool = Form(False),
    db: Session = Depends(get_db)
):
    """Generate flashcards from topic or chat history"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Generate flashcards using AI
        if generation_type == "topic" and topic:
            prompt = f"""Generate {card_count} flashcards about: {topic}
            
Difficulty level: {difficulty}

Return ONLY a valid JSON array of flashcard objects. Each object must have:
- "question": The question or front of the card
- "answer": The answer or back of the card
- "difficulty": "{difficulty}"

Example format:
[
  {{"question": "What is X?", "answer": "X is...", "difficulty": "{difficulty}"}},
  {{"question": "Explain Y", "answer": "Y is...", "difficulty": "{difficulty}"}}
]

Generate educational, clear flashcards. Return ONLY the JSON array, no other text."""
            
            content = topic
        elif generation_type == "chat_history" and chat_data:
            import json
            try:
                chat_history = json.loads(chat_data)
                chat_text = "\n".join([f"Q: {msg.get('question', '')}\nA: {msg.get('answer', '')}" for msg in chat_history[:10]])
            except:
                chat_text = chat_data
            
            prompt = f"""Based on this conversation, generate {card_count} flashcards to help remember the key concepts:

{chat_text}

Return ONLY a valid JSON array of flashcard objects. Each object must have:
- "question": The question or front of the card
- "answer": The answer or back of the card
- "difficulty": "{difficulty}"

Generate educational flashcards covering the main topics discussed. Return ONLY the JSON array, no other text."""
            
            content = chat_text
        else:
            raise HTTPException(status_code=400, detail="Provide topic or chat_data")
        
        # Call AI to generate flashcards
        ai_response = unified_ai.generate(prompt, max_tokens=2000, temperature=0.7)
        
        # Parse the response
        import json
        import re
        
        # Try to extract JSON from response
        json_match = re.search(r'\[[\s\S]*\]', ai_response)
        if json_match:
            try:
                flashcards = json.loads(json_match.group())
            except json.JSONDecodeError:
                flashcards = []
        else:
            flashcards = []
        
        if not flashcards:
            # Fallback: create simple flashcards
            flashcards = [
                {"question": f"What is {topic}?", "answer": f"A concept related to {topic}", "difficulty": difficulty}
            ]
        
        # Save to database if set_title is provided
        saved_to_set = False
        set_id = None
        
        if set_title and flashcards:
            flashcard_set = models.FlashcardSet(
                user_id=user.id,
                title=set_title,
                description=f"Generated {len(flashcards)} cards",
                source_type="ai_generated",
                is_public=is_public
            )
            db.add(flashcard_set)
            db.commit()
            db.refresh(flashcard_set)
            set_id = flashcard_set.id
            logger.info(f"Flashcard set created: id={set_id}, is_public={is_public}")
            
            for card in flashcards:
                db_card = models.Flashcard(
                    set_id=flashcard_set.id,
                    question=card.get("question", ""),
                    answer=card.get("answer", ""),
                    difficulty=card.get("difficulty", difficulty)
                )
                db.add(db_card)
            
            db.commit()
            saved_to_set = True
            
            # Add IDs to flashcards
            db_cards = db.query(models.Flashcard).filter(
                models.Flashcard.set_id == flashcard_set.id
            ).all()
            for i, card in enumerate(flashcards):
                if i < len(db_cards):
                    card["id"] = db_cards[i].id
        
        return {
            "success": True,
            "flashcards": flashcards,
            "saved_to_set": saved_to_set,
            "set_id": set_id,
            "set_title": set_title
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating flashcards: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/get_quiz_history")
def get_quiz_history(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get quiz history for a user (for activity timeline)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Try to get quiz sessions from various possible tables
        result = []
        
        # Check if QuizSession model exists
        if hasattr(models, 'QuizSession'):
            quiz_sessions = db.query(models.QuizSession).filter(
                models.QuizSession.user_id == user.id
            ).order_by(models.QuizSession.created_at.desc()).all()
            
            for session in quiz_sessions:
                result.append({
                    "id": session.id,
                    "title": getattr(session, 'title', 'Quiz Session'),
                    "score": getattr(session, 'score', 0),
                    "total_questions": getattr(session, 'total_questions', 0),
                    "correct_answers": getattr(session, 'correct_answers', 0),
                    "completed_at": (session.completed_at.isoformat() + 'Z' if hasattr(session, 'completed_at') and session.completed_at else session.created_at.isoformat() + 'Z'),
                    "created_at": session.created_at.isoformat() + 'Z'
                })
        
        # Check if QuestionSet model exists (alternative)
        elif hasattr(models, 'QuestionSet'):
            question_sets = db.query(models.QuestionSet).filter(
                models.QuestionSet.user_id == user.id
            ).order_by(models.QuestionSet.created_at.desc()).all()
            
            for qs in question_sets:
                result.append({
                    "id": qs.id,
                    "title": qs.title or 'Quiz Session',
                    "score": 0,
                    "total_questions": getattr(qs, 'question_count', 0),
                    "correct_answers": 0,
                    "completed_at": qs.created_at.isoformat() + 'Z',
                    "created_at": qs.created_at.isoformat() + 'Z'
                })
        
        logger.info(f"Retrieved {len(result)} quiz sessions for user {user.email}")
        return result
        
    except Exception as e:
        logger.error(f"Error getting quiz history: {str(e)}", exc_info=True)
        return []


# ENHANCED AI CONTENT GENERATION WITH SPECIFIC PROMPTS
@app.post("/api/generate_note_content/")
async def generate_note_content(
    user_id: str = Form(...),
    prompt: str = Form(...),
    content_type: str = Form("general"),
    existing_content: str = Form(""),
    db: Session = Depends(get_db)
):
    """
    Generate AI content for notes with specific prompt templates
    """
    try:
        logger.info(f"Generating {content_type} content for user: {user_id}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            logger.error(f"User not found: {user_id}")
            raise HTTPException(status_code=404, detail="User not found")
        
        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()
        
        user_profile = build_user_profile_dict(user, comprehensive_profile)
        
        content_context = existing_content[-500:] if existing_content else ""
        
        # OPTIMIZED SYSTEM PROMPTS FOR EACH ACTION TYPE (81% TOKEN REDUCTION)
        action_prompts = {
            "explain": f"Explain for {user_profile.get('difficulty_level', 'intermediate')} level. No greetings. Use HTML tags. Include examples.",
            
            "key_points": f"Extract 5-10 key points. Level: {user_profile.get('difficulty_level', 'intermediate')}. Use <ul><li>. No greetings.",
            
            "guide": f"Comprehensive guide. Level: {user_profile.get('difficulty_level', 'intermediate')}. Use <h2>, <h3>, examples. No greetings.",
            
            "summary": "Summarize concisely. Use <h2>Summary</h2>. No greetings.",
            
            "general": f"Educational content for {user_profile.get('difficulty_level', 'intermediate')} level. Use HTML. No greetings."
        }
        
        system_prompt = action_prompts.get(content_type, action_prompts["general"])
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt}
        ]

        logger.info("Calling AI (Gemini primary, Groq fallback)...")
        
        # Build prompt from messages
        full_prompt = "\n\n".join([f"{msg['role']}: {msg['content']}" for msg in messages])
        
        generated_content = call_ai(full_prompt, max_tokens=2048, temperature=0.7)
        logger.info(f"Generated content length: {len(generated_content)} characters")
        
        generated_content = clean_conversational_elements(generated_content)
        
        if not generated_content.strip().startswith('<'):
            generated_content = f"<h2>{prompt}</h2>\n<p>{generated_content}</p>"
        
        logger.info("Content generation successful")
        
        return {
            "status": "success",
            "content": generated_content,
            "content_type": content_type,
            "word_count": len(generated_content.split()),
            "model_used": GROQ_MODEL
        }
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error generating note content: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to generate content: {str(e)}"
        )


# ENHANCED AI WRITING ASSISTANT WITH MORE ACTIONS
@app.post("/api/ai_writing_assistant/")
async def ai_writing_assistant(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """
    AI Writing Assistant with multiple transformation actions
    """
    try:
        user_id = payload.get("user_id")
        content = payload.get("content", "")
        action = payload.get("action", "improve")
        tone = payload.get("tone", "professional")
        context = payload.get("context", "")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        user_profile = build_user_profile_dict(user)
        
        # DIFFERENT PROMPTS FOR EACH ACTION
        action_prompts = {
            "generate": f"""Write a detailed, comprehensive article explaining this topic in depth.

Topic: {content}

Write a complete article with:
1. An introduction paragraph explaining what {content} is
2. 2-3 body paragraphs with detailed explanations, examples, and facts
3. A conclusion paragraph

Use a {tone} tone. Write at least 300 words. Start writing now:""",

            "continue": f"""Continue writing this text naturally and coherently.

RULES:
- Match the existing style and tone
- Continue logically from where it ends
- Add 2-3 more sentences or paragraphs
- Do NOT repeat what's already written

Text to continue:
{content}

Continue writing:""",

            "improve": f"""Improve and enhance this text while keeping the same meaning.

RULES:
- Make it more clear and professional
- Fix awkward phrasing
- Improve word choice
- Keep the same structure and meaning
- Do NOT add new information

Original text:
{content}

Improved version:""",

            "simplify": f"""Simplify this text to make it easier to understand.

RULES:
- Use simpler words and shorter sentences
- Explain technical terms
- Break down complex ideas
- Keep all important information
- Make it accessible to a general audience

Text to simplify:
{content}

Simplified version:""",

            "expand": f"""Expand this text with more details, examples, and explanations.

RULES:
- Add relevant details and context
- Include practical examples
- Explain concepts more thoroughly
- Increase depth without redundancy
- Maintain the original tone

Text to expand:
{content}

Expanded version:""",

            "tone_change": f"""Rewrite this text in a {tone} tone.

RULES:
- Change the tone to be {tone}
- Keep the same information
- Adjust language and style appropriately
- Maintain clarity

Original text:
{content}

Rewritten in {tone} tone:""",

            "grammar": f"""Fix all grammar and spelling errors in this text.

RULES:
- Correct spelling mistakes
- Fix grammatical errors
- Improve punctuation
- Do NOT change the meaning or style
- Keep the same tone and voice

Text to correct:
{content}

Corrected version:""",

            "summarize": f"""Create a concise summary of this text.

RULES:
- Extract only the most important points
- Keep it brief but complete
- Maintain key information
- Use clear, simple language

Text to summarize:
{content}

Summary:""",

            "code": f"""Review, explain, or improve this code.

RULES:
- Explain what the code does
- Suggest improvements if any
- Point out potential issues
- Add helpful comments
- Keep the same programming language

Code:
{content}

Analysis and improved version:"""
        }
        
        prompt = action_prompts.get(action, action_prompts["improve"])
        
        if action == "generate":
            system_prompt = "Write detailed article (300+ words). No intro phrases. Start writing."
        else:
            system_prompt = "Process text only. No greetings or comments."

        # Use more tokens for generate action
        max_tokens = 4096 if action == "generate" else 2048
        
        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.8 if action == "generate" else 0.7,
            max_tokens=max_tokens,
            top_p=0.9,
        )
        
        result = chat_completion.choices[0].message.content.strip()
        # Post-process math notation
        result = process_math_in_response(result)
        result = enhance_display_math(result)
        
        # Log for debugging
        print(f"🔥 AI Action: {action}")
        print(f"🔥 Input length: {len(content)} chars")
        print(f"🔥 Output length: {len(result)} chars")
        print(f"🔥 Result preview: {result[:200]}...")
        
        # If generate action returned too short, try again with more explicit prompt
        if action == "generate" and len(result.split()) < 50:
            print(f" Result too short ({len(result.split())} words), retrying with explicit prompt...")
            
            retry_prompt = f"""Write a detailed educational article about {content}.

Requirements:
- Minimum 5 paragraphs
- Explain what it is, how it works, why it matters
- Include examples and real-world applications
- Write in {tone} tone
- Be thorough and informative

Start writing the article:"""
            
            chat_completion = groq_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "Write detailed articles."},
                    {"role": "user", "content": retry_prompt}
                ],
                model=GROQ_MODEL,
                temperature=0.9,
                max_tokens=4096,
                top_p=0.95,
            )
            
            result = chat_completion.choices[0].message.content.strip()
            # Post-process math notation
            result = process_math_in_response(result)
            result = enhance_display_math(result)
            print(f"🔥 Retry result length: {len(result)} chars")
            print(f"🔥 Retry result preview: {result[:200]}...")
        
        return {
            "status": "success",
            "action": action,
            "result": result,
            "model_used": GROQ_MODEL,
            "original_length": len(content.split()),
            "new_length": len(result.split())
        }
        
    except Exception as e:
        logger.error(f"AI writing assistant error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"AI assistant failed: {str(e)}")

@app.post("/api/create_note")
def create_note(note_data: NoteCreate, db: Session = Depends(get_db)):
    """Create a new note"""
    try:
        user = get_user_by_username(db, note_data.user_id) or get_user_by_email(db, note_data.user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        new_note = models.Note(
            user_id=user.id,
            title=note_data.title,
            content=note_data.content
        )
        db.add(new_note)
        db.commit()
        db.refresh(new_note)
        
        # Award points for note creation (20 pts)
        try:
            from gamification_system import award_points
            award_points(db, user.id, "note_created")
            db.commit()
        except Exception as gam_error:
            logger.warning(f"Failed to award note creation points: {gam_error}")
        
        return {
            "id": new_note.id,
            "title": new_note.title,
            "content": new_note.content,
            "user_id": user.id,
            "created_at": new_note.created_at.isoformat() + 'Z',
            "updated_at": new_note.updated_at.isoformat() + 'Z',
            "status": "success",
            "message": "Note created successfully"
        }
        
    except Exception as e:
        logger.error(f"Error creating note: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to create note: {str(e)}")


@app.get("/api/get_notes")
def get_notes(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all notes for a user (excluding deleted)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # SQLite stores boolean as 0/1, so check for both False and 0
        notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            (models.Note.is_deleted == False) | (models.Note.is_deleted == 0) | (models.Note.is_deleted == None)
        ).order_by(models.Note.updated_at.desc()).all()
        
        # Extra Python filter to be absolutely sure
        active_notes = [n for n in notes if not n.is_deleted and n.deleted_at is None]
        
        logger.info(f"Retrieved {len(active_notes)} active notes for user {user.email}")
        
        return [
            {
                "id": note.id,
                "title": note.title,
                "content": note.content,
                "created_at": note.created_at.isoformat() + 'Z' if note.created_at else None,
                "updated_at": note.updated_at.isoformat() + 'Z' if note.updated_at else None,
                "is_favorite": getattr(note, 'is_favorite', False),
                "folder_id": getattr(note, 'folder_id', None),
                "custom_font": getattr(note, 'custom_font', 'Inter'),
                "is_deleted": False
            }
            for note in active_notes
        ]
        
    except Exception as e:
        logger.error(f"Error getting notes: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


@app.put("/api/soft_delete_note/{note_id}")
def soft_delete_note(note_id: int, db: Session = Depends(get_db)):
    """Move note to trash (soft delete)"""
    try:
        note = db.query(models.Note).filter(models.Note.id == note_id).first()
        
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        # Set as integer 1 for SQLite
        note.is_deleted = 1
        note.deleted_at = datetime.now(timezone.utc)
        
        db.commit()
        db.refresh(note)
        
        logger.info(f"Note {note.id} moved to trash - is_deleted={note.is_deleted}")
        
        return {
            "message": "Note moved to trash successfully",
            "note_id": note.id,
            "status": "success"
        }
        
    except Exception as e:
        logger.error(f"Error deleting note: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.put("/api/update_note")
def update_note(note_data: NoteUpdate, db: Session = Depends(get_db)):
    """Update an existing note"""
    try:
        note = db.query(models.Note).filter(models.Note.id == note_data.note_id).first()
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        # Check if deleted (handle both boolean and integer)
        if note.is_deleted or note.is_deleted == 1 or note.deleted_at is not None:
            logger.warning(f"Rejected update to deleted note {note.id}")
            raise HTTPException(status_code=400, detail="Cannot update a deleted note")
        
        note.title = note_data.title
        note.content = note_data.content
        note.updated_at = datetime.now(timezone.utc)
        
        db.commit()
        db.refresh(note)
        
        logger.info(f"Note {note.id} updated successfully")
        
        return {
            "id": note.id,
            "title": note.title,
            "content": note.content,
            "updated_at": note.updated_at.isoformat() + 'Z',
            "is_deleted": False,
            "status": "success"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error updating note: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/delete_note/{note_id}")
def delete_note(note_id: int, db: Session = Depends(get_db)):
    """Delete a note"""
    note = db.query(models.Note).filter(models.Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    
    db.delete(note)
    db.commit()
    
    return {"message": "Note deleted successfully"}


@app.post("/api/generate_note_summary/")
async def generate_note_summary(
    user_id: str = Form(...),
    conversation_data: str = Form(...),
    session_titles: str = Form(...),
    import_mode: str = Form("summary"),
    db: Session = Depends(get_db)
):
    """Generate a summary for chat-to-note conversion"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        user_profile = build_user_profile_dict(user)
        
        if import_mode == "summary":
            prompt = f"""Create study notes from this conversation. Format with headers, bullet points, and key concepts.

Conversation: {conversation_data[:3000]}

Provide a clear title and well-formatted study notes."""
        elif import_mode == "exam_prep":
            prompt = f"""Create an exam preparation guide from this conversation. Include:
1. Executive Summary
2. Key Concepts
3. Study Strategy
4. Review Checklist

Conversation: {conversation_data[:3000]}

Format as a comprehensive exam prep guide with a clear title."""
        else:
            prompt = f"""Create a formatted transcript from this conversation.

Conversation: {conversation_data[:3000]}

Provide a title and formatted content."""
        
        response = await generate_ai_response(prompt, user_profile)
        
        lines = response.split('\n')
        title = lines[0].replace('#', '').strip() if lines else "Study Notes"
        content = '\n'.join(lines[1:]) if len(lines) > 1 else response
        
        return {
            "title": title[:100],
            "content": content
        }
        
    except Exception as e:
        logger.error(f"Error generating note summary: {str(e)}")
        return {
            "title": "Study Notes",
            "content": f"# Study Notes\n\n{conversation_data[:1000]}"
        }

# ==================== END NOTES ENDPOINTS ====================# Add this new endpoint to your main.py file

# Add this to your main.py (replace the existing generate_note_content endpoint)






# Optional: Add endpoint to generate content from existing chat
@app.post("/api/convert_chat_to_note_content/")
async def convert_chat_to_note_content(
    request: Request,
    db: Session = Depends(get_db)
):
    """
    Convert an existing chat session into formatted note content
    """
    try:
        # Parse request body
        try:
            body = await request.json()
            user_id = body.get('user_id')
            chat_session_id = body.get('chat_session_id')
            format_style = body.get('format_style', 'comprehensive')
        except Exception as e:
            logger.error(f"Error parsing request body: {str(e)}")
            raise HTTPException(status_code=422, detail=f"Invalid request format: {str(e)}")
        
        # Validate inputs first
        if not user_id or not chat_session_id:
            logger.warning(f"Missing parameters - user_id: {user_id}, chat_session_id: {chat_session_id}")
            raise HTTPException(status_code=400, detail="Missing required parameters: user_id and chat_session_id are required")
        
        logger.info(f"Converting chat {chat_session_id} to note for user {user_id}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get chat session
        chat_session = db.query(models.ChatSession).filter(
            models.ChatSession.id == chat_session_id,
            models.ChatSession.user_id == user.id
        ).first()
        
        if not chat_session:
            raise HTTPException(status_code=404, detail="Chat session not found")
        
        # Get all messages from the session
        messages = db.query(models.ChatMessage).filter(
            models.ChatMessage.chat_session_id == chat_session_id
        ).order_by(models.ChatMessage.timestamp.asc()).all()
        
        if not messages:
            raise HTTPException(status_code=404, detail="No messages found in chat session")
        
        # Compile conversation data
        conversation_text = []
        for msg in messages:
            conversation_text.append(f"Q: {msg.user_message}")
            conversation_text.append(f"A: {msg.ai_response}")
        
        full_conversation = "\n\n".join(conversation_text)
        
        # Generate formatted notes from conversation
        if format_style == "comprehensive":
            prompt = f"""Convert this chat conversation into comprehensive study notes.

Conversation:
{full_conversation[:4000]}

Create well-structured notes with:
- Clear headings and subheadings
- Key concepts explained
- Important points highlighted
- Examples included
- Summary at the end

Use HTML formatting. Start directly with content (NO conversational elements)."""

        elif format_style == "summary":
            prompt = f"""Create a concise summary of this chat conversation:

{full_conversation[:4000]}

Include:
- Main topics discussed
- Key takeaways (bullet points)
- Important conclusions

Use HTML formatting. Start directly with content."""

        elif format_style == "key_points":
            prompt = f"""Extract the key points from this conversation:

{full_conversation[:4000]}

Format as:
- Numbered or bulleted list of main points
- Brief explanations where needed
- Organized by topic

Use HTML formatting. Start directly with content."""
        
        user_profile = build_user_profile_dict(user)
        
        system_prompt = "Convert to notes. HTML format. No greetings."

        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=2048,
            top_p=0.9,
        )
        
        generated_content = chat_completion.choices[0].message.content
        # Post-process math notation
        generated_content = process_math_in_response(generated_content)
        generated_content = enhance_display_math(generated_content)
        generated_content = clean_conversational_elements(generated_content)
        
        # Generate a short title from the conversation
        title_prompt = f"""Based on this conversation, create a SHORT 2-3 word title that captures the main topic.

Conversation summary: {full_conversation[:500]}

Requirements:
- ONLY 2-3 words maximum
- Capitalize properly
- No punctuation
- Be specific and descriptive

Examples: "Python Basics", "Machine Learning", "Study Tips"

Return ONLY the title, nothing else."""

        title_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": "Generate 2-3 word title only."},
                {"role": "user", "content": title_prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.5,
            max_tokens=20,
        )
        
        generated_title = title_completion.choices[0].message.content.strip().strip('"').strip("'")
        
        return {
            "status": "success",
            "content": generated_content,
            "title": generated_title,
            "format_style": format_style,
            "original_message_count": len(messages),
            "word_count": len(generated_content.split())
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error converting chat to note: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to convert chat: {str(e)}")


# Optional: Endpoint to expand/elaborate on selected text
@app.post("/api/expand_note_content/")
async def expand_note_content(
    user_id: str = Form(...),
    selected_text: str = Form(...),
    expansion_type: str = Form("elaborate"),  # elaborate, explain, add_examples
    db: Session = Depends(get_db)
):
    """
    Expand or elaborate on selected text in a note
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        user_profile = build_user_profile_dict(user)
        
        if expansion_type == "elaborate":
            instruction = f"Provide a detailed elaboration of this text: {selected_text}\n\nAdd more depth, context, and explanation."
        elif expansion_type == "explain":
            instruction = f"Explain this concept in simpler terms: {selected_text}\n\nMake it easy to understand."
        elif expansion_type == "add_examples":
            instruction = f"Provide practical examples for: {selected_text}\n\nInclude 2-3 real-world examples."
        else:
            instruction = f"Expand on this text with additional information: {selected_text}"
        
        system_prompt = f"Expand for {user_profile.get('difficulty_level', 'intermediate')} level. HTML. No greetings."

        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": instruction}
            ],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=1024,
            top_p=0.9,
        )
        
        expanded_content = chat_completion.choices[0].message.content
        # Post-process math notation
        expanded_content = process_math_in_response(expanded_content)
        expanded_content = enhance_display_math(expanded_content)
        expanded_content = clean_conversational_elements(expanded_content)
        
        return {
            "status": "success",
            "expanded_content": expanded_content,
            "expansion_type": expansion_type,
            "original_length": len(selected_text),
            "expanded_length": len(expanded_content)
        }
        
    except Exception as e:
        logger.error(f"Error expanding content: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to expand content: {str(e)}")

# ==================== AI MEDIA PROCESSING ENDPOINTS ====================

from ai_media_processor import ai_media_processor
from media_models import MediaUpload, TranscriptionSegment, GeneratedNote, SpeakerSegment

@app.post("/api/media/process")
async def process_media(
    user_id: str = Form(...),
    file: UploadFile = File(None),
    youtube_url: str = Form(None),
    note_style: str = Form("detailed"),
    difficulty: str = Form("intermediate"),
    subject: str = Form("general"),
    custom_instructions: str = Form(None),
    generate_flashcards: bool = Form(False),
    generate_quiz: bool = Form(False),
    db: Session = Depends(get_db)
):
    """
    Comprehensive AI-powered media processing
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        transcript_data = None
        analysis_data = None
        source_type = None
        filename = None
        file_path = None
        
        # Debug logging
        logger.info(f"Received - File: {file.filename if file else None}, YouTube URL: '{youtube_url}'")
        
        # Process YouTube URL
        if youtube_url and youtube_url.strip():
            logger.info(f"Processing YouTube URL: {youtube_url}")
            result = await ai_media_processor.process_youtube_video(youtube_url.strip())
            
            if not result.get("success"):
                raise HTTPException(status_code=400, detail=result.get("error", "Failed to process YouTube video"))
            
            transcript_data = result
            source_type = "youtube"
            filename = result["video_info"]["title"]
            
        # Process uploaded file
        elif file:
            logger.info(f"Processing uploaded file: {file.filename}")
            
            # Save file temporarily
            upload_dir = "backend/uploads/media"
            os.makedirs(upload_dir, exist_ok=True)
            
            file_path = os.path.join(upload_dir, f"{user.id}_{int(datetime.now().timestamp())}_{file.filename}")
            
            with open(file_path, "wb") as f:
                content = await file.read()
                f.write(content)
            
            # Transcribe audio
            result = await ai_media_processor.transcribe_audio_groq(file_path)
            
            if not result.get("success"):
                raise HTTPException(status_code=400, detail=result.get("error", "Failed to transcribe audio"))
            
            transcript_data = result
            source_type = "upload"
            filename = file.filename
            
        else:
            raise HTTPException(status_code=400, detail="Please provide either a file or YouTube URL")
        
        # AI Analysis
        logger.info("Performing AI analysis...")
        analysis_result = await ai_media_processor.analyze_transcript_ai(
            transcript_data["transcript"],
            {"subject": subject, "difficulty": difficulty}
        )
        
        if not analysis_result.get("success"):
            logger.error(f"AI analysis failed: {analysis_result.get('error', 'Unknown error')}")
            raise HTTPException(status_code=500, detail=f"AI analysis failed: {analysis_result.get('error', 'Unknown error')}")
        
        analysis_data = analysis_result.get("analysis", {})
        logger.info(f"Analysis data keys: {list(analysis_data.keys())}")
        logger.info(f"Analysis summary: {analysis_data.get('summary', 'No summary')[:100]}")
        
        # Generate Notes
        logger.info(f"Generating {note_style} notes...")
        notes_result = await ai_media_processor.generate_notes_ai(
            transcript_data["transcript"],
            analysis_data,
            note_style,
            {
                "difficulty": difficulty,
                "subject": subject,
                "custom_instructions": custom_instructions
            }
        )
        
        if not notes_result.get("success"):
            raise HTTPException(status_code=500, detail="Note generation failed")
        
        # Generate Flashcards (optional)
        flashcards = []
        if generate_flashcards:
            logger.info("Generating flashcards...")
            flashcard_result = await ai_media_processor.generate_flashcards_ai(
                transcript_data["transcript"],
                analysis_data,
                count=10
            )
            if flashcard_result.get("success"):
                flashcards = flashcard_result.get("flashcards", [])
                logger.info(f"Generated {len(flashcards)} flashcards")
            else:
                logger.warning(f"Flashcard generation failed: {flashcard_result.get('error', 'Unknown')}")
        
        # Generate Quiz (optional)
        quiz_questions = []
        if generate_quiz:
            logger.info("Generating quiz...")
            quiz_result = await ai_media_processor.generate_quiz_ai(
                transcript_data["transcript"],
                analysis_data,
                count=10
            )
            if quiz_result.get("success"):
                quiz_questions = quiz_result.get("questions", [])
                logger.info(f"Generated {len(quiz_questions)} quiz questions")
            else:
                logger.warning(f"Quiz generation failed: {quiz_result.get('error', 'Unknown')}")
        
        # Extract key moments
        key_moments = []
        if transcript_data.get("has_timestamps") and transcript_data.get("segments"):
            key_moments = await ai_media_processor.extract_key_moments(
                transcript_data["segments"],
                analysis_data
            )
        
        # Prepare response
        response = {
            "success": True,
            "source_type": source_type,
            "filename": filename,
            "transcript": transcript_data["transcript"],
            "language": transcript_data.get("language", "en"),
            "language_name": ai_media_processor.get_language_name(transcript_data.get("language", "en")),
            "duration": transcript_data.get("duration", 0),
            "has_timestamps": transcript_data.get("has_timestamps", False),
            "segments": transcript_data.get("segments", [])[:50],  # Limit segments in response
            "analysis": analysis_data,
            "notes": {
                "content": notes_result["content"],
                "style": note_style
            },
            "flashcards": flashcards,
            "quiz_questions": quiz_questions,
            "key_moments": key_moments,
            "video_info": transcript_data.get("video_info") if source_type == "youtube" else None
        }
        
        return response
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Media processing error: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

@app.post("/api/media/save-notes")
async def save_media_notes(
    user_id: str = Body(...),
    title: str = Body(...),
    content: str = Body(...),
    transcript: str = Body(None),
    analysis: dict = Body(None),
    flashcards: list = Body(None),
    quiz_questions: list = Body(None),
    key_moments: list = Body(None),
    db: Session = Depends(get_db)
):
    """Save generated notes to user's notes"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        import json
        
        # Create note with special marker for media notes
        new_note = models.Note(
            user_id=user.id,
            title=title,
            content=content,
            custom_font="__MEDIA_NOTE__",  # Special marker to identify media-generated notes
            transcript=transcript,
            analysis=json.dumps(analysis) if analysis else None,
            flashcards=json.dumps(flashcards) if flashcards else None,
            quiz_questions=json.dumps(quiz_questions) if quiz_questions else None,
            key_moments=json.dumps(key_moments) if key_moments else None,
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow()
        )
        
        db.add(new_note)
        db.commit()
        db.refresh(new_note)
        
        # Create flashcards if provided
        if flashcards and len(flashcards) > 0:
            flashcard_set = models.FlashcardSet(
                user_id=user.id,
                title=f"Flashcards: {title}",
                source_type="media",
                source_id=new_note.id
            )
            db.add(flashcard_set)
            db.commit()
            db.refresh(flashcard_set)
            
            for fc in flashcards:
                flashcard = models.Flashcard(
                    set_id=flashcard_set.id,
                    question=fc.get("question", ""),
                    answer=fc.get("answer", ""),
                    difficulty=fc.get("difficulty", "medium")
                )
                db.add(flashcard)
            
            db.commit()
        
        # Track activity
        from gamification_system import award_points
        award_points(db, user.id, "note_created")
        
        return {
            "success": True,
            "note_id": new_note.id,
            "message": "Notes saved successfully"
        }
        
    except Exception as e:
        logger.error(f"Error saving notes: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to save notes: {str(e)}")

@app.post("/api/media/regenerate-notes")
async def regenerate_notes(
    transcript: str = Body(...),
    analysis: dict = Body(...),
    note_style: str = Body("detailed"),
    difficulty: str = Body("intermediate"),
    subject: str = Body("general"),
    custom_instructions: str = Body(None)
):
    """Regenerate notes with different style/settings"""
    try:
        notes_result = await ai_media_processor.generate_notes_ai(
            transcript,
            analysis,
            note_style,
            {
                "difficulty": difficulty,
                "subject": subject,
                "custom_instructions": custom_instructions
            }
        )
        
        if not notes_result.get("success"):
            raise HTTPException(status_code=500, detail="Note regeneration failed")
        
        return {
            "success": True,
            "content": notes_result["content"],
            "style": note_style
        }
        
    except Exception as e:
        logger.error(f"Error regenerating notes: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to regenerate notes: {str(e)}")

@app.get("/api/media/estimate-cost")
async def estimate_processing_cost(
    duration_seconds: int = Query(...),
    file_size_mb: float = Query(...)
):
    """Estimate processing cost (all free)"""
    cost_info = ai_media_processor.estimate_processing_cost(duration_seconds, file_size_mb)
    return cost_info

@app.post("/api/media/generate-title")
async def generate_smart_title(
    transcript: str = Body(...),
    key_concepts: list = Body([]),
    summary: str = Body("")
):
    """Generate AI-based 3-4 word title"""
    try:
        prompt = f"""Generate a concise 3-4 word title for these notes.

Summary: {summary[:200]}
Key Concepts: {', '.join(key_concepts[:5])}
Content: {transcript[:500]}

Return ONLY the title, nothing else. Make it descriptive and catchy."""

        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "Generate 3-4 word title only."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=20
        )
        
        title = response.choices[0].message.content.strip().strip('"').strip("'")
        # Ensure it's not too long
        words = title.split()
        if len(words) > 4:
            title = ' '.join(words[:4])
        
        return {"title": title}
        
    except Exception as e:
        logger.error(f"Title generation error: {str(e)}")
        return {"title": "Media Notes"}

@app.get("/api/get_note/{note_id}")
async def get_single_note(
    note_id: int,
    db: Session = Depends(get_db)
):
    """Get a single note by ID"""
    try:
        import json
        
        note = db.query(models.Note).filter(
            models.Note.id == note_id,
            models.Note.is_deleted == False
        ).first()
        
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        # Parse JSON fields
        analysis = {}
        flashcards = []
        quiz_questions = []
        key_moments = []
        
        try:
            if note.analysis:
                analysis = json.loads(note.analysis)
        except:
            pass
        
        try:
            if note.flashcards:
                flashcards = json.loads(note.flashcards)
        except:
            pass
        
        try:
            if note.quiz_questions:
                quiz_questions = json.loads(note.quiz_questions)
        except:
            pass
        
        try:
            if note.key_moments:
                key_moments = json.loads(note.key_moments)
        except:
            pass
        
        return {
            "id": note.id,
            "title": note.title,
            "content": note.content,
            "created_at": note.created_at.isoformat(),
            "updated_at": note.updated_at.isoformat(),
            "folder_id": note.folder_id,
            "is_favorite": note.is_favorite,
            "custom_font": note.custom_font,
            "transcript": note.transcript or "",
            "analysis": analysis,
            "flashcards": flashcards,
            "quiz_questions": quiz_questions,
            "key_moments": key_moments
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching note {note_id}: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/media/history")
async def get_media_history(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's media processing history - only media-generated notes"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get all recent notes
        all_notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            models.Note.is_deleted == False
        ).order_by(models.Note.created_at.desc()).limit(50).all()
        
        history = []
        for note in all_notes:
            # Check if note has the media marker OR has flashcards associated with it
            is_media_note = note.custom_font == "__MEDIA_NOTE__"
            
            # Also check if there are flashcards with this note as source
            if not is_media_note:
                flashcard_set = db.query(models.FlashcardSet).filter(
                    models.FlashcardSet.source_type == "media",
                    models.FlashcardSet.source_id == note.id
                ).first()
                if flashcard_set:
                    is_media_note = True
            
            if is_media_note:
                history.append({
                    "id": note.id,
                    "title": note.title,
                    "created_at": note.created_at.isoformat(),
                    "preview": note.content[:200] if note.content else ""
                })
        
        logger.info(f"Found {len(history)} media-generated notes for user {user_id}")
        return {"history": history}
        
    except Exception as e:
        logger.error(f"History fetch error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/create_note")
def create_note(note_data: NoteCreate, db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, note_data.user_id) or get_user_by_email(db, note_data.user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        new_note = models.Note(
            user_id=user.id,
            title=note_data.title,
            content=note_data.content
        )
        db.add(new_note)
        db.commit()
        db.refresh(new_note)
        
        # Check for note creation milestones
        total_notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            models.Note.is_deleted == False
        ).count()
        
        # Create milestone notifications
        notification = None
        if total_notes == 1:
            notification = models.Notification(
                user_id=user.id,
                title="First Note Created!",
                message="Great start! You've created your first note. Keep organizing your learning!",
                notification_type="milestone"
            )
            db.add(notification)
            db.commit()
        elif total_notes == 10:
            notification = models.Notification(
                user_id=user.id,
                title="10 Notes Milestone!",
                message="Impressive! You've created 10 notes. Your knowledge base is growing!",
                notification_type="milestone"
            )
            db.add(notification)
            db.commit()
        elif total_notes == 50:
            notification = models.Notification(
                user_id=user.id,
                title="50 Notes Achievement!",
                message="Amazing! You've created 50 notes. You're building a comprehensive knowledge library!",
                notification_type="milestone"
            )
            db.add(notification)
            db.commit()
        
        response = {
            "id": new_note.id,
            "title": new_note.title,
            "content": new_note.content,
            "user_id": user.id,
            "created_at": new_note.created_at.isoformat() + 'Z',
            "updated_at": new_note.updated_at.isoformat() + 'Z',
            "status": "success",
            "message": "Note created successfully"
        }
        
        # Include notification in response if created
        if notification:
            response["notification"] = {
                "title": notification.title,
                "message": notification.message,
                "notification_type": notification.notification_type
            }
        
        return response
        
    except Exception as e:
        logger.error(f"Error creating note: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to create note: {str(e)}")

@app.put("/api/update_note")
def update_note(note_data: NoteUpdate, db: Session = Depends(get_db)):
    """Update an existing note"""
    try:
        note = db.query(models.Note).filter(models.Note.id == note_data.note_id).first()
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        #  CRITICAL: Do not allow updating deleted notes
        if note.is_deleted:
            logger.warning(f" Attempted to update deleted note {note.id}")
            raise HTTPException(status_code=400, detail="Cannot update a deleted note")
        
        note.title = note_data.title
        note.content = note_data.content
        note.updated_at = datetime.now(timezone.utc)
        
        db.commit()
        db.refresh(note)
        
        logger.info(f" Note {note.id} updated successfully")
        
        return {
            "id": note.id,
            "title": note.title,
            "content": note.content,
            "updated_at": note.updated_at.isoformat() + 'Z',
            "is_deleted": note.is_deleted,
            "status": "success"
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error updating note: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


     
@app.delete("/api/delete_note/{note_id}")
def delete_note(note_id: int, db: Session = Depends(get_db)):
    note = db.query(models.Note).filter(models.Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    
    db.delete(note)
    db.commit()
    
    return {"message": "Note deleted successfully"}

@app.post("/api/generate_note_summary/")
async def generate_note_summary(
    user_id: str = Form(...),
    conversation_data: str = Form(...),
    session_titles: str = Form(...),
    import_mode: str = Form("summary"),
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        user_profile = build_user_profile_dict(user)
        
        if import_mode == "summary":
            prompt = f"""Create study notes from this conversation. Format with headers, bullet points, and key concepts.

Conversation: {conversation_data[:3000]}

Provide a clear title and well-formatted study notes."""
        elif import_mode == "exam_prep":
            prompt = f"""Create an exam preparation guide from this conversation. Include:
1. Executive Summary
2. Key Concepts
3. Study Strategy
4. Review Checklist

Conversation: {conversation_data[:3000]}

Format as a comprehensive exam prep guide with a clear title."""
        else:
            prompt = f"""Create a formatted transcript from this conversation.

Conversation: {conversation_data[:3000]}

Provide a title and formatted content."""
        
        response = await generate_ai_response(prompt, user_profile)
        
        lines = response.split('\n')
        title = lines[0].replace('#', '').strip() if lines else "Study Notes"
        content = '\n'.join(lines[1:]) if len(lines) > 1 else response
        
        return {
            "title": title[:100],
            "content": content
        }
        
    except Exception as e:
        logger.error(f"Error generating note summary: {str(e)}")
        return {
            "title": "Study Notes",
            "content": f"# Study Notes\n\n{conversation_data[:1000]}"
        }
# ==================== NEW MODELS FOR FEATURES ====================



# ==================== FOLDERS ENDPOINTS ====================

@app.post("/api/create_folder")
def create_folder(folder_data: FolderCreate, db: Session = Depends(get_db)):
    """Create a new folder"""
    try:
        user = get_user_by_username(db, folder_data.user_id) or get_user_by_email(db, folder_data.user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        new_folder = models.Folder(
            user_id=user.id,
            name=folder_data.name,
            color=folder_data.color,
            parent_id=folder_data.parent_id
        )
        db.add(new_folder)
        db.commit()
        db.refresh(new_folder)
        
        return {
            "id": new_folder.id,
            "name": new_folder.name,
            "color": new_folder.color,
            "parent_id": new_folder.parent_id,
            "created_at": new_folder.created_at.isoformat() + 'Z',
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error creating folder: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/delete_folder/{folder_id}")
def delete_folder(folder_id: int, db: Session = Depends(get_db)):
    """Delete a folder (moves notes to root)"""
    try:
        folder = db.query(models.Folder).filter(
            models.Folder.id == folder_id
        ).first()
        
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found")
        
        # Move all notes in this folder to root (no folder)
        db.query(models.Note).filter(
            models.Note.folder_id == folder_id
        ).update({"folder_id": None})
        
        db.delete(folder)
        db.commit()
        
        return {"message": "Folder deleted successfully", "status": "success"}
    except Exception as e:
        logger.error(f"Error deleting folder: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.put("/api/move_note_to_folder")
def move_note_to_folder(data: NoteUpdateFolder, db: Session = Depends(get_db)):
    """Move note to a folder or remove from folder"""
    try:
        note = db.query(models.Note).filter(
            models.Note.id == data.note_id
        ).first()
        
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        note.folder_id = data.folder_id
        db.commit()
        
        return {
            "message": "Note moved successfully",
            "note_id": note.id,
            "folder_id": data.folder_id,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error moving note: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== FAVORITES ENDPOINT ====================

@app.put("/api/toggle_favorite")
def toggle_favorite(data: NoteFavorite, db: Session = Depends(get_db)):
    """Toggle favorite status of a note"""
    try:
        note = db.query(models.Note).filter(
            models.Note.id == data.note_id
        ).first()
        
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        note.is_favorite = data.is_favorite
        note.updated_at = datetime.now(timezone.utc)
        db.commit()
        
        return {
            "message": "Favorite status updated",
            "note_id": note.id,
            "is_favorite": note.is_favorite,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error toggling favorite: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/get_favorite_notes")
def get_favorite_notes(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all favorite notes (excluding deleted)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            models.Note.is_favorite == True,
            models.Note.is_deleted == False  #  Exclude deleted notes
        ).order_by(models.Note.updated_at.desc()).all()
        
        return {
            "favorites": [
                {
                    "id": note.id,
                    "title": note.title,
                    "content": note.content,
                    "created_at": note.created_at.isoformat() + 'Z',
                    "updated_at": note.updated_at.isoformat() + 'Z'
                }
                for note in notes
            ]
        }
    except Exception as e:
        logger.error(f"Error getting favorites: {str(e)}")
        return {"favorites": []}
        
        
                        # ==================== TRASH/RECYCLE BIN ====================


@app.put("/api/restore_note/{note_id}")
def restore_note(note_id: int, db: Session = Depends(get_db)):
    """Restore note from trash"""
    try:
        note = db.query(models.Note).filter(
            models.Note.id == note_id
        ).first()
        
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        note.is_deleted = False
        note.deleted_at = None
        note.updated_at = datetime.now(timezone.utc)
        db.commit()
        
        return {
            "message": "Note restored successfully",
            "note_id": note.id,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error restoring note: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/permanent_delete_note/{note_id}")
def permanent_delete_note(note_id: int, db: Session = Depends(get_db)):
    """Permanently delete a note"""
    try:
        note = db.query(models.Note).filter(
            models.Note.id == note_id
        ).first()
        
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        db.delete(note)
        db.commit()
        
        return {
            "message": "Note permanently deleted",
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error permanently deleting note: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== CHAT FOLDERS ENDPOINTS ====================

@app.post("/api/create_chat_folder")
def create_chat_folder(folder_data: ChatFolderCreate, db: Session = Depends(get_db)):
    """Create a new chat folder"""
    try:
        user = get_user_by_username(db, folder_data.user_id) or get_user_by_email(db, folder_data.user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        new_folder = models.ChatFolder(
            user_id=user.id,
            name=folder_data.name,
            color=folder_data.color,
            parent_id=folder_data.parent_id
        )
        db.add(new_folder)
        db.commit()
        db.refresh(new_folder)
        
        logger.info(f" Chat folder created: {new_folder.name} for user {user.email}")
        
        return {
            "id": new_folder.id,
            "name": new_folder.name,
            "color": new_folder.color,
            "parent_id": new_folder.parent_id,
            "created_at": new_folder.created_at.isoformat() + 'Z',
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error creating chat folder: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_chat_folders")
def get_chat_folders(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all chat folders for a user with chat counts"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        folders = db.query(models.ChatFolder).filter(
            models.ChatFolder.user_id == user.id
        ).order_by(models.ChatFolder.name.asc()).all()
        
        result = []
        for folder in folders:
            chat_count = db.query(models.ChatSession).filter(
                models.ChatSession.folder_id == folder.id
            ).count()
            
            result.append({
                "id": folder.id,
                "name": folder.name,
                "color": folder.color,
                "parent_id": folder.parent_id,
                "chat_count": chat_count,
                "created_at": folder.created_at.isoformat() + 'Z',
                "updated_at": folder.updated_at.isoformat() + 'Z'
            })
        
        return result
    except Exception as e:
        logger.error(f"Error getting chat folders: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/delete_chat_folder/{folder_id}")
def delete_chat_folder(folder_id: int, db: Session = Depends(get_db)):
    """Delete a chat folder (moves chats to root)"""
    try:
        folder = db.query(models.ChatFolder).filter(
            models.ChatFolder.id == folder_id
        ).first()
        
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found")
        
        # Move all chats in this folder to root (no folder)
        db.query(models.ChatSession).filter(
            models.ChatSession.folder_id == folder_id
        ).update({"folder_id": None})
        
        db.delete(folder)
        db.commit()
        
        return {"message": "Chat folder deleted successfully", "status": "success"}
    except Exception as e:
        logger.error(f"Error deleting chat folder: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/move_chat_to_folder")
def move_chat_to_folder(data: ChatUpdateFolder, db: Session = Depends(get_db)):
    """Move chat to a folder or remove from folder"""
    try:
        chat = db.query(models.ChatSession).filter(
            models.ChatSession.id == data.chat_id
        ).first()
        
        if not chat:
            raise HTTPException(status_code=404, detail="Chat not found")
        
        chat.folder_id = data.folder_id
        chat.updated_at = datetime.now(timezone.utc)
        db.commit()
        
        return {
            "message": "Chat moved successfully",
            "chat_id": chat.id,
            "folder_id": data.folder_id,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error moving chat: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate_chat_title")
async def generate_chat_title(request: GenerateChatTitleRequest, db: Session = Depends(get_db)):
    """AI-generated chat title based on conversation content"""
    try:
        user = get_user_by_username(db, request.user_id) or get_user_by_email(db, request.user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        chat = db.query(models.ChatSession).filter(
            models.ChatSession.id == request.chat_id,
            models.ChatSession.user_id == user.id
        ).first()
        
        if not chat:
            raise HTTPException(status_code=404, detail="Chat not found")
        
        # Get first few messages to understand the topic
        messages = db.query(models.ChatMessage).filter(
            models.ChatMessage.chat_session_id == request.chat_id
        ).order_by(models.ChatMessage.timestamp.asc()).limit(3).all()
        
        if not messages:
            return {"title": "New Chat", "status": "success"}
        
        # Build context from messages
        conversation_context = "\n".join([
            f"User: {msg.user_message}\nAI: {msg.ai_response[:200]}"
            for msg in messages
        ])
        
        # Generate title using AI
        prompt = f"""Based on this conversation, generate a concise, descriptive title (max 6 words).
The title should capture the main topic or question being discussed.

Conversation:
{conversation_context}

Generate only the title, nothing else. Make it specific and informative."""

        try:
            response = groq_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "Generate concise conversation titles."},
                    {"role": "user", "content": prompt}
                ],
                model=GROQ_MODEL,
                temperature=0.7,
                max_tokens=50,
            )
            
            generated_title = response.choices[0].message.content.strip()
            # Clean up the title
            generated_title = generated_title.replace('"', '').replace("'", "")
            generated_title = generated_title[:60]  # Limit length
            
            # Update chat title
            chat.title = generated_title
            chat.updated_at = datetime.now(timezone.utc)
            db.commit()
            
            logger.info(f" Generated title for chat {chat.id}: {generated_title}")
            
            return {
                "title": generated_title,
                "chat_id": chat.id,
                "status": "success"
            }
            
        except Exception as ai_error:
            logger.error(f"AI title generation error: {str(ai_error)}")
            # Fallback to simple title
            fallback_title = messages[0].user_message[:50] + "..."
            chat.title = fallback_title
            db.commit()
            
            return {
                "title": fallback_title,
                "chat_id": chat.id,
                "status": "success"
            }
        
    except Exception as e:
        logger.error(f"Error generating chat title: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== AI WRITING ASSISTANT ====================

@app.post("/api/ai_writing_assistant/")
async def ai_writing_assistant(
    request: AIWritingAssistRequest,
    db: Session = Depends(get_db)
):
    """AI Writing Assistant - multiple actions"""
    try:
        user = get_user_by_username(db, request.user_id) or get_user_by_email(db, request.user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        user_profile = build_user_profile_dict(user)
        
        # Build prompt based on action
        prompts = {
            "continue": f"Continue writing this text naturally:\n\n{request.content}\n\nContinue from where it left off:",
            "improve": f"Improve and enhance this text while keeping the same meaning:\n\n{request.content}\n\nImproved version:",
            "simplify": f"Simplify this text to make it easier to understand:\n\n{request.content}\n\nSimplified version:",
            "expand": f"Expand this text with more details and examples:\n\n{request.content}\n\nExpanded version:",
            "tone_change": f"Rewrite this text in a {request.tone} tone:\n\n{request.content}\n\nRewritten version:",
            "grammar": f"Fix grammar and spelling errors in this text:\n\n{request.content}\n\nCorrected version:",
            "summarize": f"Summarize this text concisely:\n\n{request.content}\n\nSummary:"
        }
        
        prompt = prompts.get(request.action, prompts["improve"])
        
        # Call Groq AI
        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": f"Writing assistant for {user_profile.get('first_name', 'user')}."},
                {"role": "user", "content": prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=2048,
            top_p=0.9,
        )
        
        result = chat_completion.choices[0].message.content
        
        return {
            "status": "success",
            "action": request.action,
            "result": result,
            "model_used": GROQ_MODEL
        }
        
    except Exception as e:
        logger.error(f"AI writing assistant error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"AI assistant failed: {str(e)}")


# ==================== UPDATE GET_NOTES TO INCLUDE NEW FIELDS ====================

# ==================== PROACTIVE AI SYSTEM ====================

# Global lock to prevent duplicate notifications
_notification_lock = {}

@app.get("/api/check_proactive_message")
async def check_proactive_message(
    user_id: str = Query(...),
    is_idle: bool = Query(False),
    is_login: bool = Query(False),
    db: Session = Depends(get_db)
):
    """
    Check if AI should proactively reach out to user
    Uses ML to analyze learning patterns and determine optimal intervention
    Supports idle detection, login greetings, and personalized weak topic recommendations
    """
    try:
        # Prevent duplicate calls within 5 seconds
        import time
        current_time = time.time()
        lock_key = f"{user_id}_{is_login}_{is_idle}"
        
        if lock_key in _notification_lock:
            time_since = current_time - _notification_lock[lock_key]
            if time_since < 5:
                print(f"\n🔔 BLOCKED DUPLICATE CALL (called {time_since:.1f}s ago)\n")
                return {"should_notify": False, "message": None}
        
        _notification_lock[lock_key] = current_time
        
        print(f"\n{'='*80}")
        print(f"🔔 PROACTIVE CHECK ENDPOINT CALLED")
        print(f"🔔 user_id={user_id}, is_idle={is_idle}, is_login={is_login}")
        print(f"{'='*80}\n")
        logger.info(f"🔔 Proactive check: user_id={user_id}, is_idle={is_idle}, is_login={is_login}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        logger.info(f"🔔 User found: {user.first_name} {user.last_name}")
        
        # Get proactive AI engine
        proactive_engine = get_proactive_ai_engine(unified_ai)
        
        # Get comprehensive user profile for personalization
        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()
        
        user_profile = {
            "first_name": user.first_name or "there",
            "field_of_study": user.field_of_study or "General Studies",
            "learning_style": user.learning_style or "Mixed",
            "is_login": is_login
        }
        
        # Add comprehensive profile data if available
        if comprehensive_profile:
            user_profile.update({
                "difficulty_level": getattr(comprehensive_profile, "difficulty_level", "intermediate") or "intermediate",
                "study_goals": getattr(comprehensive_profile, "study_goals", None),
                "career_goals": getattr(comprehensive_profile, "career_goals", None)
            })
        
        # Check if we should send a proactive message (ML-based decision)
        logger.info(f"🔔 Calling ML engine with is_login={is_login}...")
        result = await proactive_engine.check_and_send_proactive_message(
            db, user.id, user_profile, is_idle, is_login
        )
        
        logger.info(f"🔔 ML engine returned: {result}")
        
        # If ML system determined we should reach out
        if result:
            # FIRST: Delete old empty check-in chats (no user replies)
            print(f"\n Checking for empty check-in chats to delete...")
            old_checkin_chats = db.query(models.ChatSession).filter(
                models.ChatSession.user_id == user.id,
                models.ChatSession.title == "AI Tutor Check-in"
            ).all()
            
            print(f" Found {len(old_checkin_chats)} check-in chats to check")
            
            deleted_count = 0
            for old_chat in old_checkin_chats:
                # Count messages where user actually typed something
                user_messages = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == old_chat.id,
                    models.ChatMessage.user_message != "",
                    models.ChatMessage.user_message.isnot(None)
                ).count()
                
                print(f" Chat ID {old_chat.id}: {user_messages} user messages")
                
                if user_messages == 0:
                    # No user replies - delete this chat
                    db.query(models.ChatMessage).filter(
                        models.ChatMessage.chat_session_id == old_chat.id
                    ).delete()
                    db.delete(old_chat)
                    deleted_count += 1
                    print(f"  Deleted empty check-in chat ID: {old_chat.id}")
                else:
                    print(f" Keeping check-in chat ID: {old_chat.id} (user replied)")
            
            if deleted_count > 0:
                db.commit()
                print(f" Deleted {deleted_count} empty check-in chats\n")
            else:
                print(f" No empty chats to delete\n")
            
            # NOW: Create new chat session for this proactive message
            new_session = models.ChatSession(
                user_id=user.id,
                title="AI Tutor Check-in",
                created_at=datetime.now(timezone.utc),
                updated_at=datetime.now(timezone.utc)
            )
            db.add(new_session)
            db.commit()
            db.refresh(new_session)
            
            # Save the ML-generated message
            new_message = models.ChatMessage(
                chat_session_id=new_session.id,
                user_id=user.id,
                user_message="",
                ai_response=result["message"],
                timestamp=datetime.now(timezone.utc)
            )
            db.add(new_message)
            
            # Check if we already created a notification in the last 30 seconds (prevent duplicates)
            recent_notif = db.query(models.Notification).filter(
                models.Notification.user_id == user.id,
                models.Notification.notification_type == "proactive_ai",
                models.Notification.created_at >= datetime.now(timezone.utc) - timedelta(seconds=30)
            ).first()
            
            if not recent_notif:
                # Delete old login notifications (keep only the latest one)
                old_login_notifs = db.query(models.Notification).filter(
                    models.Notification.user_id == user.id,
                    models.Notification.notification_type == "proactive_ai"
                ).all()
                
                for old_notif in old_login_notifs:
                    db.delete(old_notif)
                
                # Create new notification
                notification = models.Notification(
                    user_id=user.id,
                    title="Cerbyl AI",
                    message=result["message"],
                    notification_type="proactive_ai",
                    is_read=False,
                    created_at=datetime.now(timezone.utc)
                )
                db.add(notification)
                db.commit()
                db.refresh(notification)
                
                print(f"\n🔔 NOTIFICATION CREATED: ID={notification.id}, user_id={user.id}\n")
            else:
                print(f"\n🔔 SKIPPING DUPLICATE NOTIFICATION (recent one exists)\n")
                notification = recent_notif
            
            response_data = {
                "should_notify": True,
                "message": result["message"],
                "chat_id": new_session.id,
                "urgency_score": result["urgency_score"],
                "reason": result["reason"],
                "notification_id": notification.id
            }
            print(f"\n🔔 RETURNING SUCCESS: {response_data}\n")
            logger.info(f"ML-based proactive message sent to user {user.id}: {result['reason']}")
            return response_data
        
        # No intervention needed based on ML analysis
        print(f"\n🔔 RETURNING NO NOTIFICATION (result was None)\n")
        logger.info(f"ML system determined no intervention needed for user {user.id}")
        return {
            "should_notify": False,
            "message": None
        }
            
    except Exception as e:
        logger.error(f"Error checking proactive message: {str(e)}")
        import traceback
        traceback.print_exc()
        return {
            "should_notify": False,
            "message": None,
            "error": str(e)
        }
        return {
            "should_notify": False,
            "message": None,
            "error": str(e)
        }

# ==================== END PROACTIVE AI SYSTEM ====================

@app.post("/api/generate_welcome_message")
async def generate_welcome_message(payload: dict = Body(...), db: Session = Depends(get_db)):
    """Generate personalized welcome message using Gemini AI"""
    try:
        user_id = payload.get("user_id")
        first_name = payload.get("first_name", "there")
        field_of_study = payload.get("field_of_study", "learning")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get recent learning context
        recent_activities = db.query(models.Activity).filter(
            models.Activity.user_id == user.id
        ).order_by(models.Activity.timestamp.desc()).limit(10).all()
        
        context = ""
        if recent_activities:
            topics = list(set([a.topic for a in recent_activities if a.topic]))[:3]
            if topics:
                context = f"\nThey recently studied: {', '.join(topics)}"
            
            correct = sum(1 for a in recent_activities if a.correct)
            accuracy = (correct / len(recent_activities)) * 100 if recent_activities else 0
            context += f"\nRecent accuracy: {accuracy:.0f}%"
        
        prompt = f"""Welcome {first_name} studying {field_of_study}.{context}
Warm greeting (2 sentences max). Ask what they'd like to work on. Use emoji if appropriate."""
        
        message = call_ai(prompt, max_tokens=150, temperature=0.8)
        
        logger.info(f"Generated welcome message for {first_name}")
        return {"message": message}
        
    except Exception as e:
        logger.error(f"Error generating welcome message: {e}")
        first_name = payload.get("first_name", "there")
        return {"message": f"Hey {first_name}! 👋 Welcome back! What would you like to learn today?"}

@app.post("/api/generate_chat_summary")
async def generate_chat_summary(
    chat_data: str = Form(...),
    max_words: int = Form(4),
    format: str = Form("title"),
    db: Session = Depends(get_db)
):
    try:
        user_profile = {"first_name": "Student", "field_of_study": "General Studies"}
        
        prompt = f"""Create a precise 4-word title from this content:

{chat_data[:2000]}

Rules:
- EXACTLY 4 words
- Focus on main topic
- Use academic/technical terms
- Capitalize each word
- No punctuation

Generate only the 4-word title:"""

        response = await generate_ai_response(prompt, user_profile)
        
        title = response.strip().replace('"', '').replace("'", "")
        words = title.split()
        final_title = ' '.join(words[:4]) if len(words) >= 4 else title
        
        return {
            "summary": final_title,
            "word_count": len(final_title.split()),
            "status": "success"
        }
        
    except Exception as e:
        logger.error(f"Error generating summary: {str(e)}")
        return {
            "summary": "Study Session Cards",
            "word_count": 3,
            "status": "fallback"
        }





@app.get("/api/get_enhanced_user_stats")
def get_enhanced_user_stats(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        from gamification_system import get_user_stats as get_gamification_stats
        gamification_stats = get_gamification_stats(db, user.id)
        
        total_questions = gamification_stats.get("total_questions_answered", 0)
        total_notes = gamification_stats.get("total_notes_created", 0)
        total_flashcards = gamification_stats.get("total_flashcards_created", 0)
        total_quizzes = gamification_stats.get("total_quizzes_completed", 0)
        total_study_minutes = gamification_stats.get("total_study_minutes", 0)
        
        # Count actual chat sessions from database (only sessions with messages)
        total_chat_sessions = db.query(func.count(func.distinct(models.ChatSession.id))).join(
            models.ChatMessage, models.ChatMessage.chat_session_id == models.ChatSession.id
        ).filter(
            models.ChatSession.user_id == user.id
        ).scalar() or 0
        
        streak = gamification_stats.get("current_streak", 0)
        
        user_stats = db.query(models.UserStats).filter(
            models.UserStats.user_id == user.id
        ).first()
        
        if user_stats:
            user_stats.day_streak = streak
            user_stats.total_hours = total_study_minutes / 60
            db.commit()
        
        return {
            "streak": streak,
            "lessons": total_quizzes,
            "hours": round(total_study_minutes / 60, 1),
            "minutes": total_study_minutes,
            "accuracy": user_stats.accuracy_percentage if user_stats else 0,
            "totalQuestions": total_questions,
            "totalFlashcards": total_flashcards,
            "totalNotes": total_notes,
            "totalChatSessions": total_chat_sessions,
            "total_time_today": 0
        }
        
    except Exception as e:
        logger.error(f"Error getting stats: {str(e)}")
        return {
            "streak": 0, "lessons": 0, "hours": 0, "minutes": 0, "accuracy": 0,
            "totalQuestions": 0, "totalFlashcards": 0, "totalNotes": 0, "totalChatSessions": 0
        }

@app.get("/api/get_activity_heatmap")
def get_activity_heatmap(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        end_date = datetime.now(timezone.utc).date()
        start_date = end_date - timedelta(days=365)
        
        transactions = db.query(models.PointTransaction).filter(
            models.PointTransaction.user_id == user.id,
            models.PointTransaction.created_at >= datetime.combine(start_date, datetime.min.time()).replace(tzinfo=timezone.utc)
        ).all()
        
        logger.info(f"Found {len(transactions)} point transactions for user {user.id}")
        
        activity_dict = {}
        for tx in transactions:
            date_str = tx.created_at.date().isoformat()
            activity_dict[date_str] = activity_dict.get(date_str, 0) + 1
        
        current_date = start_date
        heatmap_data = []
        
        while current_date <= end_date:
            date_str = current_date.isoformat()
            count = activity_dict.get(date_str, 0)
            
            if count == 0:
                level = 0
            elif count == 1:
                level = 1
            elif count <= 3:
                level = 2
            elif count <= 5:
                level = 3
            elif count <= 8:
                level = 4
            else:
                level = 5
            
            heatmap_data.append({
                "date": date_str,
                "count": count,
                "level": level
            })
            current_date += timedelta(days=1)
        
        total_count = sum(activity_dict.values())
        
        return {
            "heatmap_data": heatmap_data,
            "total_count": total_count,
            "date_range": {
                "start": start_date.isoformat() + 'Z',
                "end": end_date.isoformat() + 'Z'
            }
        }
    except Exception as e:
        logger.error(f"Error getting heatmap: {str(e)}")
        return {"heatmap_data": [], "total_count": 0, "date_range": {"start": "", "end": ""}}

@app.get("/api/get_recent_activities")
def get_recent_activities(user_id: str = Query(...), limit: int = Query(5), db: Session = Depends(get_db)):
    return []

@app.get("/api/get_weekly_progress")
def get_weekly_progress(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get last 7 days
        end_date = datetime.now(timezone.utc).date()
        start_date = end_date - timedelta(days=6)  # 6 days ago + today = 7 days
        
        # Initialize daily data structure
        daily_data = {}
        for i in range(7):
            current_date = start_date + timedelta(days=i)
            daily_data[current_date] = {
                "date": current_date.isoformat(),
                "day": ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun'][current_date.weekday()],
                "points": 0,
                "ai_chats": 0,
                "notes": 0,
                "flashcards": 0,
                "quizzes": 0,
                "solo_quizzes": 0,
                "battles": 0,
                "study_minutes": 0
            }
        
        # Get point transactions for the week (most accurate source)
        transactions = db.query(models.PointTransaction).filter(
            models.PointTransaction.user_id == user.id,
            models.PointTransaction.created_at >= datetime.combine(start_date, datetime.min.time()).replace(tzinfo=timezone.utc)
        ).all()
        
        for t in transactions:
            date_key = t.created_at.date()
            if date_key in daily_data:
                daily_data[date_key]["points"] += t.points_earned
                
                # Categorize by activity type
                if t.activity_type == "ai_chat":
                    daily_data[date_key]["ai_chats"] += 1
                elif t.activity_type == "note_created":
                    daily_data[date_key]["notes"] += 1
                elif t.activity_type in ["flashcard_set", "flashcard_reviewed", "flashcard_mastered"]:
                    daily_data[date_key]["flashcards"] += 1
                elif t.activity_type == "quiz_completed":
                    daily_data[date_key]["quizzes"] += 1
                elif t.activity_type == "solo_quiz":
                    daily_data[date_key]["solo_quizzes"] += 1
                elif t.activity_type in ["battle_win", "battle_draw", "battle_loss"]:
                    daily_data[date_key]["battles"] += 1
                elif t.activity_type == "study_time":
                    # Extract minutes from metadata if available
                    try:
                        if t.activity_metadata:
                            import ast
                            meta = ast.literal_eval(t.activity_metadata)
                            daily_data[date_key]["study_minutes"] += meta.get("minutes", 0)
                    except:
                        pass
        
        # Convert to list format for frontend
        weekly_data = []
        daily_breakdown = []
        total_points = 0
        
        for i in range(7):
            current_date = start_date + timedelta(days=i)
            day_data = daily_data[current_date]
            weekly_data.append(day_data["points"])
            daily_breakdown.append(day_data)
            total_points += day_data["points"]
        
        average_per_day = total_points / 7 if total_points > 0 else 0
        
        # Get gamification stats for totals
        stats = db.query(models.UserGamificationStats).filter(
            models.UserGamificationStats.user_id == user.id
        ).first()
        
        return {
            "weekly_data": weekly_data,
            "daily_breakdown": daily_breakdown,
            "total_points": total_points,
            "average_per_day": round(average_per_day, 1),
            "start_date": start_date.isoformat(),
            "end_date": end_date.isoformat(),
            "weekly_stats": {
                "ai_chats": stats.weekly_ai_chats if stats else 0,
                "notes_created": stats.weekly_notes_created if stats else 0,
                "flashcards_created": stats.weekly_flashcards_created if stats else 0,
                "quizzes_completed": stats.weekly_quizzes_completed if stats else 0,
                "solo_quizzes": getattr(stats, 'weekly_solo_quizzes', 0) if stats else 0,
                "battles_won": stats.weekly_battles_won if stats else 0,
                "study_minutes": stats.weekly_study_minutes if stats else 0
            }
        }
        
    except Exception as e:
        logger.error(f"Error getting weekly progress: {str(e)}")
        return {
            "weekly_data": [0, 0, 0, 0, 0, 0, 0],
            "total_sessions": 0,
            "average_per_day": 0
        }


@app.get("/api/get_analytics_history")
def get_analytics_history(
    user_id: str = Query(...),
    period: str = Query("week"),
    db: Session = Depends(get_db)
):
    """Get historical analytics data for charts with full metric breakdown"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Determine date range and grouping based on period
        end_date = datetime.now(timezone.utc).date()
        if period == "week":
            start_date = end_date - timedelta(days=6)
            group_by = "day"
        elif period == "month":
            start_date = end_date - timedelta(days=29)
            group_by = "day"
        elif period == "year":
            start_date = end_date - timedelta(days=364)
            group_by = "week"  # Group by week for year view
        else:  # all
            # Get earliest transaction date
            earliest = db.query(func.min(models.PointTransaction.created_at)).filter(
                models.PointTransaction.user_id == user.id
            ).scalar()
            start_date = earliest.date() if earliest else end_date - timedelta(days=364)
            group_by = "month"  # Group by month for all-time view
        
        # Get point transactions for the period
        transactions = db.query(models.PointTransaction).filter(
            models.PointTransaction.user_id == user.id,
            models.PointTransaction.created_at >= datetime.combine(start_date, datetime.min.time()).replace(tzinfo=timezone.utc)
        ).order_by(models.PointTransaction.created_at.asc()).all()
        
        # Initialize data structure based on grouping
        data_points = {}
        
        if group_by == "day":
            # Daily breakdown
            current = start_date
            while current <= end_date:
                key = current.isoformat()
                day_names = ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun']
                data_points[key] = {
                    "date": key,
                    "day": day_names[current.weekday()],
                    "label": current.strftime("%b %d") if period == "month" else day_names[current.weekday()],
                    "points": 0,
                    "ai_chats": 0,
                    "notes": 0,
                    "flashcards": 0,
                    "quizzes": 0,
                    "solo_quizzes": 0,
                    "battles": 0,
                    "study_minutes": 0
                }
                current += timedelta(days=1)
        elif group_by == "week":
            # Weekly breakdown for year view
            current = start_date
            week_num = 0
            while current <= end_date:
                week_start = current
                week_end = min(current + timedelta(days=6), end_date)
                key = f"W{week_num}"
                data_points[key] = {
                    "date": week_start.isoformat(),
                    "day": key,
                    "label": week_start.strftime("%b %d"),
                    "points": 0,
                    "ai_chats": 0,
                    "notes": 0,
                    "flashcards": 0,
                    "quizzes": 0,
                    "solo_quizzes": 0,
                    "battles": 0,
                    "study_minutes": 0,
                    "_start": week_start,
                    "_end": week_end
                }
                current += timedelta(days=7)
                week_num += 1
        else:  # month
            # Monthly breakdown for all-time view
            current = start_date.replace(day=1)
            while current <= end_date:
                key = current.strftime("%Y-%m")
                data_points[key] = {
                    "date": current.isoformat(),
                    "day": current.strftime("%b"),
                    "label": current.strftime("%b %Y"),
                    "points": 0,
                    "ai_chats": 0,
                    "notes": 0,
                    "flashcards": 0,
                    "quizzes": 0,
                    "solo_quizzes": 0,
                    "battles": 0,
                    "study_minutes": 0,
                    "_month": current.month,
                    "_year": current.year
                }
                # Move to next month
                if current.month == 12:
                    current = current.replace(year=current.year + 1, month=1)
                else:
                    current = current.replace(month=current.month + 1)
        
        # Process transactions
        for t in transactions:
            t_date = t.created_at.date()
            
            # Find the right bucket
            if group_by == "day":
                key = t_date.isoformat()
            elif group_by == "week":
                # Find which week this belongs to
                key = None
                for k, v in data_points.items():
                    if "_start" in v and v["_start"] <= t_date <= v["_end"]:
                        key = k
                        break
                if not key:
                    continue
            else:  # month
                key = t_date.strftime("%Y-%m")
            
            if key not in data_points:
                continue
                
            data_points[key]["points"] += t.points_earned
            
            # Categorize by activity type
            if t.activity_type == "ai_chat":
                data_points[key]["ai_chats"] += 1
            elif t.activity_type == "note_created":
                data_points[key]["notes"] += 1
            elif t.activity_type in ["flashcard_set", "flashcard_reviewed", "flashcard_mastered"]:
                data_points[key]["flashcards"] += 1
            elif t.activity_type == "quiz_completed":
                data_points[key]["quizzes"] += 1
            elif t.activity_type == "solo_quiz":
                data_points[key]["solo_quizzes"] += 1
            elif t.activity_type in ["battle_win", "battle_draw", "battle_loss"]:
                data_points[key]["battles"] += 1
            elif t.activity_type == "study_time":
                try:
                    if t.activity_metadata:
                        import ast
                        meta = ast.literal_eval(t.activity_metadata)
                        data_points[key]["study_minutes"] += meta.get("minutes", 0)
                except:
                    pass
        
        # Convert to list and clean up internal fields
        history = []
        for key in sorted(data_points.keys()):
            item = {k: v for k, v in data_points[key].items() if not k.startswith("_")}
            history.append(item)
        
        # Calculate totals for the period
        total_points = sum(h["points"] for h in history)
        total_activities = sum(h["ai_chats"] + h["notes"] + h["flashcards"] + h["quizzes"] + h["solo_quizzes"] + h["battles"] for h in history)
        
        return {
            "history": history,
            "period": period,
            "group_by": group_by,
            "start_date": start_date.isoformat(),
            "end_date": end_date.isoformat(),
            "total_points": total_points,
            "total_activities": total_activities,
            "data_points_count": len(history)
        }
        
    except Exception as e:
        logger.error(f"Error getting analytics history: {str(e)}")
        import traceback
        traceback.print_exc()
        return {"history": [], "period": period, "error": str(e)}


# OLD endpoint removed - using newer comprehensive version at line ~9123

# OLD BROKEN ENDPOINT REMOVED - Using new centralized system at line ~8699

# OLD BROKEN ENDPOINT REMOVED - Using new one at line ~8814

@app.get("/api/get_global_leaderboard")
def get_global_leaderboard(limit: int = Query(10), db: Session = Depends(get_db)):
    try:
        top_users = db.query(models.UserGamificationStats).order_by(
            models.UserGamificationStats.total_points.desc()
        ).limit(limit).all()
        
        leaderboard = []
        for stats in top_users:
            user = db.query(models.User).filter(models.User.id == stats.user_id).first()
            if user:
                leaderboard.append({
                    "user_id": user.id,
                    "username": user.username,
                    "total_points": stats.total_points,
                    "level": stats.level
                })
        
        return {"leaderboard": leaderboard}
    except Exception as e:
        logger.error(f"Error getting leaderboard: {str(e)}")
        return {"leaderboard": []}

# OLD endpoint removed - using newer comprehensive version at line ~9198

@app.get("/api/get_learning_analytics")
def get_learning_analytics(user_id: str = Query(...), period: str = Query("week"), db: Session = Depends(get_db)):
    user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    end_date = datetime.now(timezone.utc).date()
    start_date = end_date - timedelta(days=7)
    
    daily_metrics = db.query(models.DailyLearningMetrics).filter(
        models.DailyLearningMetrics.user_id == user.id,
        models.DailyLearningMetrics.date >= start_date,
        models.DailyLearningMetrics.date <= end_date
    ).all()
    
    total_sessions = len([m for m in daily_metrics if m.questions_answered > 0 or m.time_spent_minutes > 0])
    total_time_minutes = sum(m.time_spent_minutes for m in daily_metrics)
    total_questions = sum(m.questions_answered for m in daily_metrics)
    
    return {
        "period": "week",
        "start_date": start_date.isoformat() + 'Z',
        "end_date": end_date.isoformat() + 'Z',
        "total_sessions": total_sessions,
        "total_time_minutes": total_time_minutes,
        "total_questions": total_questions,
        "accuracy_percentage": 100,
        "average_per_day": 0,
        "days_active": total_sessions,
        "daily_data": []
    }
@app.post("/api/firebase-auth")
async def firebase_authentication(request: Request, db: Session = Depends(get_db)):
    max_retries = 2
    for attempt in range(max_retries):
        try:
            data = await request.json()
            id_token = data.get('idToken')
            email = data.get('email')
            display_name = data.get('displayName')
            photo_url = data.get('photoURL')
            uid = data.get('uid')

            if not id_token or not email:
                raise HTTPException(status_code=400, detail="Missing required fields")

            user = get_user_by_email(db, email)
            is_new_user = False
            
            if not user:
                is_new_user = True
                names = display_name.split(' ') if display_name else ['', '']
                first_name = names[0] if len(names) > 0 else ''
                last_name = ' '.join(names[1:]) if len(names) > 1 else ''
                
                user = models.User(
                    first_name=first_name,
                    last_name=last_name,
                    email=email,
                    username=email,
                    hashed_password=get_password_hash(f"firebase_{uid}"),
                    picture_url=photo_url,
                    google_user=True
                )
                db.add(user)
                db.commit()
                db.refresh(user)
                
                # Create user stats
                user_stats = models.UserStats(user_id=user.id)
                db.add(user_stats)
                
                # Create gamification stats
                gamif_stats = models.UserGamificationStats(
                    user_id=user.id,
                    week_start_date=datetime.now(timezone.utc)
                )
                db.add(gamif_stats)
                
                db.commit()
                logger.info(f" New user created via Firebase: {email}")

            access_token = create_access_token(
                data={"sub": user.username, "user_id": user.id}
            )

            return {
                "access_token": access_token,
                "token_type": "bearer",
                "is_new_user": is_new_user,
                "user": {
                    "id": user.id,
                    "email": user.email,
                    "first_name": user.first_name,
                    "last_name": user.last_name,
                    "picture_url": user.picture_url,
                    "google_user": True
                }
            }

        except HTTPException:
            raise
        except Exception as e:
            db.rollback()
            error_msg = str(e).lower()
            
            # Check if it's a sequence error and retry
            if "duplicate key" in error_msg and "pkey" in error_msg and attempt < max_retries - 1:
                logger.warning(f" Sequence error on Firebase auth, fixing... (attempt {attempt + 1})")
                try:
                    # Fix the sequence
                    if "postgres" in DATABASE_URL:
                        max_id_query = text("SELECT COALESCE(MAX(id), 0) FROM users")
                        max_id = db.execute(max_id_query).scalar()
                        fix_query = text("SELECT setval('users_id_seq', :next_id)")
                        db.execute(fix_query, {"next_id": max_id + 1})
                        db.commit()
                        logger.info(f" Sequence fixed, retrying...")
                        continue  # Retry
                except Exception as fix_error:
                    logger.error(f" Failed to fix sequence: {str(fix_error)}")
            
            # If not a sequence error or retry failed
            logger.error(f"Firebase auth error: {str(e)}", exc_info=True)
            raise HTTPException(status_code=500, detail=f"Authentication failed: {str(e)}")
    
    # If all retries failed
    raise HTTPException(status_code=500, detail="Authentication failed after retries")
@app.get("/api/check_profile_quiz")
async def check_profile_quiz(user_id: str = Query(...), db: Session = Depends(get_db)):
    """
    Check if user should see the profile quiz.
    Logic: User is first-time ONLY if account was just created (within last 5 minutes) 
    AND they haven't completed or skipped the quiz yet.
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()

        # Check if user has completed the quiz flow
        has_completed_quiz = (
            comprehensive_profile is not None and 
            comprehensive_profile.primary_archetype is not None and 
            comprehensive_profile.primary_archetype != ""
        )
        
        has_skipped_quiz = (
            comprehensive_profile is not None and
            comprehensive_profile.quiz_skipped == True
        )
        
        # User has completed the onboarding if they did quiz OR skipped it
        quiz_flow_completed = has_completed_quiz or has_skipped_quiz
        
        logger.info(f" check_profile_quiz for {user_id}: completed={has_completed_quiz}, skipped={has_skipped_quiz}, flow_completed={quiz_flow_completed}")

        return {
            "completed": quiz_flow_completed,
            "quiz_completed": has_completed_quiz,
            "quiz_skipped": has_skipped_quiz,
            "user_id": user_id
        }

    except Exception as e:
        logger.error(f"Error checking quiz: {str(e)}")
        return {"completed": False}

@app.get("/api/is_first_time_user")
async def is_first_time_user(user_id: str = Query(...), db: Session = Depends(get_db)):
    """
    Check if user is truly first-time (just registered and first login).
    Returns true ONLY if this is their very first login session.
    We track this by checking if last_login is very close to created_at (within 2 minutes).
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        from datetime import datetime, timezone, timedelta
        now = datetime.now(timezone.utc)
        
        # Make timezone-aware if needed
        user_created = user.created_at
        if user_created.tzinfo is None:
            user_created = user_created.replace(tzinfo=timezone.utc)
            
        user_last_login = user.last_login
        if user_last_login and user_last_login.tzinfo is None:
            user_last_login = user_last_login.replace(tzinfo=timezone.utc)
        
        # Check if this is first login (last_login is very close to created_at)
        # This means they just registered and this is their first session
        if user_last_login:
            time_between_creation_and_login = abs((user_last_login - user_created).total_seconds())
            is_first_login = time_between_creation_and_login < 120  # Within 2 minutes
        else:
            is_first_login = True  # No last_login means first time
        
        # Also check quiz status for additional context
        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()

        has_completed_quiz = (
            comprehensive_profile is not None and 
            comprehensive_profile.primary_archetype is not None and 
            comprehensive_profile.primary_archetype != ""
        )
        
        has_skipped_quiz = (
            comprehensive_profile is not None and
            comprehensive_profile.quiz_skipped == True
        )
        
        quiz_flow_done = has_completed_quiz or has_skipped_quiz
        
        # First-time user = first login AND hasn't done quiz flow
        is_first_time = is_first_login and not quiz_flow_done
        
        time_since_creation = now - user_created

        logger.info(f" is_first_time_user for {user_id}: is_first_time={is_first_time}, is_first_login={is_first_login}, quiz_completed={has_completed_quiz}, quiz_skipped={has_skipped_quiz}, account_age_minutes={time_since_creation.total_seconds() / 60:.2f}")

        return {
            "is_first_time": is_first_time,
            "is_first_login": is_first_login,
            "account_age_minutes": time_since_creation.total_seconds() / 60,
            "quiz_completed": has_completed_quiz,
            "quiz_skipped": has_skipped_quiz,
            "user_id": user_id
        }

    except Exception as e:
        logger.error(f"Error checking first-time user: {str(e)}")
        return {"is_first_time": False}

@app.post("/api/start_session")
def start_session(
    user_id: str = Form(...),
    session_type: str = Form(...),
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        session_id = f"{user.id}_{session_type}_{int(datetime.now(timezone.utc).timestamp())}"
        
        return {
            "status": "success",
            "session_id": session_id,
            "start_time": datetime.now(timezone.utc).isoformat() + 'Z',
            "message": f"Started {session_type} session"
        }
        
    except Exception as e:
        logger.error(f"Error starting session: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to start session")

@app.post("/api/end_session")
def end_session(
    user_id: str = Form(...),
    session_id: str = Form(...),
    time_spent_minutes: float = Form(...),
    session_type: str = Form(...),
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        today = datetime.now(timezone.utc).date()
        daily_metric = db.query(models.DailyLearningMetrics).filter(
            and_(
                models.DailyLearningMetrics.user_id == user.id,
                models.DailyLearningMetrics.date == today
            )
        ).first()
        
        if not daily_metric:
            # New record - start at 1
            daily_metric = models.DailyLearningMetrics(
                user_id=user.id,
                date=today,
                sessions_completed=1,
                time_spent_minutes=time_spent_minutes,
                questions_answered=0,
                correct_answers=0,
                topics_studied="[]"
            )
            db.add(daily_metric)
        else:
            # Existing record - increment from current value (even if 0)
            daily_metric.time_spent_minutes += time_spent_minutes
            daily_metric.sessions_completed = (daily_metric.sessions_completed or 0) + 1
        
        user_stats = db.query(models.UserStats).filter(
            models.UserStats.user_id == user.id
        ).first()
        
        if not user_stats:
            user_stats = models.UserStats(user_id=user.id)
            db.add(user_stats)
        
        user_stats.total_hours += (time_spent_minutes / 60)
        user_stats.last_activity = datetime.now(timezone.utc)
        
        db.commit()
        
        logger.info(f" Session ended: user={user.email}, sessions_today={daily_metric.sessions_completed}, time={daily_metric.time_spent_minutes}")
        
        return {
            "status": "success",
            "message": f"Ended {session_type} session",
            "time_recorded": time_spent_minutes,
            "total_time_today": daily_metric.time_spent_minutes,
            "sessions_today": daily_metric.sessions_completed
        }
        
    except Exception as e:
        logger.error(f"Error ending session: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to end session")


@app.get("/api/conversation_starters")
def get_conversation_starters(user_id: str = Query(...), db: Session = Depends(get_db)):
    return {"suggestions": ["What would you like to learn today?"]}

@app.get("/api/get_learning_reviews")
def get_learning_reviews(user_id: str = Query(...), db: Session = Depends(get_db)):
    return {"reviews": []}



@app.post("/api/create_learning_review")
async def create_learning_review(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        user_id = payload.get("user_id")
        chat_session_ids = payload.get("chat_session_ids", [])
        slide_ids = payload.get("slide_ids", [])
        review_title = payload.get("review_title", "Learning Review")
        review_type = payload.get("review_type", "comprehensive")

        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        if not chat_session_ids and not slide_ids:
            raise HTTPException(status_code=400, detail="At least one source required")

        session_titles = []
        slide_filenames = []
        combined_text = ""

        if chat_session_ids:
            sessions = db.query(models.ChatSession).filter(
                models.ChatSession.id.in_(chat_session_ids),
                models.ChatSession.user_id == user.id
            ).all()

            if sessions:
                session_titles = [s.title for s in sessions]
                session_ids = [s.id for s in sessions]

                messages = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id.in_(session_ids)
                ).order_by(models.ChatMessage.timestamp.asc()).all()

                for msg in messages:
                    combined_text += f"Q: {msg.user_message}\nA: {msg.ai_response}\n\n"

        if slide_ids:
            slides = db.query(models.UploadedSlide).filter(
                models.UploadedSlide.id.in_(slide_ids),
                models.UploadedSlide.user_id == user.id
            ).all()

            if slides:
                slide_filenames = [s.original_filename for s in slides]
                
                for slide in slides:
                    if slide.extracted_text:
                        combined_text += f"\n\nSlide: {slide.original_filename}\n{slide.extracted_text}\n"

        combined_text = combined_text[:8000]

        prompt = f"""
        You are an intelligent learning assistant.
        Extract 8-12 key learning points from the following content.
        Return them strictly as a JSON list: ["Point 1", "Point 2", ...].

        Content:
        {combined_text}
        """

        user_profile = build_user_profile_dict(user)
        ai_response = await generate_ai_response(prompt, user_profile)
        
        try:
            json_match = re.search(r"\[.*\]", ai_response, re.DOTALL)
            if json_match:
                learning_points = json.loads(json_match.group())
            else:
                learning_points = ["Key concepts covered", "Important principles", "Main ideas"]
        except:
            learning_points = ["Key ideas extracted", "More points will appear after review"]

        review = models.LearningReview(
            user_id=user.id,
            title=review_title,
            source_sessions=json.dumps(chat_session_ids) if chat_session_ids else None,
            source_slides=json.dumps(slide_ids) if slide_ids else None,
            source_content=combined_text[:4000],
            expected_points=json.dumps(learning_points),
            review_type=review_type,
            total_points=len(learning_points),
            best_score=0.0,
            current_attempt=0,
            attempt_count=0,
            status="active",
            created_at=datetime.now(timezone.utc)
        )

        db.add(review)
        db.commit()
        db.refresh(review)

        return {
            "status": "success",
            "review_id": review.id,
            "id": review.id,
            "title": review.title,
            "session_titles": session_titles,
            "slide_filenames": slide_filenames,
            "learning_points": learning_points,
            "total_points": len(learning_points)
        }

    except Exception as e:
        logger.error(f"Error creating review: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to create learning review: {str(e)}")


@app.get("/api/get_learning_reviews")
def get_learning_reviews(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        reviews = db.query(models.LearningReview).filter(
            models.LearningReview.user_id == user.id
        ).order_by(models.LearningReview.created_at.desc()).all()

        result = []
        for review in reviews:
            session_titles = []
            slide_filenames = []
            
            try:
                if review.source_sessions:
                    session_ids = json.loads(review.source_sessions)
                    sessions = db.query(models.ChatSession).filter(
                        models.ChatSession.id.in_(session_ids)
                    ).all()
                    session_titles = [s.title for s in sessions]
            except:
                pass
            
            try:
                if review.source_slides:
                    slide_ids = json.loads(review.source_slides)
                    slides = db.query(models.UploadedSlide).filter(
                        models.UploadedSlide.id.in_(slide_ids)
                    ).all()
                    slide_filenames = [s.original_filename for s in slides]
            except:
                pass

            result.append({
                "id": review.id,
                "title": review.title,
                "status": review.status,
                "total_points": review.total_points,
                "best_score": round(review.best_score, 1),
                "attempt_count": review.attempt_count,
                "current_attempt": review.current_attempt,
                "session_titles": session_titles,
                "slide_filenames": slide_filenames,
                "created_at": review.created_at.isoformat() + 'Z',
                "updated_at": review.updated_at.isoformat() + 'Z',
                "completed_at": review.completed_at.isoformat() + 'Z' if review.completed_at else None,
                "can_continue": review.status == "active" and review.current_attempt < 5
            })

        return {"reviews": result}

    except Exception as e:
        logger.error(f"Error getting learning reviews: {str(e)}")
        return {"reviews": []}


UPLOAD_DIR = Path("uploads/slides")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

@app.post("/api/upload_slides")
async def upload_slides(
    user_id: str = Form(...),
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db)
):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        uploaded_slides = []
        
        for file in files:
            if not file.filename.lower().endswith(('.pdf', '.ppt', '.pptx')):
                continue
            
            file_content = await file.read()
            file_size = len(file_content)
            
            timestamp = int(datetime.now(timezone.utc).timestamp())
            safe_filename = f"{user.id}_{timestamp}_{file.filename}"
            file_path = UPLOAD_DIR / safe_filename
            
            with open(file_path, 'wb') as f:
                f.write(file_content)
            
            page_count = 0
            extracted_text = ""
            
            if file.filename.lower().endswith('.pdf'):
                try:
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    page_count = len(pdf_reader.pages)
                    
                    for page in pdf_reader.pages[:10]:
                        extracted_text += page.extract_text() + "\n"
                    
                    extracted_text = extracted_text[:10000]
                except Exception as e:
                    logger.error(f"Error extracting PDF text: {str(e)}")
            
            elif file.filename.lower().endswith(('.ppt', '.pptx')):
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(file_content))
                    page_count = len(prs.slides)
                    
                    # Extract text from first 10 slides
                    for slide_idx, ppt_slide in enumerate(prs.slides):
                        if slide_idx >= 10:
                            break
                        for shape in ppt_slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                extracted_text += shape.text.strip() + "\n"
                    
                    extracted_text = extracted_text[:10000]
                except Exception as e:
                    logger.error(f"Error extracting PowerPoint content: {str(e)}")
            
            slide = models.UploadedSlide(
                user_id=user.id,
                filename=safe_filename,
                original_filename=file.filename,
                file_path=str(file_path),
                file_size=file_size,
                file_type=file.content_type or "application/pdf",
                page_count=page_count,
                extracted_text=extracted_text,
                processing_status="completed",
                uploaded_at=datetime.now(timezone.utc),
                processed_at=datetime.now(timezone.utc)
            )
            
            db.add(slide)
            uploaded_slides.append(slide)
        
        db.commit()
        
        return {
            "status": "success",
            "uploaded_count": len(uploaded_slides),
            "message": f"Successfully uploaded {len(uploaded_slides)} slide(s)"
        }
        
    except Exception as e:
        logger.error(f"Error uploading slides: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to upload slides: {str(e)}")


@app.get("/api/get_uploaded_slides")
def get_uploaded_slides(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        slides = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.user_id == user.id
        ).order_by(models.UploadedSlide.uploaded_at.desc()).all()
        
        return {
            "slides": [
                {
                    "id": slide.id,
                    "filename": slide.original_filename,
                    "file_size": slide.file_size,
                    "file_type": slide.file_type,
                    "page_count": slide.page_count,
                    "uploaded_at": slide.uploaded_at.isoformat() + 'Z',
                    "preview_url": slide.preview_url,
                    "processing_status": slide.processing_status
                }
                for slide in slides
            ]
        }
        
    except Exception as e:
        logger.error(f"Error getting slides: {str(e)}")
        return {"slides": []}


@app.delete("/api/delete_slide/{slide_id}")
def delete_slide(slide_id: int, db: Session = Depends(get_db)):
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()
        
        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        if Path(slide.file_path).exists():
            Path(slide.file_path).unlink()
        
        db.delete(slide)
        db.commit()
        
        return {
            "status": "success",
            "message": "Slide deleted successfully"
        }
        
    except Exception as e:
        logger.error(f"Error deleting slide: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/analyze_slide/{slide_id}")
async def analyze_slide(
    slide_id: int, 
    force_reanalyze: bool = Query(False),
    db: Session = Depends(get_db)
):
    """
    Get comprehensive slide analysis with persistent storage
    Returns cached analysis if available, unless force_reanalyze=True
    """
    try:
        from comprehensive_slide_analyzer import get_or_create_analysis
        
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()
        
        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        file_path = Path(slide.file_path)
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")
        
        # Determine file type
        file_ext = slide.original_filename.lower().split('.')[-1]
        if file_ext not in ['pdf', 'ppt', 'pptx']:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        # Get or create comprehensive analysis
        analysis_result = get_or_create_analysis(
            slide_id=slide_id,
            file_path=file_path,
            file_type=file_ext,
            db=db,
            force_reanalyze=force_reanalyze
        )
        
        return {
            "status": "success",
            "filename": slide.original_filename,
            "total_slides": analysis_result["total_slides"],
            "presentation_summary": analysis_result["presentation_summary"],
            "slides": analysis_result["slides"],
            "analyzed_at": analysis_result["analyzed_at"]
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing slide: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error analyzing slide: {str(e)}")


@app.get("/api/slide_file/{slide_id}")
async def get_slide_file(slide_id: int, db: Session = Depends(get_db)):
    """Serve the original uploaded slide file for viewing/downloading"""
    from fastapi.responses import FileResponse
    
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()
        
        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        file_path = Path(slide.file_path)
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")
        
        # Determine content type based on file extension
        filename = slide.original_filename.lower()
        if filename.endswith('.pdf'):
            media_type = "application/pdf"
        elif filename.endswith('.pptx'):
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif filename.endswith('.ppt'):
            media_type = "application/vnd.ms-powerpoint"
        else:
            media_type = "application/octet-stream"
        
        return FileResponse(
            path=str(file_path),
            media_type=media_type,
            filename=slide.original_filename,
            headers={
                "Content-Disposition": f"inline; filename=\"{slide.original_filename}\""
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error serving slide file: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error serving slide file: {str(e)}")


@app.get("/api/slide_image/{slide_id}/{page_number}")
async def get_slide_image(slide_id: int, page_number: int, db: Session = Depends(get_db)):
    """Render a specific page/slide as an image"""
    from fastapi.responses import Response
    import base64
    
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id
        ).first()
        
        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        file_path = Path(slide.file_path)
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")
        
        # Handle PDF files using PyMuPDF
        if slide.original_filename.lower().endswith('.pdf'):
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(str(file_path))
                
                if page_number < 1 or page_number > len(doc):
                    doc.close()
                    raise HTTPException(status_code=400, detail=f"Invalid page number. File has {len(doc)} pages.")
                
                page = doc[page_number - 1]  # 0-indexed
                
                # Render at 2x resolution for better quality
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                
                # Convert to PNG bytes
                img_bytes = pix.tobytes("png")
                doc.close()
                
                return Response(content=img_bytes, media_type="image/png")
                
            except ImportError:
                raise HTTPException(status_code=500, detail="PyMuPDF not installed for PDF rendering")
            except Exception as e:
                logger.error(f"Error rendering PDF page: {e}")
                raise HTTPException(status_code=500, detail=f"Error rendering PDF: {str(e)}")
        
        # Handle PowerPoint files
        elif slide.original_filename.lower().endswith(('.ppt', '.pptx')):
            try:
                # First try LibreOffice conversion
                import subprocess
                import tempfile
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        result = subprocess.run([
                            'soffice', '--headless', '--convert-to', 'pdf',
                            '--outdir', temp_dir, str(file_path)
                        ], capture_output=True, timeout=60)
                        
                        converted_files = [f for f in os.listdir(temp_dir) if f.endswith('.pdf')]
                        if converted_files:
                            pdf_path = os.path.join(temp_dir, converted_files[0])
                            
                            import fitz
                            doc = fitz.open(pdf_path)
                            
                            if page_number < 1 or page_number > len(doc):
                                doc.close()
                                raise HTTPException(status_code=400, detail=f"Invalid page number")
                            
                            page = doc[page_number - 1]
                            mat = fitz.Matrix(2.0, 2.0)
                            pix = page.get_pixmap(matrix=mat)
                            img_bytes = pix.tobytes("png")
                            doc.close()
                            
                            return Response(content=img_bytes, media_type="image/png")
                    except (subprocess.TimeoutExpired, FileNotFoundError):
                        pass
                
                # If LibreOffice not available, use python-pptx to render slide content
                try:
                    from pptx import Presentation as PptxPresentation
                    from pptx.util import Inches, Pt
                    from PIL import Image, ImageDraw, ImageFont
                    
                    prs = PptxPresentation(str(file_path))
                    
                    if page_number < 1 or page_number > len(prs.slides):
                        raise HTTPException(status_code=400, detail=f"Invalid page number. File has {len(prs.slides)} slides.")
                    
                    ppt_slide = prs.slides[page_number - 1]
                    
                    # Create image with slide dimensions (16:9 aspect ratio)
                    width, height = 1280, 720
                    img = Image.new('RGB', (width, height), color='#ffffff')
                    draw = ImageDraw.Draw(img)
                    
                    # Try to load a font, fall back to default
                    try:
                        title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 36)
                        body_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
                        small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
                    except:
                        try:
                            title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
                            body_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
                            small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
                        except:
                            title_font = ImageFont.load_default()
                            body_font = ImageFont.load_default()
                            small_font = ImageFont.load_default()
                    
                    # Extract text from shapes
                    y_position = 40
                    texts_extracted = []
                    
                    for shape in ppt_slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts_extracted.append(shape.text.strip())
                    
                    # Draw title (first text block, usually the title)
                    if texts_extracted:
                        title = texts_extracted[0][:100]  # Limit title length
                        # Word wrap title
                        words = title.split()
                        lines = []
                        current_line = []
                        for word in words:
                            current_line.append(word)
                            test_line = ' '.join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=title_font)
                            if bbox[2] > width - 80:
                                current_line.pop()
                                if current_line:
                                    lines.append(' '.join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(' '.join(current_line))
                        
                        for line in lines[:2]:  # Max 2 lines for title
                            draw.text((40, y_position), line, fill='#1a1a2e', font=title_font)
                            y_position += 45
                        
                        y_position += 20
                    
                    # Draw body text
                    for text in texts_extracted[1:]:
                        if y_position > height - 60:
                            break
                        
                        # Word wrap body text
                        words = text.split()
                        lines = []
                        current_line = []
                        for word in words:
                            current_line.append(word)
                            test_line = ' '.join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=body_font)
                            if bbox[2] > width - 100:
                                current_line.pop()
                                if current_line:
                                    lines.append(' '.join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(' '.join(current_line))
                        
                        for line in lines[:4]:  # Max 4 lines per text block
                            if y_position > height - 60:
                                break
                            draw.text((50, y_position), f"• {line}", fill='#333333', font=body_font)
                            y_position += 28
                        
                        y_position += 15
                    
                    # Add slide number at bottom
                    draw.text((width - 60, height - 30), f"{page_number}", fill='#888888', font=small_font)
                    
                    # If no text was extracted, show a message
                    if not texts_extracted:
                        draw.text((width//2, height//2 - 20), f"Slide {page_number}", fill='#1a1a2e', font=title_font, anchor='mm')
                        draw.text((width//2, height//2 + 30), "(Content preview not available)", fill='#888888', font=body_font, anchor='mm')
                    
                    # Convert to bytes
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format='PNG', quality=95)
                    img_bytes = img_buffer.getvalue()
                    
                    return Response(content=img_bytes, media_type="image/png")
                    
                except ImportError as e:
                    logger.error(f"python-pptx not available: {e}")
                    # Final fallback - simple placeholder
                    from PIL import Image, ImageDraw
                    
                    img = Image.new('RGB', (800, 600), color='#1a1a2e')
                    draw = ImageDraw.Draw(img)
                    draw.text((400, 280), f"Slide {page_number}", fill='#d7b38c', anchor='mm')
                    draw.text((400, 320), "Preview not available", fill='#888888', anchor='mm')
                    
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format='PNG')
                    img_bytes = img_buffer.getvalue()
                    
                    return Response(content=img_bytes, media_type="image/png")
                    
            except Exception as e:
                logger.error(f"Error rendering PowerPoint: {e}")
                raise HTTPException(status_code=500, detail=f"Error rendering PowerPoint: {str(e)}")
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
            
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting slide image: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error getting slide image: {str(e)}")


@app.post("/api/generate_questions")
async def generate_questions(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        user_id = payload.get("user_id")
        chat_session_ids = payload.get("chat_session_ids", [])
        slide_ids = payload.get("slide_ids", [])
        question_count = payload.get("question_count", 10)
        difficulty_mix = payload.get("difficulty_mix", {"easy": 3, "medium": 5, "hard": 2})
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        if not chat_session_ids and not slide_ids:
            raise HTTPException(status_code=400, detail="At least one source required")
        
        source_content = ""
        
        if chat_session_ids:
            sessions = db.query(models.ChatSession).filter(
                models.ChatSession.id.in_(chat_session_ids),
                models.ChatSession.user_id == user.id
            ).all()
            
            for session in sessions:
                messages = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == session.id
                ).limit(20).all()
                
                for msg in messages:
                    source_content += f"Q: {msg.user_message}\nA: {msg.ai_response}\n\n"
        
        if slide_ids:
            slides = db.query(models.UploadedSlide).filter(
                models.UploadedSlide.id.in_(slide_ids),
                models.UploadedSlide.user_id == user.id
            ).all()
            
            for slide in slides:
                if slide.extracted_text:
                    source_content += f"\n\nSlide: {slide.original_filename}\n{slide.extracted_text[:2000]}\n"
        
        source_content = source_content[:6000]
        
        user_profile = build_user_profile_dict(user)
        
        prompt = f"""Generate {question_count} educational questions based on this content.

**DIFFICULTY DISTRIBUTION**:
- Easy: {difficulty_mix.get('easy', 3)} questions
- Medium: {difficulty_mix.get('medium', 5)} questions  
- Hard: {difficulty_mix.get('hard', 2)} questions

**QUESTION TYPES**: Mix of multiple choice, true/false, and short answer

**SOURCE CONTENT**:
{source_content}

**OUTPUT FORMAT** (JSON array only):
[
  {{
    "question_text": "Question here?",
    "question_type": "multiple_choice|true_false|short_answer",
    "correct_answer": "Correct answer",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "difficulty": "easy|medium|hard",
    "explanation": "Why this is correct",
    "topic": "Topic name"
  }}
]

Generate exactly {question_count} questions:"""

        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": "Generate questions as JSON array."},
                {"role": "user", "content": prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=3000,
        )
        
        response_text = chat_completion.choices[0].message.content
        # Post-process math notation in questions
        response_text = process_math_in_response(response_text)
        
        try:
            json_match = re.search(r'\[.*\]', response_text, re.DOTALL)
            if json_match:
                questions_data = json.loads(json_match.group())
            else:
                raise ValueError("No JSON array found")
        except Exception as e:
            logger.error(f"Failed to parse questions: {str(e)}")
            questions_data = []
        
        question_set = models.QuestionSet(
            user_id=user.id,
            title=f"Questions - {datetime.now().strftime('%Y-%m-%d')}",
            description="Auto-generated questions",
            source_type="mixed" if (chat_session_ids and slide_ids) else ("chat" if chat_session_ids else "slides"),
            source_chat_sessions=json.dumps(chat_session_ids) if chat_session_ids else None,
            source_slides=json.dumps(slide_ids) if slide_ids else None,
            question_count=len(questions_data),
            easy_count=difficulty_mix.get('easy', 0),
            medium_count=difficulty_mix.get('medium', 0),
            hard_count=difficulty_mix.get('hard', 0),
            status="active"
        )
        
        db.add(question_set)
        db.commit()
        db.refresh(question_set)
        
        for idx, q_data in enumerate(questions_data):
            question = models.Question(
                question_set_id=question_set.id,
                question_text=q_data.get("question_text", ""),
                question_type=q_data.get("question_type", "multiple_choice"),
                correct_answer=q_data.get("correct_answer", ""),
                options=json.dumps(q_data.get("options", [])) if q_data.get("options") else None,
                difficulty=q_data.get("difficulty", "medium"),
                explanation=q_data.get("explanation", ""),
                topic=q_data.get("topic", "General"),
                order_index=idx
            )
            db.add(question)
        
        db.commit()
        db.refresh(question_set)
        
        questions = db.query(models.Question).filter(
            models.Question.question_set_id == question_set.id
        ).order_by(models.Question.order_index).all()
        
        return {
            "status": "success",
            "question_set_id": question_set.id,
            "id": question_set.id,
            "title": question_set.title,
            "question_count": len(questions),
            "questions": [
                {
                    "id": q.id,
                    "question_text": q.question_text,
                    "question_type": q.question_type,
                    "difficulty": q.difficulty,
                    "options": json.loads(q.options) if q.options else []
                }
                for q in questions
            ]
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating questions: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to generate questions: {str(e)}")


@app.get("/api/get_question_sets")
def get_question_sets(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        question_sets = db.query(models.QuestionSet).filter(
            models.QuestionSet.user_id == user.id
        ).order_by(models.QuestionSet.created_at.desc()).all()
        
        return {
            "question_sets": [
                {
                    "id": qs.id,
                    "title": qs.title,
                    "description": qs.description,
                    "question_count": qs.question_count,
                    "easy_count": qs.easy_count,
                    "medium_count": qs.medium_count,
                    "hard_count": qs.hard_count,
                    "best_score": round(qs.best_score, 1),
                    "attempt_count": qs.attempt_count,
                    "status": qs.status,
                    "can_practice": True,
                    "created_at": qs.created_at.isoformat() + 'Z'
                }
                for qs in question_sets
            ]
        }
        
    except Exception as e:
        logger.error(f"Error getting question sets: {str(e)}")
        return {"question_sets": []}
@app.delete("/api/delete_question_set/{question_set_id}")
def delete_question_set(
    question_set_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    try:
        question_set = db.query(models.QuestionSet).filter(
            models.QuestionSet.id == question_set_id,
            models.QuestionSet.user_id == current_user.id
        ).first()
        
        if not question_set:
            raise HTTPException(status_code=404, detail="Question set not found")
        
        db.query(models.QuestionResult).filter(
            models.QuestionResult.question_id.in_(
                db.query(models.Question.id).filter(
                    models.Question.question_set_id == question_set_id
                )
            )
        ).delete(synchronize_session=False)
        
        db.query(models.QuestionAttempt).filter(
            models.QuestionAttempt.question_set_id == question_set_id
        ).delete()
        
        db.query(models.Question).filter(
            models.Question.question_set_id == question_set_id
        ).delete()
        
        db.delete(question_set)
        db.commit()
        
        return {
            "status": "success",
            "message": "Question set deleted successfully"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting question set: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to delete question set: {str(e)}")
        
@app.post("/api/submit_question_answers")
async def submit_question_answers(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        question_set_id = payload.get("question_set_id")
        answers = payload.get("answers", {})
        
        question_set = db.query(models.QuestionSet).filter(
            models.QuestionSet.id == question_set_id
        ).first()
        
        if not question_set:
            raise HTTPException(status_code=404, detail="Question set not found")
        
        questions = db.query(models.Question).filter(
            models.Question.question_set_id == question_set_id
        ).all()
        
        correct_count = 0
        incorrect_count = 0
        question_results = []
        
        for question in questions:
            #  CRITICAL FIX: Check if question was answered at all
            question_id_str = str(question.id)
            
            # If question ID not in answers dict at all = unanswered
            if question_id_str not in answers:
                is_correct = False
                display_answer = "No answer provided"
            else:
                # Question ID exists in answers, get the value
                user_answer = answers[question_id_str]
                
                # Check if the answer is None or empty string
                if user_answer is None or (isinstance(user_answer, str) and user_answer.strip() == ""):
                    is_correct = False
                    display_answer = "No answer provided"
                else:
                    # Valid answer provided - now check if it's correct
                    display_answer = user_answer
                    
                    if question.question_type == "multiple_choice":
                        is_correct = user_answer.strip().lower() == question.correct_answer.strip().lower()
                    elif question.question_type == "true_false":
                        is_correct = user_answer.strip().lower() == question.correct_answer.strip().lower()
                    else:  # short_answer
                        is_correct = user_answer.strip().lower() in question.correct_answer.strip().lower()
            
            # Count the result
            if is_correct:
                correct_count += 1
            else:
                incorrect_count += 1
            
            question_results.append({
                "question_id": question.id,
                "question_text": question.question_text,
                "user_answer": display_answer,
                "correct_answer": question.correct_answer,
                "is_correct": is_correct,
                "explanation": question.explanation
            })
        
        score = (correct_count / len(questions) * 100) if questions else 0
        
        attempt = models.QuestionAttempt(
            question_set_id=question_set_id,
            user_id=question_set.user_id,
            attempt_number=question_set.attempt_count + 1,
            answers=json.dumps(answers),
            score=score,
            correct_count=correct_count,
            incorrect_count=incorrect_count,
            total_questions=len(questions),
            submitted_at=datetime.now(timezone.utc)
        )
        
        db.add(attempt)
        
        question_set.attempt_count += 1
        if score > question_set.best_score:
            question_set.best_score = score
        
        db.commit()
        
        return {
            "status": "success",
            "score": round(score, 1),
            "correct_count": correct_count,
            "incorrect_count": incorrect_count,
            "total_questions": len(questions),
            "question_results": question_results,
            "feedback": f"You scored {round(score, 1)}%! {'Excellent work!' if score >= 80 else 'Keep practicing!'}"
        }
        
    except Exception as e:
        logger.error(f"Error submitting answers: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/update_learning_review")
async def update_learning_review(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        review_id = payload.get("review_id")
        slide_ids = payload.get("slide_ids", [])
        
        review = db.query(models.LearningReview).filter(
            models.LearningReview.id == review_id
        ).first()
        
        if not review:
            raise HTTPException(status_code=404, detail="Review not found")
        
        if slide_ids:
            slides = db.query(models.UploadedSlide).filter(
                models.UploadedSlide.id.in_(slide_ids)
            ).all()
            
            slide_content = ""
            for slide in slides:
                if slide.extracted_text:
                    slide_content += f"\n{slide.extracted_text[:1000]}\n"
            
            current_content = review.source_content or ""
            review.source_content = (current_content + "\n" + slide_content)[:8000]
            
            for slide_id in slide_ids:
                link = models.LearningReviewSlide(
                    review_id=review_id,
                    slide_id=slide_id
                )
                db.add(link)
        
        review.updated_at = datetime.now(timezone.utc)
        db.commit()
        
        return {
            "status": "success",
            "message": "Review updated with slides"
        }
        
    except Exception as e:
        logger.error(f"Error updating review: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/get_question_set_details/{question_set_id}")
def get_question_set_details(question_set_id: int, db: Session = Depends(get_db)):
    try:
        question_set = db.query(models.QuestionSet).filter(
            models.QuestionSet.id == question_set_id
        ).first()
        
        if not question_set:
            raise HTTPException(status_code=404, detail="Question set not found")
        
        questions = db.query(models.Question).filter(
            models.Question.question_set_id == question_set_id
        ).order_by(models.Question.order_index).all()
        
        return {
            "id": question_set.id,
            "title": question_set.title,
            "description": question_set.description,
            "question_count": len(questions),
            "status": question_set.status,
            "questions": [
                {
                    "id": q.id,
                    "question_text": q.question_text,
                    "question_type": q.question_type,
                    "correct_answer": q.correct_answer,
                    "options": json.loads(q.options) if q.options else [],
                    "difficulty": q.difficulty,
                    "explanation": q.explanation,
                    "topic": q.topic
                }
                for q in questions
            ]
        }
        
    except Exception as e:
        logger.error(f"Error getting question set details: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/submit_learning_response")
async def submit_learning_response(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """
    Submit and evaluate a learning review response
    """
    try:
        review_id = payload.get("review_id")
        user_response = payload.get("user_response", "")
        attempt_number = payload.get("attempt_number", 1)

        if not review_id or not user_response.strip():
            raise HTTPException(status_code=400, detail="Review ID and response are required")

        # Get the learning review
        review = db.query(models.LearningReview).filter(
            models.LearningReview.id == review_id
        ).first()

        if not review:
            raise HTTPException(status_code=404, detail="Learning review not found")

        # Parse expected learning points
        try:
            expected_points = json.loads(review.expected_points)
        except:
            expected_points = []

        if not expected_points:
            raise HTTPException(status_code=400, detail="No learning points found in review")

        # Get user profile for personalized feedback
        user = db.query(models.User).filter(models.User.id == review.user_id).first()
        user_profile = build_user_profile_dict(user)

        # ==================== OPTIMIZED EVALUATION PROMPT (77% TOKEN REDUCTION) ====================
        evaluation_prompt = f"""Evaluate student's learning retention.

**EXPECTED POINTS**:
{chr(10).join([f"{i+1}. {point}" for i, point in enumerate(expected_points)])}

**STUDENT RESPONSE**:
{user_response}

**RULES**: Point is "covered" if student shows understanding (semantic match, not exact words).

**OUTPUT JSON**:
{{
  "covered_points": ["Covered points from expected list"],
  "missing_points": ["Missing points from expected list"],
  "coverage_percentage": <0-100>,
  "understanding_quality": "<poor|fair|good|excellent>",
  "feedback": "Brief feedback on strengths/improvements",
  "next_steps": "Actionable advice"
}}"""

        # Call Groq AI for evaluation
        chat_completion = groq_client.chat.completions.create(
            messages=[
                {
                    "role": "system", 
                    "content": "Evaluate answers. Return JSON with scores, feedback, strengths, weaknesses."
                },
                {
                    "role": "user", 
                    "content": evaluation_prompt
                }
            ],
            model=GROQ_MODEL,
            temperature=0.3,  # Lower temperature for more consistent evaluation
            max_tokens=2048,
        )

        response_text = chat_completion.choices[0].message.content

        # Parse AI evaluation response
        try:
            # Extract JSON from response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                evaluation = json.loads(json_match.group())
            else:
                raise ValueError("No JSON found in response")
        except Exception as e:
            logger.error(f"Failed to parse evaluation JSON: {str(e)}")
            logger.error(f"Raw response: {response_text}")
            
            # Fallback evaluation if JSON parsing fails
            evaluation = {
                "covered_points": [],
                "missing_points": expected_points,
                "coverage_percentage": 0,
                "understanding_quality": "fair",
                "feedback": "Unable to properly evaluate your response. Please try again with more detail.",
                "next_steps": "Write a more comprehensive response covering all key concepts."
            }

        covered_points = evaluation.get("covered_points", [])
        missing_points = evaluation.get("missing_points", [])
        coverage_percentage = evaluation.get("coverage_percentage", 0)
        understanding_quality = evaluation.get("understanding_quality", "fair")
        feedback = evaluation.get("feedback", "")
        next_steps = evaluation.get("next_steps", "")

        # ==================== SAVE ATTEMPT ====================
        attempt = models.LearningReviewAttempt(
            review_id=review.id,
            attempt_number=attempt_number,
            user_response=user_response,
            covered_points=json.dumps(covered_points),
            missing_points=json.dumps(missing_points),
            completeness_percentage=coverage_percentage,
            feedback=feedback,
            submitted_at=datetime.now(timezone.utc)
        )
        db.add(attempt)

        # Update review record
        review.current_attempt = attempt_number
        review.attempt_count = max(review.attempt_count, attempt_number)
        
        if coverage_percentage > review.best_score:
            review.best_score = coverage_percentage

        # Check if review is complete (>= 80% coverage or 3+ attempts)
        is_complete = coverage_percentage >= 80.0 or attempt_number >= 3
        
        if is_complete and coverage_percentage >= 80.0:
            review.status = "completed"
            review.completed_at = datetime.now(timezone.utc)
        
        review.updated_at = datetime.now(timezone.utc)

        db.commit()
        db.refresh(attempt)

        # ==================== RETURN RESULTS ====================
        return {
            "status": "success",
            "attempt_id": attempt.id,
            "attempt_number": attempt_number,
            "covered_points": covered_points,
            "missing_points": missing_points,
            "completeness_percentage": round(coverage_percentage, 1),
            "understanding_quality": understanding_quality,
            "feedback": feedback,
            "next_steps": next_steps,
            "is_complete": is_complete,
            "can_continue": not is_complete and attempt_number < 5,
            "best_score": review.best_score,
            "points_covered_count": len(covered_points),
            "points_missing_count": len(missing_points),
            "total_points": len(expected_points)
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error submitting learning response: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to evaluate learning response: {str(e)}"
        )
#=====================Nodes endpoint=====================

# ==================== KNOWLEDGE ROADMAP SYSTEM ====================

@app.post("/api/create_knowledge_roadmap")
async def create_knowledge_roadmap(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """
    Create a new knowledge roadmap with ONLY root topic (no initial subtopics)
    """
    try:
        user_id = payload.get("user_id")
        root_topic = payload.get("root_topic")
        
        if not user_id or not root_topic:
            raise HTTPException(status_code=400, detail="user_id and root_topic required")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Create ONLY root node (no AI generation yet)
        root_node = models.KnowledgeNode(
            user_id=user.id,
            parent_node_id=None,
            topic_name=root_topic,
            description=f"Explore {root_topic} - Click 'Explore' to learn or 'Expand' to see subtopics",
            depth_level=0,
            ai_explanation=None,  # Will be generated on explore
            key_concepts=None,
            generated_subtopics=None,
            is_explored=False,
            exploration_count=0,
            expansion_status="unexpanded",  # Not expanded initially
            position_x=0.0,
            position_y=0.0
        )
        
        db.add(root_node)
        db.flush()
        
        # Create roadmap
        roadmap = models.KnowledgeRoadmap(
            user_id=user.id,
            title=f"Exploring {root_topic}",
            root_topic=root_topic,
            root_node_id=root_node.id,
            total_nodes=1,  # Only root node
            max_depth_reached=0,
            status="active",
            last_accessed=datetime.now(timezone.utc)
        )
        
        db.add(roadmap)
        db.commit()
        db.refresh(roadmap)
        db.refresh(root_node)
        
        return {
            "status": "success",
            "roadmap_id": roadmap.id,
            "root_node_id": root_node.id,
            "root_node": {
                "id": root_node.id,
                "topic_name": root_node.topic_name,
                "description": root_node.description,
                "depth_level": root_node.depth_level,
                "ai_explanation": root_node.ai_explanation,
                "key_concepts": json.loads(root_node.key_concepts) if root_node.key_concepts else [],
                "is_explored": root_node.is_explored,
                "expansion_status": root_node.expansion_status,
                "position": {"x": root_node.position_x, "y": root_node.position_y}
            },
            "child_nodes": [],  # No children initially
            "total_nodes": roadmap.total_nodes
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error creating roadmap: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to create roadmap: {str(e)}")

@app.post("/api/create_roadmap_from_chat")
async def create_roadmap_from_chat(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        user_id = payload.get("user_id")
        chat_session_id = payload.get("chat_session_id")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get chat session
        chat_session = db.query(models.ChatSession).filter(
            models.ChatSession.id == chat_session_id,
            models.ChatSession.user_id == user.id
        ).first()
        
        if not chat_session:
            raise HTTPException(status_code=404, detail="Chat session not found")
        
        # Get all messages
        messages = db.query(models.ChatMessage).filter(
            models.ChatMessage.chat_session_id == chat_session_id
        ).order_by(models.ChatMessage.timestamp.asc()).all()
        
        if not messages:
            raise HTTPException(status_code=404, detail="No messages in chat session")
        
        # Compile conversation
        conversation_text = []
        for msg in messages:
            conversation_text.append(f"Q: {msg.user_message}")
            conversation_text.append(f"A: {msg.ai_response}")
        
        full_conversation = "\n\n".join(conversation_text)[:4000]
        
        # Generate topic from conversation using AI
        prompt = f"""Analyze this conversation and extract the MAIN TOPIC in 2-4 words.

Conversation:
{full_conversation}

Return ONLY the topic name (e.g., "Machine Learning", "World War II", "Quantum Physics").
No explanation, just the topic:"""

        user_profile = build_user_profile_dict(user)
        
        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": "Extract main topic name only."},
                {"role": "user", "content": prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.3,
            max_tokens=50,
        )
        
        root_topic = chat_completion.choices[0].message.content.strip()
        root_topic = root_topic.replace('"', '').replace("'", "")
        
        # Return the extracted topic
        return {
            "status": "success",
            "root_topic": root_topic,
            "chat_title": chat_session.title,
            "message": f"Roadmap topic identified: {root_topic}"
        }
        
    except Exception as e:
        logger.error(f"Error creating roadmap from chat: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
        
@app.post("/api/expand_knowledge_node/{node_id}")
async def expand_knowledge_node(
    node_id: int,
    db: Session = Depends(get_db)
):
    """
    Expand a node by generating 4-5 subtopics (DOES NOT generate explanations)
    """
    try:
        node = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.id == node_id
        ).first()
        
        if not node:
            raise HTTPException(status_code=404, detail="Node not found")
        
        # Check if already expanded to avoid duplicate work
        if node.expansion_status == "expanded":
            children = db.query(models.KnowledgeNode).filter(
                models.KnowledgeNode.parent_node_id == node_id
            ).all()
            
            return {
                "status": "success",
                "message": "Node already expanded",
                "child_nodes": [
                    {
                        "id": child.id,
                        "parent_id": child.parent_node_id,
                        "topic_name": child.topic_name,
                        "description": child.description,
                        "depth_level": child.depth_level,
                        "is_explored": child.is_explored,
                        "expansion_status": child.expansion_status,
                        "position": {"x": child.position_x, "y": child.position_y}
                    }
                    for child in children
                ]
            }
        
        # Update expansion status immediately to prevent duplicate requests
        node.expansion_status = "expanding"
        db.commit()
        
        # Get user for personalized generation
        user = db.query(models.User).filter(models.User.id == node.user_id).first()
        user_profile = build_user_profile_dict(user)
        
        # Build context from parent nodes
        context_path = []
        current = node
        while current:
            context_path.insert(0, current.topic_name)
            if current.parent_node_id:
                current = db.query(models.KnowledgeNode).filter(
                    models.KnowledgeNode.id == current.parent_node_id
                ).first()
            else:
                current = None
        
        context_str = " → ".join(context_path)
        
        # OPTIMIZED EXPANSION PROMPT (72% TOKEN REDUCTION)
        expansion_prompt = f"""Expand "{node.topic_name}".
Context: {context_str}
Depth: {node.depth_level}
Level: {user_profile.get('difficulty_level', 'intermediate')}

Generate 4-5 specific subtopics (more specific than parent, 2-5 words each).

**JSON OUTPUT**:
{{
  "subtopics": [
    {{
      "name": "Short Name (2-5 words)",
      "description": "One-line description (<100 chars)",
      "complexity": "beginner|intermediate|advanced"
    }}
  ]
}}"""

        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": "Generate 4-5 subtopics as JSON."},
                {"role": "user", "content": expansion_prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.8,
            max_tokens=1000,
        )
        
        response_text = chat_completion.choices[0].message.content
        
        try:
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                ai_data = json.loads(json_match.group())
            else:
                raise ValueError("No JSON found")
        except Exception as e:
            logger.error(f"Failed to parse expansion JSON: {str(e)}")
            # Fallback subtopics
            ai_data = {
                "subtopics": [
                    {"name": "Fundamentals", "description": "Core concepts and basics", "complexity": "beginner"},
                    {"name": "Key Principles", "description": "Essential rules and theories", "complexity": "intermediate"},
                    {"name": "Applications", "description": "Real-world uses", "complexity": "intermediate"},
                    {"name": "Advanced Topics", "description": "Deep dive into complexity", "complexity": "advanced"}
                ]
            }
        
        # Create child nodes WITHOUT explanations
        subtopics = ai_data.get("subtopics", [])[:5]  # Max 5 subtopics
        child_nodes = []
        
        for idx, subtopic in enumerate(subtopics):
            # Calculate position in a circle around parent
            angle = (idx * (360 / len(subtopics))) * (3.14159 / 180)
            radius = 300 + (node.depth_level * 50)  # Increase radius with depth
            
            child_node = models.KnowledgeNode(
                user_id=node.user_id,
                parent_node_id=node.id,
                topic_name=subtopic.get("name", ""),
                description=subtopic.get("description", ""),
                depth_level=node.depth_level + 1,
                ai_explanation=None,  # NOT generated yet - only on explore
                key_concepts=None,
                generated_subtopics=None,
                is_explored=False,
                exploration_count=0,
                expansion_status="unexpanded",
                position_x=node.position_x + (radius * float(math.cos(angle))),
                position_y=node.position_y + (radius * float(math.sin(angle)))
            )
            db.add(child_node)
            child_nodes.append(child_node)
        
        # Update parent node status to expanded
        node.expansion_status = "expanded"
        node.generated_subtopics = json.dumps(subtopics)
        
        # Update roadmap stats
        roadmap = db.query(models.KnowledgeRoadmap).filter(
            models.KnowledgeRoadmap.user_id == node.user_id
        ).order_by(models.KnowledgeRoadmap.created_at.desc()).first()
        
        if roadmap:
            roadmap.total_nodes += len(child_nodes)
            roadmap.max_depth_reached = max(roadmap.max_depth_reached, node.depth_level + 1)
            roadmap.last_accessed = datetime.now(timezone.utc)
        
        db.commit()
        
        for child in child_nodes:
            db.refresh(child)
        
        return {
            "status": "success",
            "message": f"Expanded {node.topic_name} with {len(child_nodes)} subtopics",
            "child_nodes": [
                {
                    "id": child.id,
                    "parent_id": child.parent_node_id,
                    "topic_name": child.topic_name,
                    "description": child.description,
                    "depth_level": child.depth_level,
                    "is_explored": child.is_explored,
                    "expansion_status": child.expansion_status,
                    "position": {"x": child.position_x, "y": child.position_y}
                }
                for child in child_nodes
            ]
        }
        
    except HTTPException:
        db.rollback()
        raise
    except Exception as e:
        logger.error(f"Error expanding node: {str(e)}", exc_info=True)
        # Reset expansion status on error
        if node:
            node.expansion_status = "unexpanded"
            db.commit()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to expand node: {str(e)}")

@app.post("/api/explore_node/{node_id}")
async def explore_node(
    node_id: int,
    db: Session = Depends(get_db)
):
    """
    Generate AI explanation for a node (DOES NOT expand or create children)
    """
    try:
        node = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.id == node_id
        ).first()
        
        if not node:
            raise HTTPException(status_code=404, detail="Node not found")
        
        # If already has explanation, just return it
        if node.ai_explanation and node.key_concepts:
            node.exploration_count += 1
            node.last_explored = datetime.now(timezone.utc)
            db.commit()
            db.refresh(node)
            
            return {
                "status": "already_generated",
                "node": {
                    "id": node.id,
                    "topic_name": node.topic_name,
                    "ai_explanation": node.ai_explanation,
                    "key_concepts": json.loads(node.key_concepts) if node.key_concepts else [],
                    "is_explored": node.is_explored,
                    "exploration_count": node.exploration_count
                }
            }
        
        # Generate explanation
        user = db.query(models.User).filter(models.User.id == node.user_id).first()
        user_profile = build_user_profile_dict(user)
        
        # Build context
        context_path = []
        current = node.parent_node_id
        while current:
            parent = db.query(models.KnowledgeNode).filter(
                models.KnowledgeNode.id == current
            ).first()
            if parent:
                context_path.insert(0, parent.topic_name)
                current = parent.parent_node_id
            else:
                current = None
        
        context_str = " → ".join(context_path) if context_path else "Root level"
        
        # OPTIMIZED EXPLANATION PROMPT (75% TOKEN REDUCTION)
        explanation_prompt = f"""Explain "{node.topic_name}" for {user_profile.get('first_name', 'student')}.
Context: {context_str}
Level: {user_profile.get('difficulty_level', 'intermediate')}

**JSON OUTPUT**:
{{
  "explanation": "Clear explanation (250-400 words) with examples",
  "key_concepts": ["Concept 1", "Concept 2", "Concept 3", "Concept 4", "Concept 5"],
  "why_important": "Why this matters (2-3 sentences)",
  "real_world_examples": ["Example 1 with context", "Example 2 with context"],
  "learning_tips": "Practical mastery advice"
}}"""

        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": "Explain topic. Return JSON with explanation, examples, resources."},
                {"role": "user", "content": explanation_prompt}
            ],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=2048,
        )
        
        response_text = chat_completion.choices[0].message.content
        
        try:
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                ai_data = json.loads(json_match.group())
            else:
                raise ValueError("No JSON found")
        except Exception as e:
            logger.error(f"Failed to parse explanation JSON: {str(e)}")
            ai_data = {
                "explanation": f"An exploration of {node.topic_name}. This topic involves understanding the fundamental concepts and their applications.",
                "key_concepts": ["Core principles", "Key theories", "Practical applications", "Related fields", "Future directions"],
                "why_important": "Understanding this topic is essential for building a strong foundation in this domain.",
                "real_world_examples": ["Used in various industries", "Applied in research and development"],
                "learning_tips": "Practice regularly, connect concepts to real-world scenarios, and explore related topics."
            }
        
        # Update node with explanation (mark as explored)
        node.ai_explanation = ai_data.get("explanation", "")
        node.key_concepts = json.dumps(ai_data.get("key_concepts", []))
        node.why_important = ai_data.get("why_important", "")
        node.real_world_examples = json.dumps(ai_data.get("real_world_examples", []))
        node.learning_tips = ai_data.get("learning_tips", "")
        node.is_explored = True
        node.exploration_count += 1
        node.last_explored = datetime.now(timezone.utc)
        
        # Create exploration history record
        history = models.NodeExplorationHistory(
            node_id=node.id,
            user_id=node.user_id,
            exploration_duration=0,
            explored_at=datetime.now(timezone.utc)
        )
        db.add(history)
        
        db.commit()
        db.refresh(node)
        
        return {
            "status": "success",
            "node": {
                "id": node.id,
                "topic_name": node.topic_name,
                "description": node.description,
                "ai_explanation": node.ai_explanation,
                "key_concepts": json.loads(node.key_concepts) if node.key_concepts else [],
                "why_important": ai_data.get("why_important", ""),
                "real_world_examples": ai_data.get("real_world_examples", []),
                "learning_tips": ai_data.get("learning_tips", ""),
                "is_explored": node.is_explored,
                "exploration_count": node.exploration_count
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error exploring node: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to explore node: {str(e)}")

@app.get("/api/get_knowledge_roadmap/{roadmap_id}")
async def get_knowledge_roadmap(
    roadmap_id: int,
    db: Session = Depends(get_db)
):
    """
    Get complete roadmap structure with all nodes
    """
    try:
        roadmap = db.query(models.KnowledgeRoadmap).filter(
            models.KnowledgeRoadmap.id == roadmap_id
        ).first()
        
        if not roadmap:
            raise HTTPException(status_code=404, detail="Roadmap not found")
        
        # Get root node
        root_node = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.id == roadmap.root_node_id
        ).first()
        
        # Get flat list of all nodes in this roadmap tree
        roadmap_node_ids = set()
        
        def collect_node_ids(node_id):
            roadmap_node_ids.add(node_id)
            children = db.query(models.KnowledgeNode).filter(
                models.KnowledgeNode.parent_node_id == node_id
            ).all()
            for child in children:
                collect_node_ids(child.id)
        
        if root_node:
            collect_node_ids(root_node.id)
        
        all_nodes = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.id.in_(roadmap_node_ids)
        ).all()
        
        # Create a map of nodes for easier lookup
        node_map = {node.id: node for node in all_nodes}
        
        # Build the expansion hierarchy
        def build_expansion_hierarchy(node_id, expanded_nodes, path=[]):
            node = node_map.get(node_id)
            if not node:
                return
                
            # Add this node to the path
            current_path = path + [node_id]
            
            # If this node is expanded, add it to the expanded nodes set
            if node.expansion_status == 'expanded':
                expanded_nodes.add(node_id)
                
                # Recursively add all children of expanded nodes
                children = db.query(models.KnowledgeNode).filter(
                    models.KnowledgeNode.parent_node_id == node_id
                ).all()
                for child in children:
                    build_expansion_hierarchy(child.id, expanded_nodes, current_path)
        
        # Get all expanded nodes
        expanded_nodes = set()
        if root_node:
            build_expansion_hierarchy(root_node.id, expanded_nodes)
        
        nodes_flat = [
            {
                "id": node.id,
                "parent_id": node.parent_node_id,
                "topic_name": node.topic_name,
                "description": node.description,
                "depth_level": node.depth_level,
                "ai_explanation": node.ai_explanation,
                "key_concepts": json.loads(node.key_concepts) if node.key_concepts else [],
                "why_important": node.why_important,
                "real_world_examples": json.loads(node.real_world_examples) if node.real_world_examples else [],
                "learning_tips": node.learning_tips,
                "is_explored": node.is_explored,
                "exploration_count": node.exploration_count,
                "expansion_status": node.expansion_status,
                "user_notes": node.user_notes,
                "position": {"x": node.position_x, "y": node.position_y},
                "created_at": node.created_at.isoformat() + 'Z'
            }
            for node in all_nodes
        ]
        
        return {
            "roadmap": {
                "id": roadmap.id,
                "title": roadmap.title,
                "root_topic": roadmap.root_topic,
                "total_nodes": roadmap.total_nodes,
                "max_depth_reached": roadmap.max_depth_reached,
                "status": roadmap.status,
                "created_at": roadmap.created_at.isoformat() + 'Z',
                "last_accessed": roadmap.last_accessed.isoformat() + 'Z' if roadmap.last_accessed else None
            },
            "nodes_flat": nodes_flat,
            "expanded_nodes": list(expanded_nodes)  # Add this to track which nodes should be expanded
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting roadmap: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to get roadmap: {str(e)}")

@app.get("/api/get_user_roadmaps")
async def get_user_roadmaps(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """
    Get all knowledge roadmaps for a user
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        roadmaps = db.query(models.KnowledgeRoadmap).filter(
            models.KnowledgeRoadmap.user_id == user.id
        ).order_by(models.KnowledgeRoadmap.last_accessed.desc()).all()
        
        roadmap_data = []
        for roadmap in roadmaps:
            # Get root node info
            root_node = db.query(models.KnowledgeNode).filter(
                models.KnowledgeNode.id == roadmap.root_node_id
            ).first()
            
            roadmap_data.append({
                "id": roadmap.id,
                "title": roadmap.title,
                "root_topic": roadmap.root_topic,
                "total_nodes": roadmap.total_nodes,
                "max_depth_reached": roadmap.max_depth_reached,
                "status": roadmap.status,
                "created_at": roadmap.created_at.isoformat() + 'Z',
                "last_accessed": roadmap.last_accessed.isoformat() + 'Z' if roadmap.last_accessed else roadmap.created_at.isoformat() + 'Z'
            })
        
        return {
            "status": "success",
            "roadmaps": roadmap_data
        }
        
    except Exception as e:
        logger.error(f"Error getting user roadmaps: {str(e)}")
        return {
            "status": "error",
            "roadmaps": []
        }
 # ==================== MISSING ENDPOINTS FOR LEARNING REVIEW ====================

@app.get("/api/get_generated_questions")
def get_generated_questions(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Get all generated question sets for a user"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        question_sets = db.query(models.QuestionSet).filter(
            models.QuestionSet.user_id == user.id
        ).order_by(models.QuestionSet.created_at.desc()).all()
        
        return {
            "question_sets": [
                {
                    "id": qs.id,
                    "title": qs.title,
                    "question_count": qs.question_count,
                    "created_at": qs.created_at.isoformat() + 'Z',
                    "status": qs.status,
                    "best_score": round(qs.best_score, 1),
                    "attempt_count": qs.attempt_count
                }
                for qs in question_sets
            ]
        }
        
    except Exception as e:
        logger.error(f"Error getting generated questions: {str(e)}")
        return {"question_sets": []}

@app.get("/api/get_question_set/{question_set_id}")
def get_question_set_with_questions(question_set_id: int, db: Session = Depends(get_db)):
    """Get a specific question set with all its questions"""
    try:
        question_set = db.query(models.QuestionSet).filter(
            models.QuestionSet.id == question_set_id
        ).first()
        
        if not question_set:
            raise HTTPException(status_code=404, detail="Question set not found")
        
        questions = db.query(models.Question).filter(
            models.Question.question_set_id == question_set_id
        ).order_by(models.Question.order_index).all()
        
        return {
            "id": question_set.id,
            "title": question_set.title,
            "description": question_set.description,
            "questions": [
                {
                    "id": q.id,
                    "question_text": q.question_text,
                    "question_type": q.question_type,
                    "options": json.loads(q.options) if q.options else [],
                    "difficulty": q.difficulty,
                    "explanation": q.explanation,
                    "topic": q.topic
                }
                for q in questions
            ]
        }
        
    except Exception as e:
        logger.error(f"Error getting question set: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/submit_answers")
async def submit_answers(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Submit answers for a question set and get results"""
    try:
        user_id = payload.get("user_id")
        question_set_id = payload.get("question_set_id")
        answers = payload.get("answers", {})
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        question_set = db.query(models.QuestionSet).filter(
            models.QuestionSet.id == question_set_id
        ).first()
        
        if not question_set:
            raise HTTPException(status_code=404, detail="Question set not found")
        
        questions = db.query(models.Question).filter(
            models.Question.question_set_id == question_set_id
        ).all()
        
        results = []
        correct_count = 0
        
        for question in questions:
            user_answer = answers.get(str(question.id), "").strip()
            is_correct = False
            
            if question.question_type == "multiple_choice":
                is_correct = user_answer.lower() == question.correct_answer.lower()
            elif question.question_type == "true_false":
                is_correct = user_answer.lower() == question.correct_answer.lower()
            else:  # short_answer
                # For short answer, check if user's answer contains key phrases
                is_correct = any(keyword in user_answer.lower() 
                               for keyword in question.correct_answer.lower().split()[:3])
            
            if is_correct:
                correct_count += 1
            
            results.append({
                "question_id": question.id,
                "user_answer": user_answer,
                "correct_answer": question.correct_answer,
                "is_correct": is_correct,
                "explanation": question.explanation
            })
        
        score = (correct_count / len(questions)) * 100 if questions else 0
        
        # Update question set stats
        question_set.attempt_count += 1
        if score > question_set.best_score:
            question_set.best_score = score
        
        # Record the attempt
        attempt = models.QuestionAttempt(
            question_set_id=question_set_id,
            user_id=user.id,
            attempt_number=question_set.attempt_count,
            answers=json.dumps(answers),
            score=score,
            correct_count=correct_count,
            incorrect_count=len(questions) - correct_count,
            total_questions=len(questions),
            submitted_at=datetime.now(timezone.utc)
        )
        db.add(attempt)
        db.commit()
        
        return {
            "status": "success",
            "score": round(score, 1),
            "correct_count": correct_count,
            "total_questions": len(questions),
            "details": results
        }
        
    except Exception as e:
        logger.error(f"Error submitting answers: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_hints/{question_id}")
def get_hints_for_question(question_id: int, db: Session = Depends(get_db)):
    """Get hints for a specific question"""
    try:
        question = db.query(models.Question).filter(
            models.Question.id == question_id
        ).first()
        
        if not question:
            raise HTTPException(status_code=404, detail="Question not found")
        
        # Generate hints based on question content
        hints = []
        
        if question.question_type == "multiple_choice":
            options = json.loads(question.options) if question.options else []
            if len(options) >= 2:
                hints.append(f"Consider the options carefully. Eliminate obviously wrong answers first.")
                hints.append(f"The correct answer relates to: {question.topic}")
        
        elif question.question_type == "true_false":
            hints.append(f"Think about the fundamental concepts of {question.topic}")
            hints.append(f"Consider real-world applications of this concept")
        
        else:  # short_answer
            hints.append(f"Focus on the key terms related to {question.topic}")
            hints.append(f"Think about how this concept is applied in practice")
        
        hints.append(f"Review your notes on {question.topic} if you're unsure")
        
        return {
            "question_id": question_id,
            "hints": hints[:3]  # Return max 3 hints
        }
        
    except Exception as e:
        logger.error(f"Error getting hints: {str(e)}")
        return {"question_id": question_id, "hints": ["Think about the main concepts covered in this topic."]}

@app.post("/api/submit_review_response")
async def submit_review_response(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Submit a response for a learning review"""
    try:
        user_id = payload.get("user_id")
        review_id = payload.get("review_id")
        response_text = payload.get("response_text", "")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        review = db.query(models.LearningReview).filter(
            models.LearningReview.id == review_id
        ).first()
        
        if not review:
            raise HTTPException(status_code=404, detail="Learning review not found")
        
        # Parse expected points
        try:
            expected_points = json.loads(review.expected_points)
        except:
            expected_points = []
        
        # Simple evaluation - count how many points are mentioned
        covered_points = []
        for point in expected_points:
            if any(keyword in response_text.lower() for keyword in point.lower().split()[:3]):
                covered_points.append(point)
        
        coverage_percentage = (len(covered_points) / len(expected_points)) * 100 if expected_points else 0
        
        # Generate feedback
        if coverage_percentage >= 80:
            feedback = "Excellent! You've covered most of the key concepts thoroughly."
            strengths = ["Comprehensive understanding", "Good recall of main points"]
            improvements = ["Consider adding more specific examples"]
        elif coverage_percentage >= 60:
            feedback = "Good effort! You've captured the main ideas but missed some details."
            strengths = ["Solid grasp of core concepts", "Clear expression"]
            improvements = ["Add more specific details", "Expand on the applications"]
        else:
            feedback = "You're on the right track but missed several key concepts. Review the material and try again."
            strengths = ["Good attempt at explaining what you remember"]
            improvements = ["Review the main concepts more thoroughly", "Focus on key definitions and applications"]
        
        # Save the attempt
        attempt = models.LearningReviewAttempt(
            review_id=review_id,
            attempt_number=review.current_attempt + 1,
            user_response=response_text,
            covered_points=json.dumps(covered_points),
            missing_points=json.dumps([p for p in expected_points if p not in covered_points]),
            completeness_percentage=coverage_percentage,
            feedback=feedback,
            submitted_at=datetime.now(timezone.utc)
        )
        db.add(attempt)
        
        # Update review
        review.current_attempt += 1
        if coverage_percentage > review.best_score:
            review.best_score = coverage_percentage
        
        if coverage_percentage >= 80:
            review.status = "completed"
            review.completed_at = datetime.now(timezone.utc)
        
        db.commit()
        
        return {
            "status": "success",
            "score": round(coverage_percentage, 1),
            "feedback": feedback,
            "strengths": strengths,
            "improvements": improvements,
            "covered_points": covered_points,
            "total_points": len(expected_points)
        }
        
    except Exception as e:
        logger.error(f"Error submitting review response: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))
               
@app.get("/api/get_knowledge_roadmaps")
async def get_knowledge_roadmaps(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """
    Get all roadmaps for a user
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        roadmaps = db.query(models.KnowledgeRoadmap).filter(
            models.KnowledgeRoadmap.user_id == user.id
        ).order_by(models.KnowledgeRoadmap.created_at.desc()).all()
        
        return {
            "roadmaps": [
                {
                    "id": roadmap.id,
                    "title": roadmap.title,
                    "root_topic": roadmap.root_topic,
                    "total_nodes": roadmap.total_nodes,
                    "max_depth_reached": roadmap.max_depth_reached,
                    "status": roadmap.status,
                    "created_at": roadmap.created_at.isoformat() + 'Z',
                    "last_accessed": roadmap.last_accessed.isoformat() + 'Z' if roadmap.last_accessed else None
                }
                for roadmap in roadmaps
            ]
        }
        
    except Exception as e:
        logger.error(f"Error getting user roadmaps: {str(e)}")
        return {"roadmaps": []}


@app.post("/api/save_node_notes/{node_id}")
async def save_node_notes(
    node_id: int,
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """
    Save user notes on a node
    """
    try:
        node = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.id == node_id
        ).first()
        
        if not node:
            raise HTTPException(status_code=404, detail="Node not found")
        
        notes = payload.get("notes", "")
        node.user_notes = notes
        
        db.commit()
        
        return {
            "status": "success",
            "message": "Notes saved successfully"
        }
        
    except Exception as e:
        logger.error(f"Error saving notes: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/delete_roadmap/{roadmap_id}")
async def delete_roadmap(
    roadmap_id: int,
    db: Session = Depends(get_db)
):
    """
    Delete a roadmap and all its nodes
    """
    try:
        roadmap = db.query(models.KnowledgeRoadmap).filter(
            models.KnowledgeRoadmap.id == roadmap_id
        ).first()
        
        if not roadmap:
            raise HTTPException(status_code=404, detail="Roadmap not found")
        
        # Delete all nodes recursively
        def delete_node_tree(node_id):
            children = db.query(models.KnowledgeNode).filter(
                models.KnowledgeNode.parent_node_id == node_id
            ).all()
            
            for child in children:
                delete_node_tree(child.id)
            
            db.query(models.NodeExplorationHistory).filter(
                models.NodeExplorationHistory.node_id == node_id
            ).delete()
            
            db.query(models.KnowledgeNode).filter(
                models.KnowledgeNode.id == node_id
            ).delete()
        
        if roadmap.root_node_id:
            delete_node_tree(roadmap.root_node_id)
        
        db.delete(roadmap)
        db.commit()
        
        return {
            "status": "success",
            "message": "Roadmap deleted successfully"
        }
        
    except Exception as e:
        logger.error(f"Error deleting roadmap: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

# ==================== HINTS ENDPOINT ====================



@app.post("/api/get_learning_hints")
async def get_learning_hints(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """
    Generate hints for missing learning points
    """
    try:
        review_id = payload.get("review_id")
        missing_points = payload.get("missing_points", [])

        if not review_id or not missing_points:
            raise HTTPException(status_code=400, detail="Review ID and missing points required")

        review = db.query(models.LearningReview).filter(
            models.LearningReview.id == review_id
        ).first()

        if not review:
            raise HTTPException(status_code=404, detail="Learning review not found")

        # Get original source content for context
        source_content = review.source_content[:3000]

        # Limit to first 3 missing points
        missing_points = missing_points[:3]

        # Generate hints for each missing point
        hints_list = []
        
        for point in missing_points:
            hint_prompt = f"""You are a helpful learning assistant providing subtle hints.

**CONTEXT** (original learning material):
{source_content}

**MISSING LEARNING POINT**:
{point}

**TASK**: Create a helpful hint that guides the student toward remembering this point WITHOUT directly giving away the answer.

**OUTPUT FORMAT** (JSON only):
{{
  "missing_point": "{point}",
  "hint": "A subtle clue that prompts memory without revealing the full answer",
  "memory_trigger": "A keyword or phrase that might jog their memory",
  "guiding_question": "A question that leads them to think about this topic"
}}

Generate hint now:"""

            chat_completion = groq_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "Learning assistant. Return JSON."},
                    {"role": "user", "content": hint_prompt}
                ],
                model=GROQ_MODEL,
                temperature=0.7,
                max_tokens=512,
            )

            response_text = chat_completion.choices[0].message.content

            try:
                json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
                if json_match:
                    hint_data = json.loads(json_match.group())
                    hints_list.append(hint_data)
            except:
                # Fallback hint if parsing fails
                hints_list.append({
                    "missing_point": point,
                    "hint": f"Think about the key concepts related to: {point[:50]}...",
                    "memory_trigger": "Review your notes",
                    "guiding_question": "What do you remember about this topic?"
                })

        return {
            "status": "success",
            "hints": hints_list,
            "total_hints": len(hints_list)
        }

    except Exception as e:
        logger.error(f"Error generating hints: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate hints: {str(e)}")


# ==================== GET LEARNING REVIEWS ====================

@app.get("/api/get_learning_reviews")
def get_learning_reviews(user_id: str = Query(...), db: Session = Depends(get_db)):
    """
    Get all learning reviews for a user
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        reviews = db.query(models.LearningReview).filter(
            models.LearningReview.user_id == user.id
        ).order_by(models.LearningReview.created_at.desc()).all()

        result = []
        for review in reviews:
            # Get session titles
            try:
                session_ids = json.loads(review.source_sessions)
                sessions = db.query(models.ChatSession).filter(
                    models.ChatSession.id.in_(session_ids)
                ).all()
                session_titles = [s.title for s in sessions]
            except:
                session_titles = []

            result.append({
                "id": review.id,
                "title": review.title,
                "status": review.status,
                "total_points": review.total_points,
                "best_score": round(review.best_score, 1),
                "attempt_count": review.attempt_count,
                "current_attempt": review.current_attempt,
                "session_titles": session_titles,
                "created_at": review.created_at.isoformat() + 'Z',
                "updated_at": review.updated_at.isoformat() + 'Z',
                "completed_at": review.completed_at.isoformat() + 'Z' if review.completed_at else None,
                "can_continue": review.status == "active" and review.current_attempt < 5
            })

        return {"reviews": result}

    except Exception as e:
        logger.error(f"Error getting learning reviews: {str(e)}")
        return {"reviews": []}
    
@app.post("/api/save_archetype_profile")
async def save_archetype_profile(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        user_id = payload.get("user_id")
        primary_archetype = payload.get("primary_archetype")
        secondary_archetype = payload.get("secondary_archetype")
        archetype_scores = payload.get("archetype_scores", {})
        archetype_description = payload.get("archetype_description", "")
        quiz_responses = payload.get("quiz_responses", {})

        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()

        if not comprehensive_profile:
            comprehensive_profile = models.ComprehensiveUserProfile(
                user_id=user.id,
                primary_archetype=primary_archetype,
                secondary_archetype=secondary_archetype,
                archetype_scores=json.dumps(archetype_scores),
                archetype_description=archetype_description,
                quiz_responses=json.dumps(quiz_responses)
            )
            db.add(comprehensive_profile)
        else:
            comprehensive_profile.primary_archetype = primary_archetype
            comprehensive_profile.secondary_archetype = secondary_archetype
            comprehensive_profile.archetype_scores = json.dumps(archetype_scores)
            comprehensive_profile.archetype_description = archetype_description
            comprehensive_profile.quiz_responses = json.dumps(quiz_responses)

        db.commit()
        db.refresh(comprehensive_profile)

        return {
            "status": "success",
            "message": "Archetype profile saved successfully",
            "primary_archetype": primary_archetype,
            "secondary_archetype": secondary_archetype
        }

    except Exception as e:
        logger.error(f"Error saving archetype: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to save archetype: {str(e)}")



@app.get("/api/get_comprehensive_profile")
async def get_comprehensive_profile(user_id: str = Query(...), db: Session = Depends(get_db)):
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()

        result = {
            "firstName": user.first_name or "",
            "lastName": user.last_name or "",
            "email": user.email or "",
            "fieldOfStudy": user.field_of_study or "",
            "brainwaveGoal": "",
            "preferredSubjects": [],
            "difficultyLevel": "intermediate",
            "learningPace": "moderate",
            "primaryArchetype": "",
            "secondaryArchetype": "",
            "archetypeDescription": ""
        }

        if comprehensive_profile:
            show_insights_value = comprehensive_profile.show_study_insights
            logger.info(f"📊 Loading profile - show_study_insights from DB: {show_insights_value} (type: {type(show_insights_value)})")
            
            result.update({
                "difficultyLevel": comprehensive_profile.difficulty_level or "intermediate",
                "learningPace": comprehensive_profile.learning_pace or "moderate",
                "brainwaveGoal": comprehensive_profile.brainwave_goal or "",
                "primaryArchetype": comprehensive_profile.primary_archetype or "",
                "secondaryArchetype": comprehensive_profile.secondary_archetype or "",
                "archetypeDescription": comprehensive_profile.archetype_description or "",
                "showStudyInsights": show_insights_value if show_insights_value is not None else True
            })

            try:
                if comprehensive_profile.preferred_subjects:
                    result["preferredSubjects"] = json.loads(comprehensive_profile.preferred_subjects)
            except:
                result["preferredSubjects"] = []
        else:
            result["showStudyInsights"] = True

        return result

    except Exception as e:
        logger.error(f"Error getting profile: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/update_comprehensive_profile")
async def update_comprehensive_profile(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        logger.info(f"📊 Received profile update payload: {payload}")
        
        user_id = payload.get("user_id")
        if not user_id:
            logger.error("📊 No user_id in payload!")
            raise HTTPException(status_code=400, detail="user_id is required")
            
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        if payload.get("firstName"):
            user.first_name = payload["firstName"]
        if payload.get("lastName"):
            user.last_name = payload["lastName"]
        if payload.get("email"):
            user.email = payload["email"]
        if payload.get("fieldOfStudy"):
            user.field_of_study = payload["fieldOfStudy"]

        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()

        if not comprehensive_profile:
            comprehensive_profile = models.ComprehensiveUserProfile(user_id=user.id)
            db.add(comprehensive_profile)

        if payload.get("difficultyLevel"):
            comprehensive_profile.difficulty_level = payload["difficultyLevel"]
        if payload.get("learningPace"):
            comprehensive_profile.learning_pace = payload["learningPace"]
        if payload.get("brainwaveGoal"):
            comprehensive_profile.brainwave_goal = payload["brainwaveGoal"]

        if "preferredSubjects" in payload:
            comprehensive_profile.preferred_subjects = json.dumps(payload["preferredSubjects"])
        
        # Handle showStudyInsights setting
        if "showStudyInsights" in payload:
            value = payload["showStudyInsights"]
            # Ensure it's a boolean
            if isinstance(value, str):
                value = value.lower() == 'true'
            comprehensive_profile.show_study_insights = bool(value)
            logger.info(f"📊 Setting show_study_insights to: {comprehensive_profile.show_study_insights} (received: {payload['showStudyInsights']})")

        comprehensive_profile.updated_at = datetime.now(timezone.utc)

        db.commit()

        return {
            "status": "success",
            "message": "Profile updated successfully"
        }

    except Exception as e:
        logger.error(f"Error updating profile: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/suggest_subjects")
async def suggest_subjects(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """AI-powered subject suggestions based on user input"""
    try:
        input_text = payload.get("input", "")
        college_level = payload.get("college_level", "")
        
        if len(input_text) < 2:
            return {"suggestions": []}
        
        # Use AI to generate relevant subject suggestions
        prompt = f"""Based on the input "{input_text}" and college level "{college_level}", suggest 5 relevant academic subjects or courses.

Return ONLY a JSON array of subject names, nothing else. Format: ["Subject 1", "Subject 2", ...]

Examples:
- Input "calc" → ["Calculus I", "Calculus II", "Calculus III", "Multivariable Calculus", "Differential Calculus"]
- Input "bio" → ["Biology", "Molecular Biology", "Cell Biology", "Microbiology", "Biochemistry"]
- Input "comp" → ["Computer Science", "Computer Architecture", "Computational Theory", "Computer Networks", "Computer Graphics"]

Input: "{input_text}"
College Level: "{college_level}"
Suggestions:"""

        response = call_ai(prompt, max_tokens=200, temperature=0.7)
        
        # Parse AI response
        try:
            # Try to extract JSON array from response
            import re
            json_match = re.search(r'\[.*\]', response, re.DOTALL)
            if json_match:
                suggestions = json.loads(json_match.group())
            else:
                # Fallback: split by newlines and clean
                suggestions = [s.strip().strip('"').strip("'").strip('-').strip() 
                             for s in response.split('\n') if s.strip()]
                suggestions = [s for s in suggestions if len(s) > 2 and len(s) < 50][:5]
        except:
            suggestions = []
        
        return {"suggestions": suggestions}
        
    except Exception as e:
        logger.error(f"Error generating subject suggestions: {str(e)}")
        return {"suggestions": []}
    
@app.post("/api/save_complete_profile")
async def save_complete_profile(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    try:
        user_id = payload.get("user_id")
        logger.info(f" save_complete_profile called for {user_id}")
        logger.info(f" Payload: quiz_completed={payload.get('quiz_completed')}, quiz_skipped={payload.get('quiz_skipped')}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            logger.error(f" User not found: {user_id}")
            raise HTTPException(status_code=404, detail="User not found")

        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()

        if not comprehensive_profile:
            logger.info(f"➕ Creating new comprehensive profile for user {user.id}")
            comprehensive_profile = models.ComprehensiveUserProfile(user_id=user.id)
            db.add(comprehensive_profile)
        else:
            logger.info(f" Updating existing comprehensive profile for user {user.id}")

        # New profile quiz fields
        if "is_college_student" in payload:
            comprehensive_profile.is_college_student = payload["is_college_student"]
        
        if "college_level" in payload:
            comprehensive_profile.college_level = payload["college_level"]
        
        if "major" in payload:
            comprehensive_profile.major = payload["major"]
        
        if "preferred_subjects" in payload:
            comprehensive_profile.preferred_subjects = json.dumps(payload["preferred_subjects"])
        
        if "main_subject" in payload:
            user.field_of_study = payload["main_subject"]
            comprehensive_profile.main_subject = payload["main_subject"]
        
        if "brainwave_goal" in payload:
            comprehensive_profile.brainwave_goal = payload["brainwave_goal"]
        
        if payload.get("quiz_completed"):
            comprehensive_profile.primary_archetype = payload.get("primary_archetype", "")
            comprehensive_profile.secondary_archetype = payload.get("secondary_archetype", "")
            comprehensive_profile.archetype_scores = json.dumps(payload.get("archetype_scores", {}))
            comprehensive_profile.archetype_description = payload.get("archetype_description", "")

        comprehensive_profile.quiz_completed = payload.get("quiz_completed", False)
        comprehensive_profile.quiz_skipped = payload.get("quiz_skipped", False)
        comprehensive_profile.updated_at = datetime.now(timezone.utc)

        db.flush()  # Ensure changes are flushed to database
        db.commit()
        db.refresh(comprehensive_profile)
        
        logger.info(f" Saved profile for {user_id}: quiz_completed={comprehensive_profile.quiz_completed}, quiz_skipped={comprehensive_profile.quiz_skipped}, profile_id={comprehensive_profile.id}")
        
        # Verify it was actually saved by querying again
        verify_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()
        logger.info(f" Verification query: quiz_skipped={verify_profile.quiz_skipped if verify_profile else 'NOT FOUND'}")
        
        # Generate initial proactive greeting message
        if payload.get("quiz_completed"):
            try:
                proactive_engine = get_proactive_ai_engine(unified_ai)
                user_profile = {
                    "first_name": user.first_name or "there",
                    "field_of_study": user.field_of_study or "General Studies",
                    "is_college_student": payload.get("is_college_student", True),
                    "college_level": payload.get("college_level", ""),
                    "subjects": payload.get("preferred_subjects", []),
                    "main_subject": payload.get("main_subject", ""),
                    "goal": payload.get("brainwave_goal", "")
                }
                
                # Generate personalized greeting
                greeting_prompt = f"""You are an AI tutor greeting {user_profile['first_name']} for the first time after they completed their profile.

Profile:
- College Student: {user_profile['is_college_student']}
- Level: {user_profile['college_level']}
- Main Subject: {user_profile['main_subject']}
- Goal: {user_profile['goal']}

Generate a warm, personalized greeting message (2-3 sentences) that:
1. Welcomes them by name
2. Acknowledges their specific subject and goal
3. Expresses excitement to help them succeed
4. Sounds natural and encouraging

Keep it brief and friendly."""

                greeting_message = call_ai(greeting_prompt, max_tokens=150, temperature=0.8)
                
                # Create initial chat session with greeting
                new_session = models.ChatSession(
                    user_id=user.id,
                    title="Welcome to Cerbyl",
                    created_at=datetime.now(timezone.utc),
                    updated_at=datetime.now(timezone.utc)
                )
                db.add(new_session)
                db.commit()
                db.refresh(new_session)
                
                # Save greeting message
                greeting_chat = models.ChatMessage(
                    chat_session_id=new_session.id,
                    user_id=user.id,
                    user_message="🎓 PROFILE_COMPLETED",
                    ai_response=greeting_message,
                    timestamp=datetime.now(timezone.utc)
                )
                db.add(greeting_chat)
                db.commit()
                
            except Exception as e:
                logger.error(f"Error generating greeting: {str(e)}")

        return {
            "status": "success",
            "message": "Profile saved successfully",
            "quiz_completed": comprehensive_profile.quiz_completed,
            "quiz_skipped": comprehensive_profile.quiz_skipped
        }

    except Exception as e:
        logger.error(f"Error saving complete profile: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# ==================== FRIEND SYSTEM ENDPOINTS ====================

@app.get("/api/search_users")
async def search_users(
    query: str = Query(..., min_length=1),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Search for users by username or email"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="Current user not found")
        
        # Search for users (excluding current user)
        search_pattern = f"%{query}%"
        users = db.query(models.User).filter(
            and_(
                models.User.id != current_user.id,
                (models.User.username.ilike(search_pattern) | models.User.email.ilike(search_pattern))
            )
        ).limit(20).all()
        
        result = []
        for user in users:
            # Get comprehensive profile for preferred subjects
            comp_profile = db.query(models.ComprehensiveUserProfile).filter(
                models.ComprehensiveUserProfile.user_id == user.id
            ).first()
            
            # Get user stats
            user_stats = db.query(models.UserStats).filter(
                models.UserStats.user_id == user.id
            ).first()
            
            # Check friendship status
            friendship = db.query(models.Friendship).filter(
                and_(
                    models.Friendship.user_id == current_user.id,
                    models.Friendship.friend_id == user.id
                )
            ).first()
            
            # Check if there's a pending friend request
            pending_request_sent = db.query(models.FriendRequest).filter(
                and_(
                    models.FriendRequest.sender_id == current_user.id,
                    models.FriendRequest.receiver_id == user.id,
                    models.FriendRequest.status == "pending"
                )
            ).first()
            
            pending_request_received = db.query(models.FriendRequest).filter(
                and_(
                    models.FriendRequest.sender_id == user.id,
                    models.FriendRequest.receiver_id == current_user.id,
                    models.FriendRequest.status == "pending"
                )
            ).first()
            
            preferred_subjects = []
            if comp_profile and comp_profile.preferred_subjects:
                try:
                    preferred_subjects = json.loads(comp_profile.preferred_subjects)
                except:
                    preferred_subjects = []
            
            result.append({
                "id": user.id,
                "username": user.username,
                "email": user.email,
                "first_name": user.first_name or "",
                "last_name": user.last_name or "",
                "picture_url": user.picture_url or "",
                "field_of_study": user.field_of_study or "",
                "preferred_subjects": preferred_subjects,
                "stats": {
                    "total_lessons": user_stats.total_lessons if user_stats else 0,
                    "total_hours": round(user_stats.total_hours, 1) if user_stats else 0,
                    "day_streak": user_stats.day_streak if user_stats else 0,
                    "accuracy_percentage": round(user_stats.accuracy_percentage, 1) if user_stats else 0
                },
                "is_friend": friendship is not None,
                "request_sent": pending_request_sent is not None,
                "request_received": pending_request_received is not None
            })
        
        return {"users": result}
    
    except Exception as e:
        logger.error(f"Error searching users: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/send_friend_request")
async def send_friend_request(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Send a friend request to another user"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="Current user not found")
        
        receiver_id = payload.get("receiver_id")
        if not receiver_id:
            raise HTTPException(status_code=400, detail="receiver_id is required")
        
        # Check if receiver exists
        receiver = db.query(models.User).filter(models.User.id == receiver_id).first()
        if not receiver:
            raise HTTPException(status_code=404, detail="Receiver not found")
        
        # Check if they're already friends
        existing_friendship = db.query(models.Friendship).filter(
            and_(
                models.Friendship.user_id == current_user.id,
                models.Friendship.friend_id == receiver_id
            )
        ).first()
        
        if existing_friendship:
            raise HTTPException(status_code=400, detail="Already friends")
        
        # Check if there's already a pending request
        existing_request = db.query(models.FriendRequest).filter(
            and_(
                models.FriendRequest.sender_id == current_user.id,
                models.FriendRequest.receiver_id == receiver_id,
                models.FriendRequest.status == "pending"
            )
        ).first()
        
        if existing_request:
            raise HTTPException(status_code=400, detail="Friend request already sent")
        
        # Create friend request
        friend_request = models.FriendRequest(
            sender_id=current_user.id,
            receiver_id=receiver_id,
            status="pending"
        )
        db.add(friend_request)
        
        # Create notification for the receiver
        notification = models.Notification(
            user_id=receiver_id,
            title="New Friend Request",
            message=f"{current_user.username} wants to be your friend!",
            notification_type="friend_request",
            is_read=False
        )
        db.add(notification)
        db.commit()
        
        return {
            "status": "success",
            "message": "Friend request sent successfully"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error sending friend request: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/friend_requests")
async def get_friend_requests(
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get all pending friend requests for the current user"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="Current user not found")
        
        # Get received requests
        received_requests = db.query(models.FriendRequest).filter(
            and_(
                models.FriendRequest.receiver_id == current_user.id,
                models.FriendRequest.status == "pending"
            )
        ).all()
        
        # Get sent requests
        sent_requests = db.query(models.FriendRequest).filter(
            and_(
                models.FriendRequest.sender_id == current_user.id,
                models.FriendRequest.status == "pending"
            )
        ).all()
        
        received_result = []
        for req in received_requests:
            sender = req.sender
            received_result.append({
                "request_id": req.id,
                "user_id": sender.id,
                "username": sender.username,
                "email": sender.email,
                "first_name": sender.first_name or "",
                "last_name": sender.last_name or "",
                "picture_url": sender.picture_url or "",
                "created_at": req.created_at.isoformat() + 'Z'
            })
        
        sent_result = []
        for req in sent_requests:
            receiver = req.receiver
            sent_result.append({
                "request_id": req.id,
                "user_id": receiver.id,
                "username": receiver.username,
                "email": receiver.email,
                "first_name": receiver.first_name or "",
                "last_name": receiver.last_name or "",
                "picture_url": receiver.picture_url or "",
                "created_at": req.created_at.isoformat() + 'Z'
            })
        
        return {
            "received": received_result,
            "sent": sent_result
        }
    
    except Exception as e:
        logger.error(f"Error getting friend requests: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/respond_friend_request")
async def respond_friend_request(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Accept or reject a friend request"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="Current user not found")
        
        request_id = payload.get("request_id")
        action = payload.get("action")  # "accept" or "reject"
        
        if not request_id or not action:
            raise HTTPException(status_code=400, detail="request_id and action are required")
        
        # Get the friend request
        friend_request = db.query(models.FriendRequest).filter(
            and_(
                models.FriendRequest.id == request_id,
                models.FriendRequest.receiver_id == current_user.id,
                models.FriendRequest.status == "pending"
            )
        ).first()
        
        if not friend_request:
            raise HTTPException(status_code=404, detail="Friend request not found")
        
        if action == "accept":
            # Update request status
            friend_request.status = "accepted"
            friend_request.responded_at = datetime.now(timezone.utc)
            
            # Create friendship (bidirectional)
            friendship1 = models.Friendship(
                user_id=current_user.id,
                friend_id=friend_request.sender_id
            )
            friendship2 = models.Friendship(
                user_id=friend_request.sender_id,
                friend_id=current_user.id
            )
            
            db.add(friendship1)
            db.add(friendship2)
            
            # Create notification for the sender that their request was accepted
            notification = models.Notification(
                user_id=friend_request.sender_id,
                title="Friend Request Accepted!",
                message=f"{current_user.username} accepted your friend request. You're now friends!",
                notification_type="friend_accepted",
                is_read=False
            )
            db.add(notification)
            db.commit()
            
            return {
                "status": "success",
                "message": "Friend request accepted"
            }
        
        elif action == "reject":
            friend_request.status = "rejected"
            friend_request.responded_at = datetime.now(timezone.utc)
            db.commit()
            
            return {
                "status": "success",
                "message": "Friend request rejected"
            }
        
        else:
            raise HTTPException(status_code=400, detail="Invalid action")
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error responding to friend request: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/friends")
async def get_friends(
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get all friends of the current user"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="Current user not found")
        
        # Get all friendships
        friendships = db.query(models.Friendship).filter(
            models.Friendship.user_id == current_user.id
        ).all()
        
        result = []
        for friendship in friendships:
            friend = friendship.friend
            
            # Get comprehensive profile
            comp_profile = db.query(models.ComprehensiveUserProfile).filter(
                models.ComprehensiveUserProfile.user_id == friend.id
            ).first()
            
            # Get user stats
            user_stats = db.query(models.UserStats).filter(
                models.UserStats.user_id == friend.id
            ).first()
            
            preferred_subjects = []
            if comp_profile and comp_profile.preferred_subjects:
                try:
                    preferred_subjects = json.loads(comp_profile.preferred_subjects)
                except:
                    preferred_subjects = []
            
            result.append({
                "id": friend.id,
                "username": friend.username,
                "email": friend.email,
                "first_name": friend.first_name or "",
                "last_name": friend.last_name or "",
                "picture_url": friend.picture_url or "",
                "field_of_study": friend.field_of_study or "",
                "preferred_subjects": preferred_subjects,
                "stats": {
                    "total_lessons": user_stats.total_lessons if user_stats else 0,
                    "total_hours": round(user_stats.total_hours, 1) if user_stats else 0,
                    "day_streak": user_stats.day_streak if user_stats else 0,
                    "accuracy_percentage": round(user_stats.accuracy_percentage, 1) if user_stats else 0
                },
                "friends_since": friendship.created_at.isoformat() + 'Z'
            })
        
        return {"friends": result}
    
    except Exception as e:
        logger.error(f"Error getting friends: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/remove_friend")
async def remove_friend(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Remove a friend"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="Current user not found")
        
        friend_id = payload.get("friend_id")
        if not friend_id:
            raise HTTPException(status_code=400, detail="friend_id is required")
        
        # Delete both friendship records
        db.query(models.Friendship).filter(
            and_(
                models.Friendship.user_id == current_user.id,
                models.Friendship.friend_id == friend_id
            )
        ).delete()
        
        db.query(models.Friendship).filter(
            and_(
                models.Friendship.user_id == friend_id,
                models.Friendship.friend_id == current_user.id
            )
        ).delete()
        
        db.commit()
        
        return {
            "status": "success",
            "message": "Friend removed successfully"
        }
    
    except Exception as e:
        logger.error(f"Error removing friend: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# ==================== FRIEND ACTIVITY FEED ====================

@app.get("/api/friend_activity_feed")
async def get_friend_activity_feed(
    username: str = Depends(verify_token),
    db: Session = Depends(get_db),
    limit: int = Query(50, le=100)
):
    """Get activity feed of friends' achievements and milestones"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get friend IDs
        friendships = db.query(models.Friendship).filter(
            models.Friendship.user_id == current_user.id
        ).all()
        friend_ids = [f.friend_id for f in friendships]
        
        if not friend_ids:
            return {"activities": []}
        
        # Get friend activities
        activities = db.query(models.FriendActivity).filter(
            models.FriendActivity.user_id.in_(friend_ids)
        ).order_by(models.FriendActivity.created_at.desc()).limit(limit).all()
        
        result = []
        for activity in activities:
            # Get kudos count and check if current user gave kudos
            kudos_count = db.query(models.Kudos).filter(
                models.Kudos.activity_id == activity.id
            ).count()
            
            user_gave_kudos = db.query(models.Kudos).filter(
                and_(
                    models.Kudos.activity_id == activity.id,
                    models.Kudos.user_id == current_user.id
                )
            ).first() is not None
            
            result.append({
                "id": activity.id,
                "user": {
                    "id": activity.user.id,
                    "username": activity.user.username,
                    "first_name": activity.user.first_name or "",
                    "last_name": activity.user.last_name or "",
                    "picture_url": activity.user.picture_url or ""
                },
                "activity_type": activity.activity_type,
                "title": activity.title,
                "description": activity.description or "",
                "icon": activity.icon or "Trophy",
                "metadata": json.loads(activity.activity_data) if activity.activity_data else {},
                "kudos_count": kudos_count,
                "user_gave_kudos": user_gave_kudos,
                "created_at": activity.created_at.isoformat() + 'Z'
            })
        
        return {"activities": result}
    
    except Exception as e:
        logger.error(f"Error fetching friend activity feed: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/give_kudos")
async def give_kudos(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Give kudos to a friend's activity"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        activity_id = payload.get("activity_id")
        reaction_type = payload.get("reaction_type", "👏")
        
        if not activity_id:
            raise HTTPException(status_code=400, detail="activity_id is required")
        
        # Check if already gave kudos
        existing = db.query(models.Kudos).filter(
            and_(
                models.Kudos.activity_id == activity_id,
                models.Kudos.user_id == current_user.id
            )
        ).first()
        
        if existing:
            # Remove kudos
            db.delete(existing)
            db.commit()
            return {"status": "removed", "message": "Kudos removed"}
        else:
            # Add kudos
            kudos = models.Kudos(
                activity_id=activity_id,
                user_id=current_user.id,
                reaction_type=reaction_type
            )
            db.add(kudos)
            db.commit()
            return {"status": "added", "message": "Kudos given"}
    
    except Exception as e:
        logger.error(f"Error giving kudos: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/create_activity")
async def create_activity(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a new activity (used internally when user achieves something)"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        activity = models.FriendActivity(
            user_id=current_user.id,
            activity_type=payload.get("activity_type"),
            title=payload.get("title"),
            description=payload.get("description", ""),
            icon=payload.get("icon", "Trophy"),
            activity_data=json.dumps(payload.get("metadata", {}))
        )
        
        db.add(activity)
        db.commit()
        
        return {"status": "success", "activity_id": activity.id}
    
    except Exception as e:
        logger.error(f"Error creating activity: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# ==================== LEADERBOARDS ====================

@app.get("/api/leaderboard")
async def get_leaderboard(
    category: str = Query("global", pattern="^(global|friends|subject|archetype)$"),
    metric: str = Query("total_hours", pattern="^(total_hours|accuracy|streak|lessons)$"),
    period: str = Query("all_time", pattern="^(weekly|monthly|all_time)$"),
    subject: Optional[str] = Query(None),
    limit: int = Query(50, le=100),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get leaderboard rankings"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Determine user pool based on category
        if category == "friends":
            friendships = db.query(models.Friendship).filter(
                models.Friendship.user_id == current_user.id
            ).all()
            user_ids = [f.friend_id for f in friendships] + [current_user.id]
        else:
            user_ids = None  # All users
        
        # Get user stats
        query = db.query(models.UserStats, models.User).join(
            models.User, models.UserStats.user_id == models.User.id
        )
        
        if user_ids:
            query = query.filter(models.UserStats.user_id.in_(user_ids))
        
        # Order by metric
        if metric == "total_hours":
            query = query.order_by(models.UserStats.total_hours.desc())
            score_field = "total_hours"
        elif metric == "accuracy":
            query = query.order_by(models.UserStats.accuracy_percentage.desc())
            score_field = "accuracy_percentage"
        elif metric == "streak":
            query = query.order_by(models.UserStats.day_streak.desc())
            score_field = "day_streak"
        else:  # lessons
            query = query.order_by(models.UserStats.total_lessons.desc())
            score_field = "total_lessons"
        
        results = query.limit(limit).all()
        
        leaderboard = []
        for rank, (stats, user) in enumerate(results, start=1):
            score = getattr(stats, score_field)
            leaderboard.append({
                "rank": rank,
                "user_id": user.id,
                "username": user.username,
                "first_name": user.first_name or "",
                "last_name": user.last_name or "",
                "picture_url": user.picture_url or "",
                "score": round(score, 1) if isinstance(score, float) else score,
                "metric": metric,
                "is_current_user": user.id == current_user.id
            })
        
        # Find current user's rank if not in top results
        current_user_rank = next((item for item in leaderboard if item["is_current_user"]), None)
        
        return {
            "leaderboard": leaderboard,
            "current_user_rank": current_user_rank,
            "category": category,
            "metric": metric,
            "period": period
        }
    
    except Exception as e:
        logger.error(f"Error fetching leaderboard: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== QUIZ BATTLES ====================

@app.post("/api/create_quiz_battle")
async def create_quiz_battle(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a new quiz battle challenge and notify opponent"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        opponent_id = payload.get("opponent_id")
        subject = payload.get("subject")
        difficulty = payload.get("difficulty", "intermediate")
        question_count = payload.get("question_count", 10)
        time_limit = payload.get("time_limit_seconds", 300)
        
        if not opponent_id or not subject:
            raise HTTPException(status_code=400, detail="opponent_id and subject are required")
        
        # Check friendship
        friendship = db.query(models.Friendship).filter(
            and_(
                models.Friendship.user_id == current_user.id,
                models.Friendship.friend_id == opponent_id
            )
        ).first()
        
        if not friendship:
            raise HTTPException(status_code=400, detail="Can only battle with friends")
        
        # Create battle
        battle = models.QuizBattle(
            challenger_id=current_user.id,
            opponent_id=opponent_id,
            subject=subject,
            difficulty=difficulty,
            question_count=question_count,
            time_limit_seconds=time_limit,
            expires_at=datetime.now(timezone.utc) + timedelta(days=7)
        )
        
        db.add(battle)
        db.commit()
        db.refresh(battle)
        
        logger.info(f" Battle created: ID={battle.id}")
        
        # Create persistent notification for opponent
        battle_notification = models.Notification(
            user_id=opponent_id,
            title="Quiz Battle Challenge",
            message=f"{current_user.username} has challenged you to a quiz battle on {subject}!",
            notification_type="battle_challenge",
            is_read=False
        )
        db.add(battle_notification)
        db.commit()
        
        # Prepare notification data
        battle_data = {
            "id": battle.id,
            "subject": battle.subject,
            "difficulty": battle.difficulty,
            "question_count": battle.question_count,
            "time_limit_seconds": battle.time_limit_seconds,
            "challenger": {
                "id": current_user.id,
                "username": current_user.username,
                "first_name": current_user.first_name or "",
                "last_name": current_user.last_name or "",
                "picture_url": current_user.picture_url or ""
            },
            "is_challenger": False
        }
        
        # Send WebSocket notification
        notification_sent = await notify_battle_challenge(opponent_id, battle_data)
        
        if notification_sent:
            logger.info(f" Notification sent to opponent {opponent_id}")
        else:
            logger.warning(f" Opponent {opponent_id} not connected to WebSocket - notification not sent")
        
        # Log active connections for debugging
        logger.info(f"📊 Active WebSocket connections: {list(manager.active_connections.keys())}")
        
        return {
            "status": "success",
            "battle_id": battle.id,
            "message": "Quiz battle created",
            "notification_sent": notification_sent,
            "opponent_connected": notification_sent
        }
    
    except Exception as e:
        logger.error(f" Error creating battle: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/quiz_battles")
async def get_quiz_battles(
    username: str = Depends(verify_token),
    db: Session = Depends(get_db),
    status: str = Query("active", pattern="^(pending|active|completed|all)$")
):
    """Get user's quiz battles"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get battles where user is challenger or opponent
        query = db.query(models.QuizBattle).filter(
            (models.QuizBattle.challenger_id == current_user.id) |
            (models.QuizBattle.opponent_id == current_user.id)
        )
        
        if status != "all":
            query = query.filter(models.QuizBattle.status == status)
        
        battles = query.order_by(models.QuizBattle.created_at.desc()).all()
        
        result = []
        for battle in battles:
            is_challenger = battle.challenger_id == current_user.id
            opponent = battle.opponent if is_challenger else battle.challenger
            
            result.append({
                "id": battle.id,
                "opponent": {
                    "id": opponent.id,
                    "username": opponent.username,
                    "first_name": opponent.first_name or "",
                    "last_name": opponent.last_name or "",
                    "picture_url": opponent.picture_url or ""
                },
                "subject": battle.subject,
                "difficulty": battle.difficulty,
                "status": battle.status,
                "question_count": battle.question_count,
                "time_limit_seconds": battle.time_limit_seconds,
                "your_score": battle.challenger_score if is_challenger else battle.opponent_score,
                "opponent_score": battle.opponent_score if is_challenger else battle.challenger_score,
                "your_completed": battle.challenger_completed if is_challenger else battle.opponent_completed,
                "opponent_completed": battle.opponent_completed if is_challenger else battle.challenger_completed,
                "is_challenger": is_challenger,
                "created_at": battle.created_at.isoformat() + 'Z',
                "expires_at": battle.expires_at.isoformat() + 'Z' if battle.expires_at else None
            })
        
        return {"battles": result}
    
    except Exception as e:
        logger.error(f"Error fetching quiz battles: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/complete_quiz_battle")
async def complete_quiz_battle(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Complete a quiz battle with score"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        battle_id = payload.get("battle_id")
        score = payload.get("score")
        answers = payload.get("answers", [])
        
        if not battle_id or score is None:
            raise HTTPException(status_code=400, detail="battle_id and score are required")
        
        battle = db.query(models.QuizBattle).filter(
            models.QuizBattle.id == battle_id
        ).first()
        
        if not battle:
            raise HTTPException(status_code=404, detail="Battle not found")
        
        # Update score and store answers
        is_challenger = battle.challenger_id == current_user.id
        if is_challenger:
            battle.challenger_score = score
            battle.challenger_completed = True
            battle.challenger_answers = json.dumps(answers)
        else:
            battle.opponent_score = score
            battle.opponent_completed = True
            battle.opponent_answers = json.dumps(answers)
        
        # Get opponent ID for notification
        opponent_id = battle.opponent_id if is_challenger else battle.challenger_id
        
        # Check if both completed
        if battle.challenger_completed and battle.opponent_completed:
            battle.status = "completed"
            battle.completed_at = datetime.now(timezone.utc)
            
            # Create activity for winner
            winner_id = battle.challenger_id if battle.challenger_score > battle.opponent_score else battle.opponent_id
            winner = battle.challenger if winner_id == battle.challenger_id else battle.opponent
            loser = battle.opponent if winner_id == battle.challenger_id else battle.challenger
            
            winner_score = battle.challenger_score if winner_id == battle.challenger_id else battle.opponent_score
            loser_score = battle.opponent_score if winner_id == battle.challenger_id else battle.challenger_score
            
            activity = models.FriendActivity(
                user_id=winner_id,
                activity_type="quiz_battle_won",
                title=f"Won Quiz Battle!",
                description=f"Defeated {loser.username} in {battle.subject}",
                icon="Swords",
                activity_data=json.dumps({
                    "winner_score": winner_score,
                    "loser_score": loser_score,
                    "subject": battle.subject
                })
            )
            db.add(activity)
            
            # Create notifications for both players
            # Calculate percentage scores
            total_questions = battle.question_count or 10
            winner_percentage = round((winner_score / total_questions) * 100) if total_questions > 0 else 0
            loser_percentage = round((loser_score / total_questions) * 100) if total_questions > 0 else 0
            
            # Winner notification
            winner_notification = models.Notification(
                user_id=winner_id,
                title="Battle Victory",
                message=f"You won the quiz battle against {loser.username}! Score: {winner_score}/{total_questions} ({winner_percentage}%)",
                notification_type="battle_won"
            )
            db.add(winner_notification)
            
            # Loser notification (encouraging)
            loser_notification = models.Notification(
                user_id=loser.id,
                title="Battle Complete",
                message=f"Good effort! You scored {loser_score}/{total_questions} ({loser_percentage}%) against {winner.username}. Practice and challenge them again!",
                notification_type="battle_lost"
            )
            db.add(loser_notification)
        elif battle.status == "pending":
            battle.status = "active"
            battle.started_at = datetime.now(timezone.utc)
        
        db.commit()
        
        # Notify opponent that user completed
        await manager.send_personal_message({
            "type": "battle_opponent_completed",
            "battle_id": battle.id,
            "opponent_completed": True
        }, opponent_id)
        
        # If both completed, notify both users
        if battle.challenger_completed and battle.opponent_completed:
            await notify_battle_completed([battle.challenger_id, battle.opponent_id], battle.id, 
                                         battle.challenger_id if battle.challenger_score > battle.opponent_score else battle.opponent_id)
        
        return {
            "status": "success",
            "battle_status": battle.status,
            "message": "Score submitted",
            "both_completed": battle.challenger_completed and battle.opponent_completed
        }
    
    except Exception as e:
        logger.error(f"Error completing quiz battle: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# ==================== CHALLENGES ====================

@app.post("/api/create_challenge")
async def create_challenge(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a new challenge"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        challenge = models.Challenge(
            creator_id=current_user.id,
            title=payload.get("title"),
            description=payload.get("description", ""),
            challenge_type=payload.get("challenge_type"),
            subject=payload.get("subject"),
            target_metric=payload.get("target_metric"),
            target_value=payload.get("target_value"),
            time_limit_minutes=payload.get("time_limit_minutes"),
            starts_at=datetime.now(timezone.utc),
            ends_at=datetime.now(timezone.utc) + timedelta(minutes=payload.get("time_limit_minutes", 60))
        )
        
        db.add(challenge)
        db.commit()
        
        return {
            "status": "success",
            "challenge_id": challenge.id
        }
    
    except Exception as e:
        logger.error(f"Error creating challenge: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/challenges")
async def get_challenges(
    username: str = Depends(verify_token),
    db: Session = Depends(get_db),
    filter_type: str = Query("active", pattern="^(active|completed|my_challenges|all)$")
):
    """Get available challenges"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        query = db.query(models.Challenge)
        
        if filter_type == "my_challenges":
            query = query.filter(models.Challenge.creator_id == current_user.id)
        elif filter_type != "all":
            query = query.filter(models.Challenge.status == filter_type)
        
        challenges = query.order_by(models.Challenge.created_at.desc()).all()
        
        result = []
        for challenge in challenges:
            # Check if user is participating
            participation = db.query(models.ChallengeParticipation).filter(
                and_(
                    models.ChallengeParticipation.challenge_id == challenge.id,
                    models.ChallengeParticipation.user_id == current_user.id
                )
            ).first()
            
            result.append({
                "id": challenge.id,
                "creator": {
                    "id": challenge.creator.id,
                    "username": challenge.creator.username,
                    "first_name": challenge.creator.first_name or "",
                    "last_name": challenge.creator.last_name or ""
                },
                "title": challenge.title,
                "description": challenge.description or "",
                "challenge_type": challenge.challenge_type,
                "subject": challenge.subject or "",
                "target_metric": challenge.target_metric,
                "target_value": challenge.target_value,
                "time_limit_minutes": challenge.time_limit_minutes,
                "status": challenge.status,
                "participant_count": challenge.participant_count,
                "is_participating": participation is not None,
                "user_progress": participation.progress if participation else 0,
                "user_completed": participation.completed if participation else False,
                "created_at": challenge.created_at.isoformat() + 'Z',
                "ends_at": challenge.ends_at.isoformat() + 'Z' if challenge.ends_at else None
            })
        
        return {"challenges": result}
    
    except Exception as e:
        logger.error(f"Error fetching challenges: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/join_challenge")
async def join_challenge(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Join a challenge"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        challenge_id = payload.get("challenge_id")
        if not challenge_id:
            raise HTTPException(status_code=400, detail="challenge_id is required")
        
        # Check if already participating
        existing = db.query(models.ChallengeParticipation).filter(
            and_(
                models.ChallengeParticipation.challenge_id == challenge_id,
                models.ChallengeParticipation.user_id == current_user.id
            )
        ).first()
        
        if existing:
            raise HTTPException(status_code=400, detail="Already participating in this challenge")
        
        participation = models.ChallengeParticipation(
            challenge_id=challenge_id,
            user_id=current_user.id
        )
        
        db.add(participation)
        
        # Update participant count
        challenge = db.query(models.Challenge).filter(models.Challenge.id == challenge_id).first()
        if challenge:
            challenge.participant_count += 1
        
        db.commit()
        
        return {"status": "success", "message": "Joined challenge"}
    
    except Exception as e:
        logger.error(f"Error joining challenge: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/quiz_battle/{battle_id}")
async def get_quiz_battle_detail(
    battle_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get detailed information about a specific battle"""
    try:
        battle = db.query(models.QuizBattle).filter(
            models.QuizBattle.id == battle_id
        ).first()
        
        if not battle:
            raise HTTPException(status_code=404, detail="Battle not found")
        
        # Check if user is part of this battle
        if battle.challenger_id != current_user.id and battle.opponent_id != current_user.id:
            raise HTTPException(status_code=403, detail="Not authorized to view this battle")
        
        # Get existing questions if any
        questions = db.query(models.BattleQuestion).filter(
            models.BattleQuestion.battle_id == battle_id
        ).all()
        
        question_list = []
        for q in questions:
            question_list.append({
                "id": q.id,
                "question": q.question,
                "options": json.loads(q.options),
                "correct_answer": q.correct_answer,
                "explanation": q.explanation
            })
        
        is_challenger = battle.challenger_id == current_user.id
        
        # Get opponent user
        try:
            opponent_id = battle.opponent_id if is_challenger else battle.challenger_id
            opponent = db.query(models.User).filter(models.User.id == opponent_id).first()
            if not opponent:
                raise HTTPException(status_code=404, detail="Opponent not found")
        except Exception as e:
            logger.error(f"Error getting opponent: {str(e)}")
            raise HTTPException(status_code=500, detail="Error loading battle opponent")
        
        battle_data = {
            "id": battle.id,
            "subject": battle.subject,
            "difficulty": battle.difficulty,
            "status": battle.status,
            "question_count": battle.question_count,
            "time_limit_seconds": battle.time_limit_seconds,
            "your_score": battle.challenger_score if is_challenger else battle.opponent_score,
            "opponent_score": battle.opponent_score if is_challenger else battle.challenger_score,
            "your_completed": battle.challenger_completed if is_challenger else battle.opponent_completed,
            "opponent_completed": battle.opponent_completed if is_challenger else battle.challenger_completed,
            "is_challenger": is_challenger,
            "opponent": {
                "id": opponent.id,
                "username": opponent.username,
                "first_name": opponent.first_name or "",
                "last_name": opponent.last_name or "",
                "picture_url": opponent.picture_url or ""
            }
        }
        
        # Include answers if both completed
        if battle.challenger_completed and battle.opponent_completed:
            try:
                your_answers = json.loads(battle.challenger_answers if is_challenger else battle.opponent_answers) if (battle.challenger_answers if is_challenger else battle.opponent_answers) else []
                opponent_answers = json.loads(battle.opponent_answers if is_challenger else battle.challenger_answers) if (battle.opponent_answers if is_challenger else battle.challenger_answers) else []
                battle_data["your_answers"] = your_answers
                battle_data["opponent_answers"] = opponent_answers
            except:
                pass
        
        return {
            "battle": battle_data,
            "questions": question_list
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting battle detail: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate_battle_questions")
async def generate_battle_questions(
    payload: dict,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Generate AI questions for a quiz battle"""
    try:
        battle_id = payload.get("battle_id")
        subject = payload.get("subject")
        difficulty = payload.get("difficulty", "intermediate")
        question_count = payload.get("question_count", 10)
        
        # Verify battle exists and user is participant
        battle = db.query(models.QuizBattle).filter(
            models.QuizBattle.id == battle_id
        ).first()
        
        if not battle:
            raise HTTPException(status_code=404, detail="Battle not found")
        
        if battle.challenger_id != current_user.id and battle.opponent_id != current_user.id:
            raise HTTPException(status_code=403, detail="Not authorized")
        
        # Check if questions already exist
        existing = db.query(models.BattleQuestion).filter(
            models.BattleQuestion.battle_id == battle_id
        ).first()
        
        if existing:
            # Return existing questions
            questions = db.query(models.BattleQuestion).filter(
                models.BattleQuestion.battle_id == battle_id
            ).all()
            
            return {
                "questions": [{
                    "id": q.id,
                    "question": q.question,
                    "options": json.loads(q.options),
                    "correct_answer": q.correct_answer,
                    "explanation": q.explanation
                } for q in questions]
            }
        
        # Generate new questions using AI
        difficulty_map = {
            "beginner": "easy, suitable for beginners",
            "intermediate": "moderate difficulty",
            "advanced": "challenging, advanced level"
        }
        
        prompt = f"""Generate exactly {question_count} multiple choice questions about {subject}.
Difficulty level: {difficulty_map.get(difficulty, 'moderate')}.

Return ONLY a valid JSON array with this exact structure:
[
  {{
    "question": "Question text here?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct_answer": 0,
    "explanation": "Brief explanation of the correct answer"
  }}
]

Requirements:
- Each question must have exactly 4 options
- correct_answer must be 0, 1, 2, or 3 (index of the correct option)
- Questions should be clear and unambiguous
- Explanations should be concise (1-2 sentences)
- Make questions engaging and educational
- Return ONLY the JSON array, no additional text"""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=4000
        )
        
        content = response.choices[0].message.content.strip()
        
        # Clean the response
        if content.startswith("```json"):
            content = content[7:]
        if content.startswith("```"):
            content = content[3:]
        if content.endswith("```"):
            content = content[:-3]
        content = content.strip()
        
        questions_data = json.loads(content)
        
        # Save questions to database
        saved_questions = []
        for q_data in questions_data:
            battle_question = models.BattleQuestion(
                battle_id=battle_id,
                question=q_data["question"],
                options=json.dumps(q_data["options"]),
                correct_answer=q_data["correct_answer"],
                explanation=q_data.get("explanation", "")
            )
            db.add(battle_question)
            db.flush()
            
            saved_questions.append({
                "id": battle_question.id,
                "question": battle_question.question,
                "options": q_data["options"],
                "correct_answer": battle_question.correct_answer,
                "explanation": battle_question.explanation
            })
        
        # Update battle status to active
        if battle.status == "pending":
            battle.status = "active"
            battle.started_at = datetime.utcnow()
        
        db.commit()
        
        return {"questions": saved_questions}
    
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {str(e)}, Content: {content}")
        raise HTTPException(status_code=500, detail="Failed to parse AI response")
    except Exception as e:
        db.rollback()
        logger.error(f"Error generating battle questions: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== CHALLENGE ENDPOINTS ====================

@app.get("/api/challenge/{challenge_id}")
async def get_challenge_detail(
    challenge_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get detailed information about a specific challenge"""
    try:
        challenge = db.query(models.Challenge).filter(
            models.Challenge.id == challenge_id
        ).first()
        
        if not challenge:
            raise HTTPException(status_code=404, detail="Challenge not found")
        
        # Get user's participation
        participation = db.query(models.ChallengeParticipation).filter(
            and_(
                models.ChallengeParticipation.challenge_id == challenge_id,
                models.ChallengeParticipation.user_id == current_user.id
            )
        ).first()
        
        if not participation:
            raise HTTPException(status_code=403, detail="Not participating in this challenge")
        
        # Get existing questions if any
        questions = db.query(models.ChallengeQuestion).filter(
            models.ChallengeQuestion.challenge_id == challenge_id
        ).all()
        
        question_list = []
        for q in questions:
            question_list.append({
                "id": q.id,
                "question": q.question,
                "options": json.loads(q.options),
                "correct_answer": q.correct_answer,
                "explanation": q.explanation
            })
        
        return {
            "challenge": {
                "id": challenge.id,
                "title": challenge.title,
                "description": challenge.description,
                "challenge_type": challenge.challenge_type,
                "subject": challenge.subject,
                "target_metric": challenge.target_metric,
                "target_value": challenge.target_value,
                "time_limit_minutes": challenge.time_limit_minutes,
                "status": challenge.status,
                "progress": participation.progress,
                "completed": participation.completed
            },
            "questions": question_list
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting challenge detail: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate_challenge_questions")
async def generate_challenge_questions(
    payload: dict,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Generate AI questions for a challenge"""
    try:
        challenge_id = payload.get("challenge_id")
        subject = payload.get("subject", "General Knowledge")
        challenge_type = payload.get("challenge_type", "speed")
        question_count = payload.get("question_count", 10)
        
        # Verify challenge exists and user is participant
        challenge = db.query(models.Challenge).filter(
            models.Challenge.id == challenge_id
        ).first()
        
        if not challenge:
            raise HTTPException(status_code=404, detail="Challenge not found")
        
        participation = db.query(models.ChallengeParticipation).filter(
            and_(
                models.ChallengeParticipation.challenge_id == challenge_id,
                models.ChallengeParticipation.user_id == current_user.id
            )
        ).first()
        
        if not participation:
            raise HTTPException(status_code=403, detail="Not participating in this challenge")
        
        # Check if questions already exist
        existing = db.query(models.ChallengeQuestion).filter(
            models.ChallengeQuestion.challenge_id == challenge_id
        ).first()
        
        if existing:
            # Return existing questions
            questions = db.query(models.ChallengeQuestion).filter(
                models.ChallengeQuestion.challenge_id == challenge_id
            ).all()
            
            return {
                "questions": [{
                    "id": q.id,
                    "question": q.question,
                    "options": json.loads(q.options),
                    "correct_answer": q.correct_answer,
                    "explanation": q.explanation
                } for q in questions]
            }
        
        # Generate questions based on challenge type
        type_descriptions = {
            "speed": "fast-paced questions that can be answered quickly",
            "accuracy": "precise questions requiring careful consideration",
            "topic_mastery": "comprehensive questions testing deep understanding",
            "streak": "progressively challenging questions"
        }
        
        prompt = f"""Generate exactly {question_count} multiple choice questions about {subject}.
Challenge type: {challenge_type} - {type_descriptions.get(challenge_type, '')}.

Return ONLY a valid JSON array with this exact structure:
[
  {{
    "question": "Question text here?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct_answer": 0,
    "explanation": "Brief explanation of the correct answer"
  }}
]

Requirements:
- Each question must have exactly 4 options
- correct_answer must be 0, 1, 2, or 3 (index of the correct option)
- Questions should be clear and educational
- Explanations should be concise (1-2 sentences)
- Return ONLY the JSON array, no additional text"""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=4000
        )
        
        content = response.choices[0].message.content.strip()
        
        # Clean the response
        if content.startswith("```json"):
            content = content[7:]
        if content.startswith("```"):
            content = content[3:]
        if content.endswith("```"):
            content = content[:-3]
        content = content.strip()
        
        questions_data = json.loads(content)
        
        # Save questions to database
        saved_questions = []
        for q_data in questions_data:
            challenge_question = models.ChallengeQuestion(
                challenge_id=challenge_id,
                question=q_data["question"],
                options=json.dumps(q_data["options"]),
                correct_answer=q_data["correct_answer"],
                explanation=q_data.get("explanation", "")
            )
            db.add(challenge_question)
            db.flush()
            
            saved_questions.append({
                "id": challenge_question.id,
                "question": challenge_question.question,
                "options": q_data["options"],
                "correct_answer": challenge_question.correct_answer,
                "explanation": challenge_question.explanation
            })
        
        db.commit()
        
        return {"questions": saved_questions}
    
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {str(e)}, Content: {content}")
        raise HTTPException(status_code=500, detail="Failed to parse AI response")
    except Exception as e:
        db.rollback()
        logger.error(f"Error generating challenge questions: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/update_challenge_progress")
async def update_challenge_progress(
    payload: dict,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Update user's progress in a challenge"""
    try:
        challenge_id = payload.get("challenge_id")
        questions_answered = payload.get("questions_answered", 0)
        accuracy_percentage = payload.get("accuracy_percentage", 0)
        answers = payload.get("answers", [])
        
        # Get challenge and participation
        challenge = db.query(models.Challenge).filter(
            models.Challenge.id == challenge_id
        ).first()
        
        if not challenge:
            raise HTTPException(status_code=404, detail="Challenge not found")
        
        participation = db.query(models.ChallengeParticipation).filter(
            and_(
                models.ChallengeParticipation.challenge_id == challenge_id,
                models.ChallengeParticipation.user_id == current_user.id
            )
        ).first()
        
        if not participation:
            raise HTTPException(status_code=404, detail="Participation not found")
        
        # Save answers
        for answer_data in answers:
            battle_answer = models.ChallengeAnswer(
                challenge_id=challenge_id,
                user_id=current_user.id,
                question_id=answer_data["question_id"],
                selected_answer=answer_data["selected_answer"],
                is_correct=answer_data["is_correct"]
            )
            db.add(battle_answer)
        
        # Calculate progress based on target metric
        progress = 0
        score = questions_answered if challenge.target_metric == "questions_answered" else accuracy_percentage
        
        if challenge.target_metric == "questions_answered":
            progress = min((questions_answered / challenge.target_value) * 100, 100)
            participation.score = questions_answered
        elif challenge.target_metric == "accuracy_percentage":
            progress = min((accuracy_percentage / challenge.target_value) * 100, 100)
            participation.score = accuracy_percentage
        
        participation.progress = progress
        
        # Check if completed
        if progress >= 100:
            participation.completed = True
            participation.completed_at = datetime.utcnow()
            
            # Create activity
            activity = models.FriendActivity(
                user_id=current_user.id,
                activity_type="challenge_completed",
                title=f"Completed Challenge: {challenge.title}",
                description=f"Achieved {progress:.1f}% progress",
                icon="trophy",
                activity_data=json.dumps({
                    "challenge_id": challenge.id,
                    "score": float(score),
                    "progress": float(progress)
                })
            )
            db.add(activity)
            
            # Create notification for challenge completion
            notification = models.Notification(
                user_id=current_user.id,
                title="Challenge Completed",
                message=f"Congratulations! You've completed the challenge '{challenge.title}' with {progress:.0f}% progress!",
                notification_type="challenge_completed",
                is_read=False
            )
            db.add(notification)
        
        db.commit()
        
        return {
            "message": "Progress updated successfully",
            "progress": progress,
            "completed": participation.completed
        }
    
    except Exception as e:
        db.rollback()
        logger.error(f"Error updating challenge progress: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
from sqlalchemy import and_, or_
@app.post("/api/share_content")
async def share_content(
    share_data: ShareContentRequest,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Share a note or chat with friends"""
    try:
        logger.info(f"📤 Sharing content for user {current_user.id}")
        
        # Verify content ownership
        if share_data.content_type == 'note':
            content = db.query(models.Note).filter(
                models.Note.id == share_data.content_id,
                models.Note.user_id == current_user.id
            ).first()
            if not content:
                raise HTTPException(status_code=404, detail="Note not found or not owned by user")
        elif share_data.content_type == 'chat':
            content = db.query(models.ChatSession).filter(
                models.ChatSession.id == share_data.content_id,
                models.ChatSession.user_id == current_user.id
            ).first()
            if not content:
                raise HTTPException(status_code=404, detail="Chat not found or not owned by user")
        else:
            raise HTTPException(status_code=400, detail="Invalid content type")
        
        # Create share records for each friend
        shared_records = []
        for friend_id in share_data.friend_ids:
            #  FIXED: Check for 'active' status instead of 'accepted'
            friendship = db.query(models.Friendship).filter(
                and_(
                    or_(
                        and_(
                            models.Friendship.user_id == current_user.id,
                            models.Friendship.friend_id == friend_id
                        ),
                        and_(
                            models.Friendship.user_id == friend_id,
                            models.Friendship.friend_id == current_user.id
                        )
                    ),
                    models.Friendship.status == 'active'  #  Changed from 'accepted' to 'active'
                )
            ).first()
            
            if not friendship:
                logger.warning(f" User {friend_id} is not a friend of {current_user.id}")
                continue  # Skip if not friends
            
            # Check if already shared
            existing_share = db.query(models.SharedContent).filter(
                models.SharedContent.owner_id == current_user.id,
                models.SharedContent.shared_with_id == friend_id,
                models.SharedContent.content_type == share_data.content_type,
                models.SharedContent.content_id == share_data.content_id
            ).first()
            
            if existing_share:
                # Update existing share
                existing_share.permission = share_data.permission
                existing_share.message = share_data.message
                existing_share.shared_at = datetime.now(timezone.utc)
                shared_records.append(existing_share)
                logger.info(f"🔄 Updated existing share for friend {friend_id}")
            else:
                # Create new share
                shared_content = models.SharedContent(
                    owner_id=current_user.id,
                    shared_with_id=friend_id,
                    content_type=share_data.content_type,
                    content_id=share_data.content_id,
                    permission=share_data.permission,
                    message=share_data.message,
                    shared_at=datetime.now(timezone.utc)
                )
                db.add(shared_content)
                shared_records.append(shared_content)
                logger.info(f" Created new share for friend {friend_id}")
                
                # Create notification for the friend
                content_title = content.title if hasattr(content, 'title') else share_data.content_type
                share_notification = models.Notification(
                    user_id=friend_id,
                    title="New Shared Content",
                    message=f"{current_user.username} shared a {share_data.content_type} with you: {content_title}",
                    notification_type="content_shared",
                    is_read=False
                )
                db.add(share_notification)
        
        db.commit()
        
        return {
            "success": True,
            "message": f"Content shared with {len(shared_records)} friend(s)",
            "shared_count": len(shared_records)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f" Error sharing content: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/shared_with_me")
def get_shared_with_me(
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get all content shared with the current user"""
    try:
        logger.info(f" Starting shared_with_me endpoint for user: {current_user.id}")
        
        # Get all shared content for this user
        shared_items = db.query(models.SharedContent).filter(
            models.SharedContent.shared_with_id == current_user.id
        ).order_by(models.SharedContent.shared_at.desc()).all()
        
        logger.info(f" Found {len(shared_items)} shared items for user {current_user.id}")
        
        # Debug: Log all shared items
        for item in shared_items:
            logger.info(f"🔄 Shared item: id={item.id}, type={item.content_type}, content_id={item.content_id}, owner_id={item.owner_id}")
        
        result = []
        for item in shared_items:
            # Get owner info
            owner = db.query(models.User).filter(models.User.id == item.owner_id).first()
            if not owner:
                logger.warning(f" Owner not found for shared item {item.id}")
                continue
            
            # Get content details
            title = "Untitled"
            content_exists = False
            
            if item.content_type == 'note':
                note = db.query(models.Note).filter(models.Note.id == item.content_id).first()
                if note:
                    title = note.title or "Untitled Note"
                    content_exists = True
                    logger.info(f" Note found: {title}")
                else:
                    logger.warning(f" Note not found for ID: {item.content_id}")
                    title = "Deleted Note"
            elif item.content_type == 'chat':
                chat = db.query(models.ChatSession).filter(models.ChatSession.id == item.content_id).first()
                if chat:
                    title = chat.title or "Untitled Chat"
                    content_exists = True
                    logger.info(f" Chat found: {title}")
                else:
                    logger.warning(f" Chat not found for ID: {item.content_id}")
                    title = "Deleted Chat"
            
            # Only include if content still exists
            if content_exists:
                result.append({
                    "id": item.id,
                    "content_type": item.content_type,
                    "content_id": item.content_id,
                    "title": title,
                    "permission": item.permission,
                    "message": item.message,
                    "shared_at": item.shared_at.isoformat() + 'Z' if item.shared_at else None,
                    "shared_by": {
                        "id": owner.id,
                        "username": owner.username,
                        "email": owner.email,
                        "first_name": owner.first_name or "",
                        "last_name": owner.last_name or "",
                        "picture_url": owner.picture_url or ""
                    }
                })
            else:
                logger.info(f" Skipping deleted content: {item.content_type} {item.content_id}")
        
        logger.info(f" Returning {len(result)} valid shared items")
        return {"shared_items": result}
        
    except Exception as e:
        logger.error(f" Error getting shared content: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/debug_friendships")
def debug_friendships(
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Debug endpoint to check friendship status"""
    try:
        # Get all friendships for current user
        friendships = db.query(models.Friendship).filter(
            models.Friendship.user_id == current_user.id
        ).all()
        
        # Also get reverse friendships
        reverse_friendships = db.query(models.Friendship).filter(
            models.Friendship.friend_id == current_user.id
        ).all()
        
        result = {
            "user": {
                "id": current_user.id,
                "username": current_user.username
            },
            "friendships": [
                {
                    "id": f.id,
                    "user_id": f.user_id,
                    "friend_id": f.friend_id,
                    "status": f.status,
                    "created_at": f.created_at.isoformat() + 'Z'
                }
                for f in friendships
            ],
            "reverse_friendships": [
                {
                    "id": f.id,
                    "user_id": f.user_id,
                    "friend_id": f.friend_id,
                    "status": f.status,
                    "created_at": f.created_at.isoformat() + 'Z'
                }
                for f in reverse_friendships
            ],
            "total_friendships": len(friendships) + len(reverse_friendships)
        }
        
        return result
        
    except Exception as e:
        return {"error": str(e)}
@app.get("/api/shared/{content_type}/{content_id}")
def get_shared_content(
    content_type: str,
    content_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get details of shared content"""
    try:
        logger.info(f" Accessing shared content: {content_type} {content_id} for user: {current_user.id}")
        
        # Check if content is shared with user or owned by user
        shared = db.query(models.SharedContent).filter(
            models.SharedContent.content_type == content_type,
            models.SharedContent.content_id == content_id,
            models.SharedContent.shared_with_id == current_user.id
        ).first()
        
        # Get content
        if content_type == 'note':
            content = db.query(models.Note).filter(models.Note.id == content_id).first()
            if not content:
                raise HTTPException(status_code=404, detail="Note not found")
            
            # Check permissions
            is_owner = content.user_id == current_user.id
            if not is_owner and not shared:
                raise HTTPException(status_code=403, detail="No access to this note")
            
            owner = db.query(models.User).filter(models.User.id == content.user_id).first()
            
            return {
                "content_type": "note",
                "content_id": content.id,
                "title": content.title,
                "content": content.content,
                "created_at": content.created_at.isoformat() + 'Z',
                "updated_at": content.updated_at.isoformat() + 'Z',
                "permission": shared.permission if shared else "owner",
                "is_owner": is_owner,
                "owner": {
                    "id": owner.id,
                    "username": owner.username,
                    "first_name": owner.first_name or "",
                    "last_name": owner.last_name or "",
                    "picture_url": owner.picture_url or ""
                }
            }
            
        elif content_type == 'chat':
            content = db.query(models.ChatSession).filter(models.ChatSession.id == content_id).first()
            if not content:
                raise HTTPException(status_code=404, detail="Chat not found")
            
            # Check permissions
            is_owner = content.user_id == current_user.id
            if not is_owner and not shared:
                raise HTTPException(status_code=403, detail="No access to this chat")
            
            # Get messages
            messages = db.query(models.ChatMessage).filter(
                models.ChatMessage.chat_session_id == content_id
            ).order_by(models.ChatMessage.timestamp.asc()).all()
            
            owner = db.query(models.User).filter(models.User.id == content.user_id).first()
            
            return {
                "content_type": "chat",
                "content_id": content.id,
                "title": content.title,
                "created_at": content.created_at.isoformat() + 'Z',
                "updated_at": content.updated_at.isoformat() + 'Z',
                "permission": shared.permission if shared else "owner",
                "is_owner": is_owner,
                "owner": {
                    "id": owner.id,
                    "username": owner.username,
                    "first_name": owner.first_name or "",
                    "last_name": owner.last_name or "",
                    "picture_url": owner.picture_url or ""
                },
                "messages": [
                    {
                        "user_message": msg.user_message,
                        "ai_response": msg.ai_response,
                        "timestamp": msg.timestamp.isoformat() + 'Z'
                    }
                    for msg in messages
                ]
            }
        else:
            raise HTTPException(status_code=400, detail="Invalid content type")
            
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f" Error getting shared content: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/debug_shared_content")
def debug_shared_content(
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Debug endpoint to check shared content"""
    try:
        # Get all shared content
        shared_items = db.query(models.SharedContent).filter(
            models.SharedContent.shared_with_id == current_user.id
        ).all()
        
        # Also check if user has any shared items as owner
        owned_shares = db.query(models.SharedContent).filter(
            models.SharedContent.owner_id == current_user.id
        ).all()
        
        return {
            "user": {
                "id": current_user.id,
                "username": current_user.username,
                "email": current_user.email
            },
            "received_shares": [
                {
                    "id": item.id,
                    "owner_id": item.owner_id,
                    "shared_with_id": item.shared_with_id,
                    "content_type": item.content_type,
                    "content_id": item.content_id,
                    "permission": item.permission,
                    "message": item.message,
                    "shared_at": item.shared_at.isoformat() + 'Z' if item.shared_at else None
                }
                for item in shared_items
            ],
            "sent_shares": [
                {
                    "id": item.id,
                    "owner_id": item.owner_id,
                    "shared_with_id": item.shared_with_id,
                    "content_type": item.content_type,
                    "content_id": item.content_id,
                    "permission": item.permission,
                    "message": item.message
                }
                for item in owned_shares
            ],
            "total_received": len(shared_items),
            "total_sent": len(owned_shares)
        }
        
    except Exception as e:
        return {"error": str(e)}
@app.delete("/api/remove_shared_access/{share_id}")
def remove_shared_access(
    share_id: int,
    current_user: models.User = Depends(get_current_user),  # Use the existing dependency
    db: Session = Depends(get_db)
):
    """Remove access to shared content"""
    try:
        logger.info(f" User {current_user.id} attempting to remove shared access {share_id}")
        
        # Find shared content - user can remove if they are the recipient OR the owner
        shared = db.query(models.SharedContent).filter(
            models.SharedContent.id == share_id
        ).first()
        
        if not shared:
            logger.error(f" Shared content not found: {share_id}")
            raise HTTPException(status_code=404, detail="Shared content not found")
        
        # Check permissions: user can remove if they are the recipient OR the owner
        can_remove = (shared.shared_with_id == current_user.id) or (shared.owner_id == current_user.id)
        
        if not can_remove:
            logger.warning(f" User {current_user.id} not authorized to remove share {share_id}")
            raise HTTPException(status_code=403, detail="Not authorized to remove this shared content")
        
        logger.info(f" Removing shared access: share_id={share_id}, owner={shared.owner_id}, recipient={shared.shared_with_id}")
        
        db.delete(shared)
        db.commit()
        
        return {
            "success": True, 
            "message": "Access removed successfully",
            "removed_share_id": share_id
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f" Error removing shared access: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/update_shared_note/{note_id}")
def update_shared_note(
    note_id: int,
    note_data: dict,
    credentials: HTTPAuthorizationCredentials = Depends(security),
    db: Session = Depends(get_db)
):
    """Update a shared note (if user has edit permission)"""
    try:
        # Verify token and get user
        payload = jwt.decode(credentials.credentials, SECRET_KEY, algorithms=[ALGORITHM])
        user_email = payload.get("sub")
        user = get_user_by_email(db, user_email)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get note
        note = db.query(models.Note).filter(models.Note.id == note_id).first()
        if not note:
            raise HTTPException(status_code=404, detail="Note not found")
        
        # Check if owner
        if note.user_id == user.id:
            # Owner can always edit
            pass
        else:
            # Check if shared with edit permission
            shared = db.query(models.SharedContent).filter(
                models.SharedContent.content_type == 'note',
                models.SharedContent.content_id == note_id,
                models.SharedContent.shared_with_id == user.id,
                models.SharedContent.permission == 'edit'
            ).first()
            
            if not shared:
                raise HTTPException(status_code=403, detail="No edit permission for this note")
        
        # Update note
        if 'content' in note_data:
            note.content = note_data['content']
        if 'title' in note_data:
            note.title = note_data['title']
        
        note.updated_at = datetime.now(timezone.utc)
        db.commit()
        
        return {
            "success": True,
            "message": "Note updated successfully",
            "note": {
                "id": note.id,
                "title": note.title,
                "content": note.content,
                "updated_at": note.updated_at.isoformat() + 'Z'
            }
        }
        
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error updating shared note: {str(e)}", exc_info=True)
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

# ==================== SERVE REACT APP ====================

# Serve static files from React build
# ==================== SERVE REACT APP ====================

# ==================== SERVE REACT APP ====================

# ==================== SERVE REACT APP (OPTIONAL - NOT USED IN PRODUCTION) ====================

build_dir = Path(__file__).parent.parent / "build"

# Only serve React if build exists (for local testing)
if build_dir.exists() and os.getenv("SERVE_REACT", "false").lower() == "true":
    logger.info(f" Serving React app from {build_dir}")
    
    app.mount("/static", StaticFiles(directory=build_dir / "static"), name="static")
    
    @app.get("/{full_path:path}")
    async def serve_react_app(full_path: str):
        """Serve React app for all non-API routes (local testing only)"""
        
        # Block /api/* from being caught
        if full_path.startswith("api/"):
            raise HTTPException(status_code=404, detail="API endpoint not found")
        
        file_path = build_dir / full_path
        if file_path.is_file():
            return FileResponse(file_path)
        
        return FileResponse(build_dir / "index.html")
else:
    logger.info("️ Not serving React app (frontend on Vercel)")
    
    # Simple root endpoint
    @app.get("/")
    async def root():
        return {
            "message": "Brainwave Backend API v3.0.0",
            "status": "running",
            "frontend": "https://l1-m71fagwct-asphar0057s-projects.vercel.app"
        }
# ==================== WEBSOCKET ENDPOINT ====================

# REPLACE THE WEBSOCKET ENDPOINT in main.py (lines 7741-7827)
# with this improved version:

from typing import Optional  # Add this to your imports at the top

from typing import Optional  # Add this import at the top with other imports

@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket, token: Optional[str] = None):  # ← NEW
    """WebSocket endpoint for real-time notifications"""
    user = None
    user_id = None
    db = None
    
    try:
        logger.info(f"📥 WebSocket connection attempt")
        
        # Validate token BEFORE accepting
        if not token:
            logger.error(" No token provided")
            await websocket.close(code=1008, reason="No token")
            return
        
        db = SessionLocal()
        
        # Verify JWT token
        try:
            from jose import jwt, JWTError
            payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
            username = payload.get("sub")
            
            if not username:
                logger.error(" No username in token")
                await websocket.close(code=1008, reason="Invalid token")
                return
            
            logger.info(f"🔑 Token verified for: {username}")
            
            # Get user from database
            user = get_user_by_username(db, username) or get_user_by_email(db, username)
            if not user:
                logger.error(f" User not found: {username}")
                await websocket.close(code=1008, reason="User not found")
                return
            
            user_id = user.id
            logger.info(f" User {user_id} authenticated successfully")
            
        except JWTError as e:
            logger.error(f" JWT Error: {str(e)}")
            await websocket.close(code=1008, reason="Invalid token")
            return
        except Exception as e:
            logger.error(f" Auth error: {str(e)}")
            await websocket.close(code=1011, reason="Auth error")
            return
        
        # NOW accept the connection (after validation)
        await websocket.accept()
        logger.info(f" WebSocket accepted for user {user_id}")
        
        # Close database connection immediately after auth
        if db:
            db.close()
            db = None
            logger.info(f"🔒 Database connection closed for user {user_id}")
        
        # Store connection
        manager.active_connections[user_id] = websocket
        logger.info(f" User {user_id} connected (Total: {len(manager.active_connections)})")
        
        # Send confirmation
        await websocket.send_json({
            "type": "connected",
            "message": "Connected to battle system",
            "user_id": user_id
        })
        
        # Keep connection alive
        while True:
            try:
                data = await websocket.receive_json()
                
                if data.get("type") == "ping":
                    await websocket.send_json({"type": "pong"})
                    logger.debug(f"🏓 Ping from user {user_id}")
                
            except WebSocketDisconnect:
                logger.info(f"🔌 User {user_id} disconnected")
                break
            except Exception as e:
                logger.error(f" Error in WebSocket loop: {str(e)}")
                break
    
    except Exception as e:
        logger.error(f" WebSocket error: {str(e)}")
        try:
            await websocket.close(code=1011, reason="Error")
        except:
            pass
    
    finally:
        # Cleanup
        if user_id and user_id in manager.active_connections:
            del manager.active_connections[user_id]
            logger.info(f"👋 User {user_id} cleaned up")
        
        # Ensure database connection is closed
        if db:
            try:
                db.close()
                logger.info(f"🔒 Database connection closed in cleanup")
            except:
                pass
            
@app.post("/api/accept_quiz_battle")
async def accept_quiz_battle(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Accept a quiz battle challenge"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        battle_id = payload.get("battle_id")
        if not battle_id:
            raise HTTPException(status_code=400, detail="battle_id is required")
        
        battle = db.query(models.QuizBattle).filter(
            models.QuizBattle.id == battle_id
        ).first()
        
        if not battle:
            raise HTTPException(status_code=404, detail="Battle not found")
        
        if battle.opponent_id != current_user.id:
            raise HTTPException(status_code=403, detail="Not authorized")
        
        if battle.status != "pending":
            raise HTTPException(status_code=400, detail=f"Battle is {battle.status}")
        
        # Update battle status
        battle.status = "active"
        battle.started_at = datetime.now(timezone.utc)
        
        db.commit()
        db.refresh(battle)
        
        logger.info(f" Battle {battle_id} accepted by user {current_user.id}")
        
        # Notify challenger
        await notify_battle_accepted(battle.challenger_id, battle.id)
        await notify_battle_started([battle.challenger_id, battle.opponent_id], battle.id)
        
        return {
            "status": "success",
            "message": "Battle accepted!",
            "battle_id": battle.id,
            "redirect_to": f"/quiz-battle/{battle.id}"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f" Error accepting battle: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/decline_quiz_battle")
async def decline_quiz_battle(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Decline a quiz battle challenge"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        battle_id = payload.get("battle_id")
        if not battle_id:
            raise HTTPException(status_code=400, detail="battle_id is required")
        
        battle = db.query(models.QuizBattle).filter(
            models.QuizBattle.id == battle_id
        ).first()
        
        if not battle:
            raise HTTPException(status_code=404, detail="Battle not found")
        
        if battle.opponent_id != current_user.id:
            raise HTTPException(status_code=403, detail="Not authorized")
        
        if battle.status != "pending":
            raise HTTPException(status_code=400, detail=f"Battle is {battle.status}")
        
        # Update battle status
        battle.status = "expired"
        battle.completed_at = datetime.now(timezone.utc)
        
        db.commit()
        
        logger.info(f" Battle {battle_id} declined by user {current_user.id}")
        
        # Notify challenger
        opponent_name = f"{current_user.first_name} {current_user.last_name}" if current_user.first_name else current_user.username
        await notify_battle_declined(battle.challenger_id, battle.id, opponent_name)
        
        return {
            "status": "success",
            "message": "Battle declined",
            "battle_id": battle.id
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f" Error declining battle: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/submit_battle_answer")
async def submit_battle_answer(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Submit a single answer during battle and notify opponent"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        battle_id = payload.get("battle_id")
        question_index = payload.get("question_index")
        is_correct = payload.get("is_correct")
        
        if battle_id is None or question_index is None or is_correct is None:
            raise HTTPException(status_code=400, detail="Missing required fields")
        
        battle = db.query(models.QuizBattle).filter(
            models.QuizBattle.id == battle_id
        ).first()
        
        if not battle:
            raise HTTPException(status_code=404, detail="Battle not found")
        
        # Determine opponent
        is_challenger = battle.challenger_id == current_user.id
        opponent_id = battle.opponent_id if is_challenger else battle.challenger_id
        
        # Send live notification to opponent
        logger.info(f"📤 Sending answer notification to opponent {opponent_id}: Battle {battle_id}, Q{question_index}, Correct: {is_correct}")
        logger.info(f"📊 Active WebSocket connections: {list(manager.active_connections.keys())}")
        
        success = await manager.send_personal_message({
            "type": "battle_answer_submitted",
            "battle_id": battle_id,
            "question_index": question_index,
            "is_correct": is_correct,
            "is_opponent": True
        }, opponent_id)
        
        if success:
            logger.info(f" Answer notification delivered to opponent {opponent_id}")
        else:
            logger.warning(f" Failed to deliver notification - opponent {opponent_id} not connected")
        
        return {
            "status": "success",
            "message": "Answer submitted"
        }
    
    except Exception as e:
        logger.error(f" Error submitting answer: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/create_solo_quiz")
async def create_solo_quiz(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a new solo quiz"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        subject = payload.get("subject")
        difficulty = payload.get("difficulty", "intermediate")
        question_count = payload.get("question_count", 10)
        
        # Create quiz
        quiz = models.SoloQuiz(
            user_id=current_user.id,
            subject=subject,
            difficulty=difficulty,
            question_count=question_count,
            time_limit_seconds=300
        )
        
        db.add(quiz)
        db.commit()
        db.refresh(quiz)
        
        # Generate questions using AI
        questions = await generate_quiz_questions(subject, difficulty, question_count)
        
        # Save questions
        for q_data in questions:
            question = models.SoloQuizQuestion(
                quiz_id=quiz.id,
                question=q_data["question"],
                options=json.dumps(q_data["options"]),
                correct_answer=q_data["correct_answer"],
                explanation=q_data.get("explanation", "")
            )
            db.add(question)
        
        db.commit()
        
        return {
            "status": "success",
            "quiz_id": quiz.id
        }
    
    except Exception as e:
        logger.error(f"Error creating solo quiz: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/solo_quiz/{quiz_id}")
async def get_solo_quiz(
    quiz_id: int,
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get solo quiz details"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        quiz = db.query(models.SoloQuiz).filter(
            models.SoloQuiz.id == quiz_id,
            models.SoloQuiz.user_id == current_user.id
        ).first()
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        questions = db.query(models.SoloQuizQuestion).filter(
            models.SoloQuizQuestion.quiz_id == quiz_id
        ).all()
        
        return {
            "quiz": {
                "id": quiz.id,
                "subject": quiz.subject,
                "difficulty": quiz.difficulty,
                "question_count": quiz.question_count,
                "time_limit_seconds": quiz.time_limit_seconds
            },
            "questions": [{
                "id": q.id,
                "question": q.question,
                "options": json.loads(q.options),
                "correct_answer": q.correct_answer,
                "explanation": q.explanation
            } for q in questions]
        }
    
    except Exception as e:
        logger.error(f"Error getting solo quiz: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/complete_solo_quiz")
async def complete_solo_quiz(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Complete a solo quiz with gamification points"""
    try:
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        quiz_id = payload.get("quiz_id")
        score = payload.get("score")
        answers = payload.get("answers", [])
        
        quiz = db.query(models.SoloQuiz).filter(
            models.SoloQuiz.id == quiz_id,
            models.SoloQuiz.user_id == current_user.id
        ).first()
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        logger.info(f"Quiz completion - User: {current_user.id}, Score: {score}%")
        
        quiz.score = score
        quiz.completed = True
        quiz.status = "completed"
        quiz.answers = json.dumps(answers)
        quiz.completed_at = datetime.now(timezone.utc)
        
        db.commit()
        logger.info(f"Quiz saved successfully")
        
        # Award gamification points using the new formula
        from gamification_system import award_points, calculate_solo_quiz_points
        
        points_result = award_points(db, current_user.id, "solo_quiz", {
            "difficulty": quiz.difficulty,
            "question_count": quiz.question_count,
            "score_percentage": score
        })
        
        # Get detailed breakdown for response
        quiz_points = calculate_solo_quiz_points(quiz.difficulty, quiz.question_count, score)
        
        db.commit()
        logger.info(f"Awarded {points_result['points_earned']} points for solo quiz")
        
        # Create notifications based on quiz performance
        logger.info(f"Checking notification conditions - Score: {score}")
        notification = None
        
        if score < 50:
            # Poor performance
            logger.info(f"Creating poor performance notification (score {score} < 50)")
            notification = models.Notification(
                user_id=current_user.id,
                title="Quiz Performance Alert",
                message=f"Your recent quiz score was {score}%. Review the material and try again to improve!",
                notification_type="quiz_poor_performance"
            )
            db.add(notification)
            db.commit()
            logger.info(f" Created poor performance notification for user {current_user.id}")
        elif score >= 90:
            # Excellent performance
            logger.info(f"Creating excellent performance notification (score {score} >= 90)")
            notification = models.Notification(
                user_id=current_user.id,
                title="Excellent Work!",
                message=f"Amazing! You scored {score}% on your quiz. You earned {points_result['points_earned']} points!",
                notification_type="quiz_excellent"
            )
            db.add(notification)
            db.commit()
            logger.info(f" Created excellent performance notification for user {current_user.id}")
        else:
            logger.info(f"No notification created - score {score} is between 50-89")
        
        response = {
            "status": "success",
            "message": "Quiz completed",
            "points_earned": points_result["points_earned"],
            "total_points": points_result["total_points"],
            "level": points_result["level"],
            "points_breakdown": quiz_points
        }
        
        # Include notification in response if created
        if notification:
            response["notification"] = {
                "title": notification.title,
                "message": notification.message,
                "notification_type": notification.notification_type
            }
        
        return response
    
    except Exception as e:
        logger.error(f"Error completing solo quiz: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

async def generate_quiz_questions(subject: str, difficulty: str, count: int):
    """Generate quiz questions using AI"""
    try:
        prompt = f"""Generate {count} multiple choice questions about {subject} at {difficulty} level.

For each question provide:
- A clear question
- 4 answer options
- The index of the correct answer (0-3)
- A brief explanation of the correct answer

IMPORTANT: Return ONLY a valid JSON array, no markdown formatting, no code blocks, no extra text.
Use this exact structure:
[{{"question": "...", "options": ["A", "B", "C", "D"], "correct_answer": 0, "explanation": "..."}}]"""

        response = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=GROQ_MODEL,
            temperature=0.7,
            max_tokens=3000
        )
        
        content = response.choices[0].message.content.strip()
        
        # Remove markdown code blocks if present
        if content.startswith("```"):
            # Remove opening ```json or ```
            content = re.sub(r'^```(?:json)?\s*\n', '', content)
            # Remove closing ```
            content = re.sub(r'\n```\s*$', '', content)
            content = content.strip()
        
        logger.info(f"Cleaned content: {content[:200]}...")
        questions = json.loads(content)
        
        # Validate the structure
        if not isinstance(questions, list) or len(questions) == 0:
            raise ValueError("Invalid questions format")
        
        return questions
        
    except Exception as e:
        logger.error(f"Error generating questions: {str(e)}")
        logger.error(f"Raw content: {content if 'content' in locals() else 'No content'}")
        # Return fallback questions
        return [{
            "question": f"Sample question {i+1} about {subject}",
            "options": ["Option A", "Option B", "Option C", "Option D"],
            "correct_answer": 0,
            "explanation": "This is a sample explanation."
        } for i in range(count)]

@app.get("/api/debug/websocket-connections")
async def debug_websocket_connections(username: str = Depends(verify_token)):
    """Debug endpoint to check active WebSocket connections"""
    return {
        "active_connections": list(manager.active_connections.keys()),
        "total_connections": len(manager.active_connections),
        "requesting_user": username
    }

# ==================== GAMIFICATION SYSTEM ====================

def get_week_start():
    """Get the start of the current week (Monday)"""
    today = datetime.now(timezone.utc).date()
    return today - timedelta(days=today.weekday())

def calculate_level_from_xp(xp: int) -> int:
    """Calculate level based on XP with progressive thresholds"""
    # Level thresholds: 0, 100, 282, 500, 800, 1200, 1700, 2300, 3000...
    if xp < 100:
        return 1
    elif xp < 282:
        return 2
    elif xp < 500:
        return 3
    elif xp < 800:
        return 4
    elif xp < 1200:
        return 5
    elif xp < 1700:
        return 6
    elif xp < 2300:
        return 7
    elif xp < 3000:
        return 8
    else:
        # For higher levels: each 1000 XP = 1 level
        return 8 + int((xp - 3000) / 1000)

def get_xp_for_next_level(current_level: int) -> int:
    """Get XP required for next level"""
    thresholds = [0, 100, 282, 500, 800, 1200, 1700, 2300, 3000]
    if current_level < len(thresholds):
        return thresholds[current_level]
    else:
        return 3000 + ((current_level - 8) * 1000)

def calculate_xp_for_level(level: int) -> int:
    """Calculate XP needed for a specific level"""
    return int(100 * (level ** 1.5))

@app.post("/api/track_gamification_activity")
async def track_gamification_activity(
    payload: dict = Body(...),
    username: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Track user activity and award points - Uses centralized gamification system"""
    try:
        from gamification_system import award_points
        import traceback
        
        current_user = get_user_by_username(db, username) or get_user_by_email(db, username)
        if not current_user:
            raise HTTPException(status_code=404, detail="User not found")
        
        activity_type = payload.get("activity_type")
        metadata = payload.get("metadata", {})
        
        # Use centralized gamification system
        result = award_points(db, current_user.id, activity_type, metadata)
        db.commit()
        
        return {
            "status": "success",
            **result
        }

    
    except Exception as e:
        logger.error(f"Error tracking gamification activity: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_gamification_stats")
async def get_gamification_stats(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's gamification stats - Uses centralized gamification system"""
    try:
        from gamification_system import get_user_stats
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Use centralized gamification system
        return get_user_stats(db, user.id)
    
    except Exception as e:
        logger.error(f"Error getting gamification stats: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/get_dashboard_data")
async def get_dashboard_data(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get comprehensive dashboard data for user"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get gamification stats
        from gamification_system import get_user_stats
        gamification_stats = get_user_stats(db, user.id)
        
        # Get daily metrics
        today = datetime.now(timezone.utc).date()
        daily_metrics = db.query(models.DailyLearningMetrics).filter(
            models.DailyLearningMetrics.user_id == user.id,
            models.DailyLearningMetrics.date == today
        ).first()
        
        # Get streak
        streak = calculate_day_streak(db, user.id)
        
        return {
            "status": "success",
            "gamification": gamification_stats,
            "daily_metrics": {
                "questions_answered": daily_metrics.questions_answered if daily_metrics else 0,
                "time_spent_minutes": daily_metrics.time_spent_minutes if daily_metrics else 0,
                "accuracy_rate": daily_metrics.accuracy_rate if daily_metrics else 0
            },
            "streak": streak
        }
    
    except Exception as e:
        logger.error(f"Error getting dashboard data: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/check_missed_achievements")
async def check_missed_achievements(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Check for any achievements user may have missed"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get user's current achievements
        user_achievements = db.query(models.UserAchievement).filter(
            models.UserAchievement.user_id == user.id
        ).all()
        
        earned_achievement_ids = {ua.achievement_id for ua in user_achievements}
        
        # Get all available achievements
        all_achievements = db.query(models.Achievement).all()
        
        # Check which achievements user qualifies for but hasn't earned
        missed_achievements = []
        
        for achievement in all_achievements:
            if achievement.id not in earned_achievement_ids:
                # Check if user qualifies (simplified check)
                # In a real implementation, you'd check the criteria
                missed_achievements.append({
                    "id": achievement.id,
                    "name": achievement.name,
                    "description": achievement.description,
                    "icon": achievement.icon,
                    "points": achievement.points
                })
        
        return {
            "status": "success",
            "missed_achievements": missed_achievements[:5]  # Return top 5
        }
    
    except Exception as e:
        logger.error(f"Error checking missed achievements: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/get_weekly_bingo_stats")
async def get_weekly_bingo_stats(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's weekly bingo challenge stats"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        stats = db.query(models.UserGamificationStats).filter(
            models.UserGamificationStats.user_id == user.id
        ).first()
        
        if not stats:
            return {"stats": {}}
        
        # Return stats for bingo board
        return {
            "stats": {
                "ai_chats": stats.weekly_ai_chats,
                "questions_answered": stats.weekly_questions_answered,
                "notes_created": stats.weekly_notes_created,
                "study_hours": stats.weekly_study_minutes / 60,
                "quizzes_completed": stats.weekly_quizzes_completed,
                "flashcards_created": stats.weekly_flashcards_created,
                "streak": stats.current_streak,
                "battles_won": stats.weekly_battles_won,
                "level": stats.level,
                "solo_quizzes": getattr(stats, 'weekly_solo_quizzes', 0),
                "flashcards_reviewed": getattr(stats, 'weekly_flashcards_reviewed', 0),
                "flashcards_mastered": getattr(stats, 'weekly_flashcards_mastered', 0)
            }
        }
    
    except Exception as e:
        logger.error(f"Error getting bingo stats: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_weekly_activity_progress")
async def get_weekly_activity_progress(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's weekly activity progress"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        stats = db.query(models.UserGamificationStats).filter(
            models.UserGamificationStats.user_id == user.id
        ).first()
        
        if not stats:
            return {
                "study_minutes": 0,
                "ai_chats": 0,
                "notes_created": 0,
                "questions_answered": 0,
                "quizzes_completed": 0,
                "flashcards_created": 0,
                "solo_quizzes": 0,
                "flashcards_reviewed": 0,
                "flashcards_mastered": 0
            }
        
        return {
            "study_minutes": stats.weekly_study_minutes,
            "ai_chats": stats.weekly_ai_chats,
            "notes_created": stats.weekly_notes_created,
            "questions_answered": stats.weekly_questions_answered,
            "quizzes_completed": stats.weekly_quizzes_completed,
            "flashcards_created": stats.weekly_flashcards_created,
            "solo_quizzes": getattr(stats, 'weekly_solo_quizzes', 0),
            "flashcards_reviewed": getattr(stats, 'weekly_flashcards_reviewed', 0),
            "flashcards_mastered": getattr(stats, 'weekly_flashcards_mastered', 0)
        }
    
    except Exception as e:
        logger.error(f"Error getting weekly progress: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_recent_point_activities")
async def get_recent_point_activities(
    user_id: str = Query(...),
    limit: int = Query(10),
    db: Session = Depends(get_db)
):
    """Get user's recent point-earning activities"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        transactions = db.query(models.PointTransaction).filter(
            models.PointTransaction.user_id == user.id
        ).order_by(models.PointTransaction.created_at.desc()).limit(limit).all()
        
        activities = []
        for t in transactions:
            # Handle both timezone-aware and naive datetimes
            now = datetime.now(timezone.utc)
            created = t.created_at.replace(tzinfo=timezone.utc) if t.created_at.tzinfo is None else t.created_at
            time_diff = now - created
            
            if time_diff.days > 0:
                time_ago = f"{time_diff.days}d ago"
            elif time_diff.seconds >= 3600:
                time_ago = f"{time_diff.seconds // 3600}h ago"
            elif time_diff.seconds >= 60:
                time_ago = f"{time_diff.seconds // 60}m ago"
            else:
                time_ago = "just now"
            
            activities.append({
                "description": t.description,
                "points": t.points_earned,
                "time_ago": time_ago,
                "activity_type": t.activity_type
            })
        
        return {"activities": activities}
    
    except Exception as e:
        logger.error(f"Error getting recent activities: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_daily_challenge")
async def get_daily_challenge(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's daily challenge"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get user stats for progress tracking
        stats = db.query(models.UserGamificationStats).filter(
            models.UserGamificationStats.user_id == user.id
        ).first()
        
        # Define daily challenges (rotates based on day of year)
        challenges = [
            {"id": 1, "title": "Knowledge Sprint", "description": "Answer 15 questions correctly", "target": 15, "type": "questions_answered", "reward": 100, "icon": "target"},
            {"id": 2, "title": "Chat Master", "description": "Have 25 AI conversations", "target": 25, "type": "ai_chats", "reward": 75, "icon": "chat"},
            {"id": 3, "title": "Note Taker", "description": "Create 5 new notes", "target": 5, "type": "notes_created", "reward": 150, "icon": "note"},
            {"id": 4, "title": "Study Marathon", "description": "Study for 2 hours", "target": 120, "type": "study_minutes", "reward": 200, "icon": "clock"},
            {"id": 5, "title": "Quiz Champion", "description": "Complete 3 quizzes with 80%+", "target": 3, "type": "quizzes_completed", "reward": 175, "icon": "trophy"},
            {"id": 6, "title": "Flashcard Creator", "description": "Create 20 flashcards", "target": 20, "type": "flashcards_created", "reward": 125, "icon": "cards"},
            {"id": 7, "title": "Perfect Score", "description": "Get 100% on any quiz", "target": 1, "type": "perfect_quizzes", "reward": 250, "icon": "star"}
        ]
        
        # Select challenge based on day of year
        from datetime import datetime
        day_of_year = datetime.now().timetuple().tm_yday
        challenge_index = day_of_year % len(challenges)
        daily_challenge = challenges[challenge_index]
        
        # Get current progress
        progress = 0
        if stats:
            if daily_challenge["type"] == "questions_answered":
                progress = stats.weekly_questions_answered
            elif daily_challenge["type"] == "ai_chats":
                progress = stats.weekly_ai_chats
            elif daily_challenge["type"] == "notes_created":
                progress = stats.weekly_notes_created
            elif daily_challenge["type"] == "study_minutes":
                progress = stats.weekly_study_minutes
            elif daily_challenge["type"] == "quizzes_completed":
                progress = stats.weekly_quizzes_completed
            elif daily_challenge["type"] == "flashcards_created":
                progress = stats.weekly_flashcards_created
        
        return {
            "challenge": daily_challenge,
            "progress": progress
        }
    
    except Exception as e:
        logger.error(f"Error getting daily challenge: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_leaderboard")
async def get_leaderboard(
    category: str = Query("global"),
    limit: int = Query(50),
    db: Session = Depends(get_db)
):
    """Get leaderboard - PRIORITIZES GAMIFICATION POINTS"""
    try:
        # Get all users with gamification stats, ordered by total_points (PRIMARY METRIC)
        leaderboard_query = db.query(
            models.User,
            models.UserGamificationStats
        ).join(
            models.UserGamificationStats,
            models.User.id == models.UserGamificationStats.user_id
        ).order_by(
            models.UserGamificationStats.total_points.desc(),  # PRIMARY: Gamification points
            models.UserGamificationStats.level.desc(),          # SECONDARY: Level
            models.UserGamificationStats.experience.desc()      # TERTIARY: Experience
        ).limit(limit)
        
        results = leaderboard_query.all()
        
        leaderboard = []
        for rank, (user, stats) in enumerate(results, 1):
            leaderboard.append({
                "rank": rank,
                "user_id": user.id,
                "username": user.username,
                "first_name": user.first_name or "",
                "last_name": user.last_name or "",
                "picture_url": user.picture_url or "",
                "total_points": stats.total_points,  # PRIMARY METRIC
                "level": stats.level,
                "experience": stats.experience,
                "weekly_points": stats.weekly_points,
                "current_streak": stats.current_streak
            })
        
        return {"leaderboard": leaderboard}
    
    except Exception as e:
        logger.error(f"Error getting leaderboard: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/recalculate_gamification")
async def recalculate_gamification(db: Session = Depends(get_db)):
    """ONE-TIME: Recalculate all user stats from historical data"""
    try:
        from gamification_system import recalculate_all_stats
        count = recalculate_all_stats(db)
        return {
            "status": "success",
            "users_processed": count,
            "message": "All user stats recalculated from historical data"
        }
    except Exception as e:
        logger.error(f"Recalculation error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== END SERVE REACT APP ====================


# ==================== NOTIFICATION ENDPOINTS ====================

@app.get("/api/get_notifications")
async def get_notifications(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get all notifications for a user"""
    try:
        logger.info(f"Getting notifications for user: {user_id}")
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        
        if not user:
            logger.warning(f"User not found for notifications: {user_id}")
            # Return empty notifications instead of 404
            return {"notifications": []}
        
        logger.info(f"Found user with id: {user.id}")
        
        # Check for upcoming reminders and create notifications
        # Use local time since reminders are stored in local time
        now = datetime.now()
        upcoming_reminders = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == False,
            models.Reminder.is_notified == False,
            models.Reminder.reminder_date > now
        ).all()
        
        for reminder in upcoming_reminders:
            time_until = reminder.reminder_date - now
            minutes_until = time_until.total_seconds() / 60
            
            # Create notification if within notify_before_minutes window
            if minutes_until <= reminder.notify_before_minutes:
                # Check if notification already exists for this specific reminder
                existing = db.query(models.Notification).filter(
                    models.Notification.user_id == user.id,
                    models.Notification.notification_type == 'reminder',
                    models.Notification.title.contains(f"Reminder: {reminder.title}"),
                    models.Notification.created_at >= datetime.now() - timedelta(hours=1)  # Only check recent notifications
                ).first()
                
                if not existing:
                    notification = models.Notification(
                        user_id=user.id,
                        title=f"Reminder: {reminder.title}",
                        message=f"{reminder.description or 'Upcoming reminder'} at {reminder.reminder_date.isoformat()}",
                        notification_type='reminder'
                    )
                    db.add(notification)
                    reminder.is_notified = True
                    db.commit()
                    logger.info(f"Created reminder notification for: {reminder.title}")
        
        notifications = db.query(models.Notification).filter(
            models.Notification.user_id == user.id
        ).order_by(models.Notification.created_at.desc()).all()
        
        logger.info(f"Found {len(notifications)} notifications")
        
        return {
            "notifications": [
                {
                    "id": n.id,
                    "title": n.title,
                    "message": n.message,
                    "notification_type": n.notification_type,
                    "is_read": n.is_read,
                    "created_at": n.created_at.isoformat() + 'Z'
                }
                for n in notifications
            ]
        }
    except HTTPException as he:
        logger.error(f"HTTPException in get_notifications: {he.detail}")
        raise
    except Exception as e:
        logger.error(f"Error getting notifications: {str(e)}", exc_info=True)
        # Return empty notifications on error instead of 500
        return {"notifications": []}

@app.put("/api/mark_notification_read/{notification_id}")
async def mark_notification_read(
    notification_id: int,
    db: Session = Depends(get_db)
):
    """Mark a notification as read"""
    try:
        notification = db.query(models.Notification).filter(
            models.Notification.id == notification_id
        ).first()
        
        if not notification:
            raise HTTPException(status_code=404, detail="Notification not found")
        
        notification.is_read = True
        db.commit()
        
        return {"status": "success", "message": "Notification marked as read"}
    except Exception as e:
        logger.error(f"Error marking notification as read: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/mark_all_notifications_read")
async def mark_all_notifications_read(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Mark all notifications as read for a user"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        db.query(models.Notification).filter(
            models.Notification.user_id == user.id,
            models.Notification.is_read == False
        ).update({"is_read": True})
        
        db.commit()
        
        return {"status": "success", "message": "All notifications marked as read"}
    except Exception as e:
        logger.error(f"Error marking all notifications as read: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/create_notification")
async def create_notification(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Create a notification for a user"""
    try:
        user_id = payload.get("user_id")
        title = payload.get("title")
        message = payload.get("message")
        notification_type = payload.get("notification_type", "general")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        notification = models.Notification(
            user_id=user.id,
            title=title,
            message=message,
            notification_type=notification_type
        )
        
        db.add(notification)
        db.commit()
        db.refresh(notification)
        
        return {
            "status": "success",
            "notification_id": notification.id,
            "message": "Notification created"
        }
    except Exception as e:
        logger.error(f"Error creating notification: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/debug_notifications")
async def debug_notifications(user_id: str = Query(...), db: Session = Depends(get_db)):
    """Debug endpoint to check notifications in database"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"error": "User not found", "user_id": user_id}
        
        notifications = db.query(models.Notification).filter(
            models.Notification.user_id == user.id
        ).all()
        
        return {
            "user_id": user.id,
            "username": user.username,
            "total_notifications": len(notifications),
            "notifications": [
                {
                    "id": n.id,
                    "title": n.title,
                    "message": n.message[:100],
                    "type": n.notification_type,
                    "is_read": n.is_read,
                    "created_at": str(n.created_at)
                }
                for n in notifications
            ]
        }
    except Exception as e:
        return {"error": str(e)}

@app.delete("/api/delete_notification/{notification_id}")
async def delete_notification(
    notification_id: int,
    db: Session = Depends(get_db)
):
    """Delete a specific notification"""
    try:
        notification = db.query(models.Notification).filter(
            models.Notification.id == notification_id
        ).first()
        
        if not notification:
            raise HTTPException(status_code=404, detail="Notification not found")
        
        db.delete(notification)
        db.commit()
        
        return {"status": "success", "message": "Notification deleted"}
    except Exception as e:
        logger.error(f"Error deleting notification: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/clear_old_notifications")
async def clear_old_notifications(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Clear all old notifications from previous sessions when user logs in"""
    try:
        logger.info(f"🔔 Clearing old notifications for user: {user_id}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            logger.warning(f"🔔 User not found: {user_id}")
            return {"status": "success", "cleared": 0, "message": "User not found, nothing to clear"}
        
        # Delete all notifications for this user
        deleted = db.query(models.Notification).filter(
            models.Notification.user_id == user.id
        ).delete()
        
        db.commit()
        
        logger.info(f"🔔 Cleared {deleted} old notifications for user {user_id}")
        return {"status": "success", "cleared": deleted, "message": f"Cleared {deleted} old notifications"}
    except Exception as e:
        logger.error(f"🔔 Error clearing notifications: {str(e)}")
        return {"status": "error", "cleared": 0, "message": str(e)}

@app.delete("/api/clear_all_notifications")
async def clear_all_notifications(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Clear all notifications for a user (used on fresh login)"""
    try:
        logger.info(f"🔔 Clearing ALL notifications for user: {user_id}")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            logger.warning(f"🔔 User not found: {user_id}")
            return {"status": "success", "cleared": 0, "message": "User not found"}
        
        # Delete all notifications for this user
        deleted = db.query(models.Notification).filter(
            models.Notification.user_id == user.id
        ).delete()
        
        db.commit()
        
        logger.info(f"🔔 Cleared {deleted} notifications for user {user_id}")
        return {"status": "success", "cleared": deleted, "message": f"Cleared {deleted} notifications"}
    except Exception as e:
        logger.error(f"🔔 Error clearing all notifications: {str(e)}")
        db.rollback()
        return {"status": "error", "cleared": 0, "message": str(e)}


# ==================== STUDY INSIGHTS ENDPOINTS ====================

@app.get("/api/study_insights/session_summary")
async def get_study_session_summary(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get comprehensive summary of user's last study session"""
    try:
        from study_session_analyzer import get_study_session_analyzer
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        analyzer = get_study_session_analyzer(db, user.id, unified_ai)
        summary = analyzer.generate_session_summary()
        
        return {
            "status": "success",
            "summary": summary
        }
    except Exception as e:
        logger.error(f"Error getting study session summary: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/study_insights/ai_summary")
async def get_ai_study_summary(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get AI-generated personalized summary of last study session"""
    try:
        from study_session_analyzer import get_study_session_analyzer
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        analyzer = get_study_session_analyzer(db, user.id, unified_ai)
        ai_summary = await analyzer.generate_ai_summary()
        
        return {
            "status": "success",
            "summary": ai_summary
        }
    except Exception as e:
        logger.error(f"Error getting AI study summary: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/study_insights/strengths_weaknesses")
async def get_strengths_weaknesses(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's strengths and weaknesses based on all learning data"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get topic mastery data
        mastery_data = db.query(models.TopicMastery).filter(
            models.TopicMastery.user_id == user.id
        ).all()
        
        strengths = []
        weaknesses = []
        
        for m in mastery_data:
            topic_info = {
                "topic": m.topic_name,
                "mastery_level": round(m.mastery_level * 100, 1),
                "times_studied": m.times_studied,
                "last_studied": m.last_studied.isoformat() if m.last_studied else None
            }
            
            if m.mastery_level >= 0.7:
                strengths.append(topic_info)
            elif m.mastery_level < 0.5:
                weaknesses.append(topic_info)
        
        # Sort by mastery level
        strengths.sort(key=lambda x: x["mastery_level"], reverse=True)
        weaknesses.sort(key=lambda x: x["mastery_level"])
        
        # Get quiz performance data
        quiz_performance = db.query(models.UserPerformanceMetrics).filter(
            models.UserPerformanceMetrics.user_id == user.id
        ).all()
        
        quiz_weak_topics = []
        quiz_strong_topics = []
        
        for p in quiz_performance:
            if p.accuracy_rate < 60:
                quiz_weak_topics.append({
                    "topic": p.topic_name,
                    "accuracy": round(p.accuracy_rate, 1),
                    "questions_attempted": p.total_questions
                })
            elif p.accuracy_rate >= 80:
                quiz_strong_topics.append({
                    "topic": p.topic_name,
                    "accuracy": round(p.accuracy_rate, 1),
                    "questions_attempted": p.total_questions
                })
        
        # Get flashcard weak cards
        weak_flashcards = db.query(models.Flashcard).join(
            models.FlashcardSet
        ).filter(
            models.FlashcardSet.user_id == user.id,
            models.Flashcard.marked_for_review == True
        ).all()
        
        flashcard_weaknesses = []
        for card in weak_flashcards[:10]:
            flashcard_weaknesses.append({
                "question": card.question[:100],
                "set_id": card.set_id,
                "times_reviewed": card.times_reviewed,
                "correct_rate": round((card.correct_count / max(card.times_reviewed, 1)) * 100, 1)
            })
        
        return {
            "status": "success",
            "strengths": {
                "topics": strengths[:10],
                "quiz_topics": quiz_strong_topics[:5]
            },
            "weaknesses": {
                "topics": weaknesses[:10],
                "quiz_topics": quiz_weak_topics[:5],
                "flashcards": flashcard_weaknesses
            }
        }
    except Exception as e:
        logger.error(f"Error getting strengths/weaknesses: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/study_insights/recommendations")
async def get_study_recommendations(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get personalized study recommendations"""
    try:
        from study_session_analyzer import get_study_session_analyzer
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        analyzer = get_study_session_analyzer(db, user.id, unified_ai)
        summary = analyzer.generate_session_summary()
        
        return {
            "status": "success",
            "recommendations": summary.get("recommendations", [])
        }
    except Exception as e:
        logger.error(f"Error getting study recommendations: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/study_insights/debug_session")
async def debug_session_tracking(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Debug endpoint to verify session tracking is working correctly"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        now = datetime.utcnow()
        
        # Get last_login
        last_login = user.last_login
        session_start = last_login if last_login else (now - timedelta(hours=4))
        
        # Count messages since session start
        messages_count = db.query(models.ChatMessage).join(
            models.ChatSession
        ).filter(
            models.ChatSession.user_id == user.id,
            models.ChatMessage.timestamp >= session_start
        ).count()
        
        # Count total messages
        total_messages = db.query(models.ChatMessage).join(
            models.ChatSession
        ).filter(
            models.ChatSession.user_id == user.id
        ).count()
        
        # Get recent messages for debugging
        recent_messages = db.query(models.ChatMessage).join(
            models.ChatSession
        ).filter(
            models.ChatSession.user_id == user.id
        ).order_by(models.ChatMessage.timestamp.desc()).limit(5).all()
        
        recent_list = []
        for msg in recent_messages:
            recent_list.append({
                "timestamp": msg.timestamp.isoformat() if msg.timestamp else None,
                "user_message": msg.user_message[:100] if msg.user_message else None,
                "in_session": msg.timestamp >= session_start if msg.timestamp else False
            })
        
        return {
            "status": "success",
            "debug_info": {
                "user_id": user.id,
                "username": user.username,
                "last_login": last_login.isoformat() if last_login else None,
                "session_start": session_start.isoformat(),
                "current_time": now.isoformat(),
                "messages_in_session": messages_count,
                "total_messages": total_messages,
                "recent_messages": recent_list
            }
        }
    except Exception as e:
        logger.error(f"Error in debug session: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/study_insights/reset_stats")
async def reset_user_stats(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Reset all stats for a user - for testing purposes"""
    try:
        user_id = payload.get("user_id")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Delete chat messages and sessions
        sessions = db.query(models.ChatSession).filter(models.ChatSession.user_id == user.id).all()
        for session in sessions:
            db.query(models.ChatMessage).filter(models.ChatMessage.chat_session_id == session.id).delete()
        db.query(models.ChatSession).filter(models.ChatSession.user_id == user.id).delete()
        
        # Delete topic mastery
        db.query(models.TopicMastery).filter(models.TopicMastery.user_id == user.id).delete()
        
        # Delete flashcard study sessions
        db.query(models.FlashcardStudySession).filter(models.FlashcardStudySession.user_id == user.id).delete()
        
        # Reset last_login to force new session
        user.last_login = None
        
        db.commit()
        
        return {
            "status": "success",
            "message": "All stats reset. Please log out and log back in."
        }
    except Exception as e:
        logger.error(f"Error resetting stats: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/study_insights/generate_content")
async def generate_study_content(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Generate flashcards, notes, or quiz questions based on weak areas"""
    try:
        user_id = payload.get("user_id")
        content_type = payload.get("content_type")  # flashcards, notes, quiz
        topic = payload.get("topic")
        count = payload.get("count", 5)
        context = payload.get("context", "")  # Additional context from source question
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Build context string for more personalized generation
        context_str = ""
        if context:
            context_str = f"\n\nThe student was working on problems like: {context}\nGenerate similar problems at the same level."
        
        if content_type == "flashcards":
            # Generate flashcards for the topic
            prompt = f"""Generate {count} flashcards for studying "{topic}".{context_str}
            
Return as JSON array with this format:
[
  {{"question": "...", "answer": "...", "difficulty": "easy|medium|hard"}},
  ...
]

Make the flashcards:
- SPECIFIC to {topic} (not generic)
- Include actual formulas, equations, or specific examples
- Progressive in difficulty
- Test understanding, not just memorization
Return ONLY the JSON array, no other text."""

            response = unified_ai.generate(prompt, max_tokens=1500, temperature=0.7)
            
            # Parse the response
            try:
                # Clean up response
                response = response.strip()
                if response.startswith("```"):
                    response = response.split("```")[1]
                    if response.startswith("json"):
                        response = response[4:]
                flashcards = json.loads(response)
            except:
                flashcards = []
            
            return {
                "status": "success",
                "content_type": "flashcards",
                "topic": topic,
                "content": flashcards
            }
        
        elif content_type == "quiz":
            # Generate quiz questions
            prompt = f"""Generate {count} multiple choice quiz questions about "{topic}".{context_str}
            
Return as JSON array with this format:
[
  {{
    "question": "...",
    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
    "correct_answer": "A",
    "explanation": "..."
  }},
  ...
]

Make questions:
- SPECIFIC to {topic} (include actual problems/calculations if math-related)
- Test understanding, not just memorization
- Include step-by-step explanations
Return ONLY the JSON array, no other text."""

            response = unified_ai.generate(prompt, max_tokens=2000, temperature=0.7)
            
            try:
                response = response.strip()
                if response.startswith("```"):
                    response = response.split("```")[1]
                    if response.startswith("json"):
                        response = response[4:]
                questions = json.loads(response)
            except:
                questions = []
            
            return {
                "status": "success",
                "content_type": "quiz",
                "topic": topic,
                "content": questions
            }
        
        elif content_type == "notes":
            # Generate study notes
            prompt = f"""Create comprehensive study notes about "{topic}".
            
Include:
1. Key concepts and definitions
2. Important formulas or rules (if applicable)
3. Examples
4. Common mistakes to avoid
5. Summary points

Format with clear headings and bullet points.
Make it suitable for exam preparation."""

            response = unified_ai.generate(prompt, max_tokens=2000, temperature=0.7)
            
            return {
                "status": "success",
                "content_type": "notes",
                "topic": topic,
                "content": response
            }
        
        else:
            raise HTTPException(status_code=400, detail="Invalid content_type")
        
    except Exception as e:
        logger.error(f"Error generating study content: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/study_insights/welcome_notification")
async def get_welcome_notification(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get personalized welcome notification with session summary"""
    try:
        from study_session_analyzer import get_study_session_analyzer
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        analyzer = get_study_session_analyzer(db, user.id, unified_ai)
        
        # Get AI-generated summary
        ai_summary = await analyzer.generate_ai_summary()
        
        # Get quick stats
        summary = analyzer.generate_session_summary()
        
        # Determine if user has recent activity
        has_recent_activity = summary.get("summary", {}).get("total_activities", 0) > 0
        
        user_name = user.first_name or user.username.split('@')[0]
        
        return {
            "status": "success",
            "notification": {
                "title": "Welcome Back!" if has_recent_activity else "Welcome!",
                "message": ai_summary,
                "has_insights": has_recent_activity,
                "user_name": user_name,
                "quick_stats": {
                    "chat_messages": summary.get("summary", {}).get("chat_messages", 0),
                    "flashcards_studied": summary.get("summary", {}).get("flashcards_studied", 0),
                    "quiz_questions": summary.get("summary", {}).get("quiz_questions", 0),
                    "overall_accuracy": summary.get("summary", {}).get("overall_accuracy", 0)
                },
                "top_weakness": summary.get("weaknesses", [{}])[0] if summary.get("weaknesses") else None,
                "top_recommendation": summary.get("recommendations", [{}])[0] if summary.get("recommendations") else None
            }
        }
    except Exception as e:
        logger.error(f"Error getting welcome notification: {str(e)}")
        # Return a basic welcome message on error
        return {
            "status": "success",
            "notification": {
                "title": "Welcome Back!",
                "message": "Ready to continue learning?",
                "has_insights": False
            }
        }

@app.get("/api/check_reminder_notifications")
async def check_reminder_notifications(
    user_id: str = Query(...),
    current_time: str = Query(None),  # User's current local time as ISO string
    db: Session = Depends(get_db)
):
    """Check for upcoming reminders and create notifications for them.
    
    Reminders are stored in user's local time, so we need the user's current local time
    to compare properly. Frontend sends current_time as ISO string.
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"status": "error", "message": "User not found", "notifications_created": 0}
        
        # Parse the user's current local time from frontend
        # If not provided, fall back to server time (less accurate but works)
        if current_time:
            try:
                now = datetime.fromisoformat(current_time.replace('Z', '').replace('+00:00', ''))
                logger.info(f"Using client time: {now}")
            except:
                now = datetime.now()
                logger.info(f"Failed to parse client time, using server time: {now}")
        else:
            now = datetime.now()
            logger.info(f"No client time provided, using server time: {now}")
        
        notifications_created = []
        
        # Get reminders that haven't been notified yet and have a reminder_date
        pending_reminders = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == False,
            models.Reminder.is_notified == False,
            models.Reminder.reminder_date != None
        ).all()
        
        logger.info(f"Found {len(pending_reminders)} pending reminders for user {user_id}")
        
        for reminder in pending_reminders:
            if not reminder.reminder_date:
                continue
                
            time_until = reminder.reminder_date - now
            minutes_until = time_until.total_seconds() / 60
            
            logger.info(f"Reminder '{reminder.title}': scheduled={reminder.reminder_date}, now={now}, minutes_until={minutes_until:.1f}, notify_before={reminder.notify_before_minutes}")
            
            # Only notify if:
            # 1. Time is within the notify window (minutes_until <= notify_before_minutes)
            # 2. AND the reminder hasn't passed by more than 30 minutes (to avoid old reminders)
            # 3. AND we're actually close to the notify time (not just created)
            
            notify_window_start = reminder.notify_before_minutes
            is_in_notify_window = minutes_until <= notify_window_start and minutes_until >= -30
            
            if is_in_notify_window:
                # Check if we already created a notification for this reminder recently
                existing_notification = db.query(models.Notification).filter(
                    models.Notification.user_id == user.id,
                    models.Notification.notification_type == 'reminder',
                    models.Notification.title.contains(reminder.title),
                    models.Notification.created_at >= datetime.now() - timedelta(hours=1)
                ).first()
                
                if existing_notification:
                    logger.info(f"Skipping duplicate notification for: {reminder.title}")
                    continue
                
                # Format the time for display (already in user's local time)
                reminder_time = reminder.reminder_date.strftime('%I:%M %p')
                reminder_date_str = reminder.reminder_date.strftime('%B %d, %Y at %I:%M %p')
                
                if minutes_until <= 0:
                    # Past due / happening now
                    notification = models.Notification(
                        user_id=user.id,
                        title=f"{reminder.title} - NOW!",
                        message=f"{reminder.description or 'Your event is happening now!'} - Scheduled for {reminder_time}",
                        notification_type='reminder'
                    )
                elif minutes_until <= 5:
                    # Due very soon (within 5 minutes)
                    notification = models.Notification(
                        user_id=user.id,
                        title=f"{reminder.title} - In {int(minutes_until)} min!",
                        message=f"{reminder.description or 'Your reminder is coming up!'} - Due at {reminder_time}",
                        notification_type='reminder'
                    )
                else:
                    # Upcoming (within notify_before_minutes window)
                    notification = models.Notification(
                        user_id=user.id,
                        title=f"{reminder.title}",
                        message=f"{reminder.description or 'Your scheduled reminder'} - Due at {reminder_time} (in {int(minutes_until)} min)",
                        notification_type='reminder'
                    )
                
                db.add(notification)
                reminder.is_notified = True
                
                notifications_created.append({
                    "reminder_id": reminder.id,
                    "title": reminder.title,
                    "minutes_until": round(minutes_until),
                    "reminder_time": reminder_date_str
                })
                
                logger.info(f"Created reminder notification for: {reminder.title} at {reminder_date_str}")
        
        if notifications_created:
            db.commit()
            logger.info(f"Created {len(notifications_created)} reminder notifications")
        
        return {
            "status": "success",
            "notifications_created": len(notifications_created),
            "details": notifications_created,
            "server_time": datetime.now().isoformat(),
            "client_time_received": current_time
        }
    except Exception as e:
        logger.error(f"Error checking reminder notifications: {str(e)}")
        import traceback
        traceback.print_exc()
        db.rollback()
        return {"status": "error", "message": str(e), "notifications_created": 0}

@app.get("/api/get_concept_web")
async def get_concept_web(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get user's complete concept web (nodes and connections)"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"nodes": [], "connections": []}
        
        # Get all concept nodes
        nodes = db.query(models.ConceptNode).filter(
            models.ConceptNode.user_id == user.id
        ).all()
        
        # Get all connections
        connections = db.query(models.ConceptConnection).filter(
            models.ConceptConnection.user_id == user.id
        ).all()
        
        return {
            "nodes": [
                {
                    "id": node.id,
                    "concept_name": node.concept_name,
                    "description": node.description,
                    "category": node.category,
                    "importance_score": node.importance_score,
                    "mastery_level": node.mastery_level,
                    "position_x": node.position_x,
                    "position_y": node.position_y,
                    "notes_count": node.notes_count,
                    "quizzes_count": node.quizzes_count,
                    "flashcards_count": node.flashcards_count,
                    "created_at": node.created_at.isoformat() + 'Z'
                }
                for node in nodes
            ],
            "connections": [
                {
                    "id": conn.id,
                    "source_id": conn.source_concept_id,
                    "target_id": conn.target_concept_id,
                    "connection_type": conn.connection_type,
                    "strength": conn.strength,
                    "ai_generated": conn.ai_generated,
                    "user_confirmed": conn.user_confirmed
                }
                for conn in connections
            ]
        }
    except Exception as e:
        logger.error(f"Error getting concept web: {str(e)}")
        return {"nodes": [], "connections": []}

@app.post("/api/generate_concept_web")
async def generate_concept_web(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Generate concept web from user's actual learning content"""
    try:
        user_id = payload.get("user_id")
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        logger.info(f"Generating concept web for user: {user.username}")
        
        # Delete ALL existing concepts first
        db.query(models.ConceptConnection).filter(
            models.ConceptConnection.user_id == user.id
        ).delete()
        db.query(models.ConceptNode).filter(
            models.ConceptNode.user_id == user.id
        ).delete()
        db.commit()
        
        # Gather ACTUAL content titles/subjects
        raw_concepts = []  # List of (title, source_type, category)
        
        # From Notes - use actual note titles
        notes = db.query(models.Note).filter(
            models.Note.user_id == user.id,
            models.Note.is_deleted == False
        ).all()
        logger.info(f"Found {len(notes)} notes")
        for note in notes:
            if note.title and len(note.title.strip()) > 2:
                raw_concepts.append((note.title.strip(), "Note", "Academic"))
        
        # From Quizzes - use actual quiz subjects
        quizzes = db.query(models.SoloQuiz).filter(
            models.SoloQuiz.user_id == user.id
        ).all()
        logger.info(f"Found {len(quizzes)} quizzes")
        for quiz in quizzes:
            if quiz.subject and len(quiz.subject.strip()) > 2:
                raw_concepts.append((quiz.subject.strip(), "Quiz", "Academic"))
        
        # From Flashcard Sets - use actual set titles
        flashcard_sets = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.user_id == user.id
        ).all()
        logger.info(f"Found {len(flashcard_sets)} flashcard sets")
        for fs in flashcard_sets:
            if fs.title and len(fs.title.strip()) > 2:
                raw_concepts.append((fs.title.strip(), "Flashcards", "Academic"))
        
        # From AI Chat Sessions - use chat titles
        chat_sessions = db.query(models.ChatSession).filter(
            models.ChatSession.user_id == user.id
        ).order_by(models.ChatSession.updated_at.desc()).limit(50).all()
        logger.info(f"Found {len(chat_sessions)} chat sessions")
        for session in chat_sessions:
            if session.title and session.title != "New Chat" and len(session.title.strip()) > 2:
                raw_concepts.append((session.title.strip(), "AI Chat", "Discussion"))
        
        if not raw_concepts:
            logger.info("No content found")
            return {"status": "no_content", "message": "No learning content found"}
        
        # DEDUPLICATE: Normalize and merge similar topics
        def normalize_topic(title):
            """Normalize topic name for comparison"""
            import re
            # Convert to lowercase
            normalized = title.lower().strip()
            # Remove common prefixes/suffixes
            prefixes_to_remove = ['introduction to ', 'intro to ', 'basics of ', 'advanced ', 'fundamentals of ', 'learning ', 'study ', 'notes on ', 'flashcards on ', 'quiz on ']
            for prefix in prefixes_to_remove:
                if normalized.startswith(prefix):
                    normalized = normalized[len(prefix):]
            # Remove special characters except spaces
            normalized = re.sub(r'[^\w\s]', '', normalized)
            # Remove extra whitespace
            normalized = ' '.join(normalized.split())
            return normalized
        
        def topics_are_similar(topic1, topic2):
            """Check if two topics are semantically similar"""
            norm1 = normalize_topic(topic1)
            norm2 = normalize_topic(topic2)
            
            # Exact match after normalization
            if norm1 == norm2:
                return True
            
            # One contains the other
            if norm1 in norm2 or norm2 in norm1:
                return True
            
            # Check word overlap (Jaccard similarity)
            words1 = set(norm1.split())
            words2 = set(norm2.split())
            if not words1 or not words2:
                return False
            
            intersection = words1 & words2
            union = words1 | words2
            similarity = len(intersection) / len(union)
            
            # If more than 60% word overlap, consider similar
            return similarity > 0.6
        
        # Group similar topics together
        concepts_to_create = {}  # normalized_key -> (display_title, source_types, category)
        
        for title, source_type, category in raw_concepts:
            normalized = normalize_topic(title)
            
            # Check if this topic is similar to any existing one
            found_match = False
            for existing_key in list(concepts_to_create.keys()):
                existing_title = concepts_to_create[existing_key][0]
                if topics_are_similar(title, existing_title):
                    # Merge: keep the shorter/cleaner title, combine source types
                    existing_sources = concepts_to_create[existing_key][1]
                    if source_type not in existing_sources:
                        existing_sources.append(source_type)
                    # Keep the shorter title as display name (usually cleaner)
                    if len(title) < len(existing_title):
                        concepts_to_create[existing_key] = (title, existing_sources, category)
                    found_match = True
                    break
            
            if not found_match:
                concepts_to_create[normalized] = (title, [source_type], category)
        
        logger.info(f"Deduplicated {len(raw_concepts)} raw concepts to {len(concepts_to_create)} unique concepts")
        
        logger.info(f"Creating {len(concepts_to_create)} concepts with AI classification")
        
        # Get AI classification agent with Gemini primary, Groq fallback
        from concept_classification_agent import get_concept_agent
        agent = get_concept_agent(groq_client, GROQ_MODEL, gemini_client, GEMINI_MODEL, GEMINI_API_KEY)
        
        # BATCH CLASSIFY ALL CONCEPTS IN ONE REQUEST (avoids rate limits!)
        concept_data = list(concepts_to_create.values())  # [(display_title, source_types, category), ...]
        concept_names = [data[0] for data in concept_data]
        logger.info(f"Batch classifying {len(concept_names)} concepts in ONE AI request...")
        
        try:
            classifications = agent.ai_classify_batch_concepts(concept_names)
        except Exception as e:
            logger.error(f"Batch classification failed: {e}, falling back to basic classification")
            classifications = [
                {
                    "category": data[2],
                    "subcategory": "",
                    "advanced_topic": data[0],
                    "related_concepts": [],
                    "prerequisites": []
                }
                for data in concept_data
            ]
        
        # Create concept nodes with AI-powered classification
        concept_map = {}
        connections_to_create = []
        
        for i, (display_title, source_types, category) in enumerate(concept_data):
            # Get classification for this concept
            classification = classifications[i] if i < len(classifications) else {}
            
            # Use AI classification - prefer the MOST SPECIFIC category
            ai_category_raw = classification.get("category", category)
            subcategory = classification.get("subcategory", "")
            advanced_topic = classification.get("advanced_topic", display_title)
            
            # Smart category selection: use the most specific one
            if ai_category_raw and ai_category_raw not in ["General", "Academic", "Discussion"]:
                ai_category = ai_category_raw
                logger.info(f"✓ '{display_title}' → '{ai_category}' (from AI category)")
            elif subcategory and subcategory not in ["General", "Academic", "Discussion"]:
                ai_category = subcategory
                logger.info(f"✓ '{display_title}' → '{subcategory}' (from AI subcategory)")
            else:
                ai_category = category
                logger.info(f"✓ '{display_title}' → '{category}' (fallback)")
            
            # Enhanced description showing all source types
            source_str = ", ".join(source_types)
            description = f"{advanced_topic}"
            if subcategory and subcategory != advanced_topic and subcategory != ai_category:
                description = f"{subcategory}: {advanced_topic}"
            description += f" (from {source_str})"
            
            # Store related concepts and prerequisites for connection creation
            related = classification.get("related_concepts", [])
            prereqs = classification.get("prerequisites", [])
            
            # Create the node
            node = models.ConceptNode(
                user_id=user.id,
                concept_name=display_title,
                description=description,
                category=ai_category,
                importance_score=0.7
            )
            db.add(node)
            db.flush()
            concept_map[display_title] = node.id
            
            connections_to_create.append((node.id, display_title, related, prereqs))
        
        # Create connections between related concepts
        connections_created = 0
        
        # Get all nodes for connection creation
        all_nodes = db.query(models.ConceptNode).filter(
            models.ConceptNode.user_id == user.id
        ).all()
        
        for node_id, concept_name, related_concepts, prerequisites in connections_to_create:
            current_node = next((n for n in all_nodes if n.id == node_id), None)
            if not current_node:
                continue
            
            # Strategy 1: AI-suggested related concepts
            for related_name in related_concepts[:3]:  # Top 3 related
                for other_node in all_nodes:
                    if other_node.id != node_id and related_name.lower() in other_node.concept_name.lower():
                        conn = models.ConceptConnection(
                            user_id=user.id,
                            source_concept_id=node_id,
                            target_concept_id=other_node.id,
                            connection_type="related",
                            strength=0.7,
                            ai_generated=True
                        )
                        db.add(conn)
                        connections_created += 1
                        break
            
            # Strategy 2: AI-suggested prerequisites
            for prereq_name in prerequisites[:2]:  # Top 2 prerequisites
                for other_node in all_nodes:
                    if other_node.id != node_id and prereq_name.lower() in other_node.concept_name.lower():
                        conn = models.ConceptConnection(
                            user_id=user.id,
                            source_concept_id=other_node.id,  # Prerequisite points to this concept
                            target_concept_id=node_id,
                            connection_type="prerequisite",
                            strength=0.8,
                            ai_generated=True
                        )
                        db.add(conn)
                        connections_created += 1
                        break
            
            # Strategy 3: Same category connections (create web within categories)
            same_category_nodes = [n for n in all_nodes if n.id != node_id and n.category == current_node.category]
            for other_node in same_category_nodes[:2]:  # Connect to 2 nodes in same category
                # Check if connection already exists
                existing = db.query(models.ConceptConnection).filter(
                    models.ConceptConnection.user_id == user.id,
                    models.ConceptConnection.source_concept_id == node_id,
                    models.ConceptConnection.target_concept_id == other_node.id
                ).first()
                
                if not existing:
                    conn = models.ConceptConnection(
                        user_id=user.id,
                        source_concept_id=node_id,
                        target_concept_id=other_node.id,
                        connection_type="similar",
                        strength=0.5,
                        ai_generated=True
                    )
                    db.add(conn)
                    connections_created += 1
        
        db.commit()
        logger.info(f"Successfully created {len(concept_map)} concepts and {connections_created} connections")
        
        return {
            "status": "success",
            "concepts_created": len(concept_map),
            "connections_created": connections_created
        }
        
    except Exception as e:
        logger.error(f"Error generating concept web: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/add_concept_node")
async def add_concept_node(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Manually add a concept node with AI-powered classification"""
    try:
        user_id = payload.get("user_id")
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        concept_name = payload.get("concept_name")
        description = payload.get("description", "")
        
        # Use AI to classify the concept with maximum specificity
        from concept_classification_agent import get_concept_agent
        agent = get_concept_agent(groq_client, GROQ_MODEL, gemini_client, GEMINI_MODEL, GEMINI_API_KEY)
        
        classification = agent.ai_classify_single_concept(concept_name, description)
        
        # Use AI classification - prefer subcategory as main category if it's specific
        base_category = classification.get("category", payload.get("category", "General"))
        subcategory = classification.get("subcategory", "")
        advanced_topic = classification.get("advanced_topic", concept_name)
        
        # Use subcategory as the main category if it's more specific
        if subcategory and subcategory not in ["General", base_category]:
            category = subcategory
        else:
            category = base_category
        
        # Create enhanced description
        enhanced_description = description
        if not enhanced_description:
            if subcategory and subcategory != category:
                enhanced_description = f"{subcategory}: {advanced_topic}"
            else:
                enhanced_description = f"{advanced_topic}"
        
        node = models.ConceptNode(
            user_id=user.id,
            concept_name=concept_name,
            description=enhanced_description,
            category=category
        )
        db.add(node)
        db.flush()
        
        # Auto-create connections to related concepts
        related_concepts = classification.get("related_concepts", [])
        prerequisites = classification.get("prerequisites", [])
        
        # Find existing nodes that match related concepts
        for related_name in related_concepts[:3]:  # Top 3 related
            related_node = db.query(models.ConceptNode).filter(
                models.ConceptNode.user_id == user.id,
                models.ConceptNode.concept_name.ilike(f"%{related_name}%")
            ).first()
            
            if related_node:
                # Create bidirectional "related" connection
                conn = models.ConceptConnection(
                    user_id=user.id,
                    source_concept_id=node.id,
                    target_concept_id=related_node.id,
                    connection_type="related",
                    strength=0.7,
                    ai_generated=True
                )
                db.add(conn)
        
        # Find prerequisite nodes
        for prereq_name in prerequisites[:2]:  # Top 2 prerequisites
            prereq_node = db.query(models.ConceptNode).filter(
                models.ConceptNode.user_id == user.id,
                models.ConceptNode.concept_name.ilike(f"%{prereq_name}%")
            ).first()
            
            if prereq_node:
                # Create "prerequisite" connection (prereq -> new concept)
                conn = models.ConceptConnection(
                    user_id=user.id,
                    source_concept_id=prereq_node.id,
                    target_concept_id=node.id,
                    connection_type="prerequisite",
                    strength=0.8,
                    ai_generated=True
                )
                db.add(conn)
        
        db.commit()
        db.refresh(node)
        
        return {
            "status": "success",
            "node_id": node.id,
            "concept_name": node.concept_name,
            "classification": {
                "category": category,
                "subcategory": subcategory,
                "advanced_topic": advanced_topic,
                "difficulty": classification.get("difficulty_level", "intermediate")
            },
            "connections_created": len(related_concepts) + len(prerequisites)
        }
    except Exception as e:
        logger.error(f"Error adding concept node: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/update_node_position")
async def update_node_position(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Update node position after drag"""
    try:
        node_id = payload.get("node_id")
        x = payload.get("x")
        y = payload.get("y")
        
        node = db.query(models.ConceptNode).filter(
            models.ConceptNode.id == node_id
        ).first()
        
        if node:
            node.position_x = x
            node.position_y = y
            db.commit()
            return {"status": "success"}
        
        return {"status": "not_found"}
    except Exception as e:
        logger.error(f"Error updating node position: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/update_concept_mastery")
async def update_concept_mastery(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Update concept mastery level"""
    try:
        node_id = payload.get("node_id")
        mastery_level = payload.get("mastery_level")
        
        if mastery_level < 0 or mastery_level > 1:
            raise HTTPException(status_code=400, detail="Mastery level must be between 0 and 1")
        
        node = db.query(models.ConceptNode).filter(
            models.ConceptNode.id == node_id
        ).first()
        
        if node:
            node.mastery_level = mastery_level
            db.commit()
            return {"status": "success", "mastery_level": mastery_level}
        
        return {"status": "not_found"}
    except Exception as e:
        logger.error(f"Error updating mastery level: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate_concept_notes")
async def generate_concept_notes(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Generate study notes for a specific concept"""
    try:
        user_id = payload.get("user_id")
        concept_id = payload.get("concept_id")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        concept = db.query(models.ConceptNode).filter(
            models.ConceptNode.id == concept_id,
            models.ConceptNode.user_id == user.id
        ).first()
        
        if not concept:
            raise HTTPException(status_code=404, detail="Concept not found")
        
        # Get related chat context
        recent_chats = db.query(models.ChatMessage).join(
            models.ChatSession
        ).filter(
            models.ChatSession.user_id == user.id
        ).order_by(models.ChatMessage.timestamp.desc()).limit(10).all()
        
        chat_context = "\n".join([
            f"Q: {msg.user_message}\nA: {msg.ai_response}" 
            for msg in recent_chats
        ])[:2000]
        
        # Generate notes
        prompt = f"""Create comprehensive study notes about: {concept.concept_name}

Description: {concept.description}
Category: {concept.category}

Recent learning context:
{chat_context}

Generate detailed notes covering:
1. Key concepts and definitions
2. Important points to remember
3. Examples and applications
4. Common misconceptions

Format as clear, organized study notes."""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1500
        )
        
        content = response.choices[0].message.content
        
        # Create note
        note = models.Note(
            user_id=user.id,
            title=f"Study Notes: {concept.concept_name}",
            content=content
        )
        db.add(note)
        db.flush()
        
        concept.notes_count += 1
        concept.mastery_level = min(1.0, concept.mastery_level + 0.1)
        
        db.commit()
        
        return {
            "status": "success",
            "note_id": note.id,
            "new_mastery": concept.mastery_level,
            "message": f"Generated notes for {concept.concept_name}"
        }
        
    except Exception as e:
        logger.error(f"Error generating notes: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate_concept_flashcards")
async def generate_concept_flashcards(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Generate flashcards for a specific concept"""
    try:
        user_id = payload.get("user_id")
        concept_id = payload.get("concept_id")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        concept = db.query(models.ConceptNode).filter(
            models.ConceptNode.id == concept_id,
            models.ConceptNode.user_id == user.id
        ).first()
        
        if not concept:
            raise HTTPException(status_code=404, detail="Concept not found")
        
        # Get related chat context
        recent_chats = db.query(models.ChatMessage).join(
            models.ChatSession
        ).filter(
            models.ChatSession.user_id == user.id
        ).order_by(models.ChatMessage.timestamp.desc()).limit(10).all()
        
        chat_context = "\n".join([
            f"Q: {msg.user_message}\nA: {msg.ai_response}" 
            for msg in recent_chats
        ])[:2000]
        
        # Generate flashcards
        prompt = f"""Create 10 flashcard pairs (question/answer) about: {concept.concept_name}

Description: {concept.description}
Category: {concept.category}

Recent learning context:
{chat_context}

Return ONLY a JSON array with this format:
[
  {{"front": "Question 1", "back": "Answer 1"}},
  {{"front": "Question 2", "back": "Answer 2"}}
]

Make questions clear and answers concise."""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1500
        )
        
        content = response.choices[0].message.content.strip()
        
        # Parse JSON
        import re
        json_match = re.search(r'\[[\s\S]*\]', content)
        if json_match:
            flashcards_data = json.loads(json_match.group())
        else:
            flashcards_data = json.loads(content)
        
        # Create flashcard set
        flashcard_set = models.FlashcardSet(
            user_id=user.id,
            title=f"Flashcards: {concept.concept_name}",
            description=f"AI-generated flashcards for {concept.concept_name}"
        )
        db.add(flashcard_set)
        db.flush()
        
        # Create individual flashcards
        for card_data in flashcards_data[:10]:
            flashcard = models.Flashcard(
                set_id=flashcard_set.id,
                question=card_data.get("front", ""),
                answer=card_data.get("back", "")
            )
            db.add(flashcard)
        
        concept.flashcards_count += 1
        concept.mastery_level = min(1.0, concept.mastery_level + 0.1)
        
        db.commit()
        
        return {
            "status": "success",
            "set_id": flashcard_set.id,
            "new_mastery": concept.mastery_level,
            "message": f"Generated flashcards for {concept.concept_name}"
        }
        
    except Exception as e:
        logger.error(f"Error generating flashcards: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate_concept_quiz")
async def generate_concept_quiz(
    payload: dict = Body(...),
    db: Session = Depends(get_db)
):
    """Generate quiz for a specific concept"""
    try:
        user_id = payload.get("user_id")
        concept_id = payload.get("concept_id")
        
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        concept = db.query(models.ConceptNode).filter(
            models.ConceptNode.id == concept_id,
            models.ConceptNode.user_id == user.id
        ).first()
        
        if not concept:
            raise HTTPException(status_code=404, detail="Concept not found")
        
        # Get related chat context
        recent_chats = db.query(models.ChatMessage).join(
            models.ChatSession
        ).filter(
            models.ChatSession.user_id == user.id
        ).order_by(models.ChatMessage.timestamp.desc()).limit(10).all()
        
        chat_context = "\n".join([
            f"Q: {msg.user_message}\nA: {msg.ai_response}" 
            for msg in recent_chats
        ])[:2000]
        
        # Generate quiz
        prompt = f"""Create 5 multiple choice questions about: {concept.concept_name}

Description: {concept.description}
Category: {concept.category}

Recent learning context:
{chat_context}

Return ONLY a JSON array:
[
  {{
    "question": "Question text",
    "options": ["A", "B", "C", "D"],
    "correct": 0,
    "explanation": "Why this is correct"
  }}
]"""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1500
        )
        
        content = response.choices[0].message.content.strip()
        
        # Parse JSON
        import re
        json_match = re.search(r'\[[\s\S]*\]', content)
        if json_match:
            questions_data = json.loads(json_match.group())
        else:
            questions_data = json.loads(content)
        
        # Create quiz
        quiz = models.SoloQuiz(
            user_id=user.id,
            subject=concept.concept_name,
            difficulty="intermediate",
            question_count=len(questions_data),
            answers=json.dumps(questions_data)  # Store questions as JSON
        )
        db.add(quiz)
        db.flush()
        
        concept.quizzes_count += 1
        concept.mastery_level = min(1.0, concept.mastery_level + 0.1)
        
        db.commit()
        
        return {
            "status": "success",
            "quiz_id": quiz.id,
            "questions": questions_data,
            "new_mastery": concept.mastery_level,
            "message": f"Generated quiz for {concept.concept_name}"
        }
        
    except Exception as e:
        logger.error(f"Error generating quiz: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/delete_concept_node/{node_id}")
async def delete_concept_node(
    node_id: int,
    db: Session = Depends(get_db)
):
    """Delete a concept node and its connections"""
    try:
        # Delete connections
        db.query(models.ConceptConnection).filter(
            (models.ConceptConnection.source_concept_id == node_id) |
            (models.ConceptConnection.target_concept_id == node_id)
        ).delete()
        
        # Delete node
        db.query(models.ConceptNode).filter(
            models.ConceptNode.id == node_id
        ).delete()
        
        db.commit()
        return {"status": "success"}
    except Exception as e:
        logger.error(f"Error deleting concept node: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/delete_all_concepts")
async def delete_all_concepts(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Delete all concept nodes and connections for a user"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Delete all connections
        db.query(models.ConceptConnection).filter(
            models.ConceptConnection.user_id == user.id
        ).delete()
        
        # Delete all nodes
        db.query(models.ConceptNode).filter(
            models.ConceptNode.user_id == user.id
        ).delete()
        
        db.commit()
        return {"status": "success", "message": "All concepts deleted"}
    except Exception as e:
        logger.error(f"Error deleting all concepts: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

# ==================== REMINDER/CALENDAR EVENTS API (Apple Reminders Style) ====================

def serialize_reminder(r):
    """Helper to serialize a reminder object"""
    # Debug logging for timezone issues (only log first few to avoid spam)
    if r.reminder_date and r.id % 10 == 1:  # Only log every 10th reminder
        logger.info(f"📤 SERIALIZING REMINDER: {r.title}")
        logger.info(f"   Stored date: {r.reminder_date}")
        logger.info(f"   Serialized as: {r.reminder_date.isoformat()} (no Z suffix = local time)")
    
    return {
        "id": r.id,
        "list_id": r.list_id,
        "parent_id": r.parent_id,
        "title": r.title,
        "description": r.description,
        "notes": r.notes,
        "url": r.url,
        "reminder_date": r.reminder_date.isoformat() if r.reminder_date else None,
        "due_date": r.due_date.isoformat() if r.due_date else None,
        "reminder_type": r.reminder_type,
        "priority": r.priority,
        "color": r.color,
        "is_completed": r.is_completed,
        "completed_at": r.completed_at.isoformat() + 'Z' if r.completed_at else None,
        "is_flagged": r.is_flagged,
        "is_notified": r.is_notified,
        "notify_before_minutes": r.notify_before_minutes,
        "recurring": r.recurring,
        "recurring_interval": r.recurring_interval,
        "recurring_end_date": r.recurring_end_date.isoformat() if r.recurring_end_date else None,
        "location": r.location,
        "tags": json.loads(r.tags) if r.tags else [],
        "sort_order": r.sort_order,
        "created_at": r.created_at.isoformat() + 'Z',
        "subtasks": [serialize_reminder(s) for s in r.subtasks] if r.subtasks else []
    }

# ==================== REMINDER LISTS API ====================

@app.post("/api/create_reminder_list")
async def create_reminder_list(
    user_id: str = Form(...),
    name: str = Form(...),
    color: str = Form("#3b82f6"),
    icon: str = Form("list"),
    db: Session = Depends(get_db)
):
    """Create a new reminder list"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Get max sort order
        max_order = db.query(func.max(models.ReminderList.sort_order)).filter(
            models.ReminderList.user_id == user.id
        ).scalar() or 0
        
        reminder_list = models.ReminderList(
            user_id=user.id,
            name=name,
            color=color,
            icon=icon,
            sort_order=max_order + 1
        )
        
        db.add(reminder_list)
        db.commit()
        db.refresh(reminder_list)
        
        return {
            "id": reminder_list.id,
            "name": reminder_list.name,
            "color": reminder_list.color,
            "icon": reminder_list.icon,
            "sort_order": reminder_list.sort_order,
            "reminder_count": 0
        }
    except Exception as e:
        logger.error(f"Error creating reminder list: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_reminder_lists")
async def get_reminder_lists(
    user_id: str = Query(...),
    db: Session = Depends(get_db)
):
    """Get all reminder lists for a user with counts"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        lists = db.query(models.ReminderList).filter(
            models.ReminderList.user_id == user.id,
            models.ReminderList.is_smart_list == False
        ).order_by(models.ReminderList.sort_order).all()
        
        # Get counts for each list
        result = []
        for lst in lists:
            count = db.query(models.Reminder).filter(
                models.Reminder.list_id == lst.id,
                models.Reminder.is_completed == False,
                models.Reminder.parent_id == None
            ).count()
            
            result.append({
                "id": lst.id,
                "name": lst.name,
                "color": lst.color,
                "icon": lst.icon,
                "sort_order": lst.sort_order,
                "reminder_count": count
            })
        
        # Add smart list counts
        today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
        tomorrow = today + timedelta(days=1)
        
        today_count = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == False,
            models.Reminder.reminder_date >= today,
            models.Reminder.reminder_date < tomorrow
        ).count()
        
        scheduled_count = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == False,
            models.Reminder.reminder_date != None
        ).count()
        
        flagged_count = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == False,
            models.Reminder.is_flagged == True
        ).count()
        
        all_count = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == False,
            models.Reminder.parent_id == None
        ).count()
        
        completed_count = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.is_completed == True
        ).count()
        
        smart_lists = {
            "today": today_count,
            "scheduled": scheduled_count,
            "flagged": flagged_count,
            "all": all_count,
            "completed": completed_count
        }
        
        return {
            "lists": result,
            "smart_lists": smart_lists
        }
    except Exception as e:
        logger.error(f"Error getting reminder lists: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/update_reminder_list/{list_id}")
async def update_reminder_list(
    list_id: int,
    name: str = Form(None),
    color: str = Form(None),
    icon: str = Form(None),
    db: Session = Depends(get_db)
):
    """Update a reminder list"""
    try:
        reminder_list = db.query(models.ReminderList).filter(models.ReminderList.id == list_id).first()
        if not reminder_list:
            raise HTTPException(status_code=404, detail="List not found")
        
        if name is not None:
            reminder_list.name = name
        if color is not None:
            reminder_list.color = color
        if icon is not None:
            reminder_list.icon = icon
        
        reminder_list.updated_at = datetime.utcnow()
        db.commit()
        
        return {"status": "success", "message": "List updated"}
    except Exception as e:
        logger.error(f"Error updating reminder list: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/delete_reminder_list/{list_id}")
async def delete_reminder_list(
    list_id: int,
    db: Session = Depends(get_db)
):
    """Delete a reminder list and all its reminders"""
    try:
        reminder_list = db.query(models.ReminderList).filter(models.ReminderList.id == list_id).first()
        if not reminder_list:
            raise HTTPException(status_code=404, detail="List not found")
        
        db.delete(reminder_list)
        db.commit()
        
        return {"status": "success", "message": "List deleted"}
    except Exception as e:
        logger.error(f"Error deleting reminder list: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

# ==================== REMINDERS API ====================

@app.post("/api/create_reminder")
async def create_reminder(
    user_id: str = Form(...),
    title: str = Form(...),
    description: str = Form(None),
    notes: str = Form(None),
    url: str = Form(None),
    reminder_date: str = Form(None),
    due_date: str = Form(None),
    reminder_type: str = Form("reminder"),
    priority: str = Form("none"),
    color: str = Form("#3b82f6"),
    is_flagged: bool = Form(False),
    notify_before_minutes: int = Form(15),
    list_id: int = Form(None),
    parent_id: int = Form(None),
    recurring: str = Form("none"),
    recurring_interval: int = Form(1),
    recurring_end_date: str = Form(None),
    location: str = Form(None),
    tags: str = Form(None),
    user_timezone: str = Form("UTC"),
    timezone_offset: int = Form(0),
    db: Session = Depends(get_db)
):
    """Create a new reminder with Apple Reminders-style features
    
    Note: Reminder dates are stored as naive datetime objects representing the user's local time.
    This approach ensures that when a user sets a reminder for 9:35 AM, it triggers at 9:35 AM
    in their local timezone, regardless of server timezone or daylight saving changes.
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Parse dates - treat as local time, not UTC
        parsed_reminder_date = None
        parsed_due_date = None
        parsed_recurring_end = None
        
        if reminder_date:
            # Parse as naive datetime (local time) and store as-is
            # The datetime-local input sends time in user's local timezone
            parsed_reminder_date = datetime.fromisoformat(reminder_date.replace('Z', '').replace('+00:00', ''))
            logger.info(f"📅 REMINDER CREATION DEBUG:")
            logger.info(f"   Raw input: {reminder_date}")
            logger.info(f"   Parsed as: {parsed_reminder_date} (treated as local time)")
            logger.info(f"   User timezone: {user_timezone}")
            logger.info(f"   Timezone offset: {timezone_offset} minutes")
        if due_date:
            parsed_due_date = datetime.fromisoformat(due_date.replace('Z', '').replace('+00:00', ''))
            logger.info(f"Parsed due_date: {parsed_due_date} (treated as local time)")
        if recurring_end_date:
            parsed_recurring_end = datetime.fromisoformat(recurring_end_date.replace('Z', '').replace('+00:00', ''))
            logger.info(f"Parsed recurring_end_date: {parsed_recurring_end} (treated as local time)")
        
        # Get max sort order for the list
        max_order = db.query(func.max(models.Reminder.sort_order)).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.list_id == list_id
        ).scalar() or 0
        
        reminder = models.Reminder(
            user_id=user.id,
            list_id=list_id,
            parent_id=parent_id,
            title=title,
            description=description,
            notes=notes,
            url=url,
            reminder_date=parsed_reminder_date,
            due_date=parsed_due_date,
            reminder_type=reminder_type,
            priority=priority,
            color=color,
            is_flagged=is_flagged,
            notify_before_minutes=notify_before_minutes,
            recurring=recurring,
            recurring_interval=recurring_interval,
            recurring_end_date=parsed_recurring_end,
            location=location,
            tags=tags,
            sort_order=max_order + 1
        )
        
        db.add(reminder)
        db.commit()
        db.refresh(reminder)
        
        # Don't create immediate notification - let the reminder notification system handle it
        # when the time comes (notify_before_minutes before the reminder_date)
        
        logger.info(f"Created reminder {reminder.id} for user {user.email}")
        
        return serialize_reminder(reminder)
    except Exception as e:
        logger.error(f"Error creating reminder: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_reminders")
async def get_reminders(
    user_id: str = Query(...),
    list_id: int = Query(None),
    smart_list: str = Query(None),  # today, scheduled, flagged, all, completed
    start_date: str = Query(None),
    end_date: str = Query(None),
    include_completed: bool = Query(False),
    db: Session = Depends(get_db)
):
    """Get reminders with smart list filtering"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        query = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            models.Reminder.parent_id == None  # Only top-level reminders
        )
        
        today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
        tomorrow = today + timedelta(days=1)
        
        # Smart list filtering
        if smart_list == "today":
            query = query.filter(
                models.Reminder.is_completed == False,
                models.Reminder.reminder_date >= today,
                models.Reminder.reminder_date < tomorrow
            )
        elif smart_list == "scheduled":
            query = query.filter(
                models.Reminder.is_completed == False,
                models.Reminder.reminder_date != None
            )
        elif smart_list == "flagged":
            query = query.filter(
                models.Reminder.is_completed == False,
                models.Reminder.is_flagged == True
            )
        elif smart_list == "all":
            query = query.filter(models.Reminder.is_completed == False)
        elif smart_list == "completed":
            query = query.filter(models.Reminder.is_completed == True)
        elif list_id:
            query = query.filter(models.Reminder.list_id == list_id)
            if not include_completed:
                query = query.filter(models.Reminder.is_completed == False)
        else:
            if not include_completed:
                query = query.filter(models.Reminder.is_completed == False)
        
        # Date range filtering
        if start_date:
            query = query.filter(models.Reminder.reminder_date >= datetime.fromisoformat(start_date.replace('Z', '').replace('+00:00', '')))
        if end_date:
            query = query.filter(models.Reminder.reminder_date <= datetime.fromisoformat(end_date.replace('Z', '').replace('+00:00', '')))
        
        reminders = query.order_by(
            models.Reminder.is_completed,
            models.Reminder.reminder_date.nullslast(),
            models.Reminder.sort_order
        ).all()
        
        return [serialize_reminder(r) for r in reminders]
    except Exception as e:
        logger.error(f"Error getting reminders: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/update_reminder/{reminder_id}")
async def update_reminder(
    reminder_id: int,
    title: str = Form(None),
    description: str = Form(None),
    notes: str = Form(None),
    url: str = Form(None),
    reminder_date: str = Form(None),
    due_date: str = Form(None),
    reminder_type: str = Form(None),
    priority: str = Form(None),
    color: str = Form(None),
    is_completed: bool = Form(None),
    is_flagged: bool = Form(None),
    notify_before_minutes: int = Form(None),
    list_id: int = Form(None),
    recurring: str = Form(None),
    recurring_interval: int = Form(None),
    recurring_end_date: str = Form(None),
    location: str = Form(None),
    tags: str = Form(None),
    db: Session = Depends(get_db)
):
    """Update a reminder with all Apple Reminders features"""
    try:
        reminder = db.query(models.Reminder).filter(models.Reminder.id == reminder_id).first()
        if not reminder:
            raise HTTPException(status_code=404, detail="Reminder not found")
        
        if title is not None:
            reminder.title = title
        if description is not None:
            reminder.description = description
        if notes is not None:
            reminder.notes = notes
        if url is not None:
            reminder.url = url
        if reminder_date is not None:
            if reminder_date == '':
                reminder.reminder_date = None
            else:
                reminder.reminder_date = datetime.fromisoformat(reminder_date.replace('Z', '').replace('+00:00', ''))
        if due_date is not None:
            if due_date == '':
                reminder.due_date = None
            else:
                reminder.due_date = datetime.fromisoformat(due_date.replace('Z', '').replace('+00:00', ''))
        if reminder_type is not None:
            reminder.reminder_type = reminder_type
        if priority is not None:
            reminder.priority = priority
        if color is not None:
            reminder.color = color
        if is_completed is not None:
            reminder.is_completed = is_completed
            if is_completed:
                reminder.completed_at = datetime.utcnow()
                # Handle recurring reminders
                if reminder.recurring != "none" and reminder.recurring:
                    await create_next_recurring_reminder(db, reminder)
            else:
                reminder.completed_at = None
        if is_flagged is not None:
            reminder.is_flagged = is_flagged
        if notify_before_minutes is not None:
            reminder.notify_before_minutes = notify_before_minutes
        if list_id is not None:
            reminder.list_id = list_id if list_id > 0 else None
        if recurring is not None:
            reminder.recurring = recurring
        if recurring_interval is not None:
            reminder.recurring_interval = recurring_interval
        if recurring_end_date is not None:
            if recurring_end_date == '':
                reminder.recurring_end_date = None
            else:
                reminder.recurring_end_date = datetime.fromisoformat(recurring_end_date.replace('Z', '').replace('+00:00', ''))
        if location is not None:
            reminder.location = location
        if tags is not None:
            reminder.tags = tags
        
        reminder.updated_at = datetime.utcnow()
        db.commit()
        
        return {"status": "success", "message": "Reminder updated", "reminder": serialize_reminder(reminder)}
    except Exception as e:
        logger.error(f"Error updating reminder: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

async def create_next_recurring_reminder(db: Session, original: models.Reminder):
    """Create the next occurrence of a recurring reminder"""
    if not original.reminder_date or original.recurring == "none":
        return
    
    next_date = original.reminder_date
    interval = original.recurring_interval or 1
    
    if original.recurring == "daily":
        next_date += timedelta(days=interval)
    elif original.recurring == "weekly":
        next_date += timedelta(weeks=interval)
    elif original.recurring == "monthly":
        # Add months
        month = next_date.month + interval
        year = next_date.year + (month - 1) // 12
        month = ((month - 1) % 12) + 1
        day = min(next_date.day, [31, 29 if year % 4 == 0 else 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31][month - 1])
        next_date = next_date.replace(year=year, month=month, day=day)
    elif original.recurring == "yearly":
        next_date = next_date.replace(year=next_date.year + interval)
    
    # Check if past end date
    if original.recurring_end_date and next_date > original.recurring_end_date:
        return
    
    # Create new reminder
    new_reminder = models.Reminder(
        user_id=original.user_id,
        list_id=original.list_id,
        title=original.title,
        description=original.description,
        notes=original.notes,
        url=original.url,
        reminder_date=next_date,
        due_date=original.due_date,
        reminder_type=original.reminder_type,
        priority=original.priority,
        color=original.color,
        is_flagged=original.is_flagged,
        notify_before_minutes=original.notify_before_minutes,
        recurring=original.recurring,
        recurring_interval=original.recurring_interval,
        recurring_end_date=original.recurring_end_date,
        location=original.location,
        tags=original.tags
    )
    
    db.add(new_reminder)
    db.commit()

@app.delete("/api/delete_reminder/{reminder_id}")
async def delete_reminder(
    reminder_id: int,
    db: Session = Depends(get_db)
):
    """Delete a reminder and its subtasks"""
    try:
        reminder = db.query(models.Reminder).filter(models.Reminder.id == reminder_id).first()
        if not reminder:
            raise HTTPException(status_code=404, detail="Reminder not found")
        
        db.delete(reminder)
        db.commit()
        
        return {"status": "success", "message": "Reminder deleted"}
    except Exception as e:
        logger.error(f"Error deleting reminder: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/add_subtask/{reminder_id}")
async def add_subtask(
    reminder_id: int,
    title: str = Form(...),
    db: Session = Depends(get_db)
):
    """Add a subtask to a reminder"""
    try:
        parent = db.query(models.Reminder).filter(models.Reminder.id == reminder_id).first()
        if not parent:
            raise HTTPException(status_code=404, detail="Parent reminder not found")
        
        subtask = models.Reminder(
            user_id=parent.user_id,
            parent_id=reminder_id,
            list_id=parent.list_id,
            title=title,
            color=parent.color
        )
        
        db.add(subtask)
        db.commit()
        db.refresh(subtask)
        
        return serialize_reminder(subtask)
    except Exception as e:
        logger.error(f"Error adding subtask: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/toggle_reminder_flag/{reminder_id}")
async def toggle_reminder_flag(
    reminder_id: int,
    db: Session = Depends(get_db)
):
    """Toggle the flagged status of a reminder"""
    try:
        reminder = db.query(models.Reminder).filter(models.Reminder.id == reminder_id).first()
        if not reminder:
            raise HTTPException(status_code=404, detail="Reminder not found")
        
        reminder.is_flagged = not reminder.is_flagged
        reminder.updated_at = datetime.utcnow()
        db.commit()
        
        return {"status": "success", "is_flagged": reminder.is_flagged}
    except Exception as e:
        logger.error(f"Error toggling flag: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/get_upcoming_reminders")
async def get_upcoming_reminders(
    user_id: str = Query(...),
    hours: int = Query(24),
    db: Session = Depends(get_db)
):
    """Get upcoming reminders within the next X hours"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        now = datetime.now()
        future = now + timedelta(hours=hours)
        
        reminders = db.query(models.Reminder).filter(
            and_(
                models.Reminder.user_id == user.id,
                models.Reminder.reminder_date >= now,
                models.Reminder.reminder_date <= future,
                models.Reminder.is_completed == False
            )
        ).order_by(models.Reminder.reminder_date).all()
        
        return [serialize_reminder(r) for r in reminders]
    except Exception as e:
        logger.error(f"Error getting upcoming reminders: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/search_reminders")
async def search_reminders(
    user_id: str = Query(...),
    query: str = Query(...),
    db: Session = Depends(get_db)
):
    """Search reminders by title or description"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        reminders = db.query(models.Reminder).filter(
            models.Reminder.user_id == user.id,
            (models.Reminder.title.ilike(f"%{query}%") | models.Reminder.description.ilike(f"%{query}%"))
        ).order_by(models.Reminder.reminder_date.nullslast()).all()
        
        return [serialize_reminder(r) for r in reminders]
    except Exception as e:
        logger.error(f"Error searching reminders: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))





# ==================== FILE IMPORT API ====================

@app.post("/api/import_document")
async def import_document(
    file: UploadFile = File(...),
    user_id: str = Form(...),
    db: Session = Depends(get_db)
):
    """Import PDF or DOCX file and convert to note"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Read file content
        content = await file.read()
        file_extension = file.filename.split('.')[-1].lower()
        
        extracted_text = ""
        
        # Extract text based on file type
        if file_extension == 'pdf':
            # Extract from PDF
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() + "\n\n"
        
        elif file_extension in ['docx', 'doc']:
            # Extract from DOCX
            try:
                import docx
                docx_file = io.BytesIO(content)
                doc = docx.Document(docx_file)
                
                for paragraph in doc.paragraphs:
                    extracted_text += paragraph.text + "\n"
            except ImportError:
                raise HTTPException(
                    status_code=500, 
                    detail="python-docx not installed. Install with: pip install python-docx"
                )
        
        else:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type: {file_extension}. Supported: PDF, DOCX"
            )
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the file")
        
        # Convert to HTML blocks
        paragraphs = extracted_text.strip().split('\n')
        html_content = ""
        
        for para in paragraphs:
            para = para.strip()
            if para:
                # Check if it looks like a heading (short, all caps, or ends with colon)
                if len(para) < 50 and (para.isupper() or para.endswith(':')):
                    html_content += f"<h2>{para}</h2>"
                else:
                    html_content += f"<p>{para}</p>"
        
        # Create note
        note_title = file.filename.rsplit('.', 1)[0]
        new_note = models.Note(
            user_id=user.id,
            title=note_title,
            content=html_content
        )
        
        db.add(new_note)
        db.commit()
        db.refresh(new_note)
        
        logger.info(f"Imported {file_extension.upper()} file as note {new_note.id} for user {user.email}")
        
        return {
            "status": "success",
            "note_id": new_note.id,
            "title": note_title,
            "extracted_length": len(extracted_text),
            "message": f"Successfully imported {file.filename}"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error importing document: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to import document: {str(e)}")


@app.post("/api/upload-attachment")
async def upload_attachment(
    file: UploadFile = File(...),
    db: Session = Depends(get_db)
):
    """Upload PDF or DOCX file as attachment and return URL"""
    try:
        # Validate file type
        file_extension = file.filename.split('.')[-1].lower()
        if file_extension not in ['pdf', 'docx', 'doc']:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type: {file_extension}. Supported: PDF, DOCX"
            )
        
        # Create attachments directory if it doesn't exist
        attachments_dir = Path("backend/attachments")
        attachments_dir.mkdir(exist_ok=True)
        
        # Generate unique filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_filename = "".join(c for c in file.filename if c.isalnum() or c in (' ', '.', '_', '-')).rstrip()
        unique_filename = f"{timestamp}_{safe_filename}"
        file_path = attachments_dir / unique_filename
        
        # Save file
        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)
        
        # Return file info
        file_size = len(content)
        file_url = f"/api/attachments/{unique_filename}"
        
        logger.info(f"Uploaded attachment: {unique_filename} ({file_size} bytes)")
        
        return {
            "status": "success",
            "filename": file.filename,
            "url": file_url,
            "size": file_size,
            "type": file_extension
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading attachment: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to upload attachment: {str(e)}")


@app.get("/api/attachments/{filename}")
async def get_attachment(filename: str):
    """Serve attachment file"""
    try:
        file_path = Path("backend/attachments") / filename
        
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="File not found")
        
        # Determine content type
        file_extension = filename.split('.')[-1].lower()
        content_types = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'doc': 'application/msword'
        }
        content_type = content_types.get(file_extension, 'application/octet-stream')
        
        return FileResponse(
            path=file_path,
            media_type=content_type,
            filename=filename
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error serving attachment: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to serve attachment: {str(e)}")

# ==================== LEARNING PLAYLIST API ====================
logger.info(" Registering Learning Playlist API endpoints...")

@app.get("/api/playlists/test")
async def test_playlist_endpoint():
    """Test endpoint to verify playlist routes are working"""
    return {"message": "Playlist API is working!"}

@app.get("/api/playlists")
async def get_playlists(
    category: Optional[str] = None,
    difficulty: Optional[str] = None,
    search: Optional[str] = None,
    my_playlists: bool = False,
    following: bool = False,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get playlists with filters"""
    try:
        query = db.query(models.LearningPlaylist)
        
        if my_playlists:
            query = query.filter(models.LearningPlaylist.creator_id == current_user.id)
        elif following:
            followed_ids = db.query(models.PlaylistFollower.playlist_id).filter(
                models.PlaylistFollower.user_id == current_user.id
            ).all()
            followed_ids = [f[0] for f in followed_ids]
            query = query.filter(models.LearningPlaylist.id.in_(followed_ids))
        else:
            query = query.filter(models.LearningPlaylist.is_public == True)
        
        if category:
            query = query.filter(models.LearningPlaylist.category == category)
        if difficulty:
            query = query.filter(models.LearningPlaylist.difficulty_level == difficulty)
        if search:
            search_term = f"%{search}%"
            query = query.filter(
                or_(
                    models.LearningPlaylist.title.ilike(search_term),
                    models.LearningPlaylist.description.ilike(search_term)
                )
            )
        
        playlists = query.order_by(desc(models.LearningPlaylist.created_at)).all()
        
        result = []
        for p in playlists:
            # Check if user is following
            is_following = db.query(models.PlaylistFollower).filter(
                and_(
                    models.PlaylistFollower.playlist_id == p.id,
                    models.PlaylistFollower.user_id == current_user.id
                )
            ).first() is not None
            
            # Get user progress if following
            user_progress = None
            if is_following:
                follower = db.query(models.PlaylistFollower).filter(
                    and_(
                        models.PlaylistFollower.playlist_id == p.id,
                        models.PlaylistFollower.user_id == current_user.id
                    )
                ).first()
                if follower:
                    user_progress = {
                        "progress_percentage": follower.progress_percentage or 0,
                        "completed_items": follower.completed_items or []
                    }
            
            # Get item count
            item_count = db.query(models.PlaylistItem).filter(
                models.PlaylistItem.playlist_id == p.id
            ).count()
            
            result.append({
                "id": p.id,
                "title": p.title,
                "description": p.description,
                "category": p.category,
                "difficulty_level": p.difficulty_level,
                "estimated_hours": p.estimated_hours,
                "is_public": p.is_public,
                "cover_color": p.cover_color,
                "tags": p.tags or [],
                "fork_count": p.fork_count or 0,
                "follower_count": p.follower_count or 0,
                "completion_count": p.completion_count or 0,
                "item_count": item_count,
                "created_at": p.created_at.isoformat() if p.created_at else None,
                "creator": {
                    "id": p.creator.id,
                    "username": p.creator.username,
                    "first_name": p.creator.first_name,
                    "picture_url": p.creator.picture_url
                },
                "items": [],
                "is_owner": current_user.id == p.creator_id,
                "is_following": is_following,
                "user_progress": user_progress
            })
        
        return {"playlists": result}
    except Exception as e:
        logger.error(f"Error getting playlists: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

from pydantic import field_validator

class PlaylistCreateRequest(BaseModel):
    model_config = {"extra": "ignore"}
    
    title: str
    description: Optional[str] = None
    category: Optional[str] = None
    difficulty_level: str = "intermediate"
    estimated_hours: Optional[float] = None
    is_public: bool = True
    is_collaborative: bool = False
    cover_color: str = "#4A90E2"
    tags: Optional[List[str]] = None
    items: Optional[List] = None  # Allow items but ignore them for now
    
    @field_validator('estimated_hours', mode='before')
    @classmethod
    def empty_str_to_none(cls, v):
        if v == '' or v is None:
            return None
        return v

@app.post("/api/playlists")
async def create_playlist(
    request: Request,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Create a new learning playlist"""
    try:
        # Get raw body for debugging
        body = await request.json()
        logger.info(f" Received playlist data: {body}")
        
        # Parse into model
        playlist_data = PlaylistCreateRequest(**body)
        new_playlist = models.LearningPlaylist(
            creator_id=current_user.id,
            title=playlist_data.title,
            description=playlist_data.description,
            category=playlist_data.category,
            difficulty_level=playlist_data.difficulty_level,
            estimated_hours=playlist_data.estimated_hours,
            is_public=playlist_data.is_public,
            is_collaborative=playlist_data.is_collaborative,
            cover_color=playlist_data.cover_color,
            tags=playlist_data.tags or []
        )
        
        db.add(new_playlist)
        db.commit()
        db.refresh(new_playlist)
        
        return {
            "id": new_playlist.id,
            "title": new_playlist.title,
            "description": new_playlist.description,
            "category": new_playlist.category,
            "difficulty_level": new_playlist.difficulty_level,
            "cover_color": new_playlist.cover_color,
            "creator": {
                "id": current_user.id,
                "username": current_user.username,
                "first_name": current_user.first_name,
                "picture_url": current_user.picture_url
            },
            "items": [],
            "is_owner": True,
            "message": "Playlist created successfully"
        }
    except Exception as e:
        db.rollback()
        logger.error(f"Error creating playlist: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/playlists/{playlist_id}")
async def get_playlist_detail(
    playlist_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get detailed playlist with items"""
    try:
        playlist = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if not playlist:
            raise HTTPException(status_code=404, detail="Playlist not found")
        
        # Get items
        items = db.query(models.PlaylistItem).filter(
            models.PlaylistItem.playlist_id == playlist_id
        ).order_by(models.PlaylistItem.order_index).all()
        
        # Check if user is following
        follower = db.query(models.PlaylistFollower).filter(
            and_(
                models.PlaylistFollower.playlist_id == playlist_id,
                models.PlaylistFollower.user_id == current_user.id
            )
        ).first()
        
        user_progress = None
        if follower:
            user_progress = {
                "progress_percentage": follower.progress_percentage or 0,
                "completed_items": follower.completed_items or [],
                "is_completed": follower.is_completed or False,
                "started_at": follower.started_at.isoformat() if follower.started_at else None,
                "last_accessed": follower.last_accessed.isoformat() if follower.last_accessed else None
            }
        
        return {
            "id": playlist.id,
            "title": playlist.title,
            "description": playlist.description,
            "category": playlist.category,
            "difficulty_level": playlist.difficulty_level,
            "estimated_hours": playlist.estimated_hours,
            "is_public": playlist.is_public,
            "is_collaborative": playlist.is_collaborative,
            "cover_color": playlist.cover_color,
            "tags": playlist.tags or [],
            "fork_count": playlist.fork_count or 0,
            "follower_count": playlist.follower_count or 0,
            "completion_count": playlist.completion_count or 0,
            "created_at": playlist.created_at.isoformat() if playlist.created_at else None,
            "creator": {
                "id": playlist.creator.id,
                "username": playlist.creator.username,
                "first_name": playlist.creator.first_name,
                "last_name": playlist.creator.last_name,
                "picture_url": playlist.creator.picture_url
            },
            "items": [{
                "id": item.id,
                "order_index": item.order_index,
                "item_type": item.item_type,
                "item_id": item.item_id,
                "title": item.title,
                "url": item.url,
                "description": item.description,
                "duration_minutes": item.duration_minutes,
                "platform": item.platform,
                "is_required": item.is_required,
                "notes": item.notes
            } for item in items],
            "is_owner": current_user.id == playlist.creator_id,
            "is_following": follower is not None,
            "user_progress": user_progress
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting playlist: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

class PlaylistItemRequest(BaseModel):
    model_config = {"extra": "ignore"}
    
    item_type: str
    item_id: Optional[int] = None
    title: str
    url: Optional[str] = None
    description: Optional[str] = None
    duration_minutes: Optional[int] = None
    platform: Optional[str] = None
    is_required: bool = True
    notes: Optional[str] = None

@app.post("/api/playlists/{playlist_id}/items")
async def add_playlist_item(
    playlist_id: int,
    item_data: PlaylistItemRequest,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Add an item to a playlist"""
    try:
        playlist = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if not playlist:
            raise HTTPException(status_code=404, detail="Playlist not found")
        
        if playlist.creator_id != current_user.id:
            raise HTTPException(status_code=403, detail="Only creator can add items")
        
        # Get max order index
        max_order = db.query(func.max(models.PlaylistItem.order_index)).filter(
            models.PlaylistItem.playlist_id == playlist_id
        ).scalar() or -1
        
        # Create item
        item = models.PlaylistItem(
            playlist_id=playlist_id,
            order_index=max_order + 1,
            item_type=item_data.item_type,
            item_id=item_data.item_id,
            title=item_data.title,
            url=item_data.url,
            description=item_data.description,
            duration_minutes=item_data.duration_minutes,
            platform=item_data.platform,
            is_required=item_data.is_required,
            notes=item_data.notes
        )
        
        db.add(item)
        db.commit()
        db.refresh(item)
        
        return {
            "id": item.id,
            "message": "Item added successfully"
        }
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error adding item: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/playlists/{playlist_id}/items/{item_id}")
async def delete_playlist_item(
    playlist_id: int,
    item_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Delete an item from a playlist"""
    try:
        playlist = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if not playlist:
            raise HTTPException(status_code=404, detail="Playlist not found")
        
        if playlist.creator_id != current_user.id:
            raise HTTPException(status_code=403, detail="Only creator can delete items")
        
        item = db.query(models.PlaylistItem).filter(
            and_(
                models.PlaylistItem.id == item_id,
                models.PlaylistItem.playlist_id == playlist_id
            )
        ).first()
        
        if not item:
            raise HTTPException(status_code=404, detail="Item not found")
        
        db.delete(item)
        db.commit()
        
        return {"message": "Item deleted successfully"}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error deleting item: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/playlists/{playlist_id}/follow")
async def follow_playlist(
    playlist_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Follow a playlist"""
    try:
        playlist = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if not playlist:
            raise HTTPException(status_code=404, detail="Playlist not found")
        
        # Check if already following
        existing = db.query(models.PlaylistFollower).filter(
            and_(
                models.PlaylistFollower.playlist_id == playlist_id,
                models.PlaylistFollower.user_id == current_user.id
            )
        ).first()
        
        if existing:
            return {"message": "Already following this playlist"}
        
        # Create follower record
        follower = models.PlaylistFollower(
            playlist_id=playlist_id,
            user_id=current_user.id,
            completed_items=[]
        )
        
        db.add(follower)
        
        # Update follower count
        playlist.follower_count = (playlist.follower_count or 0) + 1
        
        db.commit()
        
        return {"message": "Successfully following playlist"}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error following playlist: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/playlists/{playlist_id}/follow")
async def unfollow_playlist(
    playlist_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Unfollow a playlist"""
    try:
        follower = db.query(models.PlaylistFollower).filter(
            and_(
                models.PlaylistFollower.playlist_id == playlist_id,
                models.PlaylistFollower.user_id == current_user.id
            )
        ).first()
        
        if not follower:
            raise HTTPException(status_code=404, detail="Not following this playlist")
        
        db.delete(follower)
        
        # Update follower count
        playlist = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if playlist:
            playlist.follower_count = max(0, (playlist.follower_count or 0) - 1)
        
        db.commit()
        
        return {"message": "Successfully unfollowed playlist"}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error unfollowing playlist: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/playlists/{playlist_id}/fork")
async def fork_playlist(
    playlist_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Fork a playlist (create a copy)"""
    try:
        original = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if not original:
            raise HTTPException(status_code=404, detail="Playlist not found")
        
        # Create forked playlist
        forked = models.LearningPlaylist(
            creator_id=current_user.id,
            title=f"{original.title} (Fork)",
            description=original.description,
            category=original.category,
            difficulty_level=original.difficulty_level,
            estimated_hours=original.estimated_hours,
            is_public=False,  # Forks start as private
            is_collaborative=original.is_collaborative,
            cover_color=original.cover_color,
            tags=original.tags
        )
        
        db.add(forked)
        db.flush()
        
        # Copy items
        original_items = db.query(models.PlaylistItem).filter(
            models.PlaylistItem.playlist_id == playlist_id
        ).order_by(models.PlaylistItem.order_index).all()
        
        for item in original_items:
            forked_item = models.PlaylistItem(
                playlist_id=forked.id,
                order_index=item.order_index,
                item_type=item.item_type,
                item_id=item.item_id,
                title=item.title,
                url=item.url,
                description=item.description,
                duration_minutes=item.duration_minutes,
                is_required=item.is_required,
                notes=item.notes
            )
            db.add(forked_item)
        
        # Create fork record
        fork_record = models.PlaylistFork(
            original_playlist_id=playlist_id,
            forked_playlist_id=forked.id,
            forked_by_id=current_user.id
        )
        db.add(fork_record)
        
        # Update fork count
        original.fork_count = (original.fork_count or 0) + 1
        
        db.commit()
        db.refresh(forked)
        
        return {
            "id": forked.id,
            "message": "Playlist forked successfully"
        }
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error forking playlist: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/playlists/{playlist_id}/items/{item_id}/view")
async def view_playlist_item(
    playlist_id: int,
    item_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """View a playlist item's content (read-only for notes/chats)"""
    try:
        # Get the playlist
        playlist = db.query(models.LearningPlaylist).filter(
            models.LearningPlaylist.id == playlist_id
        ).first()
        
        if not playlist:
            raise HTTPException(status_code=404, detail="Playlist not found")
        
        # Check if playlist is public or user has access
        if not playlist.is_public and playlist.creator_id != current_user.id:
            raise HTTPException(status_code=403, detail="Access denied")
        
        # Get the item
        item = db.query(models.PlaylistItem).filter(
            and_(
                models.PlaylistItem.id == item_id,
                models.PlaylistItem.playlist_id == playlist_id
            )
        ).first()
        
        if not item:
            raise HTTPException(status_code=404, detail="Item not found")
        
        # Return content based on item type
        if item.item_type == 'note' and item.item_id:
            note = db.query(models.Note).filter(models.Note.id == item.item_id).first()
            if note:
                return {
                    "type": "note",
                    "id": note.id,
                    "title": note.title,
                    "content": note.content,
                    "created_at": note.created_at.isoformat() if note.created_at else None,
                    "owner": {
                        "id": note.user_id,
                        "username": note.user.username if note.user else None
                    },
                    "can_edit": note.user_id == current_user.id
                }
        
        elif item.item_type == 'chat' and item.item_id:
            chat = db.query(models.ChatSession).filter(models.ChatSession.id == item.item_id).first()
            if chat:
                messages = db.query(models.ChatMessage).filter(
                    models.ChatMessage.chat_session_id == chat.id
                ).order_by(models.ChatMessage.timestamp).all()
                
                return {
                    "type": "chat",
                    "id": chat.id,
                    "title": chat.title,
                    "messages": [{
                        "id": msg.id,
                        "user_message": msg.user_message,
                        "ai_response": msg.ai_response,
                        "timestamp": msg.timestamp.isoformat() if msg.timestamp else None
                    } for msg in messages],
                    "created_at": chat.created_at.isoformat() if chat.created_at else None,
                    "owner": {
                        "id": chat.user_id,
                        "username": chat.user.username if chat.user else None
                    },
                    "can_edit": chat.user_id == current_user.id
                }
        
        # For other types, return basic info
        return {
            "type": item.item_type,
            "title": item.title,
            "url": item.url,
            "description": item.description,
            "can_edit": False
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error viewing playlist item: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/playlists/{playlist_id}/progress")
async def update_playlist_progress(
    playlist_id: int,
    item_id: int,
    completed: bool,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Update user's progress on a playlist item"""
    try:
        follower = db.query(models.PlaylistFollower).filter(
            and_(
                models.PlaylistFollower.playlist_id == playlist_id,
                models.PlaylistFollower.user_id == current_user.id
            )
        ).first()
        
        if not follower:
            raise HTTPException(status_code=404, detail="Not following this playlist")
        
        completed_items = follower.completed_items or []
        
        if completed and item_id not in completed_items:
            completed_items.append(item_id)
        elif not completed and item_id in completed_items:
            completed_items.remove(item_id)
        
        follower.completed_items = completed_items
        follower.last_accessed = datetime.utcnow()
        
        # Calculate progress percentage
        total_items = db.query(func.count(models.PlaylistItem.id)).filter(
            models.PlaylistItem.playlist_id == playlist_id
        ).scalar()
        
        if total_items > 0:
            follower.progress_percentage = (len(completed_items) / total_items) * 100
            
            # Check if completed
            if follower.progress_percentage >= 100 and not follower.is_completed:
                follower.is_completed = True
                follower.completed_at = datetime.utcnow()
                
                # Update completion count
                playlist = db.query(models.LearningPlaylist).filter(
                    models.LearningPlaylist.id == playlist_id
                ).first()
                if playlist:
                    playlist.completion_count = (playlist.completion_count or 0) + 1
        
        db.commit()
        
        return {
            "message": "Progress updated",
            "progress_percentage": follower.progress_percentage,
            "completed_items": completed_items
        }
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error updating progress: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== COMPREHENSIVE IMPORT/EXPORT API ====================

from import_export_service import ImportExportService

@app.post("/api/import_export/notes_to_flashcards")
async def convert_notes_to_flashcards(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Convert notes to flashcards"""
    try:
        note_ids = payload.get("note_ids", [])
        card_count = payload.get("card_count", 10)
        difficulty = payload.get("difficulty", "medium")
        
        service = ImportExportService(db)
        result = await service.notes_to_flashcards(
            note_ids=note_ids,
            user_id=current_user.id,
            card_count=card_count,
            difficulty=difficulty
        )
        
        if result["success"]:
            # Track in history
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="notes",
                destination_type="flashcards",
                source_ids=json.dumps(note_ids),
                destination_ids=json.dumps([result["set_id"]]),
                item_count=result["card_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in notes_to_flashcards: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/notes_to_questions")
async def convert_notes_to_questions(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Convert notes to practice questions"""
    try:
        note_ids = payload.get("note_ids", [])
        question_count = payload.get("question_count", 10)
        difficulty = payload.get("difficulty", "medium")
        
        service = ImportExportService(db)
        result = await service.notes_to_questions(
            note_ids=note_ids,
            user_id=current_user.id,
            question_count=question_count,
            difficulty=difficulty
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="notes",
                destination_type="questions",
                source_ids=json.dumps(note_ids),
                destination_ids=json.dumps([result["set_id"]]),
                item_count=result["question_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in notes_to_questions: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/flashcards_to_notes")
async def convert_flashcards_to_notes(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Convert flashcard sets to notes"""
    try:
        set_ids = payload.get("set_ids", [])
        format_style = payload.get("format_style", "structured")
        
        service = ImportExportService(db)
        result = await service.flashcards_to_notes(
            set_ids=set_ids,
            user_id=current_user.id,
            format_style=format_style
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="flashcards",
                destination_type="notes",
                source_ids=json.dumps(set_ids),
                destination_ids=json.dumps([result["note_id"]]),
                item_count=result["card_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in flashcards_to_notes: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/flashcards_to_questions")
async def convert_flashcards_to_questions(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Convert flashcards to quiz questions"""
    try:
        set_ids = payload.get("set_ids", [])
        
        service = ImportExportService(db)
        result = await service.flashcards_to_questions(
            set_ids=set_ids,
            user_id=current_user.id
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="flashcards",
                destination_type="questions",
                source_ids=json.dumps(set_ids),
                destination_ids=json.dumps([result["set_id"]]),
                item_count=result["question_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in flashcards_to_questions: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/questions_to_flashcards")
async def convert_questions_to_flashcards(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Convert question sets to flashcards"""
    try:
        set_ids = payload.get("set_ids", [])
        
        service = ImportExportService(db)
        result = await service.questions_to_flashcards(
            set_ids=set_ids,
            user_id=current_user.id
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="questions",
                destination_type="flashcards",
                source_ids=json.dumps(set_ids),
                destination_ids=json.dumps([result["set_id"]]),
                item_count=result["card_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in questions_to_flashcards: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/questions_to_notes")
async def convert_questions_to_notes(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Convert question sets to study guide notes"""
    try:
        set_ids = payload.get("set_ids", [])
        
        service = ImportExportService(db)
        result = await service.questions_to_notes(
            set_ids=set_ids,
            user_id=current_user.id
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="questions",
                destination_type="notes",
                source_ids=json.dumps(set_ids),
                destination_ids=json.dumps([result["note_id"]]),
                item_count=1,
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in questions_to_notes: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/media_to_questions")
async def convert_media_to_questions(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Generate questions from media transcripts"""
    try:
        media_ids = payload.get("media_ids", [])
        question_count = payload.get("question_count", 10)
        
        service = ImportExportService(db)
        result = await service.media_to_questions(
            media_ids=media_ids,
            user_id=current_user.id,
            question_count=question_count
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="media",
                destination_type="questions",
                source_ids=json.dumps(media_ids),
                destination_ids=json.dumps([result["set_id"]]),
                item_count=result["question_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in media_to_questions: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/playlist_to_notes")
async def convert_playlist_to_notes(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Compile playlist content into notes"""
    try:
        playlist_id = int(payload.get("playlist_id"))
        
        service = ImportExportService(db)
        result = await service.playlist_to_notes(
            playlist_id=playlist_id,
            user_id=current_user.id
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="playlist",
                destination_type="notes",
                source_ids=json.dumps([playlist_id]),
                destination_ids=json.dumps([result["note_id"]]),
                item_count=result["items_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in playlist_to_notes: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/playlist_to_flashcards")
async def convert_playlist_to_flashcards(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Generate flashcards from playlist content"""
    try:
        playlist_id = int(payload.get("playlist_id"))
        card_count = int(payload.get("card_count", 15))
        
        service = ImportExportService(db)
        result = await service.playlist_to_flashcards(
            playlist_id=playlist_id,
            user_id=current_user.id,
            card_count=card_count
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="import",
                source_type="playlist",
                destination_type="flashcards",
                source_ids=json.dumps([playlist_id]),
                destination_ids=json.dumps([result["set_id"]]),
                item_count=result["card_count"],
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in playlist_to_flashcards: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/merge_notes")
async def merge_multiple_notes(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Merge multiple notes into one"""
    try:
        note_ids = payload.get("note_ids", [])
        new_title = payload.get("title")
        
        service = ImportExportService(db)
        result = await service.merge_notes(
            note_ids=note_ids,
            user_id=current_user.id,
            new_title=new_title
        )
        
        if result["success"]:
            history = models.BatchOperation(
                user_id=current_user.id,
                operation_name="merge_notes",
                source_type="notes",
                source_ids=json.dumps(note_ids),
                result_id=result["note_id"],
                result_type="note",
                status="completed",
                progress=100,
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error in merge_notes: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/export_flashcards_csv")
async def export_flashcards_csv(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Export flashcards to CSV"""
    try:
        set_ids = payload.get("set_ids", [])
        
        service = ImportExportService(db)
        result = service.export_flashcards_to_csv(
            set_ids=set_ids,
            user_id=current_user.id
        )
        
        if result["success"]:
            # Track export
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="export",
                source_type="flashcards",
                destination_type="csv",
                source_ids=json.dumps(set_ids),
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error exporting flashcards to CSV: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/export_questions_pdf")
async def export_questions_pdf(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Export questions to PDF-ready HTML"""
    try:
        set_ids = payload.get("set_ids", [])
        
        service = ImportExportService(db)
        result = service.export_questions_to_pdf(
            set_ids=set_ids,
            user_id=current_user.id
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="export",
                source_type="questions",
                destination_type="pdf",
                source_ids=json.dumps(set_ids),
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error exporting questions to PDF: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/import_export/export_notes_markdown")
async def export_notes_markdown(
    payload: dict = Body(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Export notes to Markdown"""
    try:
        note_ids = payload.get("note_ids", [])
        
        service = ImportExportService(db)
        result = service.export_notes_to_markdown(
            note_ids=note_ids,
            user_id=current_user.id
        )
        
        if result["success"]:
            history = models.ImportExportHistory(
                user_id=current_user.id,
                operation_type="export",
                source_type="notes",
                destination_type="markdown",
                source_ids=json.dumps(note_ids),
                status="completed",
                completed_at=datetime.utcnow()
            )
            db.add(history)
            db.commit()
        
        return result
    except Exception as e:
        logger.error(f"Error exporting notes to markdown: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/import_export/history")
async def get_import_export_history(
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
    limit: int = 50
):
    """Get user's import/export history"""
    try:
        history = db.query(models.ImportExportHistory).filter(
            models.ImportExportHistory.user_id == current_user.id
        ).order_by(models.ImportExportHistory.created_at.desc()).limit(limit).all()
        
        return {
            "history": [
                {
                    "id": h.id,
                    "operation_type": h.operation_type,
                    "source_type": h.source_type,
                    "destination_type": h.destination_type,
                    "item_count": h.item_count,
                    "status": h.status,
                    "created_at": h.created_at.isoformat() if h.created_at else None,
                    "completed_at": h.completed_at.isoformat() if h.completed_at else None
                }
                for h in history
            ]
        }
    except Exception as e:
        logger.error(f"Error getting import/export history: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== SEARCH HUB ENDPOINTS ====================

async def get_expanded_search_terms(query: str) -> list:
    """
    Use AI to expand search query into semantically related terms.
    For example: "Irish Revolution" -> ["Irish Revolution", "Irish War of Independence", "Easter Rising", "Irish independence", "1916 Rising"]
    """
    try:
        # Always include the original query terms
        original_terms = [query]
        
        # Also include individual words from the query (for partial matching)
        words = [w.strip() for w in query.split() if len(w.strip()) > 2]
        
        # Use AI to generate related terms
        ai_prompt = f"""Given the search query "{query}", generate a list of semantically related terms that someone might use to find the same content.

RULES:
1. Include synonyms, alternative names, and related concepts
2. Include historical alternative names if applicable
3. Include common abbreviations or acronyms
4. Keep terms concise (1-4 words each)
5. Return ONLY a JSON array of strings, nothing else
6. Include 5-10 related terms maximum

EXAMPLES:
Query: "Irish Revolution"
["Irish War of Independence", "Easter Rising", "1916 Rising", "Irish independence", "Anglo-Irish War", "Irish rebellion"]

Query: "machine learning"
["ML", "artificial intelligence", "AI", "deep learning", "neural networks", "data science"]

Query: "World War 2"
["WWII", "WW2", "Second World War", "World War II", "1939-1945"]

NOW GENERATE RELATED TERMS FOR: "{query}"
Return ONLY the JSON array:"""

        ai_response = call_ai(ai_prompt, max_tokens=200, temperature=0.3)
        
        # Parse the response
        ai_response_clean = ai_response.strip()
        
        # Remove markdown code blocks if present
        if ai_response_clean.startswith('```'):
            lines = ai_response_clean.split('\n')
            json_lines = []
            in_code_block = False
            for line in lines:
                if line.strip().startswith('```'):
                    in_code_block = not in_code_block
                    continue
                if in_code_block or (not line.strip().startswith('```')):
                    json_lines.append(line)
            ai_response_clean = '\n'.join(json_lines).strip()
        
        related_terms = json.loads(ai_response_clean)
        
        if isinstance(related_terms, list):
            # Combine original query, individual words, and AI-generated terms
            all_terms = original_terms + words + related_terms
            # Remove duplicates while preserving order
            seen = set()
            unique_terms = []
            for term in all_terms:
                term_lower = term.lower()
                if term_lower not in seen and len(term.strip()) > 0:
                    seen.add(term_lower)
                    unique_terms.append(term)
            
            logger.info(f"Expanded '{query}' to {len(unique_terms)} terms: {unique_terms}")
            return unique_terms[:15]  # Limit to 15 terms
        
    except Exception as e:
        logger.error(f"Error expanding search terms: {str(e)}")
    
    # Fallback: return original query and its words
    fallback_terms = [query] + [w.strip() for w in query.split() if len(w.strip()) > 2]
    return list(set(fallback_terms))


async def get_spelling_suggestion(query: str, db, user_id: int) -> Optional[str]:
    """
    Generate "Did you mean..." spelling suggestions for queries with no results.
    Checks against user's actual content titles.
    """
    try:
        # Get all user's content titles for comparison
        all_titles = []
        
        # Flashcard set titles
        flashcard_sets = db.query(models.FlashcardSet.title).filter(
            models.FlashcardSet.user_id == user_id
        ).all()
        all_titles.extend([fs.title for fs in flashcard_sets if fs.title])
        
        # Note titles
        notes = db.query(models.Note.title).filter(
            models.Note.user_id == user_id,
            models.Note.is_deleted == False
        ).all()
        all_titles.extend([n.title for n in notes if n.title])
        
        # Chat session titles
        chats = db.query(models.ChatSession.title).filter(
            models.ChatSession.user_id == user_id
        ).all()
        all_titles.extend([c.title for c in chats if c.title and c.title != "New Chat"])
        
        if not all_titles:
            return None
        
        # Use AI to find the closest match
        ai_prompt = f"""The user searched for "{query}" but found no results.
Here are the titles of content they have:
{json.dumps(all_titles[:50])}

If the search query appears to be a typo or misspelling of one of these titles, suggest the correct title.
If no close match exists, return null.

Return ONLY a JSON object in this format (no markdown):
{{"suggestion": "correct title" or null, "confidence": 0.0 to 1.0}}

Examples:
Query: "quantm computing", Titles: ["Quantum Computing", "Machine Learning"]
{{"suggestion": "Quantum Computing", "confidence": 0.95}}

Query: "xyz123", Titles: ["Math", "Science"]
{{"suggestion": null, "confidence": 0.0}}
"""
        
        ai_response = call_ai(ai_prompt, max_tokens=100, temperature=0.1)
        ai_response_clean = ai_response.strip()
        
        # Remove markdown if present
        if ai_response_clean.startswith('```'):
            lines = ai_response_clean.split('\n')
            json_lines = [l for l in lines if not l.strip().startswith('```')]
            ai_response_clean = '\n'.join(json_lines).strip()
        
        result = json.loads(ai_response_clean)
        
        if result.get("suggestion") and result.get("confidence", 0) > 0.6:
            return result["suggestion"]
        
        return None
        
    except Exception as e:
        logger.error(f"Error getting spelling suggestion: {str(e)}")
        return None


async def get_related_searches(query: str, results: list) -> list:
    """
    Generate related search suggestions based on the query and results.
    """
    try:
        # Extract topics from results for context
        result_topics = [r.get("title", "") for r in results[:10]]
        
        ai_prompt = f"""Based on the search query "{query}" and these related results: {result_topics[:5]},
suggest 4-6 related searches the user might be interested in.

RULES:
1. Suggestions should be related but different from the original query
2. Include both broader and narrower topics
3. Include related concepts or prerequisites
4. Keep suggestions concise (2-5 words each)
5. Return ONLY a JSON array of strings

Examples:
Query: "machine learning"
["deep learning basics", "neural networks", "python for ML", "supervised vs unsupervised", "ML algorithms"]

Query: "Irish history"
["Easter Rising 1916", "Irish independence", "British rule in Ireland", "Irish famine", "Celtic history"]

NOW GENERATE RELATED SEARCHES FOR: "{query}"
Return ONLY the JSON array:"""

        ai_response = call_ai(ai_prompt, max_tokens=150, temperature=0.5)
        ai_response_clean = ai_response.strip()
        
        # Remove markdown if present
        if ai_response_clean.startswith('```'):
            lines = ai_response_clean.split('\n')
            json_lines = [l for l in lines if not l.strip().startswith('```')]
            ai_response_clean = '\n'.join(json_lines).strip()
        
        related = json.loads(ai_response_clean)
        
        if isinstance(related, list):
            # Filter out the original query
            related = [r for r in related if r.lower() != query.lower()]
            return related[:6]
        
        return []
        
    except Exception as e:
        logger.error(f"Error getting related searches: {str(e)}")
        return []


def get_smart_actions(result: dict) -> list:
    """
    Generate smart action buttons for each search result based on its type.
    """
    actions = []
    result_type = result.get("type", "")
    title = result.get("title", "")
    
    if result_type == "flashcard_set":
        actions = [
            {"action": "study", "label": "Study", "icon": "play"},
            {"action": "quiz", "label": "Start Quiz", "icon": "help-circle"},
            {"action": "review", "label": "Review", "icon": "refresh-cw"}
        ]
    elif result_type == "flashcard":
        actions = [
            {"action": "view_set", "label": "View Set", "icon": "layers"},
            {"action": "quiz", "label": "Quiz This", "icon": "help-circle"}
        ]
    elif result_type == "note":
        actions = [
            {"action": "edit", "label": "Edit", "icon": "edit"},
            {"action": "create_flashcards", "label": "Make Flashcards", "icon": "layers"},
            {"action": "summarize", "label": "Summarize", "icon": "file-text"}
        ]
    elif result_type == "chat":
        actions = [
            {"action": "continue", "label": "Continue Chat", "icon": "message-circle"},
            {"action": "create_flashcards", "label": "Make Flashcards", "icon": "layers"}
        ]
    elif result_type == "question_set":
        actions = [
            {"action": "start_quiz", "label": "Start Quiz", "icon": "play"},
            {"action": "practice", "label": "Practice", "icon": "target"}
        ]
    
    return actions


@app.post("/api/search_content")
async def search_content(
    user_id: str = Form(...),
    query: str = Form(...),
    content_types: str = Form("all"),  # "all" or comma-separated: "flashcard_set,flashcard,note,chat,question_set"
    sort_by: str = Form("relevance"),  # "relevance", "date_desc", "date_asc", "title_asc", "title_desc"
    date_from: Optional[str] = Form(None),  # ISO date string
    date_to: Optional[str] = Form(None),  # ISO date string
    db: Session = Depends(get_db)
):
    """
    Universal search across user's own content with advanced filters and semantic expansion
    
    Filters:
    - content_types: "all" or comma-separated list (flashcard_set,flashcard,note,chat,question_set)
    - sort_by: "relevance", "date_desc", "date_asc", "title_asc", "title_desc"
    - date_from: Filter results created after this date (ISO format)
    - date_to: Filter results created before this date (ISO format)
    """
    try:
        logger.info(f"Search request - user_id: {user_id}, query: {query}, filters: types={content_types}, sort={sort_by}")
        
        # Get actual user ID from email/username
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"results": [], "total": 0, "message": "User not found"}
        
        actual_user_id = user.id
        results = []
        
        # Generate expanded search terms using AI for semantic matching
        expanded_terms = await get_expanded_search_terms(query)
        logger.info(f"Expanded search terms: {expanded_terms}")
        
        # Parse content types filter
        if content_types == "all":
            # For "all", exclude individual flashcards - only show sets
            enabled_types = ["flashcard_set", "note", "chat", "question_set"]
        else:
            enabled_types = [t.strip() for t in content_types.split(",")]
            # Remove "flashcard" if present - we only want sets
            if "flashcard" in enabled_types:
                enabled_types.remove("flashcard")
        
        # Parse date filters
        date_from_obj = None
        date_to_obj = None
        if date_from:
            try:
                date_from_obj = datetime.fromisoformat(date_from.replace('Z', '+00:00'))
            except:
                pass
        if date_to:
            try:
                date_to_obj = datetime.fromisoformat(date_to.replace('Z', '+00:00'))
            except:
                pass
        
        # Build OR conditions for all search terms
        def build_search_conditions(column, terms):
            """Build OR conditions for multiple search terms"""
            conditions = []
            for term in terms:
                search_term = f"%{term.lower()}%"
                conditions.append(func.lower(column).like(search_term))
            return or_(*conditions) if conditions else func.lower(column).like(f"%{query.lower()}%")
        
        # Search Flashcard Sets (user's own + public from others)
        if "flashcard_set" in enabled_types:
            try:
                title_conditions = build_search_conditions(models.FlashcardSet.title, expanded_terms)
                desc_conditions = build_search_conditions(models.FlashcardSet.description, expanded_terms)
                
                # Search user's own flashcard sets OR public flashcard sets from any user
                query_builder = db.query(models.FlashcardSet).filter(
                    and_(
                        or_(title_conditions, desc_conditions),
                        or_(
                            models.FlashcardSet.user_id == actual_user_id,  # User's own sets
                            models.FlashcardSet.is_public == True  # Public sets from anyone
                        )
                    )
                )
                
                # Apply date filters
                if date_from_obj:
                    query_builder = query_builder.filter(models.FlashcardSet.created_at >= date_from_obj)
                if date_to_obj:
                    query_builder = query_builder.filter(models.FlashcardSet.created_at <= date_to_obj)
                
                flashcard_sets = query_builder.all()
                
                for fset in flashcard_sets:
                    card_count = db.query(models.Flashcard).filter(
                        models.Flashcard.set_id == fset.id
                    ).count()
                    
                    # Get author username
                    author = db.query(models.User).filter(models.User.id == fset.user_id).first()
                    author_name = author.username if author else "Unknown"
                    
                    results.append({
                        "id": fset.id,
                        "type": "flashcard_set",
                        "title": fset.title or "Untitled Set",
                        "description": fset.description or "",
                        "created_at": fset.created_at.isoformat() if fset.created_at else None,
                        "card_count": card_count,
                        "source_type": fset.source_type,
                        "author": author_name,
                        "author_id": fset.user_id,
                        "is_public": fset.is_public,
                        "is_own": fset.user_id == actual_user_id
                    })
            except Exception as e:
                logger.error(f"Error searching flashcard sets: {str(e)}")
        
        # Search Individual Flashcards (user's own + public from others)
        if "flashcard" in enabled_types:
            try:
                question_conditions = build_search_conditions(models.Flashcard.question, expanded_terms)
                answer_conditions = build_search_conditions(models.Flashcard.answer, expanded_terms)
                
                # Search flashcards from user's own sets OR public sets from any user
                query_builder = db.query(models.Flashcard).join(
                    models.FlashcardSet
                ).filter(
                    and_(
                        or_(question_conditions, answer_conditions),
                        or_(
                            models.FlashcardSet.user_id == actual_user_id,  # User's own flashcards
                            models.FlashcardSet.is_public == True  # Public flashcards from anyone
                        )
                    )
                )
                
                # Apply date filters
                if date_from_obj:
                    query_builder = query_builder.filter(models.Flashcard.created_at >= date_from_obj)
                if date_to_obj:
                    query_builder = query_builder.filter(models.Flashcard.created_at <= date_to_obj)
                
                flashcards = query_builder.limit(50).all()
                logger.info(f"Found {len(flashcards)} individual flashcards (own + public)")
                
                for card in flashcards:
                    fset = db.query(models.FlashcardSet).filter(
                        models.FlashcardSet.id == card.set_id
                    ).first()
                    
                    # Get author username
                    author = db.query(models.User).filter(models.User.id == fset.user_id).first() if fset else None
                    author_name = author.username if author else "Unknown"
                    
                    results.append({
                        "id": card.id,
                        "type": "flashcard",
                        "title": card.question[:100] if card.question else "Flashcard",
                        "description": card.answer[:200] if card.answer else "",
                        "created_at": card.created_at.isoformat() if card.created_at else None,
                        "set_name": fset.title if fset else None,
                        "set_id": card.set_id,
                        "difficulty": card.difficulty,
                        "author": author_name,
                        "author_id": fset.user_id if fset else None,
                        "is_public": fset.is_public if fset else False,
                        "is_own": fset.user_id == actual_user_id if fset else False
                    })
            except Exception as e:
                logger.error(f"Error searching flashcards: {str(e)}")
        
        # Search Notes (user's own + public from others)
        if "note" in enabled_types:
            try:
                title_conditions = build_search_conditions(models.Note.title, expanded_terms)
                content_conditions = build_search_conditions(models.Note.content, expanded_terms)
                
                query_builder = db.query(models.Note).filter(
                    and_(
                        models.Note.is_deleted == False,
                        or_(title_conditions, content_conditions),
                        or_(
                            models.Note.user_id == actual_user_id,  # User's own notes
                            models.Note.is_public == True  # Public notes from anyone
                        )
                    )
                )
                
                # Apply date filters
                if date_from_obj:
                    query_builder = query_builder.filter(models.Note.created_at >= date_from_obj)
                if date_to_obj:
                    query_builder = query_builder.filter(models.Note.created_at <= date_to_obj)
                
                notes = query_builder.limit(50).all()
                logger.info(f"Found {len(notes)} notes (own + public)")
                
                for note in notes:
                    # Get author username
                    author = db.query(models.User).filter(models.User.id == note.user_id).first()
                    author_name = author.username if author else "Unknown"
                    
                    results.append({
                        "id": note.id,
                        "type": "note",
                        "title": note.title if note.title else "Untitled Note",
                        "description": note.content[:200] if note.content else "",
                        "created_at": note.created_at.isoformat() if note.created_at else None,
                        "is_favorite": note.is_favorite,
                        "folder_id": note.folder_id,
                        "author": author_name,
                        "author_id": note.user_id,
                        "is_public": note.is_public,
                        "is_own": note.user_id == actual_user_id
                    })
            except Exception as e:
                logger.error(f"Error searching notes: {str(e)}")
        
        # Search Chat Sessions (only user's own)
        if "chat" in enabled_types:
            try:
                title_conditions = build_search_conditions(models.ChatSession.title, expanded_terms)
                
                query_builder = db.query(models.ChatSession).filter(
                    and_(
                        models.ChatSession.user_id == actual_user_id,
                        title_conditions
                    )
                )
                
                # Apply date filters
                if date_from_obj:
                    query_builder = query_builder.filter(models.ChatSession.created_at >= date_from_obj)
                if date_to_obj:
                    query_builder = query_builder.filter(models.ChatSession.created_at <= date_to_obj)
                
                chats = query_builder.limit(50).all()
                logger.info(f"Found {len(chats)} chat sessions")
                
                for chat in chats:
                    message_count = db.query(models.ChatMessage).filter(
                        models.ChatMessage.chat_session_id == chat.id
                    ).count()
                    
                    results.append({
                        "id": chat.id,
                        "type": "chat",
                        "title": chat.title or "Untitled Chat",
                        "description": f"{message_count} messages",
                        "created_at": chat.created_at.isoformat() if chat.created_at else None,
                        "updated_at": chat.updated_at.isoformat() if chat.updated_at else None,
                        "message_count": message_count,
                        "folder_id": chat.folder_id
                    })
            except Exception as e:
                logger.error(f"Error searching chats: {str(e)}")
        
        # Search Question Sets (only user's own)
        if "question_set" in enabled_types:
            try:
                title_conditions = build_search_conditions(models.QuestionSet.title, expanded_terms)
                desc_conditions = build_search_conditions(models.QuestionSet.description, expanded_terms)
                
                query_builder = db.query(models.QuestionSet).filter(
                    and_(
                        models.QuestionSet.user_id == actual_user_id,
                        or_(title_conditions, desc_conditions)
                    )
                )
                
                # Apply date filters
                if date_from_obj:
                    query_builder = query_builder.filter(models.QuestionSet.created_at >= date_from_obj)
                if date_to_obj:
                    query_builder = query_builder.filter(models.QuestionSet.created_at <= date_to_obj)
                
                question_sets = query_builder.limit(50).all()
                logger.info(f"Found {len(question_sets)} question sets")
                
                for qset in question_sets:
                    question_count = db.query(models.Question).filter(
                        models.Question.set_id == qset.id
                    ).count()
                    
                    results.append({
                        "id": qset.id,
                        "type": "question_set",
                        "title": qset.title,
                        "description": qset.description or "",
                        "created_at": qset.created_at.isoformat() if qset.created_at else None,
                        "question_count": question_count,
                        "difficulty": qset.difficulty_level,
                        "subject": qset.subject
                    })
            except Exception as e:
                logger.error(f"Error searching question sets: {str(e)}")
        
        # Apply sorting
        if sort_by == "date_desc":
            results.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        elif sort_by == "date_asc":
            results.sort(key=lambda x: x.get('created_at', ''))
        elif sort_by == "title_asc":
            results.sort(key=lambda x: x.get('title', '').lower())
        elif sort_by == "title_desc":
            results.sort(key=lambda x: x.get('title', '').lower(), reverse=True)
        else:  # relevance (default)
            # Sort by created_at desc as a simple relevance metric
            results.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        
        logger.info(f"Total results: {len(results)}")
        
        # Get filter statistics
        type_counts = {}
        for result in results:
            result_type = result['type']
            type_counts[result_type] = type_counts.get(result_type, 0) + 1
        
        # AI ENHANCEMENTS
        
        # 1. "Did you mean..." - Suggest corrections for unclear queries
        did_you_mean = None
        if len(results) == 0:
            did_you_mean = await get_spelling_suggestion(query, db, actual_user_id)
        
        # 2. Related searches - Generate related search suggestions
        related_searches = await get_related_searches(query, results)
        
        # 3. Add smart actions to results
        for result in results:
            result["smart_actions"] = get_smart_actions(result)
        
        return {
            "total_results": len(results),
            "results": results,
            "query": query,
            "filters_applied": {
                "content_types": content_types,
                "sort_by": sort_by,
                "date_from": date_from,
                "date_to": date_to
            },
            "type_counts": type_counts,
            "did_you_mean": did_you_mean,
            "related_searches": related_searches,
            "expanded_terms": expanded_terms[:5] if expanded_terms else []
        }
        
    except Exception as e:
        logger.error(f"Search error: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Search failed: {str(e)}")


@app.post("/api/autocomplete")
async def autocomplete(
    user_id: str = Form(...),
    query: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Google-style autocomplete suggestions as user types.
    Returns commands, recent searches, matching content, and smart suggestions.
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"suggestions": []}
        
        suggestions = []
        query_lower = query.lower().strip()
        
        # 1. Command suggestions - COMPREHENSIVE LIST
        commands = [
            # Creation commands
            {"text": "create flashcards on", "subtext": "Generate AI flashcards on any topic", "type": "command", "category": "create", "needs_topic": True},
            {"text": "create a quiz on", "subtext": "Test your knowledge on any topic", "type": "command", "category": "create", "needs_topic": True},
            {"text": "create a note on", "subtext": "AI writes comprehensive notes", "type": "command", "category": "create", "needs_topic": True},
            {"text": "create questions on", "subtext": "Generate practice questions", "type": "command", "category": "create", "needs_topic": True},
            {"text": "create study plan for", "subtext": "Plan your learning journey", "type": "command", "category": "create", "needs_topic": True},
            
            # Learning & Explanation
            {"text": "explain", "subtext": "Get AI explanation on any topic", "type": "command", "category": "learn", "needs_topic": True},
            {"text": "teach me about", "subtext": "Learn a new topic", "type": "command", "category": "learn", "needs_topic": True},
            {"text": "what is", "subtext": "Get definition", "type": "command", "category": "learn", "needs_topic": True},
            {"text": "how does", "subtext": "Understand how something works", "type": "command", "category": "learn", "needs_topic": True},
            {"text": "summarize", "subtext": "Get a quick summary", "type": "command", "category": "learn", "needs_topic": True},
            
            # Testing & Practice
            {"text": "quiz me on", "subtext": "Quick quiz on any topic", "type": "command", "category": "test", "needs_topic": True},
            {"text": "test me on", "subtext": "Test your knowledge", "type": "command", "category": "test", "needs_topic": True},
            
            # Progress & Analytics (no topic needed)
            {"text": "show my progress", "subtext": "View your learning statistics", "type": "command", "category": "progress", "needs_topic": False},
            {"text": "show my weak areas", "subtext": "Find knowledge gaps", "type": "command", "category": "progress", "needs_topic": False},
            {"text": "what is my learning style", "subtext": "AI analyzes your learning patterns", "type": "command", "category": "progress", "needs_topic": False},
            {"text": "show my achievements", "subtext": "View your badges and rewards", "type": "command", "category": "progress", "needs_topic": False},
            {"text": "show knowledge gaps", "subtext": "Find your blind spots", "type": "command", "category": "progress", "needs_topic": False},
            
            # Scheduling & Planning (no topic needed)
            {"text": "what should I study next", "subtext": "Get AI recommendations", "type": "command", "category": "schedule", "needs_topic": False},
            {"text": "predict what I'll forget", "subtext": "Forgetting curve analysis", "type": "command", "category": "schedule", "needs_topic": False},
            {"text": "optimize my retention", "subtext": "Spaced repetition schedule", "type": "command", "category": "schedule", "needs_topic": False},
            
            # Quick Actions (no topic needed)
            {"text": "review flashcards", "subtext": "Start flashcard review", "type": "command", "category": "quick", "needs_topic": False},
            {"text": "review weak flashcards", "subtext": "Focus on difficult cards", "type": "command", "category": "quick", "needs_topic": False},
            
            # Search
            {"text": "search for", "subtext": "Search your content", "type": "command", "category": "search", "needs_topic": True},
        ]
        
        # Check if user is typing a command
        is_typing_command = False
        matched_command = None
        remaining_topic = ""
        
        for cmd in commands:
            cmd_text = cmd["text"].lower()
            if query_lower.startswith(cmd_text):
                is_typing_command = True
                matched_command = cmd
                remaining_topic = query_lower[len(cmd_text):].strip()
                break
            elif cmd_text.startswith(query_lower):
                is_typing_command = True
                matched_command = cmd
                break
        
        # If user is typing a command that needs a topic
        if is_typing_command and matched_command:
            if matched_command.get("needs_topic"):
                if remaining_topic:
                    # User has started typing a topic - show the command with their topic
                    suggestions.append({
                        "text": f"{matched_command['text']} {remaining_topic}",
                        "subtext": matched_command["subtext"],
                        "type": "command",
                        "category": matched_command.get("category", "")
                    })
                else:
                    # User just typed the command - show it as-is, waiting for topic
                    suggestions.append({
                        "text": matched_command["text"],
                        "subtext": f"{matched_command['subtext']} (type a topic)",
                        "type": "command",
                        "category": matched_command.get("category", "")
                    })
            else:
                # Command doesn't need a topic - show it directly
                suggestions.append({
                    "text": matched_command["text"],
                    "subtext": matched_command["subtext"],
                    "type": "command",
                    "category": matched_command.get("category", "")
                })
        
        # Add other matching commands (prefix match)
        if not is_typing_command or len(suggestions) < 3:
            for cmd in commands:
                cmd_lower = cmd["text"].lower()
                if cmd_lower.startswith(query_lower) and cmd != matched_command:
                    suggestions.append({
                        "text": cmd["text"],
                        "subtext": cmd["subtext"],
                        "type": "command",
                        "category": cmd.get("category", "")
                    })
                    if len(suggestions) >= 4:
                        break
        
        # 2. Search user's content ONLY if not typing a command
        if not is_typing_command and len(query_lower) >= 2:
            search_term = f"%{query_lower}%"
            
            # Flashcard sets
            flashcard_sets = db.query(models.FlashcardSet).filter(
                models.FlashcardSet.user_id == user.id,
                func.lower(models.FlashcardSet.title).like(search_term)
            ).order_by(models.FlashcardSet.updated_at.desc()).limit(3).all()
            
            for fset in flashcard_sets:
                card_count = db.query(models.Flashcard).filter(models.Flashcard.set_id == fset.id).count()
                suggestions.append({
                    "text": fset.title,
                    "subtext": f"Flashcard Set • {card_count} cards",
                    "type": "content",
                    "contentType": "flashcard_set",
                    "id": fset.id
                })
            
            # Notes
            notes = db.query(models.Note).filter(
                models.Note.user_id == user.id,
                models.Note.is_deleted == False,
                func.lower(models.Note.title).like(search_term)
            ).order_by(models.Note.updated_at.desc()).limit(3).all()
            
            for note in notes:
                suggestions.append({
                    "text": note.title,
                    "subtext": "Note",
                    "type": "content",
                    "contentType": "note",
                    "id": note.id
                })
            
            # Chat sessions
            chats = db.query(models.ChatSession).filter(
                models.ChatSession.user_id == user.id,
                func.lower(models.ChatSession.title).like(search_term)
            ).order_by(models.ChatSession.updated_at.desc()).limit(2).all()
            
            for chat in chats:
                if chat.title and chat.title != "New Chat":
                    suggestions.append({
                        "text": chat.title,
                        "subtext": "Chat Session",
                        "type": "content",
                        "contentType": "chat",
                        "id": chat.id
                    })
        
        # 3. If very few suggestions, add a generic search option
        if len(suggestions) < 3 and len(query) >= 2:
            suggestions.append({
                "text": query,
                "subtext": f"Search for '{query}'",
                "type": "search"
            })
        
        # Deduplicate and limit
        seen = set()
        unique_suggestions = []
        for s in suggestions:
            key = s["text"].lower()
            if key not in seen:
                seen.add(key)
                unique_suggestions.append(s)
        
        return {"suggestions": unique_suggestions[:10]}
        
    except Exception as e:
        logger.error(f"Autocomplete error: {str(e)}")
        return {"suggestions": []}


@app.post("/api/natural_language_search")
async def natural_language_search(
    user_id: str = Form(...),
    query: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Handle natural language search queries with filters.
    Examples: "show me hard flashcards I haven't reviewed", "notes from last week", "easy flashcards on math"
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"results": [], "total_results": 0, "parsed_filters": {}}
        
        # Use AI to parse the natural language query into structured filters
        ai_prompt = f"""Parse this natural language search query into structured filters.

QUERY: "{query}"

Extract these filters if mentioned:
- content_type: "flashcard", "flashcard_set", "note", "chat", "question_set", or "all"
- difficulty: "easy", "medium", "hard", or null
- reviewed: true (reviewed), false (not reviewed), or null
- topic: the subject/topic being searched for, or null
- time_filter: "today", "yesterday", "last_week", "last_month", or null
- marked_for_review: true or false
- sort_by: "newest", "oldest", "difficulty", or null

Return ONLY a JSON object (no markdown):
{{
  "content_type": "string or null",
  "difficulty": "string or null",
  "reviewed": "boolean or null",
  "topic": "string or null",
  "time_filter": "string or null",
  "marked_for_review": "boolean or null",
  "sort_by": "string or null"
}}

Examples:
Query: "show me hard flashcards I haven't reviewed"
{{"content_type": "flashcard", "difficulty": "hard", "reviewed": false, "topic": null, "time_filter": null, "marked_for_review": null, "sort_by": null}}

Query: "notes from last week about physics"
{{"content_type": "note", "difficulty": null, "reviewed": null, "topic": "physics", "time_filter": "last_week", "marked_for_review": null, "sort_by": null}}

Query: "easy flashcards on math that need review"
{{"content_type": "flashcard", "difficulty": "easy", "reviewed": null, "topic": "math", "time_filter": null, "marked_for_review": true, "sort_by": null}}

NOW PARSE: "{query}"
"""
        
        ai_response = call_ai(ai_prompt, max_tokens=200, temperature=0.1)
        ai_response_clean = ai_response.strip()
        
        # Remove markdown if present
        if ai_response_clean.startswith('```'):
            lines = ai_response_clean.split('\n')
            json_lines = [l for l in lines if not l.strip().startswith('```')]
            ai_response_clean = '\n'.join(json_lines).strip()
        
        filters = json.loads(ai_response_clean)
        logger.info(f"Parsed natural language filters: {filters}")
        
        results = []
        
        # Build query based on parsed filters
        content_type = filters.get("content_type", "all")
        difficulty = filters.get("difficulty")
        reviewed = filters.get("reviewed")
        topic = filters.get("topic")
        time_filter = filters.get("time_filter")
        marked_for_review = filters.get("marked_for_review")
        
        # Calculate date filter
        date_from = None
        if time_filter:
            now = datetime.now(timezone.utc)
            if time_filter == "today":
                date_from = now.replace(hour=0, minute=0, second=0, microsecond=0)
            elif time_filter == "yesterday":
                date_from = (now - timedelta(days=1)).replace(hour=0, minute=0, second=0, microsecond=0)
            elif time_filter == "last_week":
                date_from = now - timedelta(days=7)
            elif time_filter == "last_month":
                date_from = now - timedelta(days=30)
        
        # Search flashcards
        if content_type in ["flashcard", "all"]:
            query_builder = db.query(models.Flashcard).join(models.FlashcardSet).filter(
                models.FlashcardSet.user_id == user.id
            )
            
            if difficulty:
                query_builder = query_builder.filter(models.Flashcard.difficulty == difficulty)
            
            if reviewed is not None:
                if reviewed:
                    query_builder = query_builder.filter(models.Flashcard.times_reviewed > 0)
                else:
                    query_builder = query_builder.filter(models.Flashcard.times_reviewed == 0)
            
            if marked_for_review:
                query_builder = query_builder.filter(models.Flashcard.marked_for_review == True)
            
            if topic:
                topic_term = f"%{topic.lower()}%"
                query_builder = query_builder.filter(
                    or_(
                        func.lower(models.Flashcard.question).like(topic_term),
                        func.lower(models.Flashcard.answer).like(topic_term),
                        func.lower(models.FlashcardSet.title).like(topic_term)
                    )
                )
            
            if date_from:
                query_builder = query_builder.filter(models.Flashcard.created_at >= date_from)
            
            flashcards = query_builder.limit(50).all()
            
            for card in flashcards:
                fset = db.query(models.FlashcardSet).filter(models.FlashcardSet.id == card.set_id).first()
                results.append({
                    "id": card.id,
                    "type": "flashcard",
                    "title": card.question[:100] if card.question else "Flashcard",
                    "description": card.answer[:200] if card.answer else "",
                    "created_at": card.created_at.isoformat() if card.created_at else None,
                    "set_name": fset.title if fset else None,
                    "set_id": card.set_id,
                    "difficulty": card.difficulty,
                    "times_reviewed": card.times_reviewed,
                    "marked_for_review": card.marked_for_review,
                    "smart_actions": get_smart_actions({"type": "flashcard"})
                })
        
        # Search flashcard sets
        if content_type in ["flashcard_set", "all"]:
            query_builder = db.query(models.FlashcardSet).filter(
                models.FlashcardSet.user_id == user.id
            )
            
            if topic:
                topic_term = f"%{topic.lower()}%"
                query_builder = query_builder.filter(
                    or_(
                        func.lower(models.FlashcardSet.title).like(topic_term),
                        func.lower(models.FlashcardSet.description).like(topic_term)
                    )
                )
            
            if date_from:
                query_builder = query_builder.filter(models.FlashcardSet.created_at >= date_from)
            
            flashcard_sets = query_builder.limit(50).all()
            
            for fset in flashcard_sets:
                card_count = db.query(models.Flashcard).filter(models.Flashcard.set_id == fset.id).count()
                results.append({
                    "id": fset.id,
                    "type": "flashcard_set",
                    "title": fset.title or "Untitled Set",
                    "description": fset.description or "",
                    "created_at": fset.created_at.isoformat() if fset.created_at else None,
                    "card_count": card_count,
                    "smart_actions": get_smart_actions({"type": "flashcard_set"})
                })
        
        # Search notes
        if content_type in ["note", "all"]:
            query_builder = db.query(models.Note).filter(
                models.Note.user_id == user.id,
                models.Note.is_deleted == False
            )
            
            if topic:
                topic_term = f"%{topic.lower()}%"
                query_builder = query_builder.filter(
                    or_(
                        func.lower(models.Note.title).like(topic_term),
                        func.lower(models.Note.content).like(topic_term)
                    )
                )
            
            if date_from:
                query_builder = query_builder.filter(models.Note.created_at >= date_from)
            
            notes = query_builder.limit(50).all()
            
            for note in notes:
                results.append({
                    "id": note.id,
                    "type": "note",
                    "title": note.title or "Untitled Note",
                    "description": note.content[:200] if note.content else "",
                    "created_at": note.created_at.isoformat() if note.created_at else None,
                    "smart_actions": get_smart_actions({"type": "note"})
                })
        
        # Sort results
        sort_by = filters.get("sort_by")
        if sort_by == "newest":
            results.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        elif sort_by == "oldest":
            results.sort(key=lambda x: x.get('created_at', ''))
        elif sort_by == "difficulty":
            difficulty_order = {"easy": 0, "medium": 1, "hard": 2}
            results.sort(key=lambda x: difficulty_order.get(x.get('difficulty', 'medium'), 1))
        else:
            results.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        
        # Generate filter description
        filter_desc = generate_filter_description(filters)
        
        return {
            "total_results": len(results),
            "results": results,
            "query": query,
            "parsed_filters": filters,
            "filter_description": filter_desc
        }
        
    except Exception as e:
        logger.error(f"Natural language search error: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        return {"results": [], "total_results": 0, "error": str(e)}


def generate_filter_description(filters: dict) -> str:
    """Generate a human-readable description of the applied filters."""
    parts = []
    
    if filters.get("content_type") and filters["content_type"] != "all":
        parts.append(f"{filters['content_type']}s")
    
    if filters.get("difficulty"):
        parts.append(f"{filters['difficulty']} difficulty")
    
    if filters.get("reviewed") is not None:
        parts.append("reviewed" if filters["reviewed"] else "not reviewed")
    
    if filters.get("marked_for_review"):
        parts.append("marked for review")
    
    if filters.get("topic"):
        parts.append(f"about '{filters['topic']}'")
    
    if filters.get("time_filter"):
        time_map = {
            "today": "from today",
            "yesterday": "from yesterday",
            "last_week": "from last week",
            "last_month": "from last month"
        }
        parts.append(time_map.get(filters["time_filter"], ""))
    
    if parts:
        return "Showing " + ", ".join(parts)
    return "Showing all results"


@app.post("/api/detect_search_intent")
async def detect_search_intent(
    user_id: str = Form(...),
    query: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    AI-powered intent detection - parses natural language and returns structured action
    """
    try:
        # Define all available actions with their required parameters AND API endpoints
        available_actions = {
            "create_note": {
                "description": "Create a new note",
                "parameters": {"title": "string (optional)", "content": "string (optional)"},
                "endpoint": None,  # Handled by frontend navigation
                "method": None
            },
            "create_flashcards": {
                "description": "Generate flashcards on a topic",
                "parameters": {"topic": "string (required)", "count": "integer (optional, default 10)"},
                "endpoint": None,  # Handled by frontend navigation
                "method": None
            },
            "create_quiz": {
                "description": "Create a quiz or test",
                "parameters": {"topics": "array of strings", "difficulty": "string (easy/medium/hard)", "count": "integer"},
                "endpoint": None,  # Handled by frontend navigation
                "method": None
            },
            "review_flashcards": {
                "description": "Review existing flashcards",
                "parameters": {"filter": "string (needs_review/marked_for_review/all)"},
                "endpoint": "/api/get_flashcards_for_review",
                "method": "POST"
            },
            "show_weak_areas": {
                "description": "Show topics user is struggling with",
                "parameters": {},
                "endpoint": "/api/get_weak_areas",
                "method": "POST"
            },
            "show_progress": {
                "description": "Show learning progress and statistics",
                "parameters": {},
                "endpoint": None,  # Navigate to dashboard
                "method": None
            },
            "show_achievements": {
                "description": "Show earned achievements and badges",
                "parameters": {},
                "endpoint": None,  # Navigate to dashboard
                "method": None
            },
            "start_chat": {
                "description": "Start AI chat conversation",
                "parameters": {"message": "string (the question or topic)", "topic": "string (optional)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "adapt_difficulty": {
                "description": "Adapt content difficulty to user's level",
                "parameters": {"topic": "string (optional)"},
                "endpoint": "/api/adaptive/difficulty",
                "method": "GET"
            },
            "show_learning_style": {
                "description": "Detect and show user's learning style",
                "parameters": {},
                "endpoint": "/api/adaptive/learning-style",
                "method": "GET"
            },
            "show_knowledge_gaps": {
                "description": "Find knowledge blind spots and gaps",
                "parameters": {},
                "endpoint": "/api/adaptive/knowledge-gaps",
                "method": "GET"
            },
            "create_curriculum": {
                "description": "Create personalized learning curriculum",
                "parameters": {"topic": "string (required)"},
                "endpoint": "/api/adaptive/curriculum",
                "method": "GET"
            },
            "optimize_retention": {
                "description": "Get spaced repetition schedule",
                "parameters": {},
                "endpoint": "/api/adaptive/retention",
                "method": "GET"
            },
            "predict_forgetting": {
                "description": "Predict what user will forget next",
                "parameters": {},
                "endpoint": "/api/adaptive/predict-forgetting",
                "method": "GET"
            },
            "detect_burnout": {
                "description": "Detect burnout risk",
                "parameters": {},
                "endpoint": "/api/adaptive/burnout-risk",
                "method": "GET"
            },
            "suggest_breaks": {
                "description": "Suggest optimal break schedule",
                "parameters": {},
                "endpoint": "/api/adaptive/break-schedule",
                "method": "GET"
            },
            "predict_focus": {
                "description": "Predict focus level at current time",
                "parameters": {},
                "endpoint": "/api/adaptive/focus-prediction",
                "method": "GET"
            },
            "find_study_twin": {
                "description": "Find study partner with similar learning patterns",
                "parameters": {},
                "endpoint": "/api/adaptive/study-twin",
                "method": "GET"
            },
            "find_complementary": {
                "description": "Find learners with complementary strengths",
                "parameters": {},
                "endpoint": "/api/adaptive/complementary-learners",
                "method": "GET"
            },
            "tutor_step_by_step": {
                "description": "Explain topic step-by-step",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat with mode
                "method": None
            },
            "create_analogies": {
                "description": "Create analogies to explain concept",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat with mode
                "method": None
            },
            "simplify_content": {
                "description": "Simplify content for beginners",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat with mode
                "method": None
            },
            "suggest_study_next": {
                "description": "Suggest what to study next",
                "parameters": {},
                "endpoint": "/api/suggest_study_next",
                "method": "POST"
            },
            "summarize_notes": {
                "description": "Summarize user's notes",
                "parameters": {"topic": "string (optional)"},
                "endpoint": "/api/summarize_notes",
                "method": "POST"
            },
            "create_study_plan": {
                "description": "Create a study plan",
                "parameters": {"topic": "string (required)", "duration": "integer (days)"},
                "endpoint": "/api/create_study_plan",
                "method": "POST"
            },
            "search_recent": {
                "description": "Search recent content",
                "parameters": {"timeframe": "string (yesterday/last_week/last_month/recent)"},
                "endpoint": "/api/search_recent_content",
                "method": "POST"
            },
            "find_study_buddies": {
                "description": "Find study buddies",
                "parameters": {},
                "endpoint": None,  # Navigate to social
                "method": None
            },
            "challenge_friend": {
                "description": "Challenge friend to quiz battle",
                "parameters": {"friend_name": "string (optional)"},
                "endpoint": None,  # Navigate to quiz battles
                "method": None
            },
            "show_popular_content": {
                "description": "Show trending/popular content",
                "parameters": {"topic": "string (optional)"},
                "endpoint": "/api/get_popular_content",
                "method": "POST"
            },
            "search": {
                "description": "Regular search for content",
                "parameters": {},
                "endpoint": "/api/search_content",
                "method": "POST"
            },
            # NEW NLP ACTIONS
            "compare_topics": {
                "description": "Compare two or more topics",
                "parameters": {"topics": "array of strings (2+ topics to compare)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "explain_like_im_five": {
                "description": "Explain a concept in very simple terms",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "give_examples": {
                "description": "Provide examples of a concept",
                "parameters": {"topic": "string (required)", "count": "integer (optional)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "test_me": {
                "description": "Quick test/quiz on a topic",
                "parameters": {"topic": "string (required)", "difficulty": "string (optional)"},
                "endpoint": None,  # Navigate to solo quiz
                "method": None
            },
            "define": {
                "description": "Get definition of a term",
                "parameters": {"term": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "list_prerequisites": {
                "description": "List prerequisites for learning a topic",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "suggest_resources": {
                "description": "Suggest learning resources for a topic",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "practice_problems": {
                "description": "Generate practice problems",
                "parameters": {"topic": "string (required)", "difficulty": "string (optional)", "count": "integer (optional)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "summarize_topic": {
                "description": "Get a summary of a topic",
                "parameters": {"topic": "string (required)", "length": "string (short/medium/long)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "show_statistics": {
                "description": "Show learning statistics and analytics",
                "parameters": {"timeframe": "string (optional: today/week/month/all)"},
                "endpoint": None,  # Navigate to dashboard
                "method": None
            },
            "set_goal": {
                "description": "Set a learning goal",
                "parameters": {"goal": "string (required)", "deadline": "string (optional)"},
                "endpoint": None,  # Navigate to dashboard
                "method": None
            },
            "remind_me": {
                "description": "Set a study reminder",
                "parameters": {"topic": "string (required)", "time": "string (optional)"},
                "endpoint": None,  # Create reminder
                "method": None
            },
            "export_content": {
                "description": "Export flashcards or notes",
                "parameters": {"content_type": "string (flashcards/notes)", "topic": "string (optional)"},
                "endpoint": None,  # Navigate to export
                "method": None
            },
            "how_to": {
                "description": "Get step-by-step instructions",
                "parameters": {"task": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "pros_and_cons": {
                "description": "List pros and cons of something",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "timeline": {
                "description": "Get a timeline of events or history",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "mind_map": {
                "description": "Create a mind map of a topic",
                "parameters": {"topic": "string (required)"},
                "endpoint": None,  # Navigate to AI chat
                "method": None
            },
            "flashcard_from_text": {
                "description": "Generate flashcards from pasted text",
                "parameters": {"text": "string (required)", "count": "integer (optional)"},
                "endpoint": None,  # Navigate to flashcard creator
                "method": None
            },
            "daily_review": {
                "description": "Start daily review session",
                "parameters": {},
                "endpoint": None,  # Navigate to review
                "method": None
            },
            "whats_due": {
                "description": "Show what's due for review today",
                "parameters": {},
                "endpoint": "/api/adaptive/retention",
                "method": "GET"
            },
            "random_flashcard": {
                "description": "Show a random flashcard",
                "parameters": {"topic": "string (optional)"},
                "endpoint": None,  # Navigate to flashcards
                "method": None
            }
        }
        
        # Create AI prompt with strict JSON format requirement
        ai_prompt = f"""You are an intent detection system. Analyze the user's query and determine their intent.

USER QUERY: "{query}"

AVAILABLE ACTIONS:
{json.dumps(available_actions, indent=2)}

INSTRUCTIONS:
1. Identify the user's primary intent from the available actions
2. Extract any parameters mentioned in the query
3. Return ONLY valid JSON in this EXACT format (no markdown, no explanation):

{{
  "intent": "action" or "search",
  "action": "action_name" or null,
  "parameters": {{}},
  "confidence": 0.0 to 1.0
}}

RULES:
- If the query matches an action, set intent="action" and specify the action name
- If it's a general search query, set intent="search" and action=null
- Extract parameters EXACTLY as specified in the action definition
- For topics/titles, extract the actual subject matter from the query
- confidence should be 0.8+ for clear matches, 0.5-0.8 for uncertain, <0.5 for unclear
- Return ONLY the JSON object, nothing else

IMPORTANT - SEARCH INTENT KEYWORDS:
The following words/phrases indicate the user wants to SEARCH for existing content (intent="search"):
- "fetch", "get", "get me", "bring", "bring me"
- "find", "find me", "look for", "look up", "search", "search for"
- "can I see", "show me", "display", "pull up", "retrieve"
- "where is", "where are", "locate", "give me"
- "my flashcards on", "my notes on", "my content on"
- "do I have", "have I created"

When these keywords are followed by a topic (e.g., "fetch my flashcards on irish", "get me notes on biology"), 
the intent is SEARCH, NOT create. The user wants to find EXISTING content.

EXAMPLES:
Query: "create flashcards on machine learning"
{{"intent": "action", "action": "create_flashcards", "parameters": {{"topic": "machine learning", "count": 10}}, "confidence": 0.95}}

Query: "what is my learning style"
{{"intent": "action", "action": "show_learning_style", "parameters": {{}}, "confidence": 0.98}}

Query: "adapt difficulty to my level"
{{"intent": "action", "action": "adapt_difficulty", "parameters": {{}}, "confidence": 0.95}}

Query: "explain neural networks step by step"
{{"intent": "action", "action": "tutor_step_by_step", "parameters": {{"topic": "neural networks"}}, "confidence": 0.92}}

Query: "python programming"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.85}}

Query: "fetch my flashcards on irish"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.95}}

Query: "get me notes on biology"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.95}}

Query: "can I see my flashcards on history"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.95}}

Query: "bring me my notes"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.90}}

Query: "find flashcards on chemistry"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.95}}

Query: "show me my content on math"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.95}}

Query: "do I have any notes on physics"
{{"intent": "search", "action": null, "parameters": {{}}, "confidence": 0.90}}

ADDITIONAL NLP EXAMPLES:

Query: "compare python and javascript"
{{"intent": "action", "action": "compare_topics", "parameters": {{"topics": ["python", "javascript"]}}, "confidence": 0.95}}

Query: "explain quantum physics like I'm 5"
{{"intent": "action", "action": "explain_like_im_five", "parameters": {{"topic": "quantum physics"}}, "confidence": 0.98}}

Query: "give me 5 examples of recursion"
{{"intent": "action", "action": "give_examples", "parameters": {{"topic": "recursion", "count": 5}}, "confidence": 0.95}}

Query: "test me on calculus"
{{"intent": "action", "action": "test_me", "parameters": {{"topic": "calculus"}}, "confidence": 0.95}}

Query: "what is photosynthesis"
{{"intent": "action", "action": "define", "parameters": {{"term": "photosynthesis"}}, "confidence": 0.90}}

Query: "what do I need to know before learning machine learning"
{{"intent": "action", "action": "list_prerequisites", "parameters": {{"topic": "machine learning"}}, "confidence": 0.92}}

Query: "give me practice problems on algebra"
{{"intent": "action", "action": "practice_problems", "parameters": {{"topic": "algebra"}}, "confidence": 0.95}}

Query: "summarize world war 2 in short"
{{"intent": "action", "action": "summarize_topic", "parameters": {{"topic": "world war 2", "length": "short"}}, "confidence": 0.93}}

Query: "how am I doing this week"
{{"intent": "action", "action": "show_statistics", "parameters": {{"timeframe": "week"}}, "confidence": 0.88}}

Query: "remind me to study biology tomorrow"
{{"intent": "action", "action": "remind_me", "parameters": {{"topic": "biology", "time": "tomorrow"}}, "confidence": 0.90}}

Query: "how to solve quadratic equations"
{{"intent": "action", "action": "how_to", "parameters": {{"task": "solve quadratic equations"}}, "confidence": 0.92}}

Query: "pros and cons of renewable energy"
{{"intent": "action", "action": "pros_and_cons", "parameters": {{"topic": "renewable energy"}}, "confidence": 0.95}}

Query: "timeline of the french revolution"
{{"intent": "action", "action": "timeline", "parameters": {{"topic": "french revolution"}}, "confidence": 0.95}}

Query: "what's due today"
{{"intent": "action", "action": "whats_due", "parameters": {{}}, "confidence": 0.95}}

Query: "start my daily review"
{{"intent": "action", "action": "daily_review", "parameters": {{}}, "confidence": 0.95}}

Query: "show me a random flashcard"
{{"intent": "action", "action": "random_flashcard", "parameters": {{}}, "confidence": 0.92}}

Query: "quiz me on hard chemistry questions"
{{"intent": "action", "action": "test_me", "parameters": {{"topic": "chemistry", "difficulty": "hard"}}, "confidence": 0.95}}

NOW ANALYZE THIS QUERY AND RETURN ONLY THE JSON:
"{query}"
"""
        
        # Call AI to parse intent
        logger.info(f"🤖 Calling AI for intent detection: '{query}'")
        ai_response = call_ai(ai_prompt, max_tokens=500, temperature=0.1)
        
        logger.info(f"📥 AI response: {ai_response[:200]}...")
        
        # Parse JSON response (handle markdown code blocks if present)
        ai_response_clean = ai_response.strip()
        
        # Remove markdown code blocks if present
        if ai_response_clean.startswith('```'):
            # Extract JSON from code block
            lines = ai_response_clean.split('\n')
            json_lines = []
            in_code_block = False
            for line in lines:
                if line.strip().startswith('```'):
                    in_code_block = not in_code_block
                    continue
                if in_code_block or (not line.strip().startswith('```')):
                    json_lines.append(line)
            ai_response_clean = '\n'.join(json_lines).strip()
        
        # Parse JSON
        try:
            result = json.loads(ai_response_clean)
            
            # Validate result structure
            if not isinstance(result, dict):
                raise ValueError("Response is not a JSON object")
            
            if "intent" not in result:
                raise ValueError("Missing 'intent' field")
            
            # Ensure all required fields exist
            result.setdefault("action", None)
            result.setdefault("parameters", {})
            result.setdefault("confidence", 0.5)
            
            # Log successful detection
            logger.info(f" Intent detected - Intent: {result['intent']}, Action: {result['action']}, Confidence: {result['confidence']}")
            
            # Add original query
            result["original_query"] = query
            
            return result
            
        except json.JSONDecodeError as e:
            logger.error(f" Failed to parse AI response as JSON: {e}")
            logger.error(f"Raw response: {ai_response_clean}")
            
            # Fallback: try to extract action from response text
            action_match = None
            for action_name in available_actions.keys():
                if action_name in ai_response_clean.lower():
                    action_match = action_name
                    break
            
            if action_match:
                logger.info(f" Fallback: Extracted action '{action_match}' from text")
                return {
                    "intent": "action",
                    "action": action_match,
                    "parameters": {},
                    "confidence": 0.6,
                    "original_query": query
                }
            else:
                # Default to search
                logger.info(f" Fallback: Defaulting to search")
                return {
                    "intent": "search",
                    "action": None,
                    "parameters": {},
                    "confidence": 0.5,
                    "original_query": query
                }
        
    except Exception as e:
        logger.error(f" Intent detection error: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Fallback to search on error
        return {
            "intent": "search",
            "action": None,
            "parameters": {},
            "confidence": 0.3,
            "original_query": query,
            "error": str(e)
        }


@app.post("/api/generate_topic_description")
async def generate_topic_description(
    user_id: str = Form(...),
    topic: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Generate a brief, educational AI description of a topic when no search results are found.
    This helps users understand the topic and provides context before creating study materials.
    """
    try:
        logger.info(f"🤖 Generating topic description for: '{topic}' (user: {user_id})")
        
        # Get user profile for personalization
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        
        # Build AI prompt for topic description
        ai_prompt = f"""You are an educational AI assistant. Provide a brief, clear, and engaging description of the following topic.

TOPIC: "{topic}"

INSTRUCTIONS:
1. Write 2-3 sentences that explain what this topic is about
2. Make it educational and informative
3. Use simple, accessible language
4. Focus on the core concepts and why it matters
5. Be encouraging and make the user curious to learn more
6. DO NOT use phrases like "I couldn't find" or "no results" - just describe the topic
7. Keep it concise (max 150 words)

{MATHEMATICAL_FORMATTING_INSTRUCTIONS if any(char in topic.lower() for char in ['math', 'calculus', 'algebra', 'equation', 'formula']) else ''}

EXAMPLES:

Topic: "quantum physics"
Description: "Quantum physics is the branch of physics that studies the behavior of matter and energy at the smallest scales - atoms and subatomic particles. It reveals a fascinating world where particles can exist in multiple states simultaneously and can be connected across vast distances. Understanding quantum physics opens doors to cutting-edge technologies like quantum computing and helps explain the fundamental nature of our universe."

Topic: "machine learning"
Description: "Machine learning is a field of artificial intelligence that enables computers to learn and improve from experience without being explicitly programmed. It powers many technologies you use daily, from recommendation systems to voice assistants. By studying patterns in data, machine learning algorithms can make predictions, recognize images, understand language, and solve complex problems."

Topic: "photosynthesis"
Description: "Photosynthesis is the remarkable process by which plants convert sunlight into chemical energy, producing the oxygen we breathe and the food that sustains life on Earth. During this process, plants use chlorophyll to capture light energy and transform carbon dioxide and water into glucose and oxygen. It's one of the most important biological processes on our planet."

Now generate a description for: "{topic}"

Return ONLY the description text, no labels or extra formatting."""

        # Call AI to generate description
        description = call_ai(ai_prompt, max_tokens=300, temperature=0.7)
        
        # Clean up the response
        description = description.strip()
        
        # Remove any labels if present
        if description.lower().startswith("description:"):
            description = description[12:].strip()
        
        logger.info(f" Generated description: {description[:100]}...")
        
        return {
            "success": True,
            "description": description,
            "topic": topic,
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f" Error generating topic description: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Return a generic but helpful fallback
        return {
            "success": True,
            "description": f"Let's explore {topic} together! This is a fascinating subject with many interesting aspects to discover. I can help you create flashcards, notes, or start a learning session to dive deeper into this topic.",
            "topic": topic,
            "fallback": True,
            "timestamp": datetime.now().isoformat()
        }


@app.post("/api/get_personalized_prompts")
async def get_personalized_prompts(
    user_id: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Get personalized prompt suggestions based on user's profile, activity, and weak areas.
    Returns 4 topic-based prompts from user activity + adaptive learning prompts.
    """
    try:
        # Get user
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        topic_prompts = []  # Topic-based prompts from user activity
        
        import random
        
        # 1. Get recent flashcard sets (HIGH PRIORITY - user's active topics)
        recent_flashcard_sets = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.user_id == user.id
        ).order_by(models.FlashcardSet.updated_at.desc()).limit(10).all()
        
        if recent_flashcard_sets:
            selected_sets = random.sample(recent_flashcard_sets, min(2, len(recent_flashcard_sets)))
            for fs in selected_sets:
                topic_prompts.append({
                    "text": f"create a quiz on {fs.title}",
                    "reason": "Test your flashcard knowledge",
                    "priority": "high"
                })
        
        # 2. Get recent notes (HIGH PRIORITY - user's study topics)
        recent_notes = db.query(models.Note).filter(
            models.Note.user_id == user.id
        ).order_by(models.Note.updated_at.desc()).limit(10).all()
        
        if recent_notes:
            selected_notes = random.sample(recent_notes, min(2, len(recent_notes)))
            for note in selected_notes:
                topic_prompts.append({
                    "text": f"create flashcards on {note.title}",
                    "reason": "Turn notes into active learning",
                    "priority": "high"
                })
        
        # 3. Get recent chat sessions (MEDIUM PRIORITY - topics user asked about)
        recent_chats = db.query(models.ChatSession).filter(
            models.ChatSession.user_id == user.id
        ).order_by(models.ChatSession.updated_at.desc()).limit(10).all()
        
        if recent_chats:
            # Filter out generic titles
            meaningful_chats = [c for c in recent_chats if c.title and c.title.lower() not in ['new chat', 'untitled', '']]
            if meaningful_chats:
                selected_chats = random.sample(meaningful_chats, min(2, len(meaningful_chats)))
                for chat in selected_chats:
                    topic_prompts.append({
                        "text": f"explain {chat.title} step-by-step",
                        "reason": "Continue learning this topic",
                        "priority": "medium"
                    })
        
        # 4. Get weak areas from profile (HIGH PRIORITY)
        profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()
        
        weak_topics = []
        if profile and profile.weak_areas:
            try:
                weak_topics = json.loads(profile.weak_areas) if isinstance(profile.weak_areas, str) else profile.weak_areas
            except:
                pass
        
        if weak_topics:
            selected_weak = random.sample(weak_topics, min(2, len(weak_topics)))
            for topic in selected_weak:
                topic_prompts.append({
                    "text": f"create flashcards on {topic}",
                    "reason": "Focus on weak areas",
                    "priority": "high"
                })
        
        # 5. Check for flashcards needing review
        weak_flashcard_sets = db.query(models.FlashcardSet).join(
            models.Flashcard
        ).filter(
            models.FlashcardSet.user_id == user.id,
            models.Flashcard.marked_for_review == True
        ).distinct().limit(5).all()
        
        if weak_flashcard_sets:
            topic_prompts.append({
                "text": "review weak flashcards",
                "reason": f"{len(weak_flashcard_sets)} sets need attention",
                "priority": "high"
            })
        
        # Shuffle and deduplicate topic prompts
        random.shuffle(topic_prompts)
        seen_texts = set()
        unique_prompts = []
        for p in topic_prompts:
            text_lower = p["text"].lower()
            if text_lower not in seen_texts:
                seen_texts.add(text_lower)
                unique_prompts.append(p)
        
        # Sort by priority
        priority_order = {"high": 0, "medium": 1, "low": 2}
        unique_prompts.sort(key=lambda x: (priority_order.get(x["priority"], 3), random.random()))
        
        logger.info(f"Generated {len(unique_prompts)} personalized prompts for user {user_id}")
        
        return {
            "prompts": unique_prompts[:4],  # Return top 4 topic-based prompts
            "user_id": user_id
        }
        
    except Exception as e:
        logger.error(f"Error getting personalized prompts: {str(e)}")
        # Return empty - frontend will use defaults
        return {
            "prompts": []
        }


@app.post("/api/get_flashcards_for_review")
async def get_flashcards_for_review(
    user_id: str = Form(...),
    filter: str = Form("needs_review"),  # "needs_review", "marked_for_review", "all"
    db: Session = Depends(get_db)
):
    """
    Get flashcards that need review based on user's performance
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"flashcards": [], "set_id": None}
        
        # Get flashcards based on filter
        query_builder = db.query(models.Flashcard).join(
            models.FlashcardSet
        ).filter(
            models.FlashcardSet.user_id == user.id
        )
        
        if filter == "marked_for_review":
            query_builder = query_builder.filter(models.Flashcard.marked_for_review == True)
        elif filter == "needs_review":
            # Get cards with low accuracy or not reviewed recently
            query_builder = query_builder.filter(
                or_(
                    models.Flashcard.times_reviewed == 0,
                    and_(
                        models.Flashcard.times_reviewed > 0,
                        (models.Flashcard.correct_count * 100.0 / models.Flashcard.times_reviewed) < 70
                    )
                )
            )
        
        flashcards = query_builder.limit(50).all()
        
        if flashcards:
            # Group by set and return the set with most cards needing review
            set_counts = {}
            for card in flashcards:
                set_counts[card.set_id] = set_counts.get(card.set_id, 0) + 1
            
            best_set_id = max(set_counts, key=set_counts.get)
            
            return {
                "set_id": best_set_id,
                "flashcards": [
                    {
                        "id": card.id,
                        "question": card.question,
                        "answer": card.answer,
                        "difficulty": card.difficulty,
                        "times_reviewed": card.times_reviewed,
                        "correct_count": card.correct_count
                    }
                    for card in flashcards if card.set_id == best_set_id
                ],
                "total_cards_needing_review": len(flashcards)
            }
        
        return {"flashcards": [], "set_id": None, "total_cards_needing_review": 0}
        
    except Exception as e:
        logger.error(f"Error getting flashcards for review: {str(e)}")
        return {"flashcards": [], "set_id": None, "error": str(e)}


@app.post("/api/get_weak_areas")
async def get_weak_areas(
    user_id: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Get user's weak areas based on performance data.
    ALWAYS returns topics - never empty!
    """
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"weak_areas": []}
        
        # Get comprehensive profile
        comprehensive_profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()
        
        weak_areas = []
        
        # 1. Check profile weak areas
        if comprehensive_profile and comprehensive_profile.weak_areas:
            try:
                weak_areas_list = json.loads(comprehensive_profile.weak_areas)
                for area in weak_areas_list:
                    weak_areas.append({
                        "id": len(weak_areas) + 1,
                        "type": "weak_area",
                        "title": area,
                        "description": f"You've been struggling with {area}",
                        "created_at": datetime.now(timezone.utc).isoformat()
                    })
            except:
                pass
        
        # 2. Check solo quiz performance
        solo_quizzes = db.query(models.SoloQuiz).filter(
            models.SoloQuiz.user_id == user.id,
            models.SoloQuiz.completed == True
        ).all()
        
        # Group quizzes by subject and calculate average score
        quiz_performance = {}
        for quiz in solo_quizzes:
            subject = quiz.subject
            if subject not in quiz_performance:
                quiz_performance[subject] = {"total_score": 0, "total_questions": 0, "quiz_count": 0}
            
            quiz_performance[subject]["total_score"] += quiz.score
            quiz_performance[subject]["total_questions"] += quiz.question_count
            quiz_performance[subject]["quiz_count"] += 1
        
        # Find weak subjects from quizzes (less than 60% accuracy)
        for subject, perf in quiz_performance.items():
            if perf["total_questions"] > 0:
                accuracy = (perf["total_score"] / perf["total_questions"]) * 100
                
                if accuracy < 60:  # Less than 60% accuracy
                    weak_areas.append({
                        "id": len(weak_areas) + 1,
                        "type": "quiz_subject",
                        "title": subject,
                        "description": f"Quiz accuracy: {accuracy:.1f}% across {perf['quiz_count']} quiz(es) - Needs more practice",
                        "created_at": datetime.now(timezone.utc).isoformat(),
                        "accuracy": accuracy
                    })
        
        # 3. Check flashcard performance
        flashcard_sets = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.user_id == user.id
        ).all()
        
        for fset in flashcard_sets:
            # Calculate set accuracy
            cards = db.query(models.Flashcard).filter(
                models.Flashcard.set_id == fset.id,
                models.Flashcard.times_reviewed > 0
            ).all()
            
            if cards:
                total_reviews = sum(c.times_reviewed for c in cards)
                total_correct = sum(c.correct_count for c in cards)
                
                if total_reviews > 0:
                    accuracy = (total_correct / total_reviews) * 100
                    
                    if accuracy < 60:  # Less than 60% accuracy
                        weak_areas.append({
                            "id": len(weak_areas) + 1,
                            "type": "flashcard_set",
                            "title": fset.title,
                            "description": f"Flashcard accuracy: {accuracy:.1f}% - Needs more practice",
                            "created_at": fset.created_at.isoformat() if fset.created_at else None,
                            "set_id": fset.id
                        })
        
        # 4. If still no weak areas, suggest topics based on user activity
        if len(weak_areas) == 0:
            # Get topics from flashcard sets
            recent_sets = db.query(models.FlashcardSet).filter(
                models.FlashcardSet.user_id == user.id
            ).order_by(models.FlashcardSet.created_at.desc()).limit(5).all()
            
            for fset in recent_sets:
                weak_areas.append({
                    "id": len(weak_areas) + 1,
                    "type": "suggestion",
                    "title": f"Review {fset.title}",
                    "description": "Keep your knowledge fresh with regular review",
                    "created_at": fset.created_at.isoformat() if fset.created_at else None,
                    "set_id": fset.id
                })
            
            # Get topics from notes
            recent_notes = db.query(models.Note).filter(
                models.Note.user_id == user.id
            ).order_by(models.Note.updated_at.desc()).limit(3).all()
            
            for note in recent_notes:
                weak_areas.append({
                    "id": len(weak_areas) + 1,
                    "type": "suggestion",
                    "title": f"Create flashcards on {note.title}",
                    "description": "Turn your notes into active learning",
                    "created_at": note.updated_at.isoformat() if note.updated_at else None
                })
        
        # 5. If STILL no weak areas, suggest based on field of study
        if len(weak_areas) == 0:
            field_of_study = comprehensive_profile.field_of_study if comprehensive_profile else user.field_of_study
            
            # Default suggestions based on common learning topics
            default_topics = [
                "Create your first flashcard set",
                "Take notes on a new topic",
                "Start a study session",
                "Explore AI chat for learning",
                "Set up your study goals"
            ]
            
            if field_of_study and field_of_study.lower() not in ['general studies', 'general', 'none', '']:
                default_topics = [
                    f"Core concepts in {field_of_study}",
                    f"Advanced topics in {field_of_study}",
                    f"Practice problems for {field_of_study}",
                    f"Review fundamentals of {field_of_study}",
                    f"Explore new areas in {field_of_study}"
                ]
            
            for i, topic in enumerate(default_topics[:5]):
                weak_areas.append({
                    "id": i + 1,
                    "type": "suggestion",
                    "title": topic,
                    "description": "Start building your knowledge base",
                    "created_at": datetime.now(timezone.utc).isoformat()
                })
        
        # Sort weak areas by type priority: quiz_subject and flashcard_set first, then suggestions
        type_priority = {"weak_area": 0, "quiz_subject": 1, "flashcard_set": 2, "suggestion": 3}
        weak_areas.sort(key=lambda x: type_priority.get(x.get("type", "suggestion"), 3))
        
        logger.info(f"Returning {len(weak_areas)} weak areas/suggestions for user {user_id}")
        return {"weak_areas": weak_areas[:10]}  # Return top 10
        
    except Exception as e:
        logger.error(f"Error getting weak areas: {str(e)}")
        # Even on error, return some default suggestions
        return {
            "weak_areas": [
                {
                    "id": 1,
                    "type": "suggestion",
                    "title": "Create your first flashcard set",
                    "description": "Start your learning journey",
                    "created_at": datetime.now(timezone.utc).isoformat()
                },
                {
                    "id": 2,
                    "type": "suggestion",
                    "title": "Take notes on a topic you're learning",
                    "description": "Document your knowledge",
                    "created_at": datetime.now(timezone.utc).isoformat()
                },
                {
                    "id": 3,
                    "type": "suggestion",
                    "title": "Ask AI to explain a concept",
                    "description": "Get personalized help",
                    "created_at": datetime.now(timezone.utc).isoformat()
                }
            ]
        }


# ==================== NEW SMART SEARCH HUB FEATURES ====================

@app.post("/api/suggest_study_next")
async def suggest_study_next(
    user_id: str = Form(...),
    db: Session = Depends(get_db)
):
    """AI suggests what the user should study next"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"suggestions": []}
        
        suggestions = []
        profile = db.query(models.ComprehensiveUserProfile).filter(
            models.ComprehensiveUserProfile.user_id == user.id
        ).first()
        
        if profile and profile.weak_areas:
            try:
                weak_topics = json.loads(profile.weak_areas) if isinstance(profile.weak_areas, str) else profile.weak_areas
                for topic in weak_topics[:3]:
                    suggestions.append({
                        "topic": topic,
                        "reason": "You've been struggling with this",
                        "priority": "high"
                    })
            except:
                pass
        
        return {"suggestions": suggestions[:5]}
    except Exception as e:
        logger.error(f"Error suggesting study next: {str(e)}")
        return {"suggestions": []}


@app.post("/api/summarize_notes")
async def summarize_notes(
    user_id: str = Form(...),
    topic: str = Form(None),
    db: Session = Depends(get_db)
):
    """AI summarizes user's notes on a topic"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"summary": "User not found"}
        
        notes_query = db.query(models.Note).filter(models.Note.user_id == user.id)
        if topic:
            notes_query = notes_query.filter(
                or_(models.Note.title.ilike(f"%{topic}%"), models.Note.content.ilike(f"%{topic}%"))
            )
        
        notes = notes_query.order_by(models.Note.updated_at.desc()).limit(10).all()
        if not notes:
            return {"summary": f"No notes found on {topic}" if topic else "No notes found"}
        
        combined_content = "\n\n".join([f"# {note.title}\n{note.content}" for note in notes])
        groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
        
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "Study assistant."},
                {"role": "user", "content": f"Summarize these notes:\n\n{combined_content[:4000]}"}
            ],
            temperature=0.7,
            max_tokens=1000
        )
        
        return {"summary": response.choices[0].message.content.strip(), "notes_count": len(notes)}
    except Exception as e:
        logger.error(f"Error summarizing notes: {str(e)}")
        return {"summary": f"Error: {str(e)}"}


@app.post("/api/create_study_plan")
async def create_study_plan(
    user_id: str = Form(...),
    topic: str = Form(...),
    duration: int = Form(30),
    db: Session = Depends(get_db)
):
    """AI creates a study plan"""
    try:
        groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "Create study plan."},
                {"role": "user", "content": f"Create a {duration} day study plan for {topic}"}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        return {"plan": response.choices[0].message.content.strip()}
    except Exception as e:
        return {"plan": f"Error: {str(e)}"}


@app.post("/api/search_recent_content")
async def search_recent_content(
    user_id: str = Form(...),
    timeframe: str = Form("recent"),
    db: Session = Depends(get_db)
):
    """Search content by timeframe"""
    try:
        user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
        if not user:
            return {"results": []}
        
        from datetime import timedelta
        now = datetime.now(timezone.utc)
        start_date = now - timedelta(days={"yesterday": 1, "last_week": 7, "last_month": 30}.get(timeframe, 3))
        
        results = []
        sets = db.query(models.FlashcardSet).filter(
            models.FlashcardSet.user_id == user.id,
            models.FlashcardSet.created_at >= start_date
        ).all()
        
        for fset in sets:
            results.append({"id": fset.id, "type": "flashcard_set", "title": fset.title})
        
        return {"results": results}
    except Exception as e:
        return {"results": []}


@app.post("/api/get_search_suggestion")
async def get_search_suggestion(
    user_id: str = Form(...),
    query: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Get AI-powered suggestions when no search results are found
    """
    try:
        groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
        
        prompt = f"""The user searched for "{query}" but no results were found.
        
Provide a helpful, encouraging response that:
1. Acknowledges what they're looking for
2. Suggests what type of content they could create (flashcards, notes, questions, etc.)
3. Offers to help them get started

Keep it brief, friendly, and actionable. 2-3 sentences max."""

        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "AI tutor assistant."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=200
        )
        
        suggestion = response.choices[0].message.content.strip()
        
        return {
            "description": suggestion,
            "suggestions": [
                "Create Flashcards",
                "Take Notes",
                "Ask AI"
            ]
        }
        
    except Exception as e:
        logger.error(f"AI suggestion error: {str(e)}")
        return {
            "description": f"I couldn't find anything matching \"{query}\". Would you like to create some learning materials on this topic?",
            "suggestions": [
                "Create Flashcards",
                "Take Notes",
                "Ask AI"
            ]
        }


@app.get("/api/get_trending_topics")
async def get_trending_topics(
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Get trending topics based on recent public content
    """
    try:
        # Get most popular public decks
        popular_decks = db.query(
            models.FlashcardDeck.name,
            func.count(models.FlashcardDeck.id).label('count')
        ).filter(
            and_(
                models.FlashcardDeck.is_public == True,
                models.FlashcardDeck.is_deleted == False
            )
        ).group_by(models.FlashcardDeck.name).order_by(
            func.count(models.FlashcardDeck.id).desc()
        ).limit(5).all()
        
        trending = []
        for deck_name, count in popular_decks:
            trending.append({
                "topic": deck_name,
                "count": count
            })
        
        # If not enough, add some defaults
        if len(trending) < 5:
            defaults = [
                {"topic": "Machine Learning Basics", "count": 45},
                {"topic": "Calculus Integration", "count": 38},
                {"topic": "World History Timeline", "count": 32},
                {"topic": "Python Programming", "count": 29},
                {"topic": "Chemistry Reactions", "count": 24}
            ]
            trending.extend(defaults[len(trending):])
        
        return {"trending": trending[:5]}
        
    except Exception as e:
        logger.error(f"Trending topics error: {str(e)}")
        return {
            "trending": [
                {"topic": "Machine Learning Basics", "count": 45},
                {"topic": "Calculus Integration", "count": 38},
                {"topic": "World History Timeline", "count": 32},
                {"topic": "Python Programming", "count": 29},
                {"topic": "Chemistry Reactions", "count": 24}
            ]
        }

# ==================== END SEARCH HUB ENDPOINTS ====================


if __name__ == "__main__":
    import uvicorn
    
    print("Starting Cerbyl Backend...")
    
    uvicorn.run(app, host="127.0.0.1", port=8000, reload=False)




